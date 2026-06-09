from __future__ import annotations

import ast
import json
import os
import re
import subprocess
import sys
from pathlib import Path

import pytest
from PIL import Image

from scripts.build_ppt_contact_sheet import build_contact_sheet
from scripts.build_ppt_run_html_viewer import DEFAULT_THREAD_ID, build_data
from scripts.check_ppt_layout_quality import check_layout, write_report
from scripts.validate_ppt_case_pack import validate_case_pack


ROOT = Path(__file__).resolve().parents[1]
PACK = ROOT / "docs" / "product" / "ppt-run2-data-skill-quality"
pytestmark = pytest.mark.skipif(
    os.environ.get("GITHUB_ACTIONS") == "true",
    reason=(
        "PPT Run 2 case-pack tests require local generated presentation outputs "
        "and the Codex desktop artifact-tool runtime."
    ),
)
EXPECTED_ARMS = {"prompt_only", "run1_5_skill", "run2_skill", "bad_aesthetic_memory"}
EXPECTED_SOURCE_CARD_IDS = {
    "card_cinematic_cover",
    "card_low_density_claim",
    "card_big_object_layout",
    "card_editorial_comparison",
    "card_proof_reveal",
    "card_visual_climax",
    "card_premium_closing",
    "card_appendix_absorption",
}
EXPECTED_VIDEO_CARD_IDS = {
    "video_keynote_rhythm",
    "video_transition_pacing",
}
EXPECTED_MULTIMODAL_MODALITIES = {"text", "image_reference", "video", "audio", "transcript", "interaction"}
EXPECTED_MULTIMODAL_RECORD_IDS = {
    "mm_text_visual_hierarchy_tutorial",
    "mm_text_whitespace_tutorial",
    "mm_video_audio_keynote_rhythm",
    "mm_interaction_presentation_product_reference",
    "mm_case_public_keynote_reference",
    "mm_duarte_visual_storytelling_webinar",
    "mm_duarte_slide_design_course",
    "mm_udel_design_principles_video",
}
EXPECTED_VISUAL_TARGET_IDS = {
    "target_report_to_visual_delta",
    "target_slide_mini_preview",
    "target_audio_rhythm_budget",
    "target_transcript_headline_compression",
    "target_public_demo_climax",
}
EXPECTED_VISUAL_COMPONENT_IDS = {
    "component_before_after_thumbnail",
    "component_slide_mini_preview",
    "component_rhythm_budget_strip",
    "component_transcript_headline_route",
    "component_public_demo_climax_object",
}
EXPECTED_VIDEO_BEAT_IDS = {
    "beat_opening_scale_reset",
    "beat_before_after_transformation",
    "beat_proof_reveal_cadence",
    "beat_climax_scale_up",
    "beat_close_release_handoff",
}
EXPECTED_MOTION_TARGET_IDS = {
    "motion_target_opening_attention_reset",
    "motion_target_before_after_reveal",
    "motion_target_proof_build",
    "motion_target_climax_scale_emphasis",
    "motion_target_release_gate_handoff",
}
EXPECTED_SEQUENCE_COMPONENT_IDS = {
    "sequence_component_opening_reset",
    "sequence_component_before_after_reveal",
    "sequence_component_proof_build",
    "sequence_component_climax_scale",
    "sequence_component_release_gate",
}
EXPECTED_PRODUCTION_REFERENCE_IDS = {
    "prod_ref_cinematic_cover_field",
    "prod_ref_editorial_before_after_delta",
    "prod_ref_product_system_mini_preview",
    "prod_ref_climax_handoff_sequence",
}
EXPECTED_PRODUCTION_MODULE_IDS = {
    "module_cinematic_cover_field",
    "module_editorial_before_after_delta",
    "module_proof_route_choreography",
    "module_system_mini_preview",
    "module_climax_hero_object",
    "module_release_handoff_gate",
}
EXPECTED_AESTHETIC_V2_MOVE_IDS = {
    "aesthetic_v2_cinematic_cover_field",
    "aesthetic_v2_editorial_before_after_delta",
    "aesthetic_v2_proof_route_choreography",
    "aesthetic_v2_system_mini_preview",
    "aesthetic_v2_climax_hero_object",
    "aesthetic_v2_release_handoff_gate",
}
EXPECTED_PRODUCTION_REFERENCE_FIELDS = {
    "id",
    "source_ids",
    "visual_component_ids",
    "motion_target_ids",
    "sequence_component_ids",
    "visual_observations",
    "composition_primitives",
    "typography_primitives",
    "spacing_primitives",
    "motion_or_sequence_primitives",
    "module_implications",
    "extraction_notes",
    "do_not_copy",
    "qa_probe",
    "release_boundary",
}
EXPECTED_AESTHETIC_V2_FIELDS = {
    "id",
    "production_reference_ids",
    "visual_component_ids",
    "motion_target_ids",
    "sequence_component_ids",
    "slide_roles",
    "composition_contract",
    "typography_contract",
    "spacing_contract",
    "density_budget",
    "provenance_boundary",
    "code_generation_contract",
    "trace_fields",
    "forbidden_report_patterns",
    "qa_probe",
    "release_boundary",
}
EXPECTED_VISUAL_PRODUCTION_MODULE_FIELDS = {
    "id",
    "production_reference_ids",
    "aesthetic_memory_v2_ids",
    "slide_roles",
    "native_ppt_primitives",
    "layout_recipe",
    "typography_recipe",
    "spacing_recipe",
    "fallback_policy",
    "provenance_boundary",
    "trace_fields",
    "qa_probe",
    "release_boundary",
}
EXPECTED_RUN2_6_SOURCE_IDS = {
    "figma_config_2025_platform_launch",
    "figma_config_2025_recap",
    "figma_slides_product",
    "figma_slides_help",
    "stripe_sessions_2025_product_keynote",
    "google_cloud_next_2025_wrap",
    "google_cloud_next_2025_sundar",
    "apple_liquid_glass_newsroom",
    "apple_liquid_glass_developer",
    "duarte_slide_design",
    "duarte_visual_storytelling",
    "slidemodel_visual_hierarchy",
}
EXPECTED_RUN2_6_USECASE_IDS = {
    "usecase_design_to_production_platform_launch",
    "usecase_fintech_product_keynote",
    "usecase_ai_cloud_keynote_demo",
    "usecase_design_language_public_reveal",
}
EXPECTED_RUN2_6_BENCHMARK_IDS = {
    "benchmark_design_to_production_grid_precision",
    "benchmark_visual_fidelity_interactive_slide_surface",
    "benchmark_fintech_product_keynote_breadth_without_grid",
    "benchmark_ai_platform_demo_climax",
    "benchmark_content_first_dynamic_material",
    "benchmark_glance_test_visual_hierarchy",
    "benchmark_story_driven_data_emphasis",
}
EXPECTED_RUN2_6_TRACE_FIELDS = {
    "commercial_usecase_id",
    "aesthetic_benchmark_ids",
    "theme_policy_id",
    "typography_system_id",
    "spacing_token_set_id",
    "workflow_decision_ids",
    "source_brand_sanitization",
    "benchmark_validation_probe",
    "theme_validation_probe",
}
EXPECTED_RUN2_6R_REPAIR_IDS = {
    "repair_editorial_typography_system",
    "repair_spacing_token_visibility",
    "repair_climax_editorial_spread",
    "repair_theme_differentiation_from_run2_5",
    "repair_mini_preview_fidelity",
}
EXPECTED_RUN2_6R_REPAIR_FIELDS = {
    "id",
    "target_slide_roles",
    "source_policy_ids",
    "typography_delta",
    "spacing_delta",
    "composition_delta",
    "theme_delta",
    "must_differ_from",
    "native_ppt_requirements",
    "qa_probe",
    "release_boundary",
}
EXPECTED_RUN2_6R_TRACE_FIELDS = {
    "visual_repair_policy_ids",
    "visual_delta_from_run2_5",
    "visual_repair_validation_probe",
}
EXPECTED_RUN2_7_USECASE_IDS = {
    "usecase_ai_design_to_production_platform_launch",
}
EXPECTED_RUN2_7_SOURCE_RECORD_IDS = {
    "mm_2_7_design_system_launch_case",
    "mm_2_7_video_climax_single_object",
    "mm_2_7_typography_hierarchy_tutorial",
    "mm_2_7_spacing_editorial_grid_tutorial",
    "mm_2_7_motion_demo_pacing_reference",
    "mm_2_7_product_surface_interaction_reference",
}
EXPECTED_RUN2_7_MEMORY_IDS = {
    "memory_typography_editorial_launch",
    "memory_spacing_climax_proof_grid",
    "memory_composition_single_object_climax",
    "memory_rhythm_six_slide_launch",
    "memory_source_brand_sanitization_v2",
}
EXPECTED_RUN2_7_SOURCE_RECORD_FIELDS = {
    "id",
    "source_id",
    "source_type",
    "allowed_use",
    "anchor",
    "modalities",
    "visual_observation",
    "transcript_or_teaching_claim",
    "extracted_design_rule",
    "slide_roles",
    "native_ppt_implication",
    "anti_copy_boundary",
    "qa_probe",
    "release_boundary",
}
EXPECTED_RUN2_7_MEMORY_FIELDS = {
    "id",
    "source_record_ids",
    "applicable_usecases",
    "applicable_slide_roles",
    "typography_rules",
    "spacing_rules",
    "composition_rules",
    "rhythm_rules",
    "native_ppt_generation_requirements",
    "forbidden_patterns",
    "qa_probes",
    "release_boundary",
}
EXPECTED_RUN2_7_TRACE_FIELDS = {
    "run2_7_usecase_id",
    "run2_7_source_record_ids",
    "run2_7_design_memory_ids",
    "run2_7_workflow_decision_ids",
    "run2_7_delta_from_run2_6r",
    "run2_7_quality_gate",
}
EXPECTED_RUN2_8_DECOMPOSITION_IDS = {
    "decomp_2_8_duarte_remove_nonessential_data",
    "decomp_2_8_type_hierarchy_readability_stack",
    "decomp_2_8_makeover_split_text_into_visual_steps",
    "decomp_2_8_design_first_code_second_pipeline",
    "decomp_2_8_climax_object_scale_and_pause",
    "decomp_2_8_source_brand_sanitized_case_evidence",
}
EXPECTED_RUN2_8_MEMORY_BINDING_IDS = {
    "binding_type_scale_readability",
    "binding_spacing_zone_grid",
    "binding_climax_hero_object",
    "binding_before_after_delta",
    "binding_public_gate_legibility",
}
EXPECTED_RUN2_8_TRACE_FIELDS = {
    "run2_8_decomposition_unit_ids",
    "run2_8_memory_binding_ids",
    "run2_8_gate_matrix_ids",
    "run2_8_code_binding_ids",
    "run2_8_layout_budget",
    "run2_8_visual_delta_from_run2_7",
}
EXPECTED_RUN2_73_SOURCE_AUDIT_FIELDS = {
    "source_id",
    "title",
    "role",
    "source_type",
    "modality",
    "extracted_design_signal",
        "generator_obligation",
        "anti_copy_boundary",
        "missing_decomposition",
        "generator_ready",
        "generator_use_status",
}
EXPECTED_RUN2_73_SOURCE_TYPES = {
    "true_ppt_tutorial",
    "product_keynote",
    "brand_case",
    "product_reference",
    "design_language_reference",
    "general_design_tutorial",
}
EXPECTED_RUN2_73_DECOMPOSITION_DEPTH_VALUES = {"absent", "partial", "usable"}
EXPECTED_RUN2_73_TUTORIAL_MOVE_PRINCIPLES = {
    "visual_hierarchy",
    "white_space",
    "keynote_rhythm",
    "motion_pacing",
}
EXPECTED_RUN2_73_TUTORIAL_MOVE_RULE_FIELDS = {
    "rule_id",
    "principle",
    "tutorial_rule",
    "design_move",
    "ppt_operations",
    "renderer_actions",
    "acceptance_checks",
    "source_rule_ids",
    "visual_target_ids",
    "motion_target_ids",
    "source_lineage",
}
EXPECTED_RUN2_73_RENDERER_ACTION_FIELDS = {
    "action_id",
    "action_type",
    "target_selector",
    "parameters",
    "fallback_strategy",
    "traceability",
}
EXPECTED_RUN2_74_PRODUCT_CLAIM_NODE_IDS = {
    "claim_product_identity",
    "claim_input_boundary",
    "claim_content_compiler",
    "claim_design_move_binding",
    "claim_native_ppt_output",
    "claim_inspectable_trust",
    "claim_release_boundary",
}
EXPECTED_RUN2_74_PRODUCT_CLAIM_NODE_FIELDS = {
    "claim_id",
    "claim_type",
    "public_claim",
    "product_capability_ids",
    "product_logic_relation_ids",
    "source_message_contract_ids",
    "narrative_proof_ids",
    "tutorial_design_rule_ids",
    "user_benefit",
    "visible_user_result",
    "proof_standard",
    "public_surface_route",
}
EXPECTED_RUN2_74_CAUSAL_EDGE_FIELDS = {
    "edge_id",
    "from_claim_id",
    "to_claim_id",
    "causal_relation",
    "business_logic",
    "user_value_delta",
    "proof_refs",
}
EXPECTED_RUN2_74_SLIDE_STORY_FIELDS = {
    "slide_id",
    "slide_index",
    "role",
    "story_function",
    "audience_question",
    "title",
    "thesis",
    "product_claim_ids",
    "causal_edge_ids",
    "source_message_contract_id",
    "narrative_proof_id",
    "proof_object",
    "visual_object",
    "on_canvas_copy",
    "speaker_note_or_viewer_route",
    "text_budget",
    "demo_goal",
    "content_risk",
    "handoff_to_scene_compiler",
}
EXPECTED_RUN2_74_SLIDE_STORY_ROLES = [
    "cover",
    "setup",
    "contrast",
    "proof",
    "climax",
    "close",
]
EXPECTED_RUN2_74_CONTENT_QA_DIMENSIONS = {
    "empty_claim",
    "duplicate_message",
    "claim_without_proof",
    "copy_visual_mismatch",
    "public_surface_routing",
    "text_budget_fit",
}
EXPECTED_RUN2_74_CONTENT_QA_FIELDS = {
    "audit_id",
    "slide_id",
    "role",
    "source_slide_story_id",
    "checked_claim_ids",
    "quality_checks",
    "blocking_issues",
    "watch_items",
    "approved_content_units",
    "scene_compiler_constraints",
}
EXPECTED_RUN2_74_APPROVED_CONTENT_UNIT_FIELDS = {
    "main_claim",
    "proof_object",
    "business_logic",
    "visual_role",
    "on_canvas_copy",
    "delete_or_note_route",
}
EXPECTED_RUN2_D1_SCENE_PLAN_FIELDS = {
    "scene_id",
    "slide_index",
    "role",
    "source_slide_story_id",
    "source_content_qa_id",
    "content_units",
    "proof_object",
    "visual_role",
    "layout_intent",
    "renderer_action_bindings",
    "off_canvas_routes",
    "renderer_execution_boundary",
}
EXPECTED_RUN2_D1_LAYOUT_INTENT_FIELDS = {
    "first_read",
    "main_visual",
    "information_containers",
    "chrome",
    "must_delete",
    "scene_question",
    "rhythm_intent",
}
EXPECTED_RUN2_D1_RENDERER_BINDING_FIELDS = {
    "binding_id",
    "source_rule_ids",
    "principles",
    "design_moves",
    "renderer_action_types",
    "renderer_action_ids",
    "applies_to_scene_questions",
    "acceptance_checks",
}
EXPECTED_RUN2_D1_SCHEMA_REQUIRED_FIELDS = {
    "content_units",
    "proof_object",
    "visual_role",
    "layout_intent",
    "renderer_action_bindings",
    "off_canvas_routes",
}
EXPECTED_RUN2_D2_COMPONENT_KEYS = {
    "hero_object",
    "proof_panel",
    "supporting_copy",
    "evidence_marks",
    "viewer_note_route",
}
EXPECTED_RUN2_D2_ALLOWED_ELEMENT_TYPES = {"shape", "text", "connector", "image", "svg"}
EXPECTED_RUN2_D2_SCENE_STRUCTURE_FIELDS = {
    "expansion_id",
    "source_scene_id",
    "slide_index",
    "role",
    "semantic_components",
    "visual_containers",
    "expanded_renderer_action_bindings",
    "off_canvas_contract",
    "renderer_ready_boundary",
}
EXPECTED_RUN2_D2_COMPONENT_FIELDS = {
    "component_id",
    "component_role",
    "element_type",
    "content_binding",
    "semantic_role_binding",
    "renderer_target",
    "empty_box_policy",
    "source_scene_plan_fields",
}
EXPECTED_RUN2_D2_EXPANDED_BINDING_FIELDS = {
    "binding_id",
    "target_component_key",
    "target_component_id",
    "component_element_type",
    "source_d1_binding_id",
    "source_b2_rule_ids",
    "source_b2_action_ids",
    "source_b2_action_types",
    "design_moves",
    "layout_intent_binding",
    "renderer_resolution_policy",
}
EXPECTED_RUN2_D3_VALIDATION_CHECK_IDS = {
    "required_components_present",
    "allowed_element_types_only",
    "components_bound_to_content_or_semantic_role",
    "visual_containers_bound_to_components",
    "visual_containers_not_empty",
    "expanded_bindings_target_components",
    "expanded_bindings_have_actions",
    "forbidden_renderer_fields_absent",
    "off_canvas_terms_not_on_canvas",
    "renderer_scope_not_started",
}
EXPECTED_RUN2_D3_SCENE_VALIDATION_FIELDS = {
    "validation_id",
    "role",
    "source_expansion_id",
    "status",
    "blocking_issues",
    "checked_components",
    "checked_visual_containers",
    "checked_bindings",
    "validation_checks",
    "renderer_handoff",
}
EXPECTED_RUN2_D3_BLOCKING_ISSUE_IDS = {
    "invalid_element_type",
    "unbound_visual_container",
    "forbidden_renderer_field",
    "off_canvas_term_on_canvas",
    "expanded_binding_missing_actions",
}
EXPECTED_RUN2_D3_FORBIDDEN_RUNTIME_IMPORTS = {
    "importlib",
    "subprocess",
    "pptx",
    "playwright",
    "runpy",
    "selenium",
    "webbrowser",
}
EXPECTED_RUN2_D3_FORBIDDEN_DYNAMIC_IMPORT_CALLS = {"__import__", "eval", "exec"}
EXPECTED_RUN2_E_VISUAL_GRAMMAR_MODULE_IDS = {
    "hero_field",
    "before_after_theater",
    "evidence_workspace",
    "product_reveal",
    "decision_map",
}
EXPECTED_RUN2_E_PAGE_MODULE_MAP = {
    "cover": "product_reveal",
    "setup": "hero_field",
    "contrast": "before_after_theater",
    "proof": "evidence_workspace",
    "climax": "product_reveal",
    "close": "decision_map",
}
EXPECTED_RUN2_E_FORBIDDEN_SCOPE = {
    "renderer_rerun",
    "pptx_output",
    "html_viewer",
    "public_release",
}
EXPECTED_RUN2_E2_STAGE_POLICY = (
    "part_e2_renderer_adapter_contracts_only_no_renderer_rerun_no_public_release"
)
EXPECTED_RUN2_E2_NEXT_REQUIRED_ACTION = "renderer_execute_from_d2_d3_e_adapter_manifest"
EXPECTED_RUN2_F_STAGE_POLICY = (
    "part_f_text_binding_strategy_only_no_renderer_rerun_no_public_release"
)
EXPECTED_RUN2_F_NEXT_REQUIRED_ACTION = "part_g_renderer_rerun_from_validated_text_binding_strategy"
EXPECTED_RUN2_F_REQUIRED_SOCKET_KEYS = {
    "headline_socket",
    "proof_label_sockets",
    "supporting_copy_socket",
    "callout_sockets",
    "viewer_note_socket",
}
EXPECTED_RUN2_F_BOUND_VISUAL_OBJECT_TYPES = {
    "product edge",
    "field route",
    "comparison seam",
    "evidence rail",
    "decision node",
    "negative space pocket",
    "connector endpoint",
}
EXPECTED_RUN2_F_FORBIDDEN_TEXT_PATTERNS = {
    "empty text box",
    "generic rectangle label",
    "duplicated headline/supporting copy",
    "text floating without bound visual object",
    "all slides using the same text layout",
}
EXPECTED_RUN2_G_REQUIRED_INPUTS = [
    "docs/product/ppt-run2-data-skill-quality/run2_73_scene_plan_expansion.json",
    "docs/product/ppt-run2-data-skill-quality/run2_73_renderer_input_validation.json",
    "docs/product/ppt-run2-data-skill-quality/run2_73_visual_grammar_modules.json",
    "docs/product/ppt-run2-data-skill-quality/run2_73_renderer_adapter_contracts.json",
    "docs/product/ppt-run2-data-skill-quality/run2_73_text_binding_strategy.json",
]
EXPECTED_RUN2_G_RESULT = (
    PACK / "results" / "run2_73_validated_scene_renderer_rerun_result.json"
)
EXPECTED_RUN2_G_SCRIPT = ROOT / "scripts" / "generate_ppt_run2_73_validated_scene_renderer_arms.mjs"
EXPECTED_RUN2_H_RESULT = (
    PACK / "results" / "run2_74_visual_quality_evaluation.json"
)
EXPECTED_RUN2_H_SCRIPT = ROOT / "scripts" / "audit_ppt_run2_74_visual_quality_evaluation.py"
EXPECTED_RUN2_H_QUESTIONS = {
    "is_2_73_better_than_2_72",
    "is_text_fused_with_visual_structure",
    "does_it_still_feel_like_engineering_report",
    "do_six_pages_have_distinct_visual_grammar",
    "which_pages_need_repair_and_which_layer",
}
EXPECTED_RUN2_H_ROOT_CAUSE_LAYERS = {
    "renderer",
    "visual_grammar",
    "text_binding",
    "content",
}
EXPECTED_RUN2_I_REQUIRED_INPUTS = [
    *EXPECTED_RUN2_G_REQUIRED_INPUTS,
    "docs/product/ppt-run2-data-skill-quality/results/run2_74_visual_quality_evaluation.json",
]
EXPECTED_RUN2_I_RESULT = (
    PACK / "results" / "run2_75_renderer_repair_rerun_result.json"
)
EXPECTED_RUN2_I_SCRIPT = ROOT / "scripts" / "generate_ppt_run2_75_renderer_repair_arms.mjs"
EXPECTED_RUN2_I_REPAIR_FLAGS = {
    "h_repair_instruction_consumed",
    "concrete_product_surface",
    "higher_visual_density",
    "stronger_text_visual_attachment",
    "public_polish_not_claimed",
}
EXPECTED_RUN2_J_RESULT = (
    PACK / "results" / "run2_76_visual_quality_evaluation.json"
)
EXPECTED_RUN2_J_SCRIPT = ROOT / "scripts" / "audit_ppt_run2_76_visual_quality_evaluation.py"
EXPECTED_RUN2_J_QUESTIONS = {
    "is_2_75_better_than_2_73",
    "does_2_75_have_stronger_product_feel",
    "are_page_differences_stronger_or_weaker",
    "is_text_binding_better",
    "does_2_75_reach_public_video_presentation_direction",
}
EXPECTED_RUN2_J_ROOT_CAUSE_LAYERS = {
    "renderer",
    "visual_grammar",
    "text_binding",
    "content",
}
EXPECTED_RUN2_K1_REPAIR_PLAN = (
    PACK / "run2_76_visual_grammar_renderer_repair_plan.json"
)
EXPECTED_RUN2_K1_REQUIRED_INPUTS = [
    "docs/product/ppt-run2-data-skill-quality/results/run2_76_visual_quality_evaluation.json",
    "docs/product/ppt-run2-data-skill-quality/results/run2_75_renderer_repair_rerun_result.json",
    "docs/product/ppt-run2-data-skill-quality/run2_73_visual_grammar_modules.json",
    "docs/product/ppt-run2-data-skill-quality/run2_73_renderer_adapter_contracts.json",
    "docs/product/ppt-run2-data-skill-quality/run2_73_text_binding_strategy.json",
    "docs/product/ppt-run2-data-skill-quality/run2_73_scene_plan_expansion.json",
]
EXPECTED_RUN2_K1_PAGE_REPAIR_FIELDS = {
    "current_failure",
    "target_scene_direction",
    "visual_grammar_change",
    "renderer_change",
    "text_binding_adjustment",
    "must_preserve_from_2_75",
    "must_remove_from_2_75",
    "acceptance_checks",
}
EXPECTED_RUN2_K1_RENDERER_CAPABILITIES = {
    "hero crop",
    "editorial mask",
    "asymmetric foreground/background depth",
    "larger proof object",
    "fewer but more meaningful labels",
    "scene-specific connectors",
    "non-grid evidence arrangement",
}
EXPECTED_RUN2_K1_FORBIDDEN_PATTERNS = {
    "debug-like outlines",
    "uniform small label wall",
    "same product window skeleton on every page",
    "schematic blueprint as final visual unless it is the actual proof object",
}
EXPECTED_RUN2_K2_REQUIRED_INPUTS = [
    *EXPECTED_RUN2_K1_REQUIRED_INPUTS,
    "docs/product/ppt-run2-data-skill-quality/run2_73_renderer_input_validation.json",
    "docs/product/ppt-run2-data-skill-quality/run2_76_visual_grammar_renderer_repair_plan.json",
]
EXPECTED_RUN2_K2_RESULT = (
    PACK / "results" / "run2_77_visual_grammar_renderer_repair_rerun_result.json"
)
EXPECTED_RUN2_K2_SCRIPT = ROOT / "scripts" / "generate_ppt_run2_77_visual_grammar_renderer_repair_arms.mjs"
EXPECTED_RUN2_K2_REPAIR_FLAGS = {
    "k1_repair_plan_consumed",
    "target_scene_direction_applied",
    "page_differentiation_repaired",
    "forbidden_fallbacks_removed",
    "public_polish_not_claimed",
}
EXPECTED_RUN2_9_VISUAL_PRIMITIVE_IDS = {
    "primitive_2_9_editorial_spread_composition",
    "primitive_2_9_product_surface_depth",
    "primitive_2_9_motion_storyboard_sequence",
    "primitive_2_9_climax_stage_composition",
    "primitive_2_9_typographic_field_composition",
}
EXPECTED_RUN2_9_VISUAL_MODULE_IDS = {
    "module_2_9_editorial_spread",
    "module_2_9_layered_product_surface",
    "module_2_9_motion_storyboard",
    "module_2_9_climax_stage",
    "module_2_9_typographic_field",
}
EXPECTED_RUN2_9_TRACE_FIELDS = {
    "run2_9_visual_primitive_ids",
    "run2_9_visual_module_ids",
    "run2_9_gate_matrix_ids",
    "run2_9_code_module_ids",
    "run2_9_boxiness_failure_probe",
    "run2_9_visual_delta_from_run2_8",
}
EXPECTED_RUN2_10_SOURCE_IDS = {
    "vs_source_2_10_editorial_keynote_system",
    "vs_source_2_10_product_theater_demo",
    "vs_source_2_10_typographic_launch_field",
    "vs_source_2_10_kinetic_climax_sequence",
    "vs_source_2_10_non_rectangular_proof_path",
}
EXPECTED_RUN2_10_MEMORY_IDS = {
    "visual_system_editorial_cinema",
    "visual_system_product_theater",
    "visual_system_typographic_field",
    "visual_system_kinetic_demo",
    "visual_system_non_rectangular_proof",
}
EXPECTED_RUN2_10_TRACE_FIELDS = {
    "run2_10_visual_system_source_ids",
    "run2_10_visual_system_memory_ids",
    "run2_10_gate_matrix_ids",
    "run2_10_code_module_ids",
    "run2_10_visual_delta_from_run2_9",
    "run2_10_sameness_failure_probe",
    "run2_10_public_demo_first_read_probe",
    "run2_10_shape_count_budget",
    "run2_10_asymmetry_whitespace_rule",
}
EXPECTED_RUN2_51_ROLES = {"cover", "setup", "contrast", "proof", "climax", "close"}
EXPECTED_RUN2_51_COPY_BUNDLE_KEYS = {
    "headline",
    "subline",
    "proof_nuggets",
    "annotations",
    "state_labels",
}
EXPECTED_RUN2_61_ROLES = {"cover", "setup", "contrast", "proof", "climax", "close"}
EXPECTED_RUN2_61_CARRIER_TYPES = {
    "product_surface",
    "operating_loop",
    "before_after_delta",
    "workspace_inspection",
    "climax_result_object",
    "release_handoff",
}
EXPECTED_RUN2_61_REQUIRED_COPY_UNITS = {
    "headline",
    "subhead",
    "proof_badges",
    "annotations",
    "state_labels",
    "speaker_note",
}
EXPECTED_RUN2_51_TRACE_FIELDS = {
    "run2_51_editorial_copy_memory_id",
    "run2_51_shape_text_socket_memory_id",
    "run2_51_renderer_archetype_gate_id",
    "run2_51_primary_archetype",
    "run2_51_public_surface_copy_status",
    "run2_51_text_socket_placement_status",
    "run2_51_shape_vocabulary_status",
    "run2_51_character_fit_status",
    "run2_51_forbidden_surface_terms_count",
    "run2_51_equal_card_cluster_count",
    "run2_51_semantic_primitive_count",
}
EXPECTED_RUN2_52_TRACE_FIELDS = {
    "run2_51_editorial_copy_memory_id",
    "run2_51_shape_text_socket_memory_id",
    "run2_51_renderer_archetype_gate_id",
    "run2_51_primary_archetype",
    "run2_51_public_surface_copy_status",
    "run2_51_text_socket_placement_status",
    "run2_51_shape_vocabulary_status",
    "run2_51_character_fit_status",
    "run2_51_forbidden_surface_terms_count",
    "run2_51_equal_card_cluster_count",
    "run2_51_semantic_primitive_count",
    "run2_52_code_module_ids",
    "run2_52_primary_surface_kind",
    "run2_52_socket_bound_public_text_elements",
    "run2_52_shape_primitive_count",
}
EXPECTED_RUN2_53_TRACE_FIELDS = {
    "run2_53_product_surface_scene_id",
    "run2_53_business_visual_evidence_id",
    "run2_53_scene_renderer_gate_id",
    "run2_53_primary_product_or_business_object",
    "run2_53_visual_specificity_status",
    "run2_53_forbidden_generic_geometry_count",
}
EXPECTED_RUN2_54_TRACE_FIELDS = {
    "run2_53_product_surface_scene_id",
    "run2_53_business_visual_evidence_id",
    "run2_53_scene_renderer_gate_id",
    "run2_53_primary_product_or_business_object",
    "run2_53_visual_specificity_status",
    "run2_53_forbidden_generic_geometry_count",
    "run2_54_code_module_ids",
    "run2_54_primary_surface_kind",
    "run2_54_product_surface_slots_rendered",
    "run2_54_business_visual_evidence_objects",
}
EXPECTED_RUN2_55_TRACE_FIELDS = {
    "run2_55_code_module_ids",
    "run2_55_text_shape_integration_status",
    "run2_55_primary_layout_strategy",
    "run2_55_named_text_containers_rendered",
    "run2_55_non_rectangular_shape_families_rendered",
    "run2_55_text_shape_binding_pairs",
    "run2_55_text_overflow_risk_count",
    "run2_55_equal_rectangle_cluster_count",
    "run2_55_editorial_hierarchy_levels",
}
EXPECTED_RUN2_56_TRACE_FIELDS = {
    "run2_56_code_module_ids",
    "run2_56_role_renderer_id",
    "run2_56_composition_family",
    "run2_56_layout_signature",
    "run2_56_visual_sameness_bucket",
    "run2_56_primary_anchor_region",
    "run2_56_role_specific_geometry_count",
    "run2_56_text_collision_risk_count",
    "run2_56_text_overflow_risk_count",
    "run2_56_distinct_role_surface_status",
    "run2_56_role_archetype_binding_status",
}
EXPECTED_RUN2_57_CAPABILITY_LAYERS = {
    "input_layer",
    "understanding_layer",
    "generation_layer",
    "comparison_layer",
    "verification_layer",
    "output_layer",
    "repair_loop",
}
EXPECTED_RUN2_57_TRACE_FIELDS = {
    "run2_57_product_capability_ids",
    "run2_57_slide_message_contract_id",
    "run2_57_content_workflow_gate_id",
    "run2_57_product_logic_relation_ids",
    "run2_57_competitor_boundary_ids",
    "run2_57_content_specificity_status",
    "run2_57_reader_question_answered_status",
    "run2_57_generic_claim_count",
    "run2_57_required_product_terms_rendered",
}
EXPECTED_RUN2_57_COMPETITOR_BOUNDARIES = {
    "prompt_only_slide_generator",
    "image_generator",
    "template_slide_tool",
    "design_editor",
    "office_copilot",
}
EXPECTED_RUN2_58_TRACE_FIELDS = {
    "run2_57_product_capability_ids",
    "run2_57_slide_message_contract_id",
    "run2_57_content_workflow_gate_id",
    "run2_57_product_logic_relation_ids",
    "run2_57_competitor_boundary_ids",
    "run2_57_content_specificity_status",
    "run2_57_reader_question_answered_status",
    "run2_57_generic_claim_count",
    "run2_57_required_product_terms_rendered",
    "run2_58_code_module_ids",
    "run2_58_product_content_contract_status",
    "run2_58_product_system_surface_kind",
    "run2_58_reader_question_visible",
    "run2_58_product_terms_visible_count",
    "run2_58_proof_object_count",
    "run2_58_bad_control_boundary_status",
}
EXPECTED_RUN2_58_OPEN_SOURCE_REFERENCE_IDS = {
    "pptxgenjs",
    "slidev",
    "marp",
    "revealjs",
    "slidecoder",
}
EXPECTED_RUN2_59_LAYOUT_MODULE_IDS = {
    "module_2_15_editorial_cover_field",
    "module_2_15_product_theater_stage",
    "module_2_15_before_after_route",
    "module_2_15_metric_reveal_stage",
    "module_2_15_quiet_release_handoff",
    "module_2_15_dense_evidence_compression",
}
EXPECTED_RUN2_59_TRACE_FIELDS = {
    "run2_59_content_contract_id",
    "run2_59_layout_module_id",
    "run2_59_content_burden_level",
    "run2_59_capacity_fit_status",
    "run2_59_public_visible_word_budget",
    "run2_59_trace_only_detail_count",
    "run2_59_public_surface_trace_policy_id",
    "run2_59_composition_gate_id",
}
EXPECTED_RUN2_60_TRACE_FIELDS = {
    "run2_59_content_contract_id",
    "run2_59_layout_module_id",
    "run2_59_content_burden_level",
    "run2_59_capacity_fit_status",
    "run2_59_public_visible_word_budget",
    "run2_59_trace_only_detail_count",
    "run2_59_public_surface_trace_policy_id",
    "run2_59_composition_gate_id",
    "run2_60_code_module_ids",
    "run2_60_public_trace_split_status",
    "run2_60_over_capacity_fallback_status",
    "run2_60_layout_collision_status",
    "run2_60_content_aware_composition_status",
}
EXPECTED_RUN2_61_CONSUMER_TRACE_FIELDS = {
    "run2_61_narrative_proof_id",
    "run2_61_visual_carrier_selector_id",
    "run2_61_text_socket_fusion_contract_id",
    "run2_61_public_proof_replacement_id",
    "run2_61_narrative_workflow_gate_id",
}
EXPECTED_RUN2_62_TRACE_FIELDS = {
    *EXPECTED_RUN2_61_CONSUMER_TRACE_FIELDS,
    "run2_62_code_module_ids",
    "run2_62_narrative_proof_consumption_status",
    "run2_62_socket_binding_count",
    "run2_62_public_proof_object_count",
    "run2_62_required_answer_visible",
}
EXPECTED_RUN2_63_RENDERER_BLOCKERS = {
    "static_socket_plan_repeated_on_every_slide",
    "generic_native_shape_grammar_collapses_role_specific_visual_carriers",
    "text_fit_and_wrapping_not_trace_gated",
    "semantic_diagram_labels_not_bound_to_active_slide_proof",
}
EXPECTED_RUN2_64_TRACE_FIELDS = {
    "run2_64_dynamic_socket_renderer_id",
    "run2_64_semantic_diagram_renderer_id",
    "run2_64_text_fit_gate_id",
    "run2_64_renderer_dry_run_binding_id",
    "run2_64_dynamic_socket_plan_status",
    "run2_64_semantic_diagram_binding_status",
    "run2_64_text_fit_preflight_status",
}
EXPECTED_RUN2_65_TRACE_FIELDS = {
    *EXPECTED_RUN2_64_TRACE_FIELDS,
    "run2_65_code_module_ids",
    "run2_65_renderer_composition_repair_status",
    "run2_65_dynamic_socket_count",
    "run2_65_semantic_diagram_type",
    "run2_65_text_fit_runtime_status",
    "run2_65_visual_delta_from_run2_62",
}
EXPECTED_RUN2_51_FORBIDDEN_PUBLIC_TERMS = {
    "run2",
    "memory",
    "workflow gate",
    "trace",
    "audit",
    "negative control",
    "public blocked",
}
EXPECTED_RUN2_11_AUDIT_FIELDS = {
    "schema_version",
    "status",
    "stage_policy",
    "audit_scope",
    "source_inventory",
    "workflow_inventory",
    "chain_records",
    "arm_trace_evidence",
    "negative_control_checks",
    "gate_summary",
    "next_required_action",
}
EXPECTED_RUN2_11_CHAIN_STATUSES = {"pass", "weak", "missing", "blocked"}
EXPECTED_RUN2_11_RUN_IDS = {"2.8", "2.9", "2.10"}
EXPECTED_RUN2_11_CONTROL_ARMS = {
    "prompt_only",
    "run1_5_skill",
    "negative_control",
}
EXPECTED_RUN2_12_EVIDENCE_IDS = {
    "thick_2_12_figma_platform_launch_arc",
    "thick_2_12_figma_design_to_code_surface",
    "thick_2_12_stripe_agentic_commerce_demo",
    "thick_2_12_stripe_product_keynote_metric_reframe",
    "thick_2_12_present_partners_whitespace_hierarchy",
    "thick_2_12_slidecow_powerpoint_whitespace",
}
EXPECTED_RUN2_12_MEMORY_IDS = {
    "memory_2_12_launch_arc_to_six_slide_route",
    "memory_2_12_product_demo_sequence_pacing",
    "memory_2_12_typographic_whitespace_system",
    "memory_2_12_metric_to_climax_object",
}
EXPECTED_RUN2_12_GATE_IDS = {
    "gate_2_12_evidence_source_integrity",
    "gate_2_12_memory_selection_before_generation",
    "gate_2_12_visual_density_and_whitespace",
    "gate_2_12_demo_sequence_climax",
}
EXPECTED_RUN2_12_EVIDENCE_FIELDS = {
    "id",
    "source_ids",
    "source_urls",
    "verified_accessed_on",
    "source_role",
    "modality_mix",
    "segment_locator",
    "frame_or_visual_observations",
    "spoken_or_text_claim_paraphrases",
    "derived_design_method",
    "native_ppt_code_obligations",
    "workflow_gate_obligations",
    "memory_seed_targets",
    "anti_copy_boundary",
    "qa_probe",
    "release_boundary",
}
EXPECTED_RUN2_12_MEMORY_FIELDS = {
    "id",
    "evidence_record_ids",
    "memory_type",
    "applies_to_slide_roles",
    "design_constraints",
    "native_ppt_contract",
    "required_trace_fields",
    "failure_probe",
    "qa_probe",
    "release_boundary",
}
EXPECTED_RUN2_12_GATE_FIELDS = {
    "id",
    "gate_type",
    "evidence_record_ids",
    "memory_seed_ids",
    "slide_roles",
    "pass_fail_checks",
    "required_before_next_rerun",
    "trace_fields",
    "release_boundary",
}
EXPECTED_RUN2_13_TRACE_FIELDS = {
    "run2_12_evidence_record_ids",
    "run2_12_memory_seed_ids",
    "run2_12_workflow_gate_ids",
    "run2_12_code_module_ids",
    "run2_12_density_whitespace_gate",
    "run2_12_demo_sequence_gate",
    "run2_12_climax_object_gate",
    "run2_13_visual_delta_from_run2_10",
    "run2_13_negative_control_probe",
}
EXPECTED_RUN2_14_TRACE_FIELDS = {
    "run2_12_evidence_record_ids",
    "run2_12_memory_seed_ids",
    "run2_12_workflow_gate_ids",
    "run2_14_aesthetic_source_run_id",
    "run2_14_aesthetic_shell_input_ids",
    "run2_14_surface_policy",
    "run2_14_visible_workflow_suppression",
    "run2_14_visual_delta_from_run2_13",
    "run2_14_code_module_ids",
    "run2_14_bad_control_probe",
}
EXPECTED_RUN2_15_TRACE_FIELDS = {
    "run2_15_selected_layout_module_ids",
    "run2_15_selector_gate_ids",
    "run2_15_trace_visibility_policy",
    "run2_15_text_resilience_result",
    "run2_15_trace_hidden_result",
    "run2_15_product_surface_probe",
    "run2_15_metric_reveal_climax_probe",
    "run2_15_bad_selector_control_probe",
}
EXPECTED_RUN2_16_TRACE_FIELDS = {
    "run2_16_aesthetic_source_run_id",
    "run2_16_aesthetic_shell_input_ids",
    "run2_16_selector_input_ids",
    "run2_16_selector_execution_status",
    "run2_16_surface_policy",
    "run2_16_visible_workflow_suppression",
    "run2_16_visual_delta_from_run2_13",
    "run2_16_code_module_ids",
    "run2_16_bad_control_probe",
}
EXPECTED_RUN2_6_USECASE_FIELDS = {
    "id",
    "source_ids",
    "audience",
    "business_decision",
    "deck_mission",
    "slide_arc",
    "must_show",
    "must_not_show",
    "failure_modes",
    "visual_risk",
    "workflow_implications",
    "qa_probe",
    "release_boundary",
}
EXPECTED_RUN2_6_BENCHMARK_FIELDS = {
    "id",
    "source_ids",
    "allowed_use",
    "composition_rules",
    "typography_rules",
    "spacing_rules",
    "theme_rules",
    "motion_or_interaction_rules",
    "native_ppt_implications",
    "anti_copy_rules",
    "qa_probe",
    "trace_fields",
    "release_boundary",
}
EXPECTED_CLAIM_IDS = {
    "claim_data_changes_deck_quality",
    "claim_aesthetic_memory_controls_rhythm",
    "claim_asset_memory_preserves_editability",
    "claim_bad_aesthetic_memory_is_negative_control",
    "claim_appendix_absorbs_proof_detail",
    "claim_public_release_requires_render_and_human_gate",
}
EXPECTED_ASSET_IDS = {
    "asset_launch_atmosphere_background",
    "asset_system_svg_mark",
    "asset_evidence_flow_diagram",
    "asset_comparison_delta_chart",
}
EXPECTED_MOVES = {
    "cinematic_cover",
    "low_density_claim",
    "big_object_layout",
    "editorial_comparison",
    "proof_reveal",
    "visual_climax",
    "premium_closing",
    "appendix_absorption",
}
EXPECTED_BAD_MEMORY_FIELDS = {
    "id",
    "replaces",
    "source_card_ids",
    "aesthetic_move",
    "trigger",
    "composition_rule",
    "typography_rule",
    "density_budget",
    "rhythm_role",
    "ppt_primitive",
    "negative_rules",
    "qa_signal",
    "expected_failure",
}
EXPECTED_RHYTHM_ROLES = {"cover", "setup", "contrast", "proof", "climax", "relief", "close"}
RUN2_8_FORBIDDEN_MEDIA_MARKERS = (
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
    ".mp4",
    ".mov",
    "http://",
    "https://",
    "data:image",
    "base64,",
)
RUN2_8_CODE_BINDING_TERMS = {"fontSize", "bbox", "spacing", "heroObject", "beforeAfter", "workflowGate"}
RUN2_18_EVIDENCE_FIELDS = {
    "record_id",
    "source_family",
    "commercial_usecase_ids",
    "source_ids",
    "modality_mix",
    "source_locator",
    "observed_design_method",
    "business_requirement",
    "derived_generation_constraint",
    "memory_targets",
    "workflow_gate_targets",
    "anti_copy_boundary",
    "bad_control_probe",
    "release_boundary",
}
RUN2_18_MEMORY_FIELDS = {
    "memory_id",
    "memory_family",
    "slide_roles",
    "evidence_record_ids",
    "composition_contract",
    "typography_contract",
    "spacing_contract",
    "motion_or_sequence_contract",
    "proof_object_contract",
    "code_generation_binding",
    "trace_fields_required",
    "negative_control_failure",
    "release_boundary",
}
RUN2_18_WORKFLOW_GATE_FIELDS = {
    "gate_id",
    "required_before_next_rerun",
    "evidence_record_ids",
    "memory_ids",
    "selection_rules",
    "rejection_rules",
    "trace_fields",
    "qa_probe",
    "bad_control_probe",
    "release_boundary",
}
EXPECTED_RUN2_19_TRACE_FIELDS = {
    "run2_18_selected_evidence_ids",
    "run2_18_selected_memory_ids",
    "run2_18_selected_gate_ids",
    "run2_18_bad_control_probe",
    "run2_18_release_boundary",
    "run2_19_code_module_ids",
    "run2_19_trace_thickness_status",
    "run2_19_visual_delta_from_run2_16",
}
EXPECTED_RUN2_22_TRACE_FIELDS = {
    "run2_21_visual_decision_memory_id",
    "run2_21_primary_evidence_id",
    "run2_21_secondary_evidence_ids",
    "run2_21_rejected_evidence_reasons",
    "run2_21_selector_gate_id",
    "run2_21_visual_decision_delta",
    "run2_22_selector_execution_status",
    "run2_22_code_module_ids",
}


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def normalize(value: str) -> str:
    return " ".join(re.sub(r"[^a-z0-9]+", " ", value.lower()).split())


def assert_contains(body: str, terms: list[str]) -> None:
    normalized = normalize(body)
    for term in terms:
        assert normalize(term) in normalized, f"missing term: {term!r}"


def assert_sequence(body: str, terms: list[str]) -> None:
    cursor = -1
    for term in terms:
        index = body.find(term, cursor + 1)
        assert index != -1, f"missing sequence term: {term!r}"
        cursor = index


def assert_mentions_any(body: str, terms: set[str]) -> None:
    normalized = normalize(body)
    assert any(normalize(term) in normalized for term in terms), f"missing one of: {sorted(terms)!r}"


def word_count(value: str) -> int:
    return len([part for part in re.split(r"\s+", value.strip()) if part])


def nested_keys(value: object) -> set[str]:
    if isinstance(value, dict):
        keys = set(value)
        for child in value.values():
            keys |= nested_keys(child)
        return keys
    if isinstance(value, list):
        keys: set[str] = set()
        for child in value:
            keys |= nested_keys(child)
        return keys
    return set()


def public_text_values(bundle: dict[str, object]) -> list[str]:
    values: list[str] = []
    for key in ("headline", "subline"):
        value = bundle.get(key)
        if isinstance(value, str):
            values.append(value)
    for key in ("proof_nuggets", "annotations", "state_labels"):
        items = bundle.get(key)
        if isinstance(items, list):
            values.extend(str(item) for item in items)
    return values


def iter_string_values(value: object):
    if isinstance(value, str):
        yield value
    elif isinstance(value, list):
        for item in value:
            yield from iter_string_values(item)
    elif isinstance(value, dict):
        for item in value.values():
            yield from iter_string_values(item)


def json_files(directory: str) -> list[Path]:
    return sorted((PACK / directory).glob("*.json"))


def tracked_files() -> list[str]:
    result = subprocess.run(
        ["git", "ls-files"],
        cwd=ROOT,
        check=True,
        capture_output=True,
        text=True,
    )
    return result.stdout.splitlines()


def test_run2_case_pack_is_valid() -> None:
    result = validate_case_pack(PACK, profile="run2")

    assert result.ok is True, result.errors


def test_run2_commercial_case_is_concrete_and_public_demo_oriented() -> None:
    body = (PACK / "commercial_case.md").read_text(encoding="utf-8")

    assert_contains(body, ["audience", "decision", "failed deck", "public-demo", "six slides"])
    assert "generic beautiful PPT" not in body


def test_run2_has_thick_source_and_video_cards() -> None:
    source_cards = json_files("source_cards")
    video_cards = json_files("video_cards")
    source_card_bodies = [load_json(path) for path in source_cards]
    video_card_bodies = [load_json(path) for path in video_cards]

    assert len(source_cards) >= 8
    assert len(video_cards) >= 2
    assert EXPECTED_SOURCE_CARD_IDS <= {card["card_id"] for card in source_card_bodies}
    assert EXPECTED_VIDEO_CARD_IDS <= {card["card_id"] for card in video_card_bodies}

    for card in [*source_card_bodies, *video_card_bodies]:
        assert_contains(card["do_not_copy"], ["do not copy"])
        assert card["allowed_use"] in {
            "short_analysis",
            "derived_rules_only",
            "visual_inspiration",
            "timestamped_observation_only",
        }

    for card in source_card_bodies:
        assert_mentions_any(
            card["ppt_translation"],
            {"editable", "native", "PPT", "PowerPoint", "shape", "text", "SVG", "diagram", "chart"},
        )
        assert_mentions_any(
            card["quality_risk"],
            {"generic", "report", "dashboard", "degrade", "weak", "thin", "overdesigned"},
        )


def test_run2_has_multimodal_lesson_database() -> None:
    database = load_json(PACK / "multimodal_database.json")
    sources = load_json(PACK / "sources.json")
    records = database["records"]
    source_ids = {source["id"] for source in sources["sources"]}
    record_ids = {record["id"] for record in records}
    covered_modalities = {modality for record in records for modality in record["modalities"]}

    assert database["status"] == "run2_2_multimodal_contract_public_blocked"
    assert EXPECTED_MULTIMODAL_RECORD_IDS <= record_ids
    assert {
        "duarte_persuasive_visual_storytelling",
        "duarte_slide_design_course",
        "udel_design_principles_video",
    } <= source_ids
    assert set(database["required_modalities"]) == EXPECTED_MULTIMODAL_MODALITIES
    assert covered_modalities == EXPECTED_MULTIMODAL_MODALITIES
    assert database["storage_policy"]["raw_media"] == "forbidden"
    assert_contains(
        json.dumps(database["storage_policy"]),
        ["Do not store", "audio", "video", "full transcripts", "screenshots", "source assets"],
    )

    for record in records:
        assert record["allowed_storage"] == "derived_observations_only"
        assert record["anchors"]
        assert set(record["modalities"]) <= EXPECTED_MULTIMODAL_MODALITIES
        assert record["do_not_store"]
        for anchor in record["anchors"]:
            assert anchor["modality"] in record["modalities"]
            assert anchor["allowed_use"] in {
                "short_analysis",
                "derived_rules_only",
                "visual_inspiration",
                "timestamped_observation_only",
            }
            assert_mentions_any(
                anchor["extracted_design_signal"],
                {"slide", "Generate", "Reject", "Convert", "Reduce", "Use"},
            )

    anchors_by_modality = {"transcript": [], "interaction": []}
    for record in records:
        for anchor in record["anchors"]:
            if anchor["modality"] in anchors_by_modality:
                anchors_by_modality[anchor["modality"]].append(anchor)
    assert anchors_by_modality["transcript"]
    assert anchors_by_modality["interaction"]
    assert_contains(
        " ".join(anchor["extracted_design_signal"] for anchor in anchors_by_modality["transcript"]),
        ["native headline", "do not place full transcript text"],
    )
    assert_contains(
        " ".join(anchor["extracted_design_signal"] for anchor in anchors_by_modality["interaction"]),
        ["selected memory", "density budget", "release gate"],
    )

    tasks = database["cross_modal_design_tasks"]
    assert {task["id"] for task in tasks} >= {
        "before_after_visual_delta",
        "audio_to_rhythm_budget",
        "transcript_to_native_headline",
        "interaction_to_skill_surface",
    }
    assert_contains(json.dumps(tasks), ["visible transformation", "density", "native", "generator behavior"])


def test_run2_has_visual_learning_targets_for_next_rerun() -> None:
    database = load_json(PACK / "multimodal_database.json")
    targets = load_json(PACK / "visual_learning_targets.json")
    record_ids = {record["id"] for record in database["records"]}
    anchor_ids = {anchor["anchor_id"] for record in database["records"] for anchor in record["anchors"]}

    assert targets["status"] == "run2_2_visual_learning_targets_public_blocked"
    assert targets["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert_contains(" ".join(targets["native_editable_definition"]), ["native", "editable", "must not contain"])
    assert EXPECTED_VISUAL_TARGET_IDS <= {target["id"] for target in targets["targets"]}

    for target in targets["targets"]:
        assert set(target["source_record_ids"]) <= record_ids
        assert set(target["anchor_ids"]) <= anchor_ids
        assert_contains(target["release_boundary"], ["public blocked"])
        assert_mentions_any(
            target["desired_behavior"], {"before/after", "mini-preview", "rhythm", "headline", "climax"}
        )
        requirements = " ".join(target["code_generation_requirements"])
        assert_contains(requirements, ["native"])
        assert_mentions_any(requirements, {"editable", "record", "trace"})
        assert_mentions_any(
            target["qa_probe"],
            {"contact sheet", "headline", "rhythm", "climax", "mini-preview"},
        )


def test_run2_has_visual_target_components_for_native_generation() -> None:
    targets = load_json(PACK / "visual_learning_targets.json")
    components = load_json(PACK / "visual_target_components.json")
    target_ids = {target["id"] for target in targets["targets"]}
    component_ids = {component["id"] for component in components["components"]}

    assert components["status"] == "run2_3_native_visual_components_public_blocked"
    assert components["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert EXPECTED_VISUAL_COMPONENT_IDS <= component_ids

    covered_targets = {target_id for component in components["components"] for target_id in component["target_ids"]}
    assert EXPECTED_VISUAL_TARGET_IDS <= covered_targets

    for component in components["components"]:
        assert set(component["target_ids"]) <= target_ids
        assert component["slide_roles"]
        assert_contains(" ".join(component["native_ppt_primitives"]), ["native", "editable"])
        assert_contains(component["layout_contract"], ["object", "canvas"])
        assert_mentions_any(
            component["layout_contract"],
            {"thumbnail", "mini-preview", "rhythm", "headline", "climax", "gate"},
        )
        assert_contains(component["density_contract"], ["visible words", "panels"])
        assert_mentions_any(
            " ".join(component["trace_fields"]),
            {"visual_learning_target_ids", "visual_component_ids", "density_counts", "native_ppt_checks"},
        )
        assert_mentions_any(
            component["qa_probe"],
            {"contact sheet", "thumbnail", "rhythm", "climax", "before/after"},
        )
        assert_contains(component["release_boundary"], ["public_blocked"])


def test_run2_has_video_demo_beat_map_for_motion_learning() -> None:
    database = load_json(PACK / "multimodal_database.json")
    beat_map = load_json(PACK / "video_demo_beat_map.json")
    record_ids = {record["id"] for record in database["records"]}
    anchor_ids = {anchor["anchor_id"] for record in database["records"] for anchor in record["anchors"]}
    video_card_ids = {card["card_id"] for card in (load_json(path) for path in json_files("video_cards"))}

    assert beat_map["status"] == "run2_4_video_demo_beat_map_public_blocked"
    assert beat_map["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert beat_map["storage_policy"]["raw_media"] == "forbidden"
    assert EXPECTED_VIDEO_BEAT_IDS <= {beat["id"] for beat in beat_map["beats"]}

    for beat in beat_map["beats"]:
        assert set(beat["source_record_ids"]) <= record_ids
        assert set(beat["anchor_ids"]) <= anchor_ids
        assert set(beat["video_card_ids"]) <= video_card_ids
        assert_contains(beat["locator"], [":"])
        assert_mentions_any(beat["motion_role"], {"attention_reset", "reveal", "build", "scale", "handoff"})
        assert_contains(" ".join(beat["reveal_sequence"]), ["native"])
        assert_contains(" ".join(beat["do_not_store"]), ["video", "frames", "audio", "transcript"])
        assert_contains(beat["release_boundary"], ["public_blocked"])


def test_run2_has_motion_learning_targets_and_sequence_components() -> None:
    beat_map = load_json(PACK / "video_demo_beat_map.json")
    motion_targets = load_json(PACK / "motion_learning_targets.json")
    visual_components = load_json(PACK / "visual_target_components.json")
    sequence_components = load_json(PACK / "presentation_sequence_components.json")
    beat_ids = {beat["id"] for beat in beat_map["beats"]}
    visual_component_ids = {component["id"] for component in visual_components["components"]}
    motion_target_ids = {target["id"] for target in motion_targets["targets"]}

    assert motion_targets["status"] == "run2_4_motion_learning_targets_public_blocked"
    assert sequence_components["status"] == "run2_4_sequence_components_public_blocked"
    assert EXPECTED_MOTION_TARGET_IDS <= motion_target_ids
    assert EXPECTED_SEQUENCE_COMPONENT_IDS <= {component["id"] for component in sequence_components["components"]}

    for target in motion_targets["targets"]:
        assert set(target["beat_ids"]) <= beat_ids
        assert set(target["visual_component_ids"]) <= visual_component_ids
        assert_contains(" ".join(target["code_generation_requirements"]), ["native", "metadata", "trace"])
        assert_mentions_any(target["desired_behavior"], {"reveal", "scale", "build", "handoff", "pause"})
        assert_contains(target["release_boundary"], ["public_blocked"])

    for component in sequence_components["components"]:
        assert set(component["motion_target_ids"]) <= motion_target_ids
        assert set(component["visual_component_ids"]) <= visual_component_ids
        assert component["sequence_steps"]
        assert [step["order"] for step in component["sequence_steps"]] == list(
            range(1, len(component["sequence_steps"]) + 1)
        )
        assert_contains(" ".join(component["native_ppt_primitives"]), ["native", "editable"])
        assert_contains(" ".join(component["trace_fields"]), ["motion_target_ids", "sequence_component_ids"])
        assert_contains(component["qa_probe"], ["reveal"])
        assert_contains(component["release_boundary"], ["public_blocked"])


def test_run2_5_has_production_reference_decompositions() -> None:
    references = load_json(PACK / "production_reference_decompositions.json")
    sources = load_json(PACK / "sources.json")
    source_ids = {source["id"] for source in sources["sources"]}
    visual_component_ids = {component["id"] for component in load_json(PACK / "visual_target_components.json")["components"]}
    motion_target_ids = {target["id"] for target in load_json(PACK / "motion_learning_targets.json")["targets"]}
    sequence_component_ids = {
        component["id"] for component in load_json(PACK / "presentation_sequence_components.json")["components"]
    }

    assert references["status"] == "run2_5_production_reference_decomposition_public_blocked"
    assert references["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert references["storage_policy"]["raw_media"] == "forbidden"
    assert_contains(json.dumps(references["extraction_policy"]), ["derived observation", "no automatic frame extraction"])
    assert EXPECTED_PRODUCTION_REFERENCE_IDS <= {record["id"] for record in references["references"]}

    for record in references["references"]:
        assert EXPECTED_PRODUCTION_REFERENCE_FIELDS <= set(record), record["id"]
        assert set(record["source_ids"]) <= source_ids
        assert set(record["visual_component_ids"]) <= visual_component_ids
        assert set(record["motion_target_ids"]) <= motion_target_ids
        assert set(record["sequence_component_ids"]) <= sequence_component_ids
        assert_contains(" ".join(record["visual_observations"]), ["native", "composition"])
        assert_contains(" ".join(record["composition_primitives"]), ["canvas"])
        assert_contains(" ".join(record["typography_primitives"]), ["headline", "hierarchy"])
        assert_contains(" ".join(record["spacing_primitives"]), ["whitespace", "margin"])
        assert_contains(" ".join(record["motion_or_sequence_primitives"]), ["sequence"])
        assert_contains(" ".join(record["module_implications"]), ["code", "module", "trace"])
        assert_contains(" ".join(record["do_not_copy"]), ["do not copy", "screenshot", "brand"])
        assert_contains(" ".join(record["extraction_notes"]), ["derived", "not copied"])
        assert_contains(record["qa_probe"], ["contact sheet"])
        assert_contains(record["release_boundary"], ["public_blocked"])


def test_run2_5_has_aesthetic_memory_v2_and_visual_production_modules() -> None:
    references = load_json(PACK / "production_reference_decompositions.json")
    aesthetic_v2 = load_json(PACK / "aesthetic_memory_v2.json")
    modules = load_json(PACK / "visual_production_modules.json")
    reference_ids = {record["id"] for record in references["references"]}
    visual_component_ids = {component["id"] for component in load_json(PACK / "visual_target_components.json")["components"]}
    motion_target_ids = {target["id"] for target in load_json(PACK / "motion_learning_targets.json")["targets"]}
    sequence_component_ids = {
        component["id"] for component in load_json(PACK / "presentation_sequence_components.json")["components"]
    }

    assert aesthetic_v2["status"] == "run2_5_aesthetic_memory_v2_public_blocked"
    assert modules["status"] == "run2_5_visual_production_modules_public_blocked"
    assert EXPECTED_AESTHETIC_V2_MOVE_IDS <= {move["id"] for move in aesthetic_v2["moves"]}
    assert EXPECTED_PRODUCTION_MODULE_IDS <= {module["id"] for module in modules["modules"]}

    for move in aesthetic_v2["moves"]:
        assert EXPECTED_AESTHETIC_V2_FIELDS <= set(move), move["id"]
        assert set(move["production_reference_ids"]) <= reference_ids
        assert set(move["visual_component_ids"]) <= visual_component_ids
        assert set(move["motion_target_ids"]) <= motion_target_ids
        assert set(move["sequence_component_ids"]) <= sequence_component_ids
        assert move["density_budget"]["max_words"] <= 52
        assert_contains(move["composition_contract"], ["canvas", "focal"])
        assert_contains(move["typography_contract"], ["headline", "hierarchy"])
        assert_contains(move["spacing_contract"], ["whitespace", "margin"])
        assert_contains(move["provenance_boundary"], ["derived", "do not copy"])
        assert_contains(" ".join(move["code_generation_contract"]), ["native", "PPT", "module"])
        assert_contains(" ".join(move["trace_fields"]), ["aesthetic_memory_v2_ids", "visual_production_module_ids"])
        assert_mentions_any(" ".join(move["forbidden_report_patterns"]), {"dashboard", "equal grid", "report", "tiny"})
        assert_contains(move["qa_probe"], ["contact sheet"])

    module_ids_by_move = {move_id for module in modules["modules"] for move_id in module["aesthetic_memory_v2_ids"]}
    assert EXPECTED_AESTHETIC_V2_MOVE_IDS <= module_ids_by_move
    for module in modules["modules"]:
        assert EXPECTED_VISUAL_PRODUCTION_MODULE_FIELDS <= set(module), module["id"]
        assert set(module["production_reference_ids"]) <= reference_ids
        assert set(module["aesthetic_memory_v2_ids"]) <= {move["id"] for move in aesthetic_v2["moves"]}
        assert_contains(" ".join(module["native_ppt_primitives"]), ["native", "editable"])
        assert_contains(module["layout_recipe"], ["canvas"])
        assert_contains(module["typography_recipe"], ["headline"])
        assert_contains(module["spacing_recipe"], ["margin"])
        assert_contains(" ".join(module["trace_fields"]), ["visual_production_module_ids", "aesthetic_memory_v2_ids"])
        assert_contains(module["fallback_policy"], ["native", "public_blocked"])
        assert_contains(module["provenance_boundary"], ["derived", "do not copy"])
        assert_contains(module["qa_probe"], ["contact sheet"])
        assert_contains(module["release_boundary"], ["public_blocked"])


def test_run2_6_has_commercial_usecase_bank() -> None:
    sources = load_json(PACK / "sources.json")
    usecases = load_json(PACK / "commercial_usecase_bank.json")
    source_ids = {source["id"] for source in sources["sources"]}

    assert EXPECTED_RUN2_6_SOURCE_IDS <= source_ids
    assert usecases["status"] == "run2_6_commercial_usecase_bank_public_blocked"
    assert usecases["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert usecases["storage_policy"]["raw_media"] == "forbidden"
    assert EXPECTED_RUN2_6_USECASE_IDS <= {item["id"] for item in usecases["usecases"]}

    for item in usecases["usecases"]:
        assert EXPECTED_RUN2_6_USECASE_FIELDS <= set(item), item["id"]
        assert set(item["source_ids"]) <= source_ids
        assert_contains(" ".join(item["must_not_show"]), ["copy", "brand"])
        assert_contains(" ".join(item["workflow_implications"]), ["select", "benchmark"])
        assert_contains(item["qa_probe"], ["contact sheet"])
        assert_contains(item["release_boundary"], ["public_blocked"])


def test_run2_6_has_aesthetic_benchmark_bank() -> None:
    sources = load_json(PACK / "sources.json")
    benchmarks = load_json(PACK / "aesthetic_benchmark_bank.json")
    source_ids = {source["id"] for source in sources["sources"]}

    assert benchmarks["status"] == "run2_6_aesthetic_benchmark_bank_public_blocked"
    assert benchmarks["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert benchmarks["storage_policy"]["raw_media"] == "forbidden"
    assert EXPECTED_RUN2_6_BENCHMARK_IDS <= {item["id"] for item in benchmarks["benchmarks"]}

    for item in benchmarks["benchmarks"]:
        assert EXPECTED_RUN2_6_BENCHMARK_FIELDS <= set(item), item["id"]
        assert set(item["source_ids"]) <= source_ids
        assert item["allowed_use"] == "derived_rules_only"
        assert_contains(" ".join(item["composition_rules"]), ["native"])
        assert_contains(" ".join(item["typography_rules"]), ["hierarchy"])
        assert_contains(" ".join(item["spacing_rules"]), ["spacing"])
        assert_contains(" ".join(item["anti_copy_rules"]), ["do not copy"])
        assert EXPECTED_RUN2_6_TRACE_FIELDS & set(item["trace_fields"])
        assert_contains(item["release_boundary"], ["public_blocked"])


def test_run2_6_has_workflow_decision_policy_and_trace_contract() -> None:
    usecases = load_json(PACK / "commercial_usecase_bank.json")
    benchmarks = load_json(PACK / "aesthetic_benchmark_bank.json")
    policy = load_json(PACK / "workflow_decision_policy.json")
    workflow = load_json(PACK / "skill_workflow.json")
    trace_contract = load_json(PACK / "results" / "trace_manifest_contract.json")

    usecase_ids = {item["id"] for item in usecases["usecases"]}
    benchmark_ids = {item["id"] for item in benchmarks["benchmarks"]}

    assert policy["status"] == "run2_6_workflow_decision_policy_public_blocked"
    assert policy["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert set(policy["selection_chain"]) == {
        "commercial_case",
        "usecase_id",
        "benchmark_ids",
        "theme_policy_id",
        "typography_system_id",
        "spacing_token_set_id",
        "visual_production_module_ids",
        "qa_probes",
    }
    for mapping in policy["usecase_benchmark_map"]:
        assert mapping["usecase_id"] in usecase_ids
        assert set(mapping["benchmark_ids"]) <= benchmark_ids
        assert mapping["theme_policy_id"]
        assert mapping["typography_system_id"]
        assert mapping["spacing_token_set_id"]
        assert_contains(" ".join(mapping["source_brand_sanitization"]), ["do not copy"])

    workflow_stage_ids = {stage["id"] for stage in workflow["stages"]}
    assert {
        "select_commercial_usecase",
        "select_aesthetic_benchmarks",
        "select_theme_typography_spacing_policy",
    } <= workflow_stage_ids
    assert EXPECTED_RUN2_6_TRACE_FIELDS <= set(trace_contract["per_slide_required_fields"])


def test_run2_6r_has_visual_repair_policy() -> None:
    policy = load_json(PACK / "visual_repair_policy.json")
    workflow_policy = load_json(PACK / "workflow_decision_policy.json")
    workflow_policy_ids = {
        mapping["theme_policy_id"] for mapping in workflow_policy["usecase_benchmark_map"]
    } | {
        mapping["typography_system_id"] for mapping in workflow_policy["usecase_benchmark_map"]
    } | {
        mapping["spacing_token_set_id"] for mapping in workflow_policy["usecase_benchmark_map"]
    } | {
        benchmark_id
        for mapping in workflow_policy["usecase_benchmark_map"]
        for benchmark_id in mapping["benchmark_ids"]
    }

    assert policy["status"] == "run2_6r_visual_repair_policy_public_blocked"
    assert policy["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert policy["default_visual_direction"] == "light_first_editorial_graphite_with_vivid_proof_color"
    assert EXPECTED_RUN2_6R_REPAIR_IDS <= {repair["id"] for repair in policy["repairs"]}
    policy_delta_text = " ".join(
        repair[field]
        for repair in policy["repairs"]
        for field in ("typography_delta", "spacing_delta", "composition_delta", "theme_delta")
    )
    assert_contains(policy_delta_text, ["Run 2.5", "native", "forest-green", "source-brand"])

    for repair in policy["repairs"]:
        assert EXPECTED_RUN2_6R_REPAIR_FIELDS <= set(repair), repair["id"]
        assert repair["target_slide_roles"]
        assert set(repair["target_slide_roles"]) <= {"cover", "setup", "contrast", "proof", "climax", "close"}
        assert set(repair["source_policy_ids"]) <= workflow_policy_ids | {"workflow_decision_policy.json"}
        assert_contains(repair["composition_delta"], ["native"])
        assert_contains(repair["theme_delta"], ["forest-green", "source-brand"])
        assert "ppt-run2-5-full-vulca" in repair["must_differ_from"]
        assert "ppt-run2-6-full-vulca" in repair["must_differ_from"]
        assert_contains(" ".join(repair["native_ppt_requirements"]), ["native", "editable"])
        assert_contains(repair["qa_probe"], ["contact sheet"])
        assert_contains(repair["release_boundary"], ["public_blocked"])


def test_run2_6r_workflow_and_trace_contract_include_visual_repair() -> None:
    workflow = load_json(PACK / "skill_workflow.json")
    trace_contract = load_json(PACK / "results" / "trace_manifest_contract.json")

    workflow_stage_ids = {stage["id"] for stage in workflow["stages"]}
    assert "select_visual_repair_policy" in workflow_stage_ids
    workflow_text = json.dumps(workflow)
    assert_contains(workflow_text, ["visual_repair_policy.json", "visual repair", "not Run 3.0"])
    assert EXPECTED_RUN2_6R_TRACE_FIELDS <= set(trace_contract["per_slide_required_fields"])


def test_run2_7_has_real_usecase_and_multimodal_source_records() -> None:
    usecase = load_json(PACK / "run2_7_commercial_usecase.json")
    source_records = load_json(PACK / "run2_7_multimodal_source_records.json")
    sources = load_json(PACK / "sources.json")
    source_ids = {source["id"] for source in sources["sources"]}

    assert usecase["status"] == "run2_7_commercial_usecase_public_blocked"
    assert usecase["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert EXPECTED_RUN2_7_USECASE_IDS <= {usecase["id"]}
    assert usecase["primary_usecase"] == "AI design-to-production platform launch deck"
    assert len(usecase["six_slide_arc"]) == 6
    assert [slide["rhythm_role"] for slide in usecase["six_slide_arc"]] == [
        "cover",
        "setup",
        "contrast",
        "proof",
        "climax",
        "close",
    ]
    assert_contains(json.dumps(usecase), ["AI product builders", "design engineering leaders", "technical founders"])
    assert_contains(json.dumps(usecase["business_job"]), ["product-system learning", "not one-shot prompting"])
    assert_contains(" ".join(usecase["must_not_show"]), ["copy", "source brand", "full-slide raster"])
    assert_contains(" ".join(usecase["proof_questions"]), ["data", "memory", "workflow", "PPT"])
    assert_contains(usecase["release_boundary"], ["public_blocked"])

    assert source_records["status"] == "run2_7_multimodal_source_records_public_blocked"
    assert source_records["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert source_records["storage_policy"]["raw_media"] == "forbidden"
    assert EXPECTED_RUN2_7_SOURCE_RECORD_IDS <= {record["id"] for record in source_records["records"]}

    covered_modalities = {modality for record in source_records["records"] for modality in record["modalities"]}
    assert {"text", "image_reference", "video", "audio", "transcript", "interaction"} <= covered_modalities

    for record in source_records["records"]:
        assert EXPECTED_RUN2_7_SOURCE_RECORD_FIELDS <= set(record), record["id"]
        assert record["source_id"] in source_ids
        assert record["allowed_use"] == "derived_rules_only"
        assert set(record["slide_roles"]) <= EXPECTED_RHYTHM_ROLES
        assert_contains(record["native_ppt_implication"], ["native", "editable"])
        assert_contains(record["anti_copy_boundary"], ["do not copy"])
        assert_contains(record["qa_probe"], ["contact sheet"])
        assert_contains(record["release_boundary"], ["public_blocked"])


def test_run2_73_source_quality_audit_covers_every_source() -> None:
    sources = load_json(PACK / "sources.json")
    audit = load_json(PACK / "run2_73_source_quality_audit.json")
    source_ids = {source["id"] for source in sources["sources"]}
    audit_source_ids = {record["source_id"] for record in audit["source_audit"]}

    assert audit["status"] == "run2_73_source_quality_audit_public_blocked"
    assert audit["stage_policy"] == "part_a_source_tutorial_audit_only"
    assert audit["storage_policy"]["raw_media"] == "forbidden"
    assert audit_source_ids == source_ids
    assert audit["summary"]["total_sources"] == len(source_ids)
    assert audit["summary"]["full_source_count"] == len(audit["source_audit"])
    assert {
        "true_ppt_tutorial",
        "product_keynote",
        "brand_case",
    } <= set(audit["summary"]["source_type_counts"])
    assert audit["summary"]["sources_with_extracted_design_signal"] < len(source_ids)
    assert audit["summary"]["sources_missing_screenshot_or_shot_or_layout_decomposition"] > 0

    grouped = audit["source_groups"]
    assert grouped["true_ppt_tutorial"]
    assert grouped["product_keynote"]
    assert grouped["brand_case"]
    assert grouped["missing_screenshot_level_or_shot_level_decomposition"]

    for record in audit["source_audit"]:
        assert EXPECTED_RUN2_73_SOURCE_AUDIT_FIELDS <= set(record), record["source_id"]
        assert record["source_type"] in EXPECTED_RUN2_73_SOURCE_TYPES
        assert record["modality"]
        assert isinstance(record["extracted_design_signal"], list)
        assert isinstance(record["generator_obligation"], list)
        assert isinstance(record["generator_ready"], bool)
        assert record["generator_obligation"]
        assert_contains(record["anti_copy_boundary"], ["do not copy"])
        assert "screenshot_level" in record["decomposition_depth"]
        assert "shot_or_timestamp_level" in record["decomposition_depth"]
        assert "layout_level" in record["decomposition_depth"]
        assert set(record["decomposition_depth"].values()) <= EXPECTED_RUN2_73_DECOMPOSITION_DEPTH_VALUES

    missing_signal_records = [
        record for record in audit["source_audit"] if record["signal_status"] == "missing_in_focused_inputs"
    ]
    assert missing_signal_records
    assert all(record["extracted_design_signal"] == [] for record in missing_signal_records)
    assert all(
        "blocked_from_generator_until_decomposed" in record["generator_obligation"]
        for record in missing_signal_records
    )
    assert all(record["generator_ready"] is False for record in missing_signal_records)
    assert all(
        record["generator_use_status"] == "blocked_from_generator_until_decomposed"
        for record in missing_signal_records
    )


def test_run2_73_tutorial_to_design_moves_reconciles_renderer_actions() -> None:
    tutorial = load_json(PACK / "run2_8_tutorial_decomposition.json")
    visual_targets = load_json(PACK / "visual_learning_targets.json")
    motion_targets = load_json(PACK / "motion_learning_targets.json")
    moves = load_json(PACK / "run2_73_tutorial_to_design_moves.json")

    source_rule_ids = {unit["id"] for unit in tutorial["units"]}
    visual_target_ids = {target["id"] for target in visual_targets["targets"]}
    motion_target_ids = {target["id"] for target in motion_targets["targets"]}
    rule_mappings = moves["tutorial_rule_mappings"]
    renderer_actions = [action for rule in rule_mappings for action in rule["renderer_actions"]]
    action_types = set(moves["renderer_action_vocabulary"])

    assert moves["status"] == "run2_73_tutorial_to_design_moves_public_blocked"
    assert moves["stage_policy"] == "part_b_reconciled_tutorial_to_renderer_actions_only"
    assert moves["artifact_scope"]["does_not_start"] == [
        "part_c_slide_story",
        "part_d_scene_compiler",
        "public_release",
    ]
    assert len(rule_mappings) == 20
    assert len(renderer_actions) == 43
    assert {rule["principle"] for rule in rule_mappings} == EXPECTED_RUN2_73_TUTORIAL_MOVE_PRINCIPLES

    for source_input in moves["source_inputs"]:
        source_path = ROOT / source_input["path"]
        assert source_input["available"] is True
        assert source_path.exists(), source_input["path"]
        assert source_input["record_count"] > 0
        assert_contains(source_input["use_in_this_artifact"], ["ids"])

    for rule in rule_mappings:
        assert EXPECTED_RUN2_73_TUTORIAL_MOVE_RULE_FIELDS <= set(rule), rule["rule_id"]
        assert rule["source_rule_ids"]
        assert rule["visual_target_ids"]
        assert rule["motion_target_ids"]
        assert set(rule["source_rule_ids"]) <= source_rule_ids
        assert set(rule["visual_target_ids"]) <= visual_target_ids
        assert set(rule["motion_target_ids"]) <= motion_target_ids
        assert rule["ppt_operations"]
        assert rule["renderer_actions"]
        assert rule["acceptance_checks"]

        lineage = rule["source_lineage"]
        assert {item["id"] for item in lineage["tutorial_rules"]} == set(rule["source_rule_ids"])
        assert {item["id"] for item in lineage["visual_targets"]} == set(rule["visual_target_ids"])
        assert {item["id"] for item in lineage["motion_targets"]} == set(rule["motion_target_ids"])

        for action in rule["renderer_actions"]:
            assert EXPECTED_RUN2_73_RENDERER_ACTION_FIELDS <= set(action), action["action_id"]
            assert action["action_type"] in action_types
            assert action["target_selector"]
            assert action["parameters"]
            assert action["fallback_strategy"]
            assert action["traceability"]["source_rule_ids"] == rule["source_rule_ids"]
            assert action["traceability"]["visual_target_ids"] == rule["visual_target_ids"]
            assert action["traceability"]["motion_target_ids"] == rule["motion_target_ids"]

    summary = moves["traceability_summary"]
    assert summary["rule_count"] == 20
    assert summary["renderer_action_count"] == 43
    assert set(summary["source_rule_ids_covered"]) == source_rule_ids
    assert set(summary["visual_target_ids_covered"]) == visual_target_ids
    assert set(summary["motion_target_ids_covered"]) == motion_target_ids
    assert all(count > 0 for count in summary["coverage_counts_by_input_id"].values())


def test_run2_74_product_claim_graph_defines_c1_product_logic() -> None:
    capability_memory = load_json(PACK / "run2_57_product_capability_memory.json")
    message_contracts = load_json(PACK / "run2_57_slide_message_contracts.json")
    narrative_proof = load_json(PACK / "run2_61_narrative_proof_dataset.json")
    tutorial_moves = load_json(PACK / "run2_73_tutorial_to_design_moves.json")
    graph = load_json(PACK / "run2_74_product_claim_graph.json")

    capability_ids = {record["id"] for record in capability_memory["product_capability_records"]}
    logic_relation_ids = {
        record["id"] for record in capability_memory["product_logic_relation_records"]
    }
    message_contract_ids = {
        record["contract_id"] for record in message_contracts["slide_message_contracts"]
    }
    narrative_proof_ids = {
        record["narrative_proof_id"] for record in narrative_proof["narrative_proof_records"]
    }
    tutorial_rule_ids = {rule["rule_id"] for rule in tutorial_moves["tutorial_rule_mappings"]}

    assert graph["artifact_id"] == "run2_74_product_claim_graph"
    assert graph["part"] == "Part C1"
    assert graph["status"] == "run2_74_product_claim_graph_public_blocked"
    assert graph["stage_policy"] == "part_c1_product_claim_graph_only"
    assert graph["selected_usecase_id"] == "usecase_design_to_production_platform_launch"
    assert graph["artifact_scope"]["does_not_start"] == [
        "part_c2_slide_story",
        "part_c3_content_quality_audit",
        "part_d_scene_compiler",
        "renderer_rerun",
        "public_release",
    ]
    forbidden_c1_root_keys = {
        "slide_story",
        "slide_stories",
        "slide_scenes",
        "scene_compiler",
        "renderer_actions",
        "layout_geometry",
        "pptx_output",
    }
    assert not forbidden_c1_root_keys & set(graph)

    for source_input in graph["source_inputs"]:
        source_path = ROOT / source_input["path"]
        assert source_input["available"] is True
        assert source_path.exists(), source_input["path"]
        assert source_input["use_in_this_artifact"]

    product_definition = graph["product_definition"]
    assert product_definition["product_category"] == "source-bound editable PPT generation system"
    assert product_definition["primary_output"] == "code-generated editable PPT deck"
    assert_contains(product_definition["one_sentence"], ["commercial brief", "editable PPT"])
    assert "image-only slide generator" in product_definition["not_product"]
    assert "generic template library" in product_definition["not_product"]

    nodes = graph["product_claim_nodes"]
    node_ids = {node["claim_id"] for node in nodes}
    assert node_ids == EXPECTED_RUN2_74_PRODUCT_CLAIM_NODE_IDS

    forbidden_empty_claims = {
        claim.lower() for claim in capability_memory["forbidden_empty_claims"]
    }
    route_terms = {term.lower() for term in graph["public_surface_routing"]["viewer_or_note_only_terms"]}
    for node in nodes:
        assert EXPECTED_RUN2_74_PRODUCT_CLAIM_NODE_FIELDS <= set(node), node["claim_id"]
        assert node["public_claim"]
        assert not any(claim in node["public_claim"].lower() for claim in forbidden_empty_claims)
        assert set(node["product_capability_ids"]) <= capability_ids
        assert set(node["product_logic_relation_ids"]) <= logic_relation_ids
        assert set(node["source_message_contract_ids"]) <= message_contract_ids
        assert set(node["narrative_proof_ids"]) <= narrative_proof_ids
        assert set(node["tutorial_design_rule_ids"]) <= tutorial_rule_ids
        assert node["user_benefit"]
        assert node["visible_user_result"]
        assert node["proof_standard"]
        assert node["public_surface_route"]["on_canvas"]
        assert node["public_surface_route"]["speaker_note_or_viewer"]
        assert "slide_index" not in node
        assert "slide_title" not in node
        assert "scene_layout" not in node
        assert "renderer_actions" not in node
        on_canvas_text = " ".join(node["public_surface_route"]["on_canvas"]).lower()
        assert not any(term in on_canvas_text for term in route_terms)

    edges = graph["causal_edges"]
    assert len(edges) == 6
    for edge in edges:
        assert EXPECTED_RUN2_74_CAUSAL_EDGE_FIELDS <= set(edge), edge["edge_id"]
        assert edge["from_claim_id"] in node_ids
        assert edge["to_claim_id"] in node_ids
        assert edge["causal_relation"]
        assert edge["business_logic"]
        assert edge["user_value_delta"]
        assert edge["proof_refs"]

    trust_model = graph["trust_model"]
    assert len(trust_model["why_credible"]) >= 4
    assert "public release remains blocked until review passes" in trust_model["release_truth"]
    assert graph["traceability_summary"]["claim_node_count"] == 7
    assert graph["traceability_summary"]["causal_edge_count"] == 6
    assert graph["next_required_action"] == "run2_74_c2_create_slide_story_from_product_claim_graph"


def test_run2_74_slide_story_maps_claim_graph_to_six_slide_story() -> None:
    claim_graph = load_json(PACK / "run2_74_product_claim_graph.json")
    message_contracts = load_json(PACK / "run2_57_slide_message_contracts.json")
    narrative_proof = load_json(PACK / "run2_61_narrative_proof_dataset.json")
    story = load_json(PACK / "run2_74_slide_story.json")

    claim_ids = {node["claim_id"] for node in claim_graph["product_claim_nodes"]}
    edge_ids = {edge["edge_id"] for edge in claim_graph["causal_edges"]}
    message_contract_by_role = {
        record["role"]: record for record in message_contracts["slide_message_contracts"]
    }
    narrative_proof_by_role = {
        record["role"]: record for record in narrative_proof["narrative_proof_records"]
    }

    assert story["artifact_id"] == "run2_74_slide_story"
    assert story["part"] == "Part C2"
    assert story["status"] == "run2_74_slide_story_public_blocked"
    assert story["stage_policy"] == "part_c2_slide_story_only"
    assert story["source_product_claim_graph"] == "run2_74_product_claim_graph.json"
    assert story["artifact_scope"]["does_not_start"] == [
        "part_c3_content_quality_audit",
        "part_d_scene_compiler",
        "renderer_rerun",
        "public_release",
    ]
    forbidden_root_keys = {
        "scene_compiler",
        "renderer_actions",
        "layout_geometry",
        "pptx_output",
        "html_viewer",
    }
    assert not forbidden_root_keys & set(story)

    slides = story["slides"]
    assert [slide["role"] for slide in slides] == EXPECTED_RUN2_74_SLIDE_STORY_ROLES
    assert [slide["slide_index"] for slide in slides] == [1, 2, 3, 4, 5, 6]

    route_terms = {
        term.lower()
        for term in claim_graph["public_surface_routing"]["viewer_or_note_only_terms"]
    }
    for slide in slides:
        assert EXPECTED_RUN2_74_SLIDE_STORY_FIELDS <= set(slide), slide["role"]
        assert slide["slide_id"] == f"slide_story_2_74_{slide['role']}"
        assert set(slide["product_claim_ids"]) <= claim_ids
        assert set(slide["causal_edge_ids"]) <= edge_ids
        assert slide["source_message_contract_id"] == message_contract_by_role[slide["role"]]["contract_id"]
        assert slide["narrative_proof_id"] == narrative_proof_by_role[slide["role"]]["narrative_proof_id"]
        assert slide["audience_question"] == message_contract_by_role[slide["role"]]["reader_question"]
        assert slide["title"]
        assert slide["thesis"]
        assert slide["proof_object"]["primary"]
        assert slide["proof_object"]["evidence"]
        assert slide["visual_object"]["role"]
        assert slide["visual_object"]["must_show"]
        assert not {"x", "y", "w", "h", "width", "height", "left", "top"} & set(
            slide["visual_object"]
        )
        assert slide["on_canvas_copy"]["headline"]
        assert slide["on_canvas_copy"]["supporting_line"]
        assert slide["on_canvas_copy"]["proof_labels"]
        assert slide["speaker_note_or_viewer_route"]["speaker_note"]
        assert slide["speaker_note_or_viewer_route"]["viewer_only"]
        assert 35 <= slide["text_budget"]["max_public_words"] <= 90
        assert slide["text_budget"]["max_proof_labels"] <= 4
        assert slide["demo_goal"]
        assert slide["content_risk"]["failure_mode"]
        assert slide["content_risk"]["avoid"]
        assert slide["handoff_to_scene_compiler"]["required_scene_question"]
        assert "layout_geometry" not in slide
        assert "renderer_actions" not in slide
        on_canvas_text = " ".join(
            [
                slide["on_canvas_copy"]["headline"],
                slide["on_canvas_copy"]["supporting_line"],
                " ".join(slide["on_canvas_copy"]["proof_labels"]),
            ]
        ).lower()
        assert not any(term in on_canvas_text for term in route_terms)

    assert story["traceability_summary"]["slide_count"] == 6
    assert story["traceability_summary"]["claim_ids_covered"] == sorted(claim_ids)
    assert story["next_required_action"] == "run2_74_c3_audit_slide_story_content_quality"


def test_run2_74_content_quality_audit_checks_slide_story_quality() -> None:
    claim_graph = load_json(PACK / "run2_74_product_claim_graph.json")
    story = load_json(PACK / "run2_74_slide_story.json")
    audit = load_json(PACK / "run2_74_content_quality_audit.json")

    claim_ids = {node["claim_id"] for node in claim_graph["product_claim_nodes"]}
    route_terms = {
        term.lower()
        for term in claim_graph["public_surface_routing"]["viewer_or_note_only_terms"]
    }
    story_by_role = {slide["role"]: slide for slide in story["slides"]}

    assert audit["artifact_id"] == "run2_74_content_quality_audit"
    assert audit["part"] == "Part C3"
    assert audit["status"] == "run2_74_content_quality_audit_public_blocked"
    assert audit["stage_policy"] == "part_c3_content_quality_audit_only"
    assert audit["source_product_claim_graph"] == "run2_74_product_claim_graph.json"
    assert audit["source_slide_story"] == "run2_74_slide_story.json"
    assert audit["artifact_scope"]["does_not_start"] == [
        "part_d_scene_compiler",
        "renderer_rerun",
        "public_release",
    ]
    forbidden_root_keys = {
        "scene_compiler",
        "renderer_actions",
        "layout_geometry",
        "pptx_output",
        "html_viewer",
    }
    assert not forbidden_root_keys & set(audit)

    for source_input in audit["source_inputs"]:
        source_path = ROOT / source_input["path"]
        assert source_input["available"] is True
        assert source_path.exists(), source_input["path"]

    assert set(audit["quality_dimensions"]) == EXPECTED_RUN2_74_CONTENT_QA_DIMENSIONS

    slide_audits = audit["slide_quality_audits"]
    assert [record["role"] for record in slide_audits] == EXPECTED_RUN2_74_SLIDE_STORY_ROLES
    assert len(slide_audits) == 6

    forbidden_constraint_keys = {"x", "y", "w", "h", "width", "height", "left", "top"}
    for record in slide_audits:
        role = record["role"]
        source_slide = story_by_role[role]
        assert EXPECTED_RUN2_74_CONTENT_QA_FIELDS <= set(record), role
        assert record["audit_id"] == f"content_qa_2_74_{role}"
        assert record["source_slide_story_id"] == source_slide["slide_id"]
        assert record["checked_claim_ids"] == source_slide["product_claim_ids"]
        assert set(record["checked_claim_ids"]) <= claim_ids
        assert record["blocking_issues"] == []
        assert record["watch_items"]

        checks = record["quality_checks"]
        assert {check["dimension"] for check in checks} == EXPECTED_RUN2_74_CONTENT_QA_DIMENSIONS
        for check in checks:
            assert check["status"] in {"pass", "watch"}
            assert check["evidence"]
            assert check["route_decision"]

        units = record["approved_content_units"]
        assert EXPECTED_RUN2_74_APPROVED_CONTENT_UNIT_FIELDS <= set(units), role
        assert units["main_claim"] == source_slide["thesis"]
        assert units["proof_object"] == source_slide["proof_object"]["primary"]
        assert units["visual_role"] == source_slide["visual_object"]["role"]
        assert units["business_logic"]
        assert units["on_canvas_copy"] == source_slide["on_canvas_copy"]
        assert units["delete_or_note_route"]["speaker_note"] == source_slide[
            "speaker_note_or_viewer_route"
        ]["speaker_note"]
        assert units["delete_or_note_route"]["viewer_only"] == source_slide[
            "speaker_note_or_viewer_route"
        ]["viewer_only"]
        assert units["delete_or_note_route"]["must_delete_from_canvas"]

        on_canvas_text = json.dumps(units["on_canvas_copy"], ensure_ascii=False).lower()
        assert not any(term in on_canvas_text for term in route_terms)
        assert word_count(units["on_canvas_copy"]["headline"]) + word_count(
            units["on_canvas_copy"]["supporting_line"]
        ) <= source_slide["text_budget"]["max_public_words"]

        constraints = record["scene_compiler_constraints"]
        assert constraints["content_only"] is True
        assert constraints["must_keep_main_claim"] == units["main_claim"]
        assert constraints["must_render_proof_object"] == units["proof_object"]
        assert constraints["must_preserve_visual_role"] == units["visual_role"]
        assert constraints["must_route_off_canvas"] == units["delete_or_note_route"]["viewer_only"]
        assert constraints["must_not_add"]
        assert not forbidden_constraint_keys & set(constraints)

    summary = audit["traceability_summary"]
    assert summary["slide_count"] == 6
    assert summary["quality_dimensions"] == sorted(EXPECTED_RUN2_74_CONTENT_QA_DIMENSIONS)
    assert summary["blocking_issue_count"] == 0
    assert summary["approved_for_next_stage"] == "part_d_scene_compiler_inputs_only"
    assert audit["next_required_action"] == "part_d_scene_compiler_from_c1_c2_c3_and_b2"


def test_run2_73_scene_compiler_contracts_merge_b2_and_c_content_units() -> None:
    design_moves = load_json(PACK / "run2_73_tutorial_to_design_moves.json")
    claim_graph = load_json(PACK / "run2_74_product_claim_graph.json")
    story = load_json(PACK / "run2_74_slide_story.json")
    audit = load_json(PACK / "run2_74_content_quality_audit.json")
    scene_contract = load_json(PACK / "run2_73_scene_compiler_contracts.json")

    story_by_role = {slide["role"]: slide for slide in story["slides"]}
    audit_by_role = {record["role"]: record for record in audit["slide_quality_audits"]}
    design_rule_by_id = {rule["rule_id"]: rule for rule in design_moves["tutorial_rule_mappings"]}
    renderer_action_types = set(design_moves["renderer_action_vocabulary"])
    route_terms = {
        term.lower()
        for term in claim_graph["public_surface_routing"]["viewer_or_note_only_terms"]
    }

    assert scene_contract["artifact_id"] == "run2_73_scene_compiler_contracts"
    assert scene_contract["part"] == "Part D1"
    assert scene_contract["status"] == "run2_73_scene_compiler_contracts_public_blocked"
    assert scene_contract["stage_policy"] == "part_d1_scene_compiler_contract_only"
    assert scene_contract["source_design_moves"] == "run2_73_tutorial_to_design_moves.json"
    assert scene_contract["source_product_claim_graph"] == "run2_74_product_claim_graph.json"
    assert scene_contract["source_slide_story"] == "run2_74_slide_story.json"
    assert scene_contract["source_content_quality_audit"] == "run2_74_content_quality_audit.json"
    assert scene_contract["artifact_scope"]["does_not_start"] == [
        "renderer_rerun",
        "pptx_output",
        "html_viewer",
        "public_release",
    ]
    assert scene_contract["scene_structure_schema"]["required_scene_fields"] == sorted(
        EXPECTED_RUN2_D1_SCHEMA_REQUIRED_FIELDS
    )
    assert (
        scene_contract["scene_structure_schema"]["content_unit_source"]
        == "slide_quality_audits[*].approved_content_units"
    )
    assert scene_contract["renderer_execution_contract"][
        "must_execute_scene_plan_before_shape_generation"
    ] is True
    assert scene_contract["renderer_execution_contract"][
        "must_not_directly_draw_placeholder_boxes"
    ] is True

    forbidden_root_keys = {
        "layout_geometry",
        "pptx_output",
        "html_viewer",
        "slide_shapes",
        "shape_rectangles",
        "direct_shape_plan",
    }
    assert not forbidden_root_keys & set(scene_contract)

    forbidden_scene_keys = {
        "layout_geometry",
        "pptx_output",
        "html_viewer",
        "slide_shapes",
        "shape_rectangles",
        "direct_shape_plan",
        "x",
        "y",
        "w",
        "h",
        "width",
        "height",
        "left",
        "top",
        "x_pct",
        "y_pct",
        "width_pct",
        "height_pct",
    }
    scenes = scene_contract["scene_plans"]
    assert [scene["role"] for scene in scenes] == EXPECTED_RUN2_74_SLIDE_STORY_ROLES
    assert [scene["slide_index"] for scene in scenes] == [1, 2, 3, 4, 5, 6]

    for scene in scenes:
        role = scene["role"]
        source_slide = story_by_role[role]
        source_audit = audit_by_role[role]
        units = source_audit["approved_content_units"]

        assert EXPECTED_RUN2_D1_SCENE_PLAN_FIELDS <= set(scene), role
        assert not forbidden_scene_keys & nested_keys(scene)
        assert scene["scene_id"] == f"scene_compiler_2_73_{role}"
        assert scene["source_slide_story_id"] == source_slide["slide_id"]
        assert scene["source_content_qa_id"] == source_audit["audit_id"]
        assert "main_claim" not in source_audit
        assert "proof_object" not in source_audit
        assert scene["content_units"] == units

        assert scene["proof_object"]["primary"] == units["proof_object"]
        assert scene["proof_object"]["evidence"] == source_slide["proof_object"]["evidence"]
        assert scene["proof_object"]["business_logic"] == units["business_logic"]
        assert scene["visual_role"] == units["visual_role"]

        layout = scene["layout_intent"]
        assert EXPECTED_RUN2_D1_LAYOUT_INTENT_FIELDS <= set(layout), role
        assert layout["first_read"] == units["on_canvas_copy"]["headline"]
        assert layout["main_visual"] == units["proof_object"]
        assert layout["scene_question"] == source_slide["handoff_to_scene_compiler"][
            "required_scene_question"
        ]
        assert layout["information_containers"]
        assert layout["chrome"]
        assert layout["must_delete"] == units["delete_or_note_route"]["must_delete_from_canvas"]
        assert {item["container_role"] for item in layout["information_containers"]} >= {
            "claim_container",
            "proof_container",
        }
        assert all(item["chrome_role"] for item in layout["chrome"])

        assert scene["off_canvas_routes"] == units["delete_or_note_route"]
        on_canvas_text = json.dumps(units["on_canvas_copy"], ensure_ascii=False).lower()
        assert not any(term in on_canvas_text for term in route_terms)

        assert scene["renderer_action_bindings"]
        for binding in scene["renderer_action_bindings"]:
            assert EXPECTED_RUN2_D1_RENDERER_BINDING_FIELDS <= set(binding), binding[
                "binding_id"
            ]
            assert set(binding["source_rule_ids"]) <= set(design_rule_by_id)
            assert set(binding["renderer_action_types"]) <= renderer_action_types
            assert binding["design_moves"]
            assert binding["applies_to_scene_questions"]
            assert binding["acceptance_checks"]

            source_actions = [
                action
                for rule_id in binding["source_rule_ids"]
                for action in design_rule_by_id[rule_id]["renderer_actions"]
            ]
            assert set(binding["renderer_action_ids"]) <= {
                action["action_id"] for action in source_actions
            }
            assert set(binding["renderer_action_types"]) <= {
                action["action_type"] for action in source_actions
            }

        boundary = scene["renderer_execution_boundary"]
        assert boundary["must_use_scene_plan"] is True
        assert boundary["must_not_read_content_from_c3_record_top_level"] is True
        assert boundary["must_not_draw_generic_boxes"] is True

    summary = scene_contract["traceability_summary"]
    assert summary["slide_count"] == 6
    assert summary["scene_count"] == 6
    assert summary["approved_content_units_source"] == (
        "run2_74_content_quality_audit.json:slide_quality_audits[*].approved_content_units"
    )
    assert set(summary["design_rule_ids_bound"]) <= set(design_rule_by_id)
    assert summary["renderer_action_type_count"] >= 6
    assert scene_contract["next_required_action"] == (
        "part_d2_renderer_execute_scene_plan_not_direct_boxes"
    )


def test_run2_73_scene_plan_expansion_builds_renderer_ready_structures() -> None:
    scene_contract = load_json(PACK / "run2_73_scene_compiler_contracts.json")
    expansion = load_json(PACK / "run2_73_scene_plan_expansion.json")

    scenes_by_role = {scene["role"]: scene for scene in scene_contract["scene_plans"]}
    allowed_types = EXPECTED_RUN2_D2_ALLOWED_ELEMENT_TYPES

    assert expansion["artifact_id"] == "run2_73_scene_plan_expansion"
    assert expansion["part"] == "Part D2"
    assert expansion["status"] == "run2_73_scene_plan_expansion_public_blocked"
    assert expansion["stage_policy"] == "part_d2_scene_plan_expansion_only"
    assert expansion["source_scene_compiler_contracts"] == "run2_73_scene_compiler_contracts.json"
    assert expansion["artifact_scope"]["does_not_start"] == [
        "part_d3_renderer_input_validation",
        "renderer_rerun",
        "pptx_output",
        "html_viewer",
        "public_release",
    ]
    assert expansion["component_schema"]["required_components"] == sorted(
        EXPECTED_RUN2_D2_COMPONENT_KEYS
    )
    assert set(expansion["component_schema"]["allowed_element_types"]) == allowed_types
    assert expansion["renderer_ready_contract"]["no_empty_boxes"] is True
    assert expansion["renderer_ready_contract"]["every_visual_container_bound"] is True
    assert expansion["renderer_ready_contract"]["geometry_resolution_deferred"] is True

    forbidden_keys = {
        "layout_geometry",
        "pptx_output",
        "html_viewer",
        "slide_shapes",
        "shape_rectangles",
        "direct_shape_plan",
        "x",
        "y",
        "w",
        "h",
        "width",
        "height",
        "left",
        "top",
        "x_pct",
        "y_pct",
        "width_pct",
        "height_pct",
    }
    assert not forbidden_keys & nested_keys(expansion)

    structures = expansion["scene_structures"]
    assert [structure["role"] for structure in structures] == EXPECTED_RUN2_74_SLIDE_STORY_ROLES
    assert [structure["slide_index"] for structure in structures] == [1, 2, 3, 4, 5, 6]

    for structure in structures:
        role = structure["role"]
        source_scene = scenes_by_role[role]
        source_binding = source_scene["renderer_action_bindings"][0]
        source_action_ids = set(source_binding["renderer_action_ids"])
        source_action_types = set(source_binding["renderer_action_types"])

        assert EXPECTED_RUN2_D2_SCENE_STRUCTURE_FIELDS <= set(structure), role
        assert structure["expansion_id"] == f"scene_expansion_2_73_{role}"
        assert structure["source_scene_id"] == source_scene["scene_id"]
        assert structure["off_canvas_contract"] == source_scene["off_canvas_routes"]

        components = structure["semantic_components"]
        assert set(components) == EXPECTED_RUN2_D2_COMPONENT_KEYS
        component_ids = {component["component_id"] for component in components.values()}
        element_types = {component["element_type"] for component in components.values()}
        assert element_types <= allowed_types
        assert {"shape", "text", "connector"} <= element_types

        for key, component in components.items():
            assert EXPECTED_RUN2_D2_COMPONENT_FIELDS <= set(component), key
            assert component["component_id"].startswith(f"{role}_")
            assert component["component_role"] == key
            assert component["empty_box_policy"] == "forbidden"
            assert component["source_scene_plan_fields"]
            assert component["renderer_target"]["target_component_id"] == component["component_id"]
            assert component["renderer_target"]["target_element_type"] == component["element_type"]
            assert component["content_binding"]["source"]
            assert component["content_binding"]["value"] or component["semantic_role_binding"]["value"]

        assert components["hero_object"]["content_binding"]["value"] == source_scene[
            "proof_object"
        ]["primary"]
        assert components["supporting_copy"]["content_binding"]["value"] == source_scene[
            "content_units"
        ]["on_canvas_copy"]
        assert components["viewer_note_route"]["content_binding"]["value"] == source_scene[
            "off_canvas_routes"
        ]

        for container in structure["visual_containers"]:
            assert container["container_id"].startswith(f"{role}_")
            assert container["bound_component_id"] in component_ids
            assert container["container_type"] in {
                "component_container",
                "information_container",
                "chrome",
            }
            assert container["empty_box_policy"] == "forbidden"
            assert container["content_binding"]["value"] or container["semantic_role_binding"]["value"]

        bindings = structure["expanded_renderer_action_bindings"]
        assert len(bindings) >= len(EXPECTED_RUN2_D2_COMPONENT_KEYS)
        assert {binding["target_component_key"] for binding in bindings} == (
            EXPECTED_RUN2_D2_COMPONENT_KEYS
        )

        covered_action_ids = set()
        for binding in bindings:
            assert EXPECTED_RUN2_D2_EXPANDED_BINDING_FIELDS <= set(binding), binding[
                "binding_id"
            ]
            assert binding["target_component_id"] in component_ids
            assert (
                binding["component_element_type"]
                == components[binding["target_component_key"]]["element_type"]
            )
            assert binding["source_d1_binding_id"] == source_binding["binding_id"]
            assert set(binding["source_b2_action_ids"]) <= source_action_ids
            assert set(binding["source_b2_action_types"]) <= source_action_types
            assert binding["source_b2_action_ids"]
            assert binding["design_moves"]
            assert binding["layout_intent_binding"] in {
                "first_read",
                "main_visual",
                "information_containers",
                "chrome",
                "must_delete",
                "off_canvas_routes",
            }
            assert binding["renderer_resolution_policy"]["geometry_resolution"] == (
                "defer_to_renderer"
            )
            assert binding["renderer_resolution_policy"]["must_not_create_empty_box"] is True
            covered_action_ids |= set(binding["source_b2_action_ids"])

        assert source_action_ids <= covered_action_ids
        boundary = structure["renderer_ready_boundary"]
        assert boundary["is_renderer_ready_json"] is True
        assert boundary["must_not_render_ppt"] is True
        assert boundary["must_validate_before_d3"] is True

    summary = expansion["traceability_summary"]
    assert summary["scene_count"] == 6
    assert summary["component_count"] == 30
    assert summary["source_scene_compiler_contracts"] == "run2_73_scene_compiler_contracts.json"
    assert summary["empty_visual_container_count"] == 0
    assert expansion["next_required_action"] == "part_d3_renderer_input_validation"


def test_run2_73_renderer_input_validation_blocks_bad_scene_inputs() -> None:
    from scripts.build_ppt_run2_73_renderer_input_validation import validate_renderer_input

    expansion = load_json(PACK / "run2_73_scene_plan_expansion.json")
    validation = load_json(PACK / "run2_73_renderer_input_validation.json")

    assert validation["artifact_id"] == "run2_73_renderer_input_validation"
    assert validation["part"] == "Part D3"
    assert validation["status"] == "run2_73_renderer_input_validation_passed_public_blocked"
    assert validation["stage_policy"] == "part_d3_renderer_input_validation_only"
    assert validation["source_scene_plan_expansion"] == "run2_73_scene_plan_expansion.json"
    assert validation["artifact_scope"]["does_not_start"] == [
        "renderer_rerun",
        "pptx_output",
        "html_viewer",
        "public_release",
    ]
    guard = validation["execution_guard"]
    assert guard["mode"] == "validation_only"
    assert guard["rendering_subprocesses_allowed"] is False
    assert guard["allowed_side_effects"] == [
        "read_run2_73_scene_plan_expansion_json",
        "write_run2_73_renderer_input_validation_json",
    ]
    assert guard["forbidden_invocations"] == validation["artifact_scope"]["does_not_start"]
    assert set(guard["forbidden_runtime_imports"]) == (
        EXPECTED_RUN2_D3_FORBIDDEN_RUNTIME_IMPORTS
    )
    assert set(guard["forbidden_dynamic_import_calls"]) == (
        EXPECTED_RUN2_D3_FORBIDDEN_DYNAMIC_IMPORT_CALLS
    )

    script_path = ROOT / "scripts" / "build_ppt_run2_73_renderer_input_validation.py"
    script_tree = ast.parse(script_path.read_text())
    script_imports = set()
    script_calls = set()
    for node in ast.walk(script_tree):
        if isinstance(node, ast.Import):
            script_imports |= {alias.name.split(".")[0] for alias in node.names}
        elif isinstance(node, ast.ImportFrom) and node.module:
            script_imports.add(node.module.split(".")[0])
        elif isinstance(node, ast.Call) and isinstance(node.func, ast.Name):
            script_calls.add(node.func.id)
    assert not EXPECTED_RUN2_D3_FORBIDDEN_RUNTIME_IMPORTS & script_imports
    assert not EXPECTED_RUN2_D3_FORBIDDEN_DYNAMIC_IMPORT_CALLS & script_calls

    policy = validation["validation_policy"]
    assert set(policy["required_check_ids"]) == EXPECTED_RUN2_D3_VALIDATION_CHECK_IDS
    assert set(policy["allowed_element_types"]) == EXPECTED_RUN2_D2_ALLOWED_ELEMENT_TYPES
    assert policy["failure_behavior"] == "block_renderer_handoff"
    assert policy["off_canvas_term_sources"] == [
        "scene_structures[*].off_canvas_contract.viewer_only",
        "scene_structures[*].off_canvas_contract.must_delete_from_canvas",
    ]
    assert {
        "layout_geometry",
        "pptx_output",
        "html_viewer",
        "x_pct",
        "width_pct",
    } <= set(policy["forbidden_renderer_field_keys"])

    forbidden_keys = {
        "layout_geometry",
        "pptx_output",
        "html_viewer",
        "slide_shapes",
        "shape_rectangles",
        "direct_shape_plan",
        "x",
        "y",
        "w",
        "h",
        "width",
        "height",
        "left",
        "top",
        "x_pct",
        "y_pct",
        "width_pct",
        "height_pct",
    }
    assert not forbidden_keys & nested_keys(validation["renderer_handoff_contract"])
    assert validation["renderer_handoff_contract"]["input_status"] == "validated_renderer_ready_json"
    assert validation["renderer_handoff_contract"]["input_contract_scope"] == (
        "strictly_d2_renderer_ready_json_not_ppt_or_html_schema"
    )
    assert validation["renderer_handoff_contract"]["must_not_render_in_d3"] is True
    assert validation["renderer_handoff_contract"]["must_consume_source_artifact"] == (
        "run2_73_scene_plan_expansion.json"
    )

    scene_results = validation["scene_validation_results"]
    assert [record["role"] for record in scene_results] == EXPECTED_RUN2_74_SLIDE_STORY_ROLES
    assert [record["status"] for record in scene_results] == ["pass"] * 6
    assert len(scene_results) == 6

    expansion_by_role = {structure["role"]: structure for structure in expansion["scene_structures"]}
    for record in scene_results:
        role = record["role"]
        source_structure = expansion_by_role[role]
        assert EXPECTED_RUN2_D3_SCENE_VALIDATION_FIELDS <= set(record), role
        assert record["validation_id"] == f"renderer_input_validation_2_73_{role}"
        assert record["source_expansion_id"] == source_structure["expansion_id"]
        assert record["blocking_issues"] == []
        assert record["checked_components"] == len(source_structure["semantic_components"])
        assert record["checked_visual_containers"] == len(source_structure["visual_containers"])
        assert record["checked_bindings"] == len(
            source_structure["expanded_renderer_action_bindings"]
        )
        assert {check["check_id"] for check in record["validation_checks"]} == (
            EXPECTED_RUN2_D3_VALIDATION_CHECK_IDS
        )
        assert all(check["status"] == "pass" for check in record["validation_checks"])
        assert record["renderer_handoff"]["can_handoff_to_renderer"] is True
        assert record["renderer_handoff"]["must_not_render_in_d3"] is True
        assert record["renderer_handoff"]["component_manifest_count"] == 5
        assert record["renderer_handoff"]["action_binding_count"] == 5

    summary = validation["traceability_summary"]
    assert summary["scene_count"] == 6
    assert summary["component_count"] == 30
    assert summary["visual_container_count"] == 60
    assert summary["expanded_renderer_binding_count"] == 30
    assert summary["blocking_issue_count"] == 0
    assert summary["renderer_handoff_approved"] == "validated_renderer_ready_json_only"
    assert validation["next_required_action"] == "part_e2_renderer_adapter_contracts"

    bad = json.loads(json.dumps(expansion))
    cover = bad["scene_structures"][0]
    cover["semantic_components"]["hero_object"]["element_type"] = "placeholder_box"
    cover["visual_containers"][0]["content_binding"]["value"] = ""
    cover["visual_containers"][0]["semantic_role_binding"]["value"] = ""
    cover["layout_geometry"] = {"x_pct": 1}
    cover["pptx_output"] = "forbidden ppt schema"
    cover["semantic_components"]["supporting_copy"]["content_binding"]["value"][
        "proof_labels"
    ].append(cover["off_canvas_contract"]["viewer_only"][0])
    cover["expanded_renderer_action_bindings"][0]["source_b2_action_ids"] = []

    bad_result = validate_renderer_input(bad)
    issue_ids = {issue["issue_id"] for issue in bad_result["blocking_issues"]}
    assert EXPECTED_RUN2_D3_BLOCKING_ISSUE_IDS <= issue_ids
    assert any(
        issue["issue_id"] == "forbidden_renderer_field" and "pptx_output" in issue["location"]
        for issue in bad_result["blocking_issues"]
    )
    assert bad_result["status"] == "fail"
    assert bad_result["traceability_summary"]["blocking_issue_count"] >= len(
        EXPECTED_RUN2_D3_BLOCKING_ISSUE_IDS
    )


def test_run2_73_visual_grammar_modules_define_part_e_non_card_layer() -> None:
    modules = load_json(PACK / "run2_73_visual_grammar_modules.json")
    design_grammar = load_json(PACK / "run2_66_reference_first_design_grammar.json")
    semantic_assets = load_json(PACK / "run2_43_semantic_visual_asset_memory.json")
    object_grammar = load_json(PACK / "run2_46_visual_object_grammar_memory.json")

    assert modules["artifact_id"] == "run2_73_visual_grammar_modules"
    assert modules["part"] == "Part E"
    assert modules["status"] == "run2_73_visual_grammar_modules_ready_public_blocked"
    assert modules["stage_policy"] == (
        "part_e_visual_grammar_modules_only_no_renderer_rerun_no_public_release"
    )
    assert EXPECTED_RUN2_E_FORBIDDEN_SCOPE <= set(modules["artifact_scope"]["does_not_start"])

    source_paths = {source["path"] for source in modules["source_inputs"]}
    assert {
        "docs/product/ppt-run2-data-skill-quality/run2_66_reference_first_design_grammar.json",
        "docs/product/ppt-run2-data-skill-quality/run2_43_semantic_visual_asset_memory.json",
        "docs/product/ppt-run2-data-skill-quality/run2_46_visual_object_grammar_memory.json",
    } <= source_paths

    page_mappings = modules["page_type_to_visual_grammar"]
    assert [record["page_type"] for record in page_mappings] == EXPECTED_RUN2_74_SLIDE_STORY_ROLES
    assert {
        record["page_type"]: record["primary_visual_grammar_module"]
        for record in page_mappings
    } == EXPECTED_RUN2_E_PAGE_MODULE_MAP

    source_reference_ids = {
        record["reference_archetype_id"]
        for record in design_grammar["role_design_grammar_records"]
    }
    source_object_grammar_ids = {
        record["visual_object_grammar_id"]
        for record in object_grammar["visual_object_grammar_records"]
    }
    source_semantic_asset_ids = {
        record["semantic_asset_id"]
        for record in semantic_assets["semantic_visual_asset_records"]
    }

    for page in page_mappings:
        role = page["page_type"]
        assert page["source_reference_archetype_id"] in source_reference_ids
        assert page["source_visual_object_grammar_id"] in source_object_grammar_ids
        assert set(page["source_semantic_asset_ids"]) <= source_semantic_asset_ids
        assert len(page["source_semantic_asset_ids"]) == 3
        assert page["main_structure"]["name"]
        assert len(page["main_structure"]["non_rectangular_or_non_card_basis"]) >= 3
        assert page["main_structure"]["serves_content_by"]
        assert len(page["draw_order"]) >= 3
        assert page["success_probe"]
        assert page["primary_visual_grammar_module"] in EXPECTED_RUN2_E_VISUAL_GRAMMAR_MODULE_IDS
        assert role in page["source_reference_archetype_id"]

    visual_modules = modules["visual_grammar_modules"]
    assert {module["module_id"] for module in visual_modules} == (
        EXPECTED_RUN2_E_VISUAL_GRAMMAR_MODULE_IDS
    )
    for module in visual_modules:
        assert len(module["how_to_draw"]) >= 3
        assert module["content_service"]
        assert module["native_ppt_primitives"]
        assert module["avoid"]

    blueprints = modules["module_geometry_blueprints"]
    assert {blueprint["module_id"] for blueprint in blueprints} == (
        EXPECTED_RUN2_E_VISUAL_GRAMMAR_MODULE_IDS
    )
    for blueprint in blueprints:
        assert blueprint["coordinate_system"] == "normalized_16_9_canvas_0_100"
        assert blueprint["primary_structure_is_not"]
        assert blueprint["primary_structure_is"]
        assert len(blueprint["native_ppt_shape_plan"]) >= 3
        assert any("why_not_card" in shape for shape in blueprint["native_ppt_shape_plan"])
        assert all(shape["semantic_role"] for shape in blueprint["native_ppt_shape_plan"])
        assert blueprint["content_attachment_points"]

    selection_rules = modules["module_selection_rules"]
    assert len(selection_rules) == 5
    assert {rule["select_module"] for rule in selection_rules} == (
        EXPECTED_RUN2_E_VISUAL_GRAMMAR_MODULE_IDS
    )
    selected_module_by_role = {
        role: rule["select_module"]
        for rule in selection_rules
        for role in rule["applies_to_page_types"]
    }
    assert selected_module_by_role == EXPECTED_RUN2_E_PAGE_MODULE_MAP
    assert set(modules["coverage_matrix"]["page_roles_covered"]) == set(
        EXPECTED_RUN2_74_SLIDE_STORY_ROLES
    )
    assert set(modules["coverage_matrix"]["modules_covered"]) == (
        EXPECTED_RUN2_E_VISUAL_GRAMMAR_MODULE_IDS
    )
    assert modules["success_criteria_check"] == {
        "every_page_has_non_rectangular_or_non_card_main_structure": True,
        "every_main_structure_serves_content": True,
        "all_requested_modules_defined": True,
        "no_module_depends_on_copied_source_media": True,
        "public_surface_trace_terms_hidden": True,
    }
    assert modules["traceability_summary"]["page_type_count"] == 6
    assert modules["traceability_summary"]["visual_grammar_module_count"] == 5


def test_run2_73_renderer_adapter_contracts_bind_d2_d3_and_part_e() -> None:
    from scripts.build_ppt_run2_73_renderer_adapter_contracts import (
        FORBIDDEN_DYNAMIC_IMPORT_CALLS as ADAPTER_FORBIDDEN_DYNAMIC_IMPORT_CALLS,
        FORBIDDEN_RUNTIME_IMPORTS as ADAPTER_FORBIDDEN_RUNTIME_IMPORTS,
    )

    expansion = load_json(PACK / "run2_73_scene_plan_expansion.json")
    validation = load_json(PACK / "run2_73_renderer_input_validation.json")
    grammar = load_json(PACK / "run2_73_visual_grammar_modules.json")
    adapter = load_json(PACK / "run2_73_renderer_adapter_contracts.json")

    assert adapter["artifact_id"] == "run2_73_renderer_adapter_contracts"
    assert adapter["part"] == "Part E2"
    assert adapter["status"] == "run2_73_renderer_adapter_contracts_ready_public_blocked"
    assert adapter["stage_policy"] == EXPECTED_RUN2_E2_STAGE_POLICY
    assert adapter["source_scene_plan_expansion"] == "run2_73_scene_plan_expansion.json"
    assert adapter["source_renderer_input_validation"] == "run2_73_renderer_input_validation.json"
    assert adapter["source_visual_grammar_modules"] == "run2_73_visual_grammar_modules.json"
    assert EXPECTED_RUN2_E_FORBIDDEN_SCOPE <= set(adapter["artifact_scope"]["does_not_start"])

    guard = adapter["execution_guard"]
    assert guard["mode"] == "adapter_contract_only"
    assert guard["rendering_subprocesses_allowed"] is False
    assert guard["allowed_side_effects"] == [
        "read_run2_73_scene_plan_expansion_json",
        "read_run2_73_renderer_input_validation_json",
        "read_run2_73_visual_grammar_modules_json",
        "write_run2_73_renderer_adapter_contracts_json",
    ]
    assert set(guard["forbidden_invocations"]) == EXPECTED_RUN2_E_FORBIDDEN_SCOPE
    assert set(guard["forbidden_runtime_imports"]) == ADAPTER_FORBIDDEN_RUNTIME_IMPORTS
    assert set(guard["forbidden_dynamic_import_calls"]) == ADAPTER_FORBIDDEN_DYNAMIC_IMPORT_CALLS

    script_path = ROOT / "scripts" / "build_ppt_run2_73_renderer_adapter_contracts.py"
    script_tree = ast.parse(script_path.read_text())
    script_imports = set()
    script_calls = set()
    for node in ast.walk(script_tree):
        if isinstance(node, ast.Import):
            script_imports |= {alias.name.split(".")[0] for alias in node.names}
        elif isinstance(node, ast.ImportFrom) and node.module:
            script_imports.add(node.module.split(".")[0])
        elif isinstance(node, ast.Call) and isinstance(node.func, ast.Name):
            script_calls.add(node.func.id)
    assert not ADAPTER_FORBIDDEN_RUNTIME_IMPORTS & script_imports
    assert not ADAPTER_FORBIDDEN_DYNAMIC_IMPORT_CALLS & script_calls

    source_paths = {source["path"] for source in adapter["source_inputs"]}
    assert {
        "docs/product/ppt-run2-data-skill-quality/run2_73_scene_plan_expansion.json",
        "docs/product/ppt-run2-data-skill-quality/run2_73_renderer_input_validation.json",
        "docs/product/ppt-run2-data-skill-quality/run2_73_visual_grammar_modules.json",
    } <= source_paths

    d2_by_role = {record["role"]: record for record in expansion["scene_structures"]}
    d3_by_role = {record["role"]: record for record in validation["scene_validation_results"]}
    e_by_role = {record["page_type"]: record for record in grammar["page_type_to_visual_grammar"]}
    blueprint_by_module = {record["module_id"]: record for record in grammar["module_geometry_blueprints"]}

    records = adapter["adapter_scene_records"]
    assert [record["role"] for record in records] == EXPECTED_RUN2_74_SLIDE_STORY_ROLES
    assert [record["slide_index"] for record in records] == [1, 2, 3, 4, 5, 6]

    for record in records:
        role = record["role"]
        d2 = d2_by_role[role]
        d3 = d3_by_role[role]
        e = e_by_role[role]
        blueprint = blueprint_by_module[e["primary_visual_grammar_module"]]

        assert record["adapter_scene_id"] == f"renderer_adapter_2_73_{role}"
        assert record["source_expansion_id"] == d2["expansion_id"]
        assert record["source_validation_id"] == d3["validation_id"]
        assert record["validation_status"] == "pass"
        assert record["source_visual_grammar_page_type"] == role
        assert record["visual_grammar_binding"]["module_id"] == e["primary_visual_grammar_module"]
        assert record["visual_grammar_binding"]["module_variant"] == e["module_variant"]
        assert record["visual_grammar_binding"]["main_structure"] == e["main_structure"]
        assert record["visual_grammar_binding"]["draw_order"] == e["draw_order"]
        assert record["geometry_blueprint_binding"]["module_id"] == blueprint["module_id"]
        assert record["geometry_blueprint_binding"]["coordinate_system"] == "normalized_16_9_canvas_0_100"
        assert record["geometry_blueprint_binding"]["native_ppt_shape_plan"] == blueprint["native_ppt_shape_plan"]

        manifest = record["renderer_adapter_manifest"]
        assert manifest["semantic_component_ids"] == [
            component["component_id"] for component in d2["semantic_components"].values()
        ]
        assert manifest["visual_container_ids"] == [
            container["container_id"] for container in d2["visual_containers"]
        ]
        assert manifest["expanded_renderer_binding_ids"] == [
            binding["binding_id"] for binding in d2["expanded_renderer_action_bindings"]
        ]
        assert manifest["d3_renderer_handoff"] == d3["renderer_handoff"]

        instructions = record["adapter_renderer_instructions"]
        assert instructions["draw_primary_structure_before_components"] is True
        assert instructions["apply_geometry_blueprint_before_component_layout"] is True
        assert instructions["bind_semantic_components_before_geometry"] is True
        assert instructions["preserve_off_canvas_contract"] is True
        assert instructions["renderer_execution_allowed_in_this_artifact"] is False
        assert instructions["public_release_allowed_in_this_artifact"] is False

    summary = adapter["traceability_summary"]
    assert summary["scene_count"] == 6
    assert summary["validated_scene_count"] == 6
    assert summary["visual_grammar_module_count"] == 5
    assert summary["geometry_blueprint_count"] == 5
    assert summary["adapter_blocking_issue_count"] == 0
    assert summary["sources_consumed"] == [
        "run2_73_scene_plan_expansion.json",
        "run2_73_renderer_input_validation.json",
        "run2_73_visual_grammar_modules.json",
    ]
    assert adapter["next_required_action"] == EXPECTED_RUN2_E2_NEXT_REQUIRED_ACTION


def test_run2_73_text_binding_strategy_binds_copy_to_visual_sockets() -> None:
    from scripts.build_ppt_run2_73_text_binding_strategy import (
        FORBIDDEN_DYNAMIC_IMPORT_CALLS as TEXT_FORBIDDEN_DYNAMIC_IMPORT_CALLS,
        FORBIDDEN_RUNTIME_IMPORTS as TEXT_FORBIDDEN_RUNTIME_IMPORTS,
    )

    expansion = load_json(PACK / "run2_73_scene_plan_expansion.json")
    adapter = load_json(PACK / "run2_73_renderer_adapter_contracts.json")
    story = load_json(PACK / "run2_74_slide_story.json")
    audit = load_json(PACK / "run2_74_content_quality_audit.json")
    strategy = load_json(PACK / "run2_73_text_binding_strategy.json")

    assert strategy["artifact_id"] == "run2_73_text_binding_strategy"
    assert strategy["part"] == "Part F"
    assert strategy["status"] == "run2_73_text_binding_strategy_ready_public_blocked"
    assert strategy["stage_policy"] == EXPECTED_RUN2_F_STAGE_POLICY
    assert strategy["source_scene_plan_expansion"] == "run2_73_scene_plan_expansion.json"
    assert strategy["source_renderer_adapter_contracts"] == "run2_73_renderer_adapter_contracts.json"
    assert strategy["source_slide_story"] == "run2_74_slide_story.json"
    assert strategy["source_content_quality_audit"] == "run2_74_content_quality_audit.json"
    assert EXPECTED_RUN2_E_FORBIDDEN_SCOPE <= set(strategy["artifact_scope"]["does_not_start"])
    assert EXPECTED_RUN2_F_FORBIDDEN_TEXT_PATTERNS <= set(
        strategy["global_forbidden_text_patterns"]
    )

    guard = strategy["execution_guard"]
    assert guard["mode"] == "text_binding_strategy_only"
    assert guard["rendering_subprocesses_allowed"] is False
    assert set(guard["forbidden_invocations"]) == EXPECTED_RUN2_E_FORBIDDEN_SCOPE
    assert set(guard["forbidden_runtime_imports"]) == TEXT_FORBIDDEN_RUNTIME_IMPORTS
    assert set(guard["forbidden_dynamic_import_calls"]) == TEXT_FORBIDDEN_DYNAMIC_IMPORT_CALLS

    script_path = ROOT / "scripts" / "build_ppt_run2_73_text_binding_strategy.py"
    script_tree = ast.parse(script_path.read_text())
    script_imports = set()
    script_calls = set()
    for node in ast.walk(script_tree):
        if isinstance(node, ast.Import):
            script_imports |= {alias.name.split(".")[0] for alias in node.names}
        elif isinstance(node, ast.ImportFrom) and node.module:
            script_imports.add(node.module.split(".")[0])
        elif isinstance(node, ast.Call) and isinstance(node.func, ast.Name):
            script_calls.add(node.func.id)
    assert not TEXT_FORBIDDEN_RUNTIME_IMPORTS & script_imports
    assert not TEXT_FORBIDDEN_DYNAMIC_IMPORT_CALLS & script_calls

    source_paths = {source["path"] for source in strategy["source_inputs"]}
    assert {
        "docs/product/ppt-run2-data-skill-quality/run2_73_scene_plan_expansion.json",
        "docs/product/ppt-run2-data-skill-quality/run2_73_renderer_adapter_contracts.json",
        "docs/product/ppt-run2-data-skill-quality/run2_73_visual_grammar_modules.json",
        "docs/product/ppt-run2-data-skill-quality/run2_74_slide_story.json",
        "docs/product/ppt-run2-data-skill-quality/run2_74_content_quality_audit.json",
    } <= source_paths

    d2_by_role = {record["role"]: record for record in expansion["scene_structures"]}
    e2_by_role = {record["role"]: record for record in adapter["adapter_scene_records"]}
    story_by_role = {record["role"]: record for record in story["slides"]}
    audit_by_role = {record["role"]: record for record in audit["slide_quality_audits"]}

    records = strategy["page_text_binding_records"]
    assert [record["role"] for record in records] == EXPECTED_RUN2_74_SLIDE_STORY_ROLES
    assert [record["slide_index"] for record in records] == [1, 2, 3, 4, 5, 6]
    assert len({record["layout_signature"] for record in records}) == 6

    for record in records:
        role = record["role"]
        d2 = d2_by_role[role]
        e2 = e2_by_role[role]
        story_record = story_by_role[role]
        audit_record = audit_by_role[role]
        known_binding_ids = {
            e2["adapter_scene_id"],
            d2["expansion_id"],
            *e2["renderer_adapter_manifest"]["semantic_component_ids"],
            *e2["renderer_adapter_manifest"]["visual_container_ids"],
            *e2["renderer_adapter_manifest"]["expanded_renderer_binding_ids"],
        }

        assert record["text_binding_id"] == f"text_binding_2_73_{role}"
        assert record["source_expansion_id"] == d2["expansion_id"]
        assert record["source_adapter_scene_id"] == e2["adapter_scene_id"]
        assert record["source_slide_story_id"] == story_record["slide_id"]
        assert record["source_content_audit_id"] == audit_record["audit_id"]
        assert record["source_visual_grammar_module"] == e2["visual_grammar_binding"]["module_id"]
        assert EXPECTED_RUN2_F_REQUIRED_SOCKET_KEYS <= set(record["text_socket_strategy"])

        socket_objects = [
            record["text_socket_strategy"]["headline_socket"],
            record["text_socket_strategy"]["supporting_copy_socket"],
            record["text_socket_strategy"]["viewer_note_socket"],
            *record["text_socket_strategy"]["proof_label_sockets"],
            *record["text_socket_strategy"]["callout_sockets"],
        ]
        assert len(record["text_socket_strategy"]["proof_label_sockets"]) >= 2
        assert len(record["text_socket_strategy"]["callout_sockets"]) >= 2

        for socket in socket_objects:
            assert socket["bound_visual_object_type"] in EXPECTED_RUN2_F_BOUND_VISUAL_OBJECT_TYPES
            assert socket["bound_source_id"] in known_binding_ids
            assert socket["bound_source_artifact"] in {
                "run2_73_scene_plan_expansion",
                "run2_73_renderer_adapter_contracts",
            }
            assert socket["binding_role"] in {
                "headline",
                "proof_label",
                "supporting_copy",
                "callout",
                "viewer_note",
            }
            capacity = socket["capacity"]
            assert type(capacity["max_words"]) is int and capacity["max_words"] > 0
            assert type(capacity["max_lines"]) is int and capacity["max_lines"] > 0
            assert capacity["hierarchy_level"] in {
                "h1",
                "h2",
                "proof_label",
                "supporting_copy",
                "callout",
                "viewer_note",
            }
            assert set(capacity["allowed_font_scale"]) == {"min", "max"}
            assert capacity["allowed_font_scale"]["min"] > 0
            assert capacity["allowed_font_scale"]["max"] >= capacity["allowed_font_scale"]["min"]
            assert capacity["overflow_behavior"] in {
                "truncate_with_route_to_viewer",
                "route_to_speaker_note",
                "route_to_html_viewer_metadata",
            }

        routing = record["text_routing"]
        assert set(routing) == {"canvas_text", "speaker_note_text", "html_viewer_metadata"}
        assert "headline_socket" in routing["canvas_text"]
        assert "supporting_copy_socket" in routing["canvas_text"]
        assert "viewer_note_socket" in routing["speaker_note_text"]
        assert "overflow_payload" in routing["html_viewer_metadata"]

        overflow = record["overflow_policy"]
        assert overflow["max_canvas_words"] <= story_record["text_budget"]["max_public_words"]
        assert overflow["max_proof_labels_on_canvas"] <= story_record["text_budget"]["max_proof_labels"]
        assert overflow["route_excess_to"] == ["speaker_note", "html_viewer_metadata"]
        assert overflow["never_create_empty_text_box"] is True
        assert EXPECTED_RUN2_F_FORBIDDEN_TEXT_PATTERNS <= set(record["forbidden_text_patterns"])

    summary = strategy["traceability_summary"]
    assert summary["page_text_binding_count"] == 6
    assert summary["socket_count"] >= 36
    assert summary["layout_signature_count"] == 6
    assert summary["sources_consumed"] == [
        "run2_73_scene_plan_expansion.json",
        "run2_73_renderer_adapter_contracts.json",
        "run2_73_visual_grammar_modules.json",
        "run2_74_slide_story.json",
        "run2_74_content_quality_audit.json",
    ]
    assert strategy["next_required_action"] == EXPECTED_RUN2_F_NEXT_REQUIRED_ACTION


def test_run2_73_validated_scene_renderer_rerun_consumes_a_f_and_updates_viewer() -> None:
    script = EXPECTED_RUN2_G_SCRIPT.read_text(encoding="utf-8")
    assert_contains(
        script,
        [
            *EXPECTED_RUN2_G_REQUIRED_INPUTS,
            "run2_73_validated_scene_renderer_rerun_result.json",
            "run2_73_validated_scene_renderer_rerun_result.md",
            "build_ppt_run_html_viewer.py",
        ],
    )

    expansion = load_json(PACK / "run2_73_scene_plan_expansion.json")
    adapter = load_json(PACK / "run2_73_renderer_adapter_contracts.json")
    result = load_json(EXPECTED_RUN2_G_RESULT)

    assert result["artifact_id"] == "run2_73_validated_scene_renderer_rerun_result"
    assert result["part"] == "Part G"
    assert result["run_id"] == "2.73"
    assert result["status"] == "run2_73_validated_scene_renderer_rerun_generated_public_blocked"
    assert result["public_ready"] is False
    assert result["public_release_started"] is False
    assert result["quality_claim_boundary"] == "generated_viewer_check_only_no_part_h_quality_verdict"
    assert result["consumed_sources"] == EXPECTED_RUN2_G_REQUIRED_INPUTS

    manifest = result["rerun_manifest"]
    assert manifest["generator"] == "scripts/generate_ppt_run2_73_validated_scene_renderer_arms.mjs"
    assert manifest["consumed_sources"] == EXPECTED_RUN2_G_REQUIRED_INPUTS
    assert manifest["best_internal_arm"] == "run2_73_full_validated_scene_renderer"
    assert manifest["viewer_update"]["latest_run_id"] == "2.73"
    assert manifest["viewer_update"]["viewer_can_reference_new_run"] is True

    outputs = manifest["outputs"]
    html_output = ROOT / outputs["html_viewer"]
    pptx_output = ROOT / outputs["pptx"]
    viewer_output = ROOT / outputs["ppt_run_viewer"]
    assert html_output.exists()
    assert pptx_output.exists()
    assert viewer_output.exists()

    d2_by_role = {record["role"]: record for record in expansion["scene_structures"]}
    e2_by_role = {record["role"]: record for record in adapter["adapter_scene_records"]}
    pages = result["rendered_pages"]
    assert [page["role"] for page in pages] == EXPECTED_RUN2_74_SLIDE_STORY_ROLES
    assert [page["slide_index"] for page in pages] == [1, 2, 3, 4, 5, 6]
    assert len({page["layout_signature"] for page in pages}) == 6

    for page in pages:
        role = page["role"]
        assert page["visual_grammar_module"] == EXPECTED_RUN2_E_PAGE_MODULE_MAP[role]
        assert page["source_text_binding_id"] == f"text_binding_2_73_{role}"
        assert EXPECTED_RUN2_F_REQUIRED_SOCKET_KEYS <= set(page["text_sockets_used"])
        assert page["source_trace_terms_visible_on_canvas"] == []
        assert set(page["forbidden_text_patterns_absent"]) >= EXPECTED_RUN2_F_FORBIDDEN_TEXT_PATTERNS

        known_binding_ids = {
            d2_by_role[role]["expansion_id"],
            e2_by_role[role]["adapter_scene_id"],
            *e2_by_role[role]["renderer_adapter_manifest"]["semantic_component_ids"],
            *e2_by_role[role]["renderer_adapter_manifest"]["visual_container_ids"],
            *e2_by_role[role]["renderer_adapter_manifest"]["expanded_renderer_binding_ids"],
        }
        socket_keys = {binding["socket_key"] for binding in page["text_socket_bindings"]}
        assert EXPECTED_RUN2_F_REQUIRED_SOCKET_KEYS <= socket_keys
        for binding in page["text_socket_bindings"]:
            assert binding["bound_visual_object_type"] in EXPECTED_RUN2_F_BOUND_VISUAL_OBJECT_TYPES
            assert binding["bound_source_id"] in known_binding_ids
            capacity = binding["capacity"]
            assert capacity["max_words"] > 0
            assert capacity["max_lines"] > 0
            assert capacity["hierarchy_level"]
            assert {"min", "max"} <= set(capacity["allowed_font_scale"])
            assert capacity["overflow_behavior"]

        assert page["visual_containers"]
        for container in page["visual_containers"]:
            assert container["empty"] is False
            assert container["bound_text_socket_ids"]
            assert container["visual_object_type"] in EXPECTED_RUN2_F_BOUND_VISUAL_OBJECT_TYPES

    checks = result["render_quality_checks"]
    assert checks["empty_visual_container_count"] == 0
    assert checks["floating_text_without_bound_visual_object_count"] == 0
    assert checks["generic_rectangle_label_count"] == 0
    assert checks["source_trace_terms_visible_on_canvas_count"] == 0
    assert checks["distinct_text_layout_signatures"] == 6
    assert checks["pages_using_expected_visual_grammar"] == 6
    assert checks["pages_using_required_text_sockets"] == 6

    viewer_html = viewer_output.read_text(encoding="utf-8")
    assert_contains(
        viewer_html,
        [
            '"latestRunId": "2.77"',
            "Run 2.73",
            "ppt-run2-73-full-vulca",
            "run2_73_validated_scene_renderer_rerun_result.json",
        ],
    )
    assert "data-run-id=\"2.73\"" in html_output.read_text(encoding="utf-8")

    viewer_data = build_data(ROOT / "outputs" / DEFAULT_THREAD_ID / "presentations", viewer_output)
    run = next(run for run in viewer_data["runs"] if run["id"] == "2.73")
    assert run["fullArm"]["id"] == "run2_73_full_validated_scene_renderer"
    assert len(run["fullArm"]["slides"]) == 6


def test_run2_74_visual_quality_evaluation_closes_viewer_comparison_loop(tmp_path: Path) -> None:
    assert EXPECTED_RUN2_H_SCRIPT.exists(), "missing Part H visual quality evaluation script"
    result_json = tmp_path / "run2_74_visual_quality_evaluation.json"
    result_md = tmp_path / "run2_74_visual_quality_evaluation.md"

    completed = subprocess.run(
        [
            sys.executable,
            str(EXPECTED_RUN2_H_SCRIPT),
            "--result-json",
            str(result_json),
            "--result-md",
            str(result_md),
        ],
        cwd=ROOT,
        check=True,
        capture_output=True,
        text=True,
    )

    audit = load_json(result_json)
    report = result_md.read_text(encoding="utf-8")
    assert "run2_74_visual_quality_evaluation_public_blocked" in completed.stdout
    assert audit["artifact_id"] == "run2_74_visual_quality_evaluation"
    assert audit["part"] == "Part H"
    assert audit["run_id"] == "2.74"
    assert audit["status"] == "run2_74_visual_quality_evaluation_public_blocked"
    assert audit["creates_new_ppt_deck"] is False
    assert audit["starts_renderer_rerun"] is False
    assert audit["public_ready"] is False
    assert audit["quality_claim_boundary"] == "part_h_evaluation_only_no_public_release_no_renderer_rerun"
    assert audit["source_runs"] == {
        "comparison_baseline": "2.72",
        "evaluated_run": "2.73",
    }

    input_chain = audit["input_chain"]
    assert input_chain["run2_72_full_trace_manifest"].endswith(
        "ppt-run2-72-full-vulca/trace_manifest.json"
    )
    assert input_chain["run2_73_full_trace_manifest"].endswith(
        "ppt-run2-73-full-vulca/trace_manifest.json"
    )
    assert input_chain["run2_72_result"].endswith("run2_72_shape_bound_text_rerun_result.json")
    assert input_chain["run2_73_result"].endswith(
        "run2_73_validated_scene_renderer_rerun_result.json"
    )
    assert input_chain["ppt_run_viewer"].endswith("ppt-run-viewer.html")

    viewer = audit["viewer_comparison_closure"]
    assert viewer["viewer_latest_run_id"] == "2.73"
    assert viewer["viewer_can_compare_2_72_and_2_73"] is True
    assert viewer["run2_72_full_preview_count"] == 6
    assert viewer["run2_73_full_preview_count"] == 6
    assert viewer["browser_check_required_for_handoff"] is True

    questions = audit["evaluation_questions"]
    assert set(questions) == EXPECTED_RUN2_H_QUESTIONS
    assert questions["is_2_73_better_than_2_72"]["answer"] == "mixed_not_public_quality_pass"
    assert questions["do_six_pages_have_distinct_visual_grammar"]["answer"] == "yes_trace_and_thumbnail"
    assert questions["does_it_still_feel_like_engineering_report"]["answer"] == "yes_but_different_failure_mode"

    assessment = audit["visual_quality_assessment"]
    assert assessment["data_workflow_entry_gate"] == "pass_internal_only"
    assert assessment["viewer_comparison_gate"] == "pass_internal_only"
    assert assessment["design_quality_gate"] == "blocked"
    assert assessment["public_video_readiness"] == "blocked"
    assert assessment["global_delta_vs_2_72"] == "structural_variety_up_public_polish_down"
    assert assessment["top_blocker"] == "thin_abstract_renderer_placeholders_do_not_read_as_product_presentation"
    assert assessment["next_layer_to_fix"] == "renderer"

    assert len(audit["role_assessments"]) == 6
    assert [record["role"] for record in audit["role_assessments"]] == EXPECTED_RUN2_74_SLIDE_STORY_ROLES
    assert [record["slide_index"] for record in audit["role_assessments"]] == [1, 2, 3, 4, 5, 6]
    assert len({record["visual_grammar_module"] for record in audit["role_assessments"]}) == 5
    needs_repair = [record for record in audit["role_assessments"] if record["repair_required"]]
    assert len(needs_repair) >= 4
    for record in audit["role_assessments"]:
        assert record["visual_grammar_module"] == EXPECTED_RUN2_E_PAGE_MODULE_MAP[record["role"]]
        assert record["root_cause_layer"] in EXPECTED_RUN2_H_ROOT_CAUSE_LAYERS
        assert record["text_visual_fusion"] in {"improved_but_weak", "partial", "weak"}
        assert record["report_like_risk"] in {"low", "medium", "high"}
        assert record["repair_instruction"]

    assert audit["root_cause_summary"]["primary_layer"] == "renderer"
    assert audit["root_cause_summary"]["not_primary_layer"] == "data_absence"
    assert audit["next_required_action"] == "part_i_renderer_repair_from_visual_quality_evaluation"
    assert_contains(
        report,
        [
            "Run 2.74 Visual Quality Evaluation",
            "2.73 vs 2.72",
            "mixed",
            "public polish down",
            "Part I",
        ],
    )


def test_run2_74_records_visual_quality_evaluation_result() -> None:
    audit = load_json(EXPECTED_RUN2_H_RESULT)
    report = (
        PACK / "results" / "run2_74_visual_quality_evaluation.md"
    ).read_text(encoding="utf-8")

    assert audit["status"] == "run2_74_visual_quality_evaluation_public_blocked"
    assert audit["public_ready"] is False
    assert audit["visual_quality_assessment"]["design_quality_gate"] == "blocked"
    assert audit["visual_quality_assessment"]["global_delta_vs_2_72"] == (
        "structural_variety_up_public_polish_down"
    )
    assert audit["viewer_comparison_closure"]["viewer_latest_run_id"] == "2.73"
    assert len(audit["role_assessments"]) == 6
    assert "thin abstract renderer placeholders" in report


def test_run2_75_renderer_repair_rerun_consumes_h_and_updates_viewer() -> None:
    assert EXPECTED_RUN2_I_SCRIPT.exists(), "missing Part I renderer repair script"
    script = EXPECTED_RUN2_I_SCRIPT.read_text(encoding="utf-8")
    assert_contains(
        script,
        [
            *EXPECTED_RUN2_I_REQUIRED_INPUTS,
            "run2_75_renderer_repair_rerun_result.json",
            "run2_75_renderer_repair_rerun_result.md",
            "build_ppt_run_html_viewer.py",
        ],
    )

    h_audit = load_json(EXPECTED_RUN2_H_RESULT)
    result = load_json(EXPECTED_RUN2_I_RESULT)

    assert result["artifact_id"] == "run2_75_renderer_repair_rerun_result"
    assert result["part"] == "Part I"
    assert result["run_id"] == "2.75"
    assert result["status"] == "run2_75_renderer_repair_rerun_generated_public_blocked"
    assert result["public_ready"] is False
    assert result["public_release_started"] is False
    assert result["quality_claim_boundary"] == (
        "renderer_repair_generated_viewer_check_only_no_part_j_quality_verdict"
    )
    assert result["consumed_sources"] == EXPECTED_RUN2_I_REQUIRED_INPUTS
    assert result["source_h_evaluation"]["status"] == h_audit["status"]
    assert result["source_h_evaluation"]["top_blocker"] == (
        "thin_abstract_renderer_placeholders_do_not_read_as_product_presentation"
    )

    manifest = result["renderer_repair_manifest"]
    assert manifest["generator"] == "scripts/generate_ppt_run2_75_renderer_repair_arms.mjs"
    assert manifest["consumed_sources"] == EXPECTED_RUN2_I_REQUIRED_INPUTS
    assert manifest["best_internal_arm"] == "run2_75_full_renderer_repair"
    assert manifest["viewer_update"]["latest_run_id"] == "2.75"
    assert manifest["viewer_update"]["viewer_can_reference_new_run"] is True

    outputs = manifest["outputs"]
    html_output = ROOT / outputs["html_viewer"]
    pptx_output = ROOT / outputs["pptx"]
    viewer_output = ROOT / outputs["ppt_run_viewer"]
    assert html_output.exists()
    assert pptx_output.exists()
    assert viewer_output.exists()

    pages = result["rendered_pages"]
    assert [page["role"] for page in pages] == EXPECTED_RUN2_74_SLIDE_STORY_ROLES
    assert [page["slide_index"] for page in pages] == [1, 2, 3, 4, 5, 6]
    assert len({page["visual_grammar_module"] for page in pages}) == 5
    assert len({page["visual_density_profile"] for page in pages}) == 6

    h_by_role = {record["role"]: record for record in h_audit["role_assessments"]}
    for page in pages:
        role = page["role"]
        assert page["visual_grammar_module"] == EXPECTED_RUN2_E_PAGE_MODULE_MAP[role]
        assert page["source_text_binding_id"] == f"text_binding_2_73_{role}"
        assert EXPECTED_RUN2_F_REQUIRED_SOCKET_KEYS <= set(page["text_sockets_used"])
        assert page["source_trace_terms_visible_on_canvas"] == []
        assert set(page["forbidden_text_patterns_absent"]) >= EXPECTED_RUN2_F_FORBIDDEN_TEXT_PATTERNS
        assert page["h_repair_source"]["repair_instruction"] == h_by_role[role]["repair_instruction"]
        assert page["h_repair_source"]["root_cause_layer"] == h_by_role[role]["root_cause_layer"]
        assert EXPECTED_RUN2_I_REPAIR_FLAGS <= set(page["renderer_repair_directives_applied"])
        assert page["product_surface_detail_count"] >= 5
        assert page["connector_or_edge_binding_count"] >= 3

    checks = result["renderer_repair_checks"]
    assert checks["empty_visual_container_count"] == 0
    assert checks["floating_text_without_bound_visual_object_count"] == 0
    assert checks["generic_rectangle_label_count"] == 0
    assert checks["source_trace_terms_visible_on_canvas_count"] == 0
    assert checks["pages_using_expected_visual_grammar"] == 6
    assert checks["pages_using_required_text_sockets"] == 6
    assert checks["pages_with_h_repair_directive_consumed"] == 6
    assert checks["pages_with_concrete_product_surface"] == 6
    assert checks["pages_with_stronger_connector_or_edge_binding"] == 6
    assert checks["distinct_visual_density_profiles"] == 6
    assert checks["public_quality_verdict_started"] is False

    viewer_html = viewer_output.read_text(encoding="utf-8")
    assert_contains(
        viewer_html,
        [
            '"latestRunId": "2.77"',
            "Run 2.75",
            "ppt-run2-75-full-vulca",
            "run2_75_renderer_repair_rerun_result.json",
        ],
    )
    assert "data-run-id=\"2.75\"" in html_output.read_text(encoding="utf-8")

    viewer_data = build_data(ROOT / "outputs" / DEFAULT_THREAD_ID / "presentations", viewer_output)
    run = next(run for run in viewer_data["runs"] if run["id"] == "2.75")
    assert run["fullArm"]["id"] == "run2_75_full_renderer_repair"
    assert len(run["fullArm"]["slides"]) == 6
    assert result["next_required_action"] == "part_j_visual_quality_evaluation_for_run2_75"


def test_run2_76_visual_quality_evaluation_compares_2_75_against_2_73_with_gemini(
    tmp_path: Path,
) -> None:
    assert EXPECTED_RUN2_J_SCRIPT.exists(), "missing Part J visual quality evaluation script"
    script = EXPECTED_RUN2_J_SCRIPT.read_text(encoding="utf-8")
    assert_contains(
        script,
        [
            "run2_73_validated_scene_renderer_rerun_result.json",
            "run2_75_renderer_repair_rerun_result.json",
            "run2_74_visual_quality_evaluation.json",
            "gemini-3.5-flash",
            "ppt-run2-73-full-vulca/preview/contact-sheet.png",
            "ppt-run2-75-full-vulca/preview/contact-sheet.png",
        ],
    )
    result_json = tmp_path / "run2_76_visual_quality_evaluation.json"
    result_md = tmp_path / "run2_76_visual_quality_evaluation.md"

    completed = subprocess.run(
        [
            sys.executable,
            str(EXPECTED_RUN2_J_SCRIPT),
            "--result-json",
            str(result_json),
            "--result-md",
            str(result_md),
        ],
        cwd=ROOT,
        check=True,
        capture_output=True,
        text=True,
    )

    audit = load_json(result_json)
    report = result_md.read_text(encoding="utf-8")
    assert "run2_76_visual_quality_evaluation_public_blocked" in completed.stdout
    assert audit["artifact_id"] == "run2_76_visual_quality_evaluation"
    assert audit["part"] == "Part J"
    assert audit["run_id"] == "2.76"
    assert audit["status"] == "run2_76_visual_quality_evaluation_public_blocked"
    assert audit["creates_new_ppt_deck"] is False
    assert audit["starts_renderer_rerun"] is False
    assert audit["public_ready"] is False
    assert audit["quality_claim_boundary"] == "part_j_evaluation_only_no_public_release_no_renderer_rerun"
    assert audit["source_runs"] == {
        "comparison_baseline": "2.73",
        "evaluated_run": "2.75",
    }

    input_chain = audit["input_chain"]
    assert input_chain["run2_73_result"].endswith("run2_73_validated_scene_renderer_rerun_result.json")
    assert input_chain["run2_75_result"].endswith("run2_75_renderer_repair_rerun_result.json")
    assert input_chain["run2_74_h_evaluation"].endswith("run2_74_visual_quality_evaluation.json")
    assert input_chain["run2_73_full_contact_sheet"].endswith("ppt-run2-73-full-vulca/preview/contact-sheet.png")
    assert input_chain["run2_75_full_contact_sheet"].endswith("ppt-run2-75-full-vulca/preview/contact-sheet.png")

    gemini = audit["gemini_agent_review_summary"]
    assert gemini["tool"] == "mcp__gemini_agent.gemini_artifact_review"
    assert gemini["model"] == "gemini-3.5-flash"
    assert gemini["review_count"] == 2
    assert gemini["used_for_verdict"] is True
    assert "engineering blueprint" in " ".join(gemini["shared_risks"])
    assert "page_differentiation_regression_for_01_02_04_05" in gemini["run2_75_findings"]

    questions = audit["evaluation_questions"]
    assert set(questions) == EXPECTED_RUN2_J_QUESTIONS
    assert questions["is_2_75_better_than_2_73"]["answer"] == (
        "mixed_product_surface_up_page_differentiation_down_public_blocked"
    )
    assert questions["does_2_75_have_stronger_product_feel"]["answer"] == "yes_but_still_wireframe"
    assert questions["are_page_differences_stronger_or_weaker"]["answer"] == (
        "weaker_for_core_product_surface_pages"
    )
    assert questions["does_2_75_reach_public_video_presentation_direction"]["answer"] == (
        "no_internal_blueprint_risk_remains"
    )

    assessment = audit["visual_quality_assessment"]
    assert assessment["design_quality_gate"] == "blocked"
    assert assessment["public_video_readiness"] == "blocked"
    assert assessment["global_delta_vs_2_73"] == (
        "product_surface_up_page_differentiation_down_public_readiness_still_blocked"
    )
    assert assessment["top_blocker"] == (
        "wireframe_blueprint_aesthetic_and_repeated_product_surfaces_still_read_as_internal_engineering_diagrams"
    )
    assert assessment["next_layer_to_fix"] == "visual_grammar_and_renderer"

    assert len(audit["role_assessments"]) == 6
    assert [record["role"] for record in audit["role_assessments"]] == EXPECTED_RUN2_74_SLIDE_STORY_ROLES
    repairs = [record for record in audit["role_assessments"] if record["repair_required"]]
    assert len(repairs) >= 5
    for record in audit["role_assessments"]:
        assert record["visual_grammar_module"] == EXPECTED_RUN2_E_PAGE_MODULE_MAP[record["role"]]
        assert record["root_cause_layer"] in EXPECTED_RUN2_J_ROOT_CAUSE_LAYERS
        assert record["product_feel_delta"] in {"improved_but_wireframe", "partial", "not_improved"}
        assert record["page_differentiation_delta"] in {"improved", "same", "weaker"}
        assert record["text_binding_delta"] in {"slightly_stronger_but_small", "partial", "weak"}
        assert record["engineering_report_risk"] in {"low", "medium", "high"}
        assert record["public_video_direction"] in {"no", "partial", "yes"}
        assert record["repair_instruction"]

    assert audit["root_cause_summary"]["primary_layer"] == "visual_grammar_and_renderer"
    assert audit["root_cause_summary"]["not_primary_layer"] == "data_absence"
    assert audit["next_required_action"] == "part_k_visual_grammar_and_renderer_repair_from_j_evaluation"
    assert_contains(
        report,
        [
            "Run 2.76 Visual Quality Evaluation",
            "2.75 vs 2.73",
            "product surface up",
            "page differentiation down",
            "public blocked",
            "Gemini",
        ],
    )


def test_run2_76_records_visual_quality_evaluation_result() -> None:
    audit = load_json(EXPECTED_RUN2_J_RESULT)
    report = (
        PACK / "results" / "run2_76_visual_quality_evaluation.md"
    ).read_text(encoding="utf-8")

    assert audit["status"] == "run2_76_visual_quality_evaluation_public_blocked"
    assert audit["public_ready"] is False
    assert audit["visual_quality_assessment"]["design_quality_gate"] == "blocked"
    assert audit["visual_quality_assessment"]["global_delta_vs_2_73"] == (
        "product_surface_up_page_differentiation_down_public_readiness_still_blocked"
    )
    assert audit["viewer_comparison_closure"]["viewer_latest_run_id"] == "2.75"
    assert len(audit["role_assessments"]) == 6
    assert "wireframe blueprint aesthetic" in report


def test_run2_76_visual_grammar_renderer_repair_plan_consumes_j_without_rerun() -> None:
    plan = load_json(EXPECTED_RUN2_K1_REPAIR_PLAN)
    j_evaluation = load_json(EXPECTED_RUN2_J_RESULT)
    run275 = load_json(EXPECTED_RUN2_I_RESULT)

    assert plan["artifact_id"] == "run2_76_visual_grammar_renderer_repair_plan"
    assert plan["part"] == "Part K1"
    assert plan["status"] == "run2_76_visual_grammar_renderer_repair_plan_ready_public_blocked"
    assert plan["creates_new_ppt_deck"] is False
    assert plan["starts_renderer_rerun"] is False
    assert plan["updates_html_viewer"] is False
    assert plan["public_release_started"] is False
    assert plan["public_ready"] is False
    assert plan["quality_claim_boundary"] == "part_k1_repair_contract_only_no_renderer_rerun_no_public_release"
    assert plan["consumed_sources"] == EXPECTED_RUN2_K1_REQUIRED_INPUTS

    source_j = plan["source_j_evaluation"]
    assert source_j["source_path"].endswith("run2_76_visual_quality_evaluation.json")
    assert source_j["status"] == j_evaluation["status"]
    assert source_j["top_blocker"] == j_evaluation["visual_quality_assessment"]["top_blocker"]
    assert source_j["primary_layer_to_fix"] == j_evaluation["root_cause_summary"]["primary_layer"]
    assert source_j["secondary_layers"] == j_evaluation["root_cause_summary"]["secondary_layers"]
    assert source_j["next_required_action"] == "part_k_visual_grammar_and_renderer_repair_from_j_evaluation"

    source_run275 = plan["source_run2_75_renderer_repair"]
    assert source_run275["source_path"].endswith("run2_75_renderer_repair_rerun_result.json")
    assert source_run275["status"] == run275["status"]
    assert source_run275["public_ready"] is False

    strategy = plan["global_repair_strategy"]
    assert "wireframe blueprint aesthetic" in strategy["from_failure_mode"]
    assert "public-video presentation scene" in strategy["target_shift"]
    assert strategy["primary_layers_to_fix"] == ["visual_grammar", "renderer"]
    assert strategy["secondary_layers_to_fix"] == ["text_binding"]
    assert "data" in strategy["not_data_layer_reason"]
    assert EXPECTED_RUN2_K1_RENDERER_CAPABILITIES <= set(strategy["renderer_capabilities_required"])
    assert EXPECTED_RUN2_K1_FORBIDDEN_PATTERNS <= set(strategy["forbidden_renderer_fallbacks"])

    page_plans = plan["page_repair_plans"]
    assert [record["role"] for record in page_plans] == EXPECTED_RUN2_74_SLIDE_STORY_ROLES
    assert [record["slide_index"] for record in page_plans] == [1, 2, 3, 4, 5, 6]
    assert len({record["target_scene_direction"] for record in page_plans}) == 6
    j_by_role = {record["role"]: record for record in j_evaluation["role_assessments"]}

    for record in page_plans:
        role = record["role"]
        assert record["visual_grammar_module"] == EXPECTED_RUN2_E_PAGE_MODULE_MAP[role]
        assert EXPECTED_RUN2_K1_PAGE_REPAIR_FIELDS <= set(record)
        assert record["source_j_assessment"]["root_cause_layer"] == j_by_role[role]["root_cause_layer"]
        assert record["source_j_assessment"]["repair_instruction"] == j_by_role[role]["repair_instruction"]
        assert record["current_failure"]
        assert record["target_scene_direction"]
        assert record["visual_grammar_change"]
        assert record["renderer_change"]
        assert record["text_binding_adjustment"]
        assert record["must_preserve_from_2_75"]
        assert record["must_remove_from_2_75"]
        assert EXPECTED_RUN2_K1_FORBIDDEN_PATTERNS & set(record["must_remove_from_2_75"])
        checks = record["acceptance_checks"]
        assert checks["page_differentiation_check"]
        assert checks["wireframe_reduction_check"]
        assert checks["renderer_capability_check"]
        assert checks["no_renderer_rerun_in_k1"] is True

    by_role = {record["role"]: record for record in page_plans}
    assert by_role["cover"]["target_scene_direction"] != by_role["climax"]["target_scene_direction"]
    assert "finished product hero" in by_role["cover"]["target_scene_direction"]
    assert "completion reveal" in by_role["climax"]["target_scene_direction"]
    assert "proof workspace" in " ".join(by_role["setup"]["must_remove_from_2_75"])
    assert "evidence inspection scene" in by_role["proof"]["target_scene_direction"]
    assert "non-grid evidence arrangement" in by_role["proof"]["renderer_change"]
    assert "decision / release gate" in by_role["close"]["target_scene_direction"]
    assert plan["next_required_action"] == "part_k2_renderer_rerun_from_visual_grammar_renderer_repair_plan"


def test_run2_77_visual_grammar_renderer_repair_rerun_consumes_k1_and_updates_viewer() -> None:
    assert EXPECTED_RUN2_K2_SCRIPT.exists(), "missing Part K2 renderer rerun script"
    script = EXPECTED_RUN2_K2_SCRIPT.read_text(encoding="utf-8")
    assert_contains(
        script,
        [
            *EXPECTED_RUN2_K2_REQUIRED_INPUTS,
            "run2_77_visual_grammar_renderer_repair_rerun_result.json",
            "run2_77_visual_grammar_renderer_repair_rerun_result.md",
            "build_ppt_run_html_viewer.py",
        ],
    )

    k1_plan = load_json(EXPECTED_RUN2_K1_REPAIR_PLAN)
    result = load_json(EXPECTED_RUN2_K2_RESULT)
    plan_by_role = {record["role"]: record for record in k1_plan["page_repair_plans"]}

    assert result["artifact_id"] == "run2_77_visual_grammar_renderer_repair_rerun_result"
    assert result["part"] == "Part K2"
    assert result["run_id"] == "2.77"
    assert result["status"] == "run2_77_visual_grammar_renderer_repair_rerun_generated_public_blocked"
    assert result["public_ready"] is False
    assert result["public_release_started"] is False
    assert result["quality_claim_boundary"] == (
        "visual_grammar_renderer_repair_generated_viewer_check_only_no_part_l_quality_verdict"
    )
    assert result["consumed_sources"] == EXPECTED_RUN2_K2_REQUIRED_INPUTS
    assert result["source_k1_repair_plan"]["status"] == k1_plan["status"]
    assert result["source_k1_repair_plan"]["top_blocker"] == k1_plan["source_j_evaluation"]["top_blocker"]
    assert result["source_k1_repair_plan"]["next_required_action"] == (
        "part_k2_renderer_rerun_from_visual_grammar_renderer_repair_plan"
    )

    manifest = result["visual_grammar_renderer_repair_manifest"]
    assert manifest["generator"] == "scripts/generate_ppt_run2_77_visual_grammar_renderer_repair_arms.mjs"
    assert manifest["consumed_sources"] == EXPECTED_RUN2_K2_REQUIRED_INPUTS
    assert manifest["best_internal_arm"] == "run2_77_full_visual_grammar_renderer_repair"
    assert manifest["viewer_update"]["latest_run_id"] == "2.77"
    assert manifest["viewer_update"]["viewer_can_reference_new_run"] is True

    outputs = manifest["outputs"]
    html_output = ROOT / outputs["html_viewer"]
    pptx_output = ROOT / outputs["pptx"]
    viewer_output = ROOT / outputs["ppt_run_viewer"]
    assert html_output.exists()
    assert pptx_output.exists()
    assert viewer_output.exists()

    pages = result["rendered_pages"]
    assert [page["role"] for page in pages] == EXPECTED_RUN2_74_SLIDE_STORY_ROLES
    assert [page["slide_index"] for page in pages] == [1, 2, 3, 4, 5, 6]
    assert len({page["target_scene_direction"] for page in pages}) == 6
    all_capabilities = set()
    for page in pages:
        role = page["role"]
        plan = plan_by_role[role]
        assert page["visual_grammar_module"] == EXPECTED_RUN2_E_PAGE_MODULE_MAP[role]
        assert page["source_k1_repair_plan"]["target_scene_direction"] == plan["target_scene_direction"]
        assert page["target_scene_direction"] == plan["target_scene_direction"]
        assert set(page["renderer_repair_directives_applied"]) >= EXPECTED_RUN2_K2_REPAIR_FLAGS
        assert set(page["renderer_capabilities_applied"]) <= EXPECTED_RUN2_K1_RENDERER_CAPABILITIES
        assert page["renderer_capabilities_applied"]
        all_capabilities.update(page["renderer_capabilities_applied"])
        assert set(page["forbidden_renderer_fallbacks_absent"]) >= EXPECTED_RUN2_K1_FORBIDDEN_PATTERNS
        assert page["label_count"] <= 5
        assert page["source_trace_terms_visible_on_canvas"] == []
        assert page["visual_containers"]
        assert page["k1_acceptance_checks"]["page_differentiation_check"]
        assert page["k1_acceptance_checks"]["wireframe_reduction_check"]
        assert page["k1_acceptance_checks"]["no_renderer_rerun_in_k1"] is True
    assert EXPECTED_RUN2_K1_RENDERER_CAPABILITIES <= all_capabilities

    by_role = {page["role"]: page for page in pages}
    assert by_role["cover"]["target_scene_direction"] != by_role["climax"]["target_scene_direction"]
    assert "finished product hero" in by_role["cover"]["target_scene_direction"]
    assert "completion reveal" in by_role["climax"]["target_scene_direction"]
    assert "evidence inspection scene" in by_role["proof"]["target_scene_direction"]
    assert "decision / release gate" in by_role["close"]["target_scene_direction"]

    checks = result["visual_grammar_renderer_repair_checks"]
    assert checks["pages_with_k1_repair_plan_consumed"] == 6
    assert checks["pages_using_expected_visual_grammar"] == 6
    assert checks["distinct_target_scene_directions"] == 6
    assert checks["pages_with_forbidden_fallbacks_absent"] == 6
    assert checks["pages_with_reduced_label_count"] == 6
    assert checks["required_renderer_capabilities_covered"] == sorted(EXPECTED_RUN2_K1_RENDERER_CAPABILITIES)
    assert checks["public_quality_verdict_started"] is False

    viewer_html = viewer_output.read_text(encoding="utf-8")
    assert_contains(
        viewer_html,
        [
            '"latestRunId": "2.77"',
            "Run 2.77",
            "ppt-run2-77-full-vulca",
            "run2_77_visual_grammar_renderer_repair_rerun_result.json",
        ],
    )
    assert "data-run-id=\"2.77\"" in html_output.read_text(encoding="utf-8")

    viewer_data = build_data(ROOT / "outputs" / DEFAULT_THREAD_ID / "presentations", viewer_output)
    run = next(run for run in viewer_data["runs"] if run["id"] == "2.77")
    assert run["fullArm"]["id"] == "run2_77_full_visual_grammar_renderer_repair"
    assert len(run["fullArm"]["slides"]) == 6
    assert result["next_required_action"] == "part_l_visual_quality_evaluation_for_run2_77"


def test_run2_7_has_serializable_design_memory() -> None:
    usecase = load_json(PACK / "run2_7_commercial_usecase.json")
    source_records = load_json(PACK / "run2_7_multimodal_source_records.json")
    memory = load_json(PACK / "run2_7_design_memory.json")
    source_record_ids = {record["id"] for record in source_records["records"]}

    assert memory["status"] == "run2_7_design_memory_public_blocked"
    assert memory["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert memory["memory_type"] == "deterministic_serializable_rules"
    assert EXPECTED_RUN2_7_MEMORY_IDS <= {record["id"] for record in memory["memories"]}

    for record in memory["memories"]:
        assert EXPECTED_RUN2_7_MEMORY_FIELDS <= set(record), record["id"]
        assert set(record["source_record_ids"]) <= source_record_ids
        assert usecase["id"] in record["applicable_usecases"]
        assert set(record["applicable_slide_roles"]) <= EXPECTED_RHYTHM_ROLES
        assert record["typography_rules"]
        assert record["spacing_rules"]
        assert record["composition_rules"]
        assert record["rhythm_rules"]
        assert_contains(" ".join(record["native_ppt_generation_requirements"]), ["native", "editable", "trace"])
        assert_mentions_any(
            " ".join(record["forbidden_patterns"]),
            {"report", "dashboard", "equal card", "source brand", "full-slide raster"},
        )
        assert_contains(" ".join(record["qa_probes"]), ["contact sheet"])
        assert_contains(record["release_boundary"], ["public_blocked"])

    climax_memory = next(
        record for record in memory["memories"] if record["id"] == "memory_composition_single_object_climax"
    )
    assert "climax" in climax_memory["applicable_slide_roles"]
    assert_contains(" ".join(climax_memory["composition_rules"]), ["40-55%", "one native proof object"])


def test_run2_7_workflow_and_trace_contract_include_memory_selection() -> None:
    usecase = load_json(PACK / "run2_7_commercial_usecase.json")
    source_records = load_json(PACK / "run2_7_multimodal_source_records.json")
    memory = load_json(PACK / "run2_7_design_memory.json")
    policy = load_json(PACK / "run2_7_workflow_policy.json")
    workflow = load_json(PACK / "skill_workflow.json")
    trace_contract = load_json(PACK / "results" / "trace_manifest_contract.json")

    source_record_ids = {record["id"] for record in source_records["records"]}
    memory_ids = {record["id"] for record in memory["memories"]}

    assert policy["status"] == "run2_7_workflow_policy_public_blocked"
    assert policy["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert policy["commercial_usecase_id"] == usecase["id"]
    assert set(policy["selection_chain"]) == {
        "commercial_usecase",
        "source_record_ids",
        "typography_memory_id",
        "spacing_memory_id",
        "composition_memory_id",
        "rhythm_memory_id",
        "brand_sanitization_memory_id",
        "visual_repair_policy_ids",
        "native_ppt_generation",
        "qa_gate",
    }

    for mapping in policy["slide_role_memory_map"]:
        assert mapping["commercial_usecase_id"] == usecase["id"]
        assert set(mapping["source_record_ids"]) <= source_record_ids
        assert set(mapping["design_memory_ids"]) <= memory_ids
        assert EXPECTED_RUN2_7_TRACE_FIELDS <= set(mapping["trace_fields"])
        assert_contains(" ".join(mapping["workflow_gates"]), ["public_blocked", "native", "source-brand"])

    workflow_stage_ids = [stage["id"] for stage in workflow["stages"]]
    assert "select_run2_7_design_memory" in workflow_stage_ids
    assert workflow_stage_ids.index("select_run2_7_design_memory") < workflow_stage_ids.index(
        "select_visual_production_modules"
    )
    workflow_text = json.dumps(workflow)
    assert_contains(
        workflow_text,
        [
            "run2_7_commercial_usecase.json",
            "run2_7_multimodal_source_records.json",
            "run2_7_design_memory.json",
            "run2_7_workflow_policy.json",
        ],
    )
    assert EXPECTED_RUN2_7_TRACE_FIELDS <= set(trace_contract["per_slide_required_fields"])


def test_run2_8_has_tutorial_video_decomposition_units() -> None:
    decomposition_path = PACK / "run2_8_tutorial_decomposition.json"
    assert decomposition_path.exists(), "missing Run 2.8 tutorial decomposition data"
    decomposition = load_json(decomposition_path)
    source_records = load_json(PACK / "run2_7_multimodal_source_records.json")
    sources = load_json(PACK / "sources.json")
    source_record_ids = {record["id"] for record in source_records["records"]}
    source_ids = {source["id"] for source in sources["sources"]}

    assert decomposition["status"] == "run2_8_tutorial_decomposition_public_blocked"
    assert decomposition["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert decomposition["storage_policy"]["raw_media"] == "forbidden"
    assert EXPECTED_RUN2_8_DECOMPOSITION_IDS <= {unit["id"] for unit in decomposition["units"]}

    required = {
        "id",
        "source_record_ids",
        "source_ids",
        "modality_mix",
        "tutorial_anchor",
        "observed_design_move",
        "derived_rule",
        "code_generation_binding",
        "native_ppt_obligation",
        "layout_budget",
        "failure_probe",
        "anti_copy_boundary",
        "qa_probe",
        "release_boundary",
    }
    derived_modalities = {"video", "audio", "transcript", "image_reference", "interaction"}
    for unit in decomposition["units"]:
        assert required <= set(unit), unit["id"]
        assert unit["source_record_ids"]
        assert unit["source_ids"]
        assert unit["modality_mix"]
        assert set(unit["source_record_ids"]) <= source_record_ids
        assert set(unit["source_ids"]) <= source_ids
        assert set(unit["modality_mix"]) & derived_modalities
        assert_contains(json.dumps(unit["code_generation_binding"]), ["native"])
        assert isinstance(unit["layout_budget"], dict) and unit["layout_budget"]
        for field_value in iter_string_values(unit):
            lowered = field_value.lower()
            for marker in RUN2_8_FORBIDDEN_MEDIA_MARKERS:
                assert marker not in lowered, f"{unit['id']} contains copied media marker {marker!r}"


def test_run2_8_has_executable_design_memory_bindings() -> None:
    memory_path = PACK / "run2_8_executable_design_memory.json"
    decomposition_path = PACK / "run2_8_tutorial_decomposition.json"
    assert memory_path.exists(), "missing Run 2.8 executable design memory"
    assert decomposition_path.exists(), "missing Run 2.8 tutorial decomposition data"
    memory = load_json(memory_path)
    decomposition = load_json(decomposition_path)
    decomposition_unit_ids = {unit["id"] for unit in decomposition["units"]}

    assert memory["status"] == "run2_8_executable_design_memory_public_blocked"
    assert memory["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert memory["memory_type"] == "executable_schema_bindings"
    assert EXPECTED_RUN2_8_MEMORY_BINDING_IDS <= {binding["id"] for binding in memory["bindings"]}

    required = {
        "id",
        "decomposition_unit_ids",
        "applies_to_slide_roles",
        "design_token",
        "code_binding",
        "native_ppt_constraints",
        "typography_constraints",
        "spacing_constraints",
        "composition_constraints",
        "negative_control_failure",
        "qa_probe",
        "release_boundary",
    }
    for binding in memory["bindings"]:
        assert required <= set(binding), binding["id"]
        assert binding["decomposition_unit_ids"]
        assert binding["applies_to_slide_roles"]
        assert set(binding["decomposition_unit_ids"]) <= decomposition_unit_ids
        assert set(binding["applies_to_slide_roles"]) <= EXPECTED_RHYTHM_ROLES
        assert isinstance(binding["code_binding"], dict), binding["id"]
        assert {"function_name", "params", "layout_budget"} <= set(binding["code_binding"]), binding["id"]
        assert isinstance(binding["code_binding"]["params"], dict) and binding["code_binding"]["params"]
        assert isinstance(binding["code_binding"]["layout_budget"], dict) and binding["code_binding"]["layout_budget"]
        code_binding_text = json.dumps(binding["code_binding"])
        assert any(term in code_binding_text for term in RUN2_8_CODE_BINDING_TERMS), binding["id"]


def test_run2_8_workflow_gate_matrix_connects_schema_to_trace() -> None:
    matrix_path = PACK / "run2_8_workflow_gate_matrix.json"
    decomposition_path = PACK / "run2_8_tutorial_decomposition.json"
    memory_path = PACK / "run2_8_executable_design_memory.json"
    assert matrix_path.exists(), "missing Run 2.8 workflow gate matrix"
    assert decomposition_path.exists(), "missing Run 2.8 tutorial decomposition data"
    assert memory_path.exists(), "missing Run 2.8 executable design memory"
    matrix = load_json(matrix_path)
    decomposition = load_json(decomposition_path)
    memory = load_json(memory_path)
    workflow = load_json(PACK / "skill_workflow.json")
    trace_contract = load_json(PACK / "results" / "trace_manifest_contract.json")
    decomposition_unit_ids = {unit["id"] for unit in decomposition["units"]}
    memory_binding_ids = {binding["id"] for binding in memory["bindings"]}
    code_binding_ids = {binding["code_binding"]["function_name"] for binding in memory["bindings"]}

    assert matrix["status"] == "run2_8_workflow_gate_matrix_public_blocked"
    assert matrix["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert set(matrix["selection_chain"]) == {
        "commercial_usecase",
        "run2_8_decomposition_units",
        "run2_8_executable_memory_bindings",
        "run2_8_gate_matrix",
        "native_ppt_code_generation",
        "layout_quality_gate",
        "delivery_gate",
        "visual_qa_gate",
    }
    assert EXPECTED_RUN2_8_TRACE_FIELDS <= set(trace_contract["per_slide_required_fields"])

    workflow_stage_ids = [stage["id"] for stage in workflow["stages"]]
    run2_8_stages = [
        "decompose_run2_8_tutorial_video_units",
        "select_run2_8_executable_design_memory",
        "apply_run2_8_workflow_gate_matrix",
    ]
    for stage_id in run2_8_stages:
        assert stage_id in workflow_stage_ids
        assert workflow_stage_ids.index(stage_id) < workflow_stage_ids.index("generate_code_first_ppt")
    assert workflow_stage_ids.index(run2_8_stages[0]) < workflow_stage_ids.index(run2_8_stages[1])
    assert workflow_stage_ids.index(run2_8_stages[1]) < workflow_stage_ids.index(run2_8_stages[2])
    workflow_text = json.dumps(workflow)
    assert_contains(
        workflow_text,
        [
            "run2_8_tutorial_decomposition.json",
            "run2_8_executable_design_memory.json",
            "run2_8_workflow_gate_matrix.json",
        ],
    )

    required = {
        "id",
        "slide_role",
        "decomposition_unit_ids",
        "memory_binding_ids",
        "required_code_bindings",
        "layout_budget",
        "pass_fail_checks",
        "trace_fields",
        "public_release_gate",
    }
    covered_trace_fields = set()
    covered_decomposition_unit_ids = set()
    covered_memory_binding_ids = set()
    covered_code_binding_ids = set()
    for gate in matrix["gates"]:
        assert required <= set(gate), gate.get("id", gate["slide_role"])
        assert gate["slide_role"] in EXPECTED_RHYTHM_ROLES
        assert gate["decomposition_unit_ids"]
        assert gate["memory_binding_ids"]
        assert gate["required_code_bindings"]
        assert set(gate["decomposition_unit_ids"]) <= decomposition_unit_ids
        assert set(gate["memory_binding_ids"]) <= memory_binding_ids
        assert set(gate["required_code_bindings"]) <= code_binding_ids
        assert set(gate["trace_fields"]) <= set(trace_contract["per_slide_required_fields"])
        covered_decomposition_unit_ids.update(gate["decomposition_unit_ids"])
        covered_memory_binding_ids.update(gate["memory_binding_ids"])
        covered_code_binding_ids.update(gate["required_code_bindings"])
        covered_trace_fields.update(gate["trace_fields"])
    assert EXPECTED_RUN2_8_DECOMPOSITION_IDS <= covered_decomposition_unit_ids
    assert EXPECTED_RUN2_8_MEMORY_BINDING_IDS <= covered_memory_binding_ids
    assert code_binding_ids <= covered_code_binding_ids
    assert EXPECTED_RUN2_8_TRACE_FIELDS <= covered_trace_fields


def test_run2_9_has_visual_primitive_repair_data() -> None:
    repair_path = PACK / "run2_9_visual_primitive_repair.json"
    assert repair_path.exists(), "missing Run 2.9 visual primitive repair data"
    repair = load_json(repair_path)
    sources = load_json(PACK / "sources.json")
    run2_8_decomposition = load_json(PACK / "run2_8_tutorial_decomposition.json")
    source_ids = {source["id"] for source in sources["sources"]}
    decomposition_ids = {unit["id"] for unit in run2_8_decomposition["units"]}

    assert repair["status"] == "run2_9_visual_primitive_repair_public_blocked"
    assert repair["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert repair["storage_policy"]["raw_media"] == "forbidden"
    assert EXPECTED_RUN2_9_VISUAL_PRIMITIVE_IDS <= {item["id"] for item in repair["primitive_repairs"]}

    required = {
        "id",
        "source_ids",
        "decomposition_unit_ids",
        "target_slide_roles",
        "visual_problem",
        "reference_method",
        "extracted_visual_primitive",
        "native_ppt_translation",
        "code_module_obligation",
        "layout_recipe",
        "typography_recipe",
        "motion_or_sequence_recipe",
        "forbidden_box_patterns",
        "boxiness_failure_probe",
        "qa_probe",
        "anti_copy_boundary",
        "release_boundary",
    }
    for primitive in repair["primitive_repairs"]:
        assert required <= set(primitive), primitive["id"]
        assert set(primitive["source_ids"]) <= source_ids
        assert set(primitive["decomposition_unit_ids"]) <= decomposition_ids
        assert set(primitive["target_slide_roles"]) <= EXPECTED_RHYTHM_ROLES
        assert_contains(primitive["native_ppt_translation"], ["native", "editable"])
        assert_contains(primitive["code_module_obligation"], ["function", "trace"])
        assert_mentions_any(
            " ".join(primitive["forbidden_box_patterns"]),
            {"dashboard", "card", "equal panel", "rectangle", "grid"},
        )
        assert_contains(primitive["boxiness_failure_probe"], ["box"])
        for field_value in iter_string_values(primitive):
            lowered = field_value.lower()
            for marker in RUN2_8_FORBIDDEN_MEDIA_MARKERS:
                assert marker not in lowered, f"{primitive['id']} contains copied media marker {marker!r}"


def test_run2_9_has_executable_visual_modules_and_gate_matrix() -> None:
    repair_path = PACK / "run2_9_visual_primitive_repair.json"
    modules_path = PACK / "run2_9_executable_visual_modules.json"
    matrix_path = PACK / "run2_9_visual_gate_matrix.json"
    trace_contract_path = PACK / "results" / "trace_manifest_contract.json"
    workflow_path = PACK / "skill_workflow.json"
    assert repair_path.exists(), "missing Run 2.9 visual primitive repair data"
    assert modules_path.exists(), "missing Run 2.9 executable visual modules"
    assert matrix_path.exists(), "missing Run 2.9 visual gate matrix"
    modules = load_json(modules_path)
    matrix = load_json(matrix_path)
    repair = load_json(repair_path)
    trace_contract = load_json(trace_contract_path)
    workflow = load_json(workflow_path)
    primitive_ids = {item["id"] for item in repair["primitive_repairs"]}
    module_ids = {module["id"] for module in modules["modules"]}
    code_module_ids = {module["code_binding"]["function_name"] for module in modules["modules"]}

    assert modules["status"] == "run2_9_executable_visual_modules_public_blocked"
    assert modules["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert modules["memory_type"] == "visual_primitive_code_modules"
    assert EXPECTED_RUN2_9_VISUAL_MODULE_IDS <= module_ids
    assert EXPECTED_RUN2_9_TRACE_FIELDS <= set(trace_contract["per_slide_required_fields"])

    module_required = {
        "id",
        "primitive_repair_ids",
        "applies_to_slide_roles",
        "code_binding",
        "composition_contract",
        "native_ppt_primitives",
        "negative_control_failure",
        "qa_probe",
        "release_boundary",
    }
    for module in modules["modules"]:
        assert module_required <= set(module), module["id"]
        assert set(module["primitive_repair_ids"]) <= primitive_ids
        assert set(module["applies_to_slide_roles"]) <= EXPECTED_RHYTHM_ROLES
        assert {"function_name", "params", "layout_budget"} <= set(module["code_binding"]), module["id"]
        assert module["code_binding"]["function_name"].startswith("drawRun29")
        assert_contains(" ".join(module["native_ppt_primitives"]), ["native", "editable"])
        assert_mentions_any(
            module["negative_control_failure"],
            {"box", "dashboard", "card", "grid", "report"},
        )

    assert matrix["status"] == "run2_9_visual_gate_matrix_public_blocked"
    assert matrix["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert set(matrix["selection_chain"]) == {
        "run2_8_visual_diagnosis",
        "run2_9_visual_primitive_repairs",
        "run2_9_executable_visual_modules",
        "run2_9_visual_gate_matrix",
        "native_ppt_code_generation",
        "visual_delta_gate",
        "delivery_gate",
    }
    gate_required = {
        "id",
        "slide_role",
        "visual_primitive_ids",
        "visual_module_ids",
        "required_code_modules",
        "boxiness_failure_probe",
        "pass_fail_checks",
        "trace_fields",
        "public_release_gate",
    }
    covered_trace_fields = set()
    for gate in matrix["gates"]:
        assert gate_required <= set(gate), gate["id"]
        assert gate["slide_role"] in EXPECTED_RHYTHM_ROLES
        assert set(gate["visual_primitive_ids"]) <= primitive_ids
        assert set(gate["visual_module_ids"]) <= module_ids
        assert set(gate["required_code_modules"]) <= code_module_ids
        assert set(gate["trace_fields"]) <= set(trace_contract["per_slide_required_fields"])
        assert_contains(gate["boxiness_failure_probe"], ["box"])
        covered_trace_fields |= set(gate["trace_fields"])
    assert EXPECTED_RUN2_9_TRACE_FIELDS <= covered_trace_fields

    workflow_stage_ids = [stage["id"] for stage in workflow["stages"]]
    run2_9_stages = [
        "repair_run2_9_visual_primitives",
        "select_run2_9_executable_visual_modules",
        "apply_run2_9_visual_gate_matrix",
    ]
    for stage_id in run2_9_stages:
        assert stage_id in workflow_stage_ids
        assert workflow_stage_ids.index(stage_id) < workflow_stage_ids.index("generate_code_first_ppt")
    assert workflow_stage_ids.index(run2_9_stages[0]) < workflow_stage_ids.index(run2_9_stages[1])
    assert workflow_stage_ids.index(run2_9_stages[1]) < workflow_stage_ids.index(run2_9_stages[2])
    assert_contains(
        json.dumps(workflow),
        [
            "run2_9_visual_primitive_repair.json",
            "run2_9_executable_visual_modules.json",
            "run2_9_visual_gate_matrix.json",
            "visual primitive",
            "boxiness",
        ],
    )


def test_run2_10_has_visual_system_sources_memory_and_gate_matrix() -> None:
    sources_path = PACK / "run2_10_visual_system_sources.json"
    memory_path = PACK / "run2_10_visual_system_memory.json"
    matrix_path = PACK / "run2_10_visual_system_gate_matrix.json"
    trace_contract_path = PACK / "results" / "trace_manifest_contract.json"
    workflow_path = PACK / "skill_workflow.json"

    assert sources_path.exists(), "missing Run 2.10 visual-system sources"
    assert memory_path.exists(), "missing Run 2.10 visual-system memory"
    assert matrix_path.exists(), "missing Run 2.10 visual-system gate matrix"

    sources = load_json(sources_path)
    memory = load_json(memory_path)
    matrix = load_json(matrix_path)
    trace_contract = load_json(trace_contract_path)
    workflow = load_json(workflow_path)

    assert sources["status"] == "run2_10_visual_system_sources_public_blocked"
    assert memory["status"] == "run2_10_visual_system_memory_public_blocked"
    assert matrix["status"] == "run2_10_visual_system_gate_matrix_public_blocked"
    assert sources["storage_policy"]["raw_media"] == "forbidden"
    assert sources["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert memory["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert matrix["stage_policy"] == "repeat_same_five_layers_not_run3"

    source_ids = {item["id"] for item in sources["sources"]}
    memory_ids = {item["visual_system_id"] for item in memory["visual_systems"]}
    assert EXPECTED_RUN2_10_SOURCE_IDS <= source_ids
    assert EXPECTED_RUN2_10_MEMORY_IDS <= memory_ids
    assert EXPECTED_RUN2_10_TRACE_FIELDS <= set(trace_contract["per_slide_required_fields"])

    required_source_fields = {
        "id",
        "source_ids",
        "reference_type",
        "allowed_use",
        "visual_system_direction",
        "typography_observation",
        "spatial_composition_observation",
        "asset_strategy_observation",
        "motion_or_sequence_observation",
        "climax_grammar_observation",
        "native_ppt_implication",
        "anti_copy_boundary",
        "public_demo_probe",
        "release_boundary",
    }
    for source in sources["sources"]:
        assert required_source_fields <= set(source), source["id"]
        assert source["allowed_use"] == "derived_observations_only"
        assert_contains(source["native_ppt_implication"], ["native", "editable"])
        for field_value in iter_string_values(source):
            lowered = field_value.lower()
            for marker in RUN2_8_FORBIDDEN_MEDIA_MARKERS:
                assert marker not in lowered, f"{source['id']} contains copied media marker {marker!r}"

    required_memory_fields = {
        "visual_system_id",
        "source_record_ids",
        "applicable_slide_roles",
        "typography_contract",
        "composition_contract",
        "asset_strategy_contract",
        "motion_sequence_contract",
        "native_ppt_module_implications",
        "forbidden_sameness_patterns",
        "public_demo_first_read_probe",
        "anti_copy_boundary",
        "release_boundary",
    }
    for entry in memory["visual_systems"]:
        assert required_memory_fields <= set(entry), entry["visual_system_id"]
        assert set(entry["source_record_ids"]) <= source_ids
        assert set(entry["applicable_slide_roles"]) <= EXPECTED_RHYTHM_ROLES
        assert_mentions_any(
            " ".join(entry["forbidden_sameness_patterns"]),
            {"rectangle", "same visual family", "palette-only", "dashboard"},
        )

    gate_required_fields = {
        "id",
        "slide_role",
        "visual_system_source_ids",
        "visual_system_memory_ids",
        "required_code_modules",
        "visual_delta_from_run2_9",
        "sameness_failure_probe",
        "public_demo_first_read_probe",
        "shape_count_budget",
        "asymmetry_whitespace_rule",
        "trace_fields",
        "public_release_gate",
    }
    covered_trace_fields = set()
    for gate in matrix["gates"]:
        assert gate_required_fields <= set(gate), gate["id"]
        assert gate["slide_role"] in EXPECTED_RHYTHM_ROLES
        assert set(gate["visual_system_source_ids"]) <= source_ids
        assert set(gate["visual_system_memory_ids"]) <= memory_ids
        assert set(gate["trace_fields"]) <= set(trace_contract["per_slide_required_fields"])
        assert_contains(gate["sameness_failure_probe"], ["Run 2.9"])
        assert_contains(gate["asymmetry_whitespace_rule"], ["asymmetry", "whitespace"])
        assert gate["shape_count_budget"]["max_native_shapes"] <= 72
        covered_trace_fields |= set(gate["trace_fields"])
    assert EXPECTED_RUN2_10_TRACE_FIELDS <= covered_trace_fields

    workflow_stage_ids = [stage["id"] for stage in workflow["stages"]]
    run2_10_stages = [
        "select_run2_10_visual_system_sources",
        "compile_run2_10_visual_system_memory",
        "apply_run2_10_visual_system_gate_matrix",
    ]
    for stage_id in run2_10_stages:
        assert stage_id in workflow_stage_ids
        assert workflow_stage_ids.index(stage_id) < workflow_stage_ids.index("generate_code_first_ppt")
    assert workflow_stage_ids.index(run2_10_stages[0]) < workflow_stage_ids.index(run2_10_stages[1])
    assert workflow_stage_ids.index(run2_10_stages[1]) < workflow_stage_ids.index(run2_10_stages[2])


def test_run2_11_data_workflow_audit_artifact_has_required_chains() -> None:
    audit_path = PACK / "results" / "run2_11_data_workflow_audit.json"
    assert audit_path.exists(), "missing Run 2.11 data/workflow audit"

    audit = load_json(audit_path)
    assert EXPECTED_RUN2_11_AUDIT_FIELDS <= set(audit)
    assert audit["status"] == "run2_11_data_workflow_audit_public_blocked"
    assert audit["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert audit["audit_scope"]["creates_new_ppt_deck"] is False
    assert audit["audit_scope"]["primary_quality_target"] == "data_workflow_evidence"

    chains = audit["chain_records"]
    assert chains, "audit must contain source-to-trace chains"

    run28_decomposition_units_by_id = {
        unit["id"]: unit for unit in load_json(PACK / "run2_8_tutorial_decomposition.json")["units"]
    }
    trace_fields_by_run = {
        "2.8": {
            "decomposition_ids": "run2_8_decomposition_unit_ids",
            "memory_ids": "run2_8_memory_binding_ids",
            "gate_ids": "run2_8_gate_matrix_ids",
            "code_ids": "run2_8_code_binding_ids",
        },
        "2.9": {
            "decomposition_ids": "run2_9_visual_primitive_ids",
            "memory_ids": "run2_9_visual_module_ids",
            "gate_ids": "run2_9_gate_matrix_ids",
            "code_ids": "run2_9_code_module_ids",
        },
        "2.10": {
            "source_ids": "run2_10_visual_system_source_ids",
            "memory_ids": "run2_10_visual_system_memory_ids",
            "gate_ids": "run2_10_gate_matrix_ids",
            "code_ids": "run2_10_code_module_ids",
        },
    }
    required_non_empty_chain_fields_by_run = {
        "2.8": {"source_ids", "decomposition_ids", "memory_ids", "gate_ids"},
        "2.9": {"decomposition_ids", "memory_ids", "gate_ids"},
        "2.10": {"source_ids", "memory_ids", "gate_ids"},
    }
    collected_run_ids = set()
    for chain in chains:
        assert {
            "chain_id",
            "run_id",
            "layer",
            "source_ids",
            "decomposition_ids",
            "memory_ids",
            "gate_ids",
            "required_code_module_ids",
            "actual_code_module_ids",
            "slide_roles",
            "trace_manifest_paths",
            "control_boundary",
            "status",
            "reason",
            "next_fix",
        } <= set(chain), chain.get("chain_id")
        assert chain["status"] in EXPECTED_RUN2_11_CHAIN_STATUSES
        assert chain["run_id"] in EXPECTED_RUN2_11_RUN_IDS
        collected_run_ids.add(chain["run_id"])
        if chain["status"] in {"weak", "missing", "blocked"}:
            assert chain["next_fix"], chain["chain_id"]
        if chain["status"] == "pass":
            assert chain["required_code_module_ids"], chain["chain_id"]
            assert chain["actual_code_module_ids"], chain["chain_id"]
            assert set(chain["required_code_module_ids"]) <= set(chain["actual_code_module_ids"])
            for field in required_non_empty_chain_fields_by_run[chain["run_id"]]:
                assert chain[field], f"{chain['chain_id']} missing {field}"
            assert chain["trace_manifest_paths"], chain["chain_id"]
            trace_values = {
                chain_field: set()
                for chain_field in trace_fields_by_run[chain["run_id"]]
            }
            for relative_path in chain["trace_manifest_paths"]:
                trace_path = ROOT / relative_path
                assert trace_path.exists(), relative_path
                trace = load_json(trace_path)
                for slide in trace.get("slides", []):
                    for chain_field, trace_field in trace_fields_by_run[chain["run_id"]].items():
                        trace_values[chain_field].update(slide.get(trace_field, []))
            for chain_field, actual_ids_from_trace in trace_values.items():
                if chain_field == "code_ids":
                    continue
                assert set(chain[chain_field]) <= actual_ids_from_trace
            if chain["run_id"] == "2.8":
                source_record_ids_from_decomposition = {
                    source_record_id
                    for decomposition_id in chain["decomposition_ids"]
                    for source_record_id in run28_decomposition_units_by_id[decomposition_id][
                        "source_record_ids"
                    ]
                }
                assert set(chain["source_ids"]) <= source_record_ids_from_decomposition
            assert set(chain["actual_code_module_ids"]) <= trace_values["code_ids"]
            assert set(chain["required_code_module_ids"]) <= trace_values["code_ids"]

    assert EXPECTED_RUN2_11_RUN_IDS <= collected_run_ids
    assert any(chain["run_id"] == "2.8" and chain["status"] == "pass" for chain in chains)
    assert any(chain["run_id"] == "2.9" and chain["status"] == "pass" for chain in chains)
    assert any(chain["run_id"] == "2.10" and chain["status"] == "pass" for chain in chains)


def test_run2_11_audit_references_existing_data_memory_gate_and_trace_fields() -> None:
    audit = load_json(PACK / "results" / "run2_11_data_workflow_audit.json")
    workflow = load_json(PACK / "skill_workflow.json")
    run27_sources = load_json(PACK / "run2_7_multimodal_source_records.json")
    run28_decomp = load_json(PACK / "run2_8_tutorial_decomposition.json")
    run28_memory = load_json(PACK / "run2_8_executable_design_memory.json")
    run28_gate = load_json(PACK / "run2_8_workflow_gate_matrix.json")
    run29_primitives = load_json(PACK / "run2_9_visual_primitive_repair.json")
    run29_modules = load_json(PACK / "run2_9_executable_visual_modules.json")
    run29_gate = load_json(PACK / "run2_9_visual_gate_matrix.json")
    run210_sources = load_json(PACK / "run2_10_visual_system_sources.json")
    run210_memory = load_json(PACK / "run2_10_visual_system_memory.json")
    run210_gate = load_json(PACK / "run2_10_visual_system_gate_matrix.json")

    assert audit["source_inventory"] == {
        "run2_7_source_records": len(run27_sources["records"]),
        "run2_8_tutorial_decomposition_units": len(run28_decomp["units"]),
        "run2_9_visual_primitives": len(run29_primitives["primitive_repairs"]),
        "run2_10_visual_system_sources": len(run210_sources["sources"]),
        "run2_10_visual_system_memory_records": len(run210_memory["visual_systems"]),
    }
    assert audit["workflow_inventory"] == {
        "skill_workflow_stages": 31,
        "run2_8_workflow_gates": len(run28_gate["gates"]),
        "run2_9_visual_gates": len(run29_gate["gates"]),
        "run2_10_visual_system_gates": len(run210_gate["gates"]),
    }
    assert audit["workflow_inventory"]["skill_workflow_stages"] <= len(workflow["stages"])

    known_source_ids = {record["id"] for record in run27_sources["records"]}
    known_decomposition_ids = {unit["id"] for unit in run28_decomp["units"]}
    known_memory_ids = {binding["id"] for binding in run28_memory["bindings"]}
    known_gate_ids = {gate["id"] for gate in run28_gate["gates"]}
    known_run29_primitive_ids = {item["id"] for item in run29_primitives["primitive_repairs"]}
    known_run29_module_ids = {item["id"] for item in run29_modules["modules"]}
    known_run29_gate_ids = {gate["id"] for gate in run29_gate["gates"]}
    known_run210_source_ids = {item["id"] for item in run210_sources["sources"]}
    known_run210_memory_ids = {item["visual_system_id"] for item in run210_memory["visual_systems"]}
    known_run210_gate_ids = {gate["id"] for gate in run210_gate["gates"]}

    for chain in audit["chain_records"]:
        if chain["run_id"] == "2.8":
            assert set(chain["source_ids"]) <= known_source_ids
            assert set(chain["decomposition_ids"]) <= known_decomposition_ids
            assert set(chain["memory_ids"]) <= known_memory_ids
            assert set(chain["gate_ids"]) <= known_gate_ids
        elif chain["run_id"] == "2.9":
            assert chain["source_ids"] == []
            assert set(chain["decomposition_ids"]) <= known_run29_primitive_ids
            assert set(chain["memory_ids"]) <= known_run29_module_ids
            assert set(chain["gate_ids"]) <= known_run29_gate_ids
        elif chain["run_id"] == "2.10":
            assert set(chain["source_ids"]) <= known_run210_source_ids
            assert chain["decomposition_ids"] == []
            assert set(chain["memory_ids"]) <= known_run210_memory_ids
            assert set(chain["gate_ids"]) <= known_run210_gate_ids


def test_run2_11_audit_records_trace_evidence_and_control_boundaries() -> None:
    audit = load_json(PACK / "results" / "run2_11_data_workflow_audit.json")

    trace_evidence = audit["arm_trace_evidence"]
    assert {"run2_8_full_skill", "run2_9_full_skill", "run2_10_full_skill"} <= set(trace_evidence)

    assert trace_evidence["run2_8_full_skill"]["trace_origin"] == "actual_native_module_calls"
    assert trace_evidence["run2_9_full_skill"]["trace_origin"] == "actual_native_visual_module_calls"
    assert trace_evidence["run2_10_full_skill"]["trace_origin"] == "actual_native_visual_module_calls"

    controls = audit["negative_control_checks"]
    control_arms = {item["arm_role"] for item in controls}
    assert EXPECTED_RUN2_11_CONTROL_ARMS <= control_arms
    for check in controls:
        assert check["status"] in {"pass", "weak", "missing", "blocked"}
        assert check["forbidden_fields"]
        assert check["observed_boundary"]


def test_run2_11_audit_is_embedded_in_html_viewer(tmp_path: Path) -> None:
    presentations_dir = tmp_path / "presentations"
    presentations_dir.mkdir()
    out = presentations_dir / "ppt-run-viewer.html"
    completed = subprocess.run(
        [
            sys.executable,
            "scripts/build_ppt_run_html_viewer.py",
            "--presentations-dir",
            str(presentations_dir),
            "--out",
            str(out),
        ],
        cwd=ROOT,
        capture_output=True,
        text=True,
        timeout=20,
    )

    assert completed.returncode == 0, completed.stderr
    body = out.read_text(encoding="utf-8")
    assert "Data/Workflow Audit" in body
    assert 'data-view="audit"' in body
    assert "run2_11_data_workflow_audit_public_blocked" in body
    assert "Source-to-slide chains" in body


def test_run2_12_has_thick_multimodal_evidence_records() -> None:
    evidence_path = PACK / "run2_12_thick_multimodal_evidence.json"
    assert evidence_path.exists(), "missing Run 2.12 thick multimodal evidence"

    evidence = load_json(evidence_path)
    sources = load_json(PACK / "sources.json")
    source_ids = {source["id"] for source in sources["sources"]}
    source_urls = {source["id"]: source["url"] for source in sources["sources"]}

    assert evidence["status"] == "run2_12_thick_multimodal_evidence_public_blocked"
    assert evidence["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert evidence["audit_response_to"] == "run2_11_data_workflow_chain_gate"
    assert evidence["storage_policy"]["raw_media"] == "forbidden"
    assert evidence["run_scope"]["creates_new_ppt_deck"] is False

    records = evidence["records"]
    assert EXPECTED_RUN2_12_EVIDENCE_IDS <= {record["id"] for record in records}
    covered_modalities = {modality for record in records for modality in record["modality_mix"]}
    assert EXPECTED_MULTIMODAL_MODALITIES <= covered_modalities

    for record in records:
        assert EXPECTED_RUN2_12_EVIDENCE_FIELDS <= set(record), record["id"]
        assert record["source_ids"]
        assert set(record["source_ids"]) <= source_ids
        assert record["source_urls"]
        assert record["verified_accessed_on"] == "2026-06-03"
        assert set(record["modality_mix"]) <= EXPECTED_MULTIMODAL_MODALITIES
        assert len(record["frame_or_visual_observations"]) >= 2
        assert len(record["spoken_or_text_claim_paraphrases"]) >= 1
        assert len(record["native_ppt_code_obligations"]) >= 2
        assert len(record["workflow_gate_obligations"]) >= 2
        assert record["memory_seed_targets"]
        assert_contains(record["anti_copy_boundary"], ["do not copy"])
        assert_contains(record["release_boundary"], ["public_blocked"])
        for source_url in record["source_urls"]:
            assert source_url["source_id"] in record["source_ids"]
            assert source_url["url"] == source_urls[source_url["source_id"]]
            assert source_url["retrieval_status"] == "verified_public_page"
        locator = record["segment_locator"]
        assert {"locator_type", "public_locator", "derived_only_note"} <= set(locator)
        assert_contains(locator["derived_only_note"], ["derived", "no raw media"])


def test_run2_12_memory_and_workflow_gate_seeds_are_referentially_integrated() -> None:
    evidence = load_json(PACK / "run2_12_thick_multimodal_evidence.json")
    memory = load_json(PACK / "run2_12_design_memory_seed.json")
    gate_seed = load_json(PACK / "run2_12_workflow_gate_seed.json")
    workflow = load_json(PACK / "skill_workflow.json")
    trace_contract = load_json(PACK / "results" / "trace_manifest_contract.json")

    evidence_ids = {record["id"] for record in evidence["records"]}
    memory_ids = {record["id"] for record in memory["memory_seeds"]}
    trace_fields = set(trace_contract["per_slide_required_fields"])

    assert memory["status"] == "run2_12_design_memory_seed_public_blocked"
    assert memory["derived_from"] == "run2_12_thick_multimodal_evidence.json"
    assert memory["run_scope"]["creates_new_ppt_deck"] is False
    assert EXPECTED_RUN2_12_MEMORY_IDS <= memory_ids
    for record in memory["memory_seeds"]:
        assert EXPECTED_RUN2_12_MEMORY_FIELDS <= set(record), record["id"]
        assert set(record["evidence_record_ids"]) <= evidence_ids
        assert set(record["applies_to_slide_roles"]) <= EXPECTED_RHYTHM_ROLES
        assert record["design_constraints"]
        assert_contains(json.dumps(record["native_ppt_contract"]), ["native", "editable"])
        assert set(record["required_trace_fields"]) <= trace_fields
        assert_contains(record["release_boundary"], ["public_blocked"])

    assert gate_seed["status"] == "run2_12_workflow_gate_seed_public_blocked"
    assert gate_seed["gate_policy"] == "required_before_next_four_arm_rerun"
    assert gate_seed["run_scope"]["creates_new_ppt_deck"] is False
    assert EXPECTED_RUN2_12_GATE_IDS <= {gate["id"] for gate in gate_seed["gates"]}
    for gate in gate_seed["gates"]:
        assert EXPECTED_RUN2_12_GATE_FIELDS <= set(gate), gate["id"]
        assert set(gate["evidence_record_ids"]) <= evidence_ids
        assert set(gate["memory_seed_ids"]) <= memory_ids
        assert set(gate["slide_roles"]) <= EXPECTED_RHYTHM_ROLES
        assert gate["pass_fail_checks"]
        assert gate["required_before_next_rerun"] is True
        assert set(gate["trace_fields"]) <= trace_fields
        assert_contains(gate["release_boundary"], ["public_blocked"])

    workflow_text = json.dumps(workflow)
    assert_contains(
        workflow_text,
        [
            "run2_12_thick_multimodal_evidence.json",
            "run2_12_design_memory_seed.json",
            "run2_12_workflow_gate_seed.json",
            "required before next four-arm rerun",
        ],
    )


def test_run2_four_arm_isolation_mentions_multimodal_and_target_boundaries() -> None:
    prompt_only = (PACK / "generation_briefs" / "prompt_only.md").read_text(encoding="utf-8")
    run1_5 = (PACK / "generation_briefs" / "run1_5_skill.md").read_text(encoding="utf-8")
    run2 = (PACK / "generation_briefs" / "run2_skill.md").read_text(encoding="utf-8")
    bad = (PACK / "generation_briefs" / "bad_aesthetic_memory.md").read_text(encoding="utf-8")

    assert_contains(
        prompt_only, ["Do not use multimodal database", "visual learning targets", "visual target components"]
    )
    assert_contains(prompt_only, ["video demo beat map", "motion learning targets", "presentation sequence components"])
    assert_contains(
        run1_5,
        ["Do not use Run 2.2 multimodal database", "visual learning targets", "visual target components"],
    )
    assert_contains(
        run1_5, ["video_demo_beat_map.json", "motion_learning_targets.json", "presentation_sequence_components.json"]
    )
    assert_contains(
        run2,
        [
            "multimodal_database.json",
            "visual_learning_targets.json",
            "visual_target_components.json",
            "video_demo_beat_map.json",
            "motion_learning_targets.json",
            "presentation_sequence_components.json",
            "production_reference_decompositions.json",
            "aesthetic_memory_v2.json",
            "visual_production_modules.json",
        ],
    )
    assert_contains(
        bad,
        [
            "multimodal_database.json",
            "visual_learning_targets.json",
            "visual_target_components.json",
            "video_demo_beat_map.json",
            "motion_learning_targets.json",
            "presentation_sequence_components.json",
        ],
    )
    assert_contains(bad, ["Good aesthetic_memory.json", "Good slide_archetypes.json"])


def test_run2_cards_have_executable_extraction_units() -> None:
    required = {"unit_id", "source_anchor", "derived_rule", "slide_role", "execution_guard", "qa_probe"}
    allowed_roles = {"cover", "setup", "contrast", "proof", "climax", "relief", "close"}

    cards = [load_json(path) for path in [*json_files("source_cards"), *json_files("video_cards")]]

    for card in cards:
        units = card.get("extraction_units")
        assert isinstance(units, list) and len(units) >= 2, card["card_id"]
        for unit in units:
            assert required <= set(unit), card["card_id"]
            assert unit["slide_role"] in allowed_roles
            for key in required - {"slide_role"}:
                assert isinstance(unit[key], str) and unit[key].strip(), f"{card['card_id']} {key}"


def test_run2_evidence_memory_has_required_claims() -> None:
    memory = load_json(PACK / "evidence_memory.json")
    claims = memory["claims"]

    assert len(claims) >= 6
    assert EXPECTED_CLAIM_IDS <= {claim["id"] for claim in claims}
    for claim in claims:
        assert claim["source_card_ids"]
        assert claim["qa_checks"]


def test_run2_asset_memory_has_required_editable_assets() -> None:
    memory = load_json(PACK / "asset_memory.json")
    assets = memory["assets"]

    assert len(assets) >= 4
    assert EXPECTED_ASSET_IDS <= {asset["id"] for asset in assets}
    for asset in assets:
        assert asset["provenance_state"]
        assert asset["text_editability"]
        assert asset["render_risks"]
        assert asset["accessibility_risks"]


def test_run2_aesthetic_memory_has_required_public_demo_moves() -> None:
    memory = load_json(PACK / "aesthetic_memory.json")
    moves = {move["aesthetic_move"] for move in memory["moves"]}

    assert EXPECTED_MOVES <= moves
    assert len(memory["moves"]) >= 8
    for move in memory["moves"]:
        assert move["source_card_ids"]
        assert move["negative_rules"]
        assert move["density_budget"]["max_words"] <= 80
        assert move["rhythm_role"] in {"cover", "setup", "contrast", "proof", "climax", "relief", "close"}


def test_run2_narrative_spine_is_six_slide_public_demo() -> None:
    spine = load_json(PACK / "narrative_spine.json")

    assert spine["deck_length"] == 6
    assert [slide["rhythm_role"] for slide in spine["slides"]] == [
        "cover",
        "setup",
        "contrast",
        "proof",
        "climax",
        "close",
    ]


def test_run2_skill_is_a_staged_deck_director() -> None:
    body = (PACK / "vulca_ppt_skill.md").read_text(encoding="utf-8")

    assert_contains(
        body,
        [
            "compile the multimodal database",
            "select the narrative spine",
            "compile evidence memory",
            "compile aesthetic memory",
            "select assets",
            "generate code-first PPT modules",
            "run structural QA",
            "run aesthetic QA",
            "repair",
            "release decision",
        ],
    )
    assert_contains(body, ["move detail to appendix", "not compress the same content into smaller text"])
    assert_contains(
        body,
        [
            "trace manifest",
            "per slide",
            "Selected multimodal record ids",
            "density counts",
            "deleted or routed content",
            "asset provenance",
            "QA outcomes",
        ],
    )
    assert_contains(body, ["text, image-reference, video, audio, transcript, and interaction anchors"])
    assert_contains(body, ["visual_learning_targets.json", "before/after visual deltas", "slide mini-previews"])
    assert_contains(body, ["visual_target_components.json", "before/after thumbnail", "visual component ids"])
    assert_contains(
        body,
        [
            "video_demo_beat_map.json",
            "motion_learning_targets.json",
            "presentation_sequence_components.json",
            "motion target ids",
            "sequence component ids",
        ],
    )
    assert_contains(
        body,
        [
            "production_reference_decompositions.json",
            "aesthetic_memory_v2.json",
            "visual_production_modules.json",
            "production modules",
            "run2_10_visual_system_sources.json",
            "run2_10_visual_system_memory.json",
            "run2_10_visual_system_gate_matrix.json",
            "sameness failure probes",
        ],
    )


def test_run2_skill_workflow_is_declarative_and_gated() -> None:
    workflow = load_json(PACK / "skill_workflow.json")
    stage_ids = [stage["id"] for stage in workflow["stages"]]

    assert workflow["workflow_type"] == "declarative_skill_director"
    assert workflow["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert workflow["five_layer_loop"] == [
        "real_commercial_case",
        "multimodal_tutorial_case_data",
        "evidence_aesthetic_asset_memory",
        "skill_workflow",
        "rerun_and_evaluation",
    ]
    assert stage_ids == [
        "read_commercial_case",
        "compile_multimodal_database",
        "compile_evidence_memory",
        "compile_aesthetic_memory",
        "compile_production_reference_decompositions",
        "compile_aesthetic_memory_v2",
        "select_slide_archetypes",
        "select_commercial_usecase",
        "select_aesthetic_benchmarks",
        "select_theme_typography_spacing_policy",
        "select_visual_repair_policy",
        "select_run2_7_design_memory",
        "select_visual_production_modules",
        "decompose_run2_8_tutorial_video_units",
        "select_run2_8_executable_design_memory",
        "apply_run2_8_workflow_gate_matrix",
        "repair_run2_9_visual_primitives",
        "select_run2_9_executable_visual_modules",
        "apply_run2_9_visual_gate_matrix",
        "select_run2_10_visual_system_sources",
        "compile_run2_10_visual_system_memory",
        "apply_run2_10_visual_system_gate_matrix",
        "apply_run2_15_layout_selector_gate_matrix",
        "generate_code_first_ppt",
        "run_structural_and_aesthetic_qa",
        "recommend_repairs",
        "refresh_trace_qa_outcomes",
        "emit_release_decision",
        "expand_run2_18_multimodal_evidence",
        "expand_run2_18_design_memory",
        "apply_run2_18_workflow_gate_expansion",
        "lock_run2_24_single_usecase_content_memory",
        "compile_run2_24_visual_evidence_asset_memory",
        "apply_run2_24_content_visual_workflow_gates",
        "compile_run2_35_visual_evidence_asset_realism_memory",
        "compile_run2_35_editorial_composition_memory",
        "apply_run2_35_visual_evidence_workflow_gates",
        "compile_run2_38_public_video_slide_direction_memory",
        "compile_run2_38_per_slide_visual_recipe_memory",
        "apply_run2_38_public_video_workflow_gates",
        "compile_run2_43_semantic_visual_asset_memory",
        "compile_run2_43_editorial_composition_typography_memory",
        "apply_run2_43_visual_asset_semantics_workflow_gates",
        "compile_run2_49_readability_memory",
        "compile_run2_49_content_evidence_density_memory",
        "apply_run2_49_editorial_renderer_workflow_gates",
        "compile_run2_51_editorial_copy_memory",
        "compile_run2_51_shape_text_socket_memory",
        "apply_run2_51_renderer_archetype_workflow_gates",
        "compile_run2_53_product_surface_scene_memory",
        "compile_run2_53_business_visual_evidence_memory",
        "apply_run2_53_scene_renderer_workflow_gates",
    ]
    assert [stage["order"] for stage in workflow["stages"]] == list(range(1, len(workflow["stages"]) + 1))
    assert workflow["repair_triggers"]
    workflow_text = json.dumps(workflow)
    assert "multimodal_database.json" in workflow_text
    assert "visual_learning_targets.json" in workflow_text
    assert "visual_target_components.json" in workflow_text
    assert "video_demo_beat_map.json" in workflow_text
    assert "motion_learning_targets.json" in workflow_text
    assert "presentation_sequence_components.json" in workflow_text
    assert "production_reference_decompositions.json" in workflow_text
    assert "aesthetic_memory_v2.json" in workflow_text
    assert "visual_production_modules.json" in workflow_text
    assert "commercial_usecase_bank.json" in workflow_text
    assert "run2_7_commercial_usecase.json" in workflow_text
    assert "run2_43_semantic_visual_asset_memory.json" in workflow_text
    assert "run2_43_editorial_composition_typography_memory.json" in workflow_text
    assert "run2_43_visual_asset_semantics_workflow_gates.json" in workflow_text
    assert "run2_7_multimodal_source_records.json" in workflow_text
    assert "run2_7_design_memory.json" in workflow_text
    assert "run2_7_workflow_policy.json" in workflow_text
    assert "run2_8_tutorial_decomposition.json" in workflow_text
    assert "run2_8_executable_design_memory.json" in workflow_text
    assert "run2_8_workflow_gate_matrix.json" in workflow_text
    assert "run2_9_visual_primitive_repair.json" in workflow_text
    assert "run2_9_executable_visual_modules.json" in workflow_text
    assert "run2_9_visual_gate_matrix.json" in workflow_text
    assert "run2_10_visual_system_sources.json" in workflow_text
    assert "run2_10_visual_system_memory.json" in workflow_text
    assert "run2_10_visual_system_gate_matrix.json" in workflow_text
    assert "run2_24_single_usecase_content_memory.json" in workflow_text
    assert "run2_24_visual_evidence_asset_memory.json" in workflow_text
    assert "run2_24_content_visual_workflow_gates.json" in workflow_text
    assert "run2_35_visual_evidence_asset_realism_memory.json" in workflow_text
    assert "run2_35_editorial_composition_memory.json" in workflow_text
    assert "run2_35_visual_evidence_workflow_gates.json" in workflow_text
    assert "run2_15_layout_selector_sources.json" in workflow_text
    assert "run2_15_layout_module_memory.json" in workflow_text
    assert "run2_15_layout_selector_gate_matrix.json" in workflow_text
    assert "aesthetic_benchmark_bank.json" in workflow_text
    assert "workflow_decision_policy.json" in workflow_text
    assert "visual_repair_policy.json" in workflow_text
    assert_contains(workflow_text, ["all required modalities covered", "no copied source media stored"])
    assert_contains(
        workflow_text, ["visual targets reference valid multimodal anchors", "selected visual learning targets"]
    )
    assert_contains(workflow_text, ["native visual components", "selected visual target components"])
    assert_contains(workflow_text, ["motion grammar", "selected motion targets", "sequence components"])
    assert_contains(workflow_text, ["production reference", "aesthetic memory v2", "visual production modules"])
    assert_contains(workflow_text, ["commercial usecase", "aesthetic benchmark", "theme policy"])
    assert_contains(workflow_text, ["visual repair", "Run 2.6R"])
    assert_contains(workflow_text, ["fallback", "visual validation"])
    assert_contains(workflow_text, ["schema validation", "aesthetic memory constraints"])
    assert_contains(workflow_text, ["provenance metadata", "copyright boundary"])
    assert_contains(workflow_text, ["visual system", "sameness failure", "asymmetry", "whitespace"])
    assert_contains(workflow_text, ["layout module selector", "text resilience", "trace visibility", "product theater realism"])
    assert "public_blocked" in workflow["release_decisions"]
    assert "public_ready" not in workflow["release_decisions"]
    assert "automated repair without human gate" not in normalize(json.dumps(workflow))


def test_run2_generation_briefs_define_four_arms() -> None:
    briefs = {path.stem for path in (PACK / "generation_briefs").glob("*.md") if path.name != "README.md"}

    assert briefs == EXPECTED_ARMS
    readme = (PACK / "generation_briefs" / "README.md").read_text(encoding="utf-8")
    assert_contains(readme, ["runtime isolation", "fresh prompt", "separate output directory", "forbidden input"])
    for arm in EXPECTED_ARMS:
        body = (PACK / "generation_briefs" / f"{arm}.md").read_text(encoding="utf-8")
        assert_contains(body, ["Allowed Inputs", "Forbidden Inputs", "Trace Output"])

    assert_contains(
        (PACK / "generation_briefs" / "prompt_only.md").read_text(encoding="utf-8"), ["commercial brief only"]
    )
    assert_contains((PACK / "generation_briefs" / "run1_5_skill.md").read_text(encoding="utf-8"), ["evidence-heavy"])
    assert_contains(
        (PACK / "generation_briefs" / "run2_skill.md").read_text(encoding="utf-8"),
        [
            "multimodal database",
            "visual learning targets",
            "visual target components",
            "video demo beat map",
            "motion learning targets",
            "presentation sequence components",
            "production reference decompositions",
            "aesthetic memory v2",
            "visual production modules",
            "aesthetic memory",
        ],
    )
    assert_contains(
        (PACK / "generation_briefs" / "bad_aesthetic_memory.md").read_text(encoding="utf-8"), ["negative control"]
    )


def test_run2_4_generator_preserves_motion_trace_boundaries() -> None:
    body = (ROOT / "scripts" / "generate_ppt_run2_4_arms.mjs").read_text(encoding="utf-8")

    assert_contains(
        body,
        [
            "prompt_only",
            "run1_5_skill",
            "run2_4_full_skill",
            "bad_aesthetic_memory",
            "no_cross_arm_reuse",
            "layout JSON",
            "contact sheets",
        ],
    )
    assert 'if (!["run2_4_full_skill", "bad_aesthetic_memory"].includes(arm.armId)) return [];' in body
    assert 'const goodOrBad = ["run2_4_full_skill", "bad_aesthetic_memory"].includes(arm.armId);' in body
    assert re.search(r"multimodal_record_ids:\s*goodOrBad\s*\?", body)
    assert re.search(r"visual_learning_target_ids:\s*goodOrBad\s*\?", body)
    assert re.search(r"visual_component_ids:\s*goodOrBad\s*\?", body)
    assert "video_beat_ids: motion.video_beat_ids ?? []" in body
    assert "motion_target_ids: motion.motion_target_ids ?? []" in body
    assert "sequence_component_ids: motion.sequence_component_ids ?? []" in body
    assert "ordered_reveal_steps: sequenceSteps.map" in body
    assert_contains(
        body,
        [
            "beat_opening_scale_reset",
            "motion_target_climax_scale_emphasis",
            "sequence_component_release_gate",
        ],
    )


def test_run2_5_generator_consumes_production_data_and_preserves_control_boundaries() -> None:
    body = (ROOT / "scripts" / "generate_ppt_run2_5_arms.mjs").read_text(encoding="utf-8")
    arm_order = ["prompt_only", "run1_5_skill", "run2_5_full_skill", "bad_aesthetic_memory"]

    def arm_block(arm_id: str) -> str:
        start = body.index(f'armId: "{arm_id}"')
        next_starts = [body.find(f'armId: "{next_arm}"', start + 1) for next_arm in arm_order]
        next_starts = [index for index in next_starts if index > start]
        end = min(next_starts) if next_starts else body.index("function sequenceStepsForSlide", start)
        return body[start:end]

    def section(block: str, start_marker: str, end_marker: str) -> str:
        start = block.index(start_marker)
        end = block.index(end_marker, start)
        return block[start:end]

    restricted_production_inputs = [
        "production_reference_decompositions.json",
        "aesthetic_memory_v2.json",
        "visual_production_modules.json",
    ]

    assert_contains(
        body,
        [
            "prompt_only",
            "run1_5_skill",
            "run2_5_full_skill",
            "bad_aesthetic_memory",
            "production_reference_decompositions.json",
            "aesthetic_memory_v2.json",
            "visual_production_modules.json",
            "no_cross_arm_reuse",
        ],
    )
    prompt_allowed = section(arm_block("prompt_only"), "allowed:", "forbidden:")
    prompt_forbidden = section(arm_block("prompt_only"), "forbidden:", "palette:")
    run1_allowed = section(arm_block("run1_5_skill"), "allowed:", "forbidden:")
    run1_forbidden = section(arm_block("run1_5_skill"), "forbidden:", "palette:")
    full_allowed = section(arm_block("run2_5_full_skill"), "allowed:", "forbidden:")
    full_forbidden = section(arm_block("run2_5_full_skill"), "forbidden:", "palette:")
    bad_allowed = section(arm_block("bad_aesthetic_memory"), "allowed:", "forbidden:")
    bad_forbidden = section(arm_block("bad_aesthetic_memory"), "forbidden:", "palette:")

    for term in restricted_production_inputs:
        assert term not in prompt_allowed
        assert term in prompt_forbidden
        assert term not in run1_allowed
        assert term in run1_forbidden
        assert term in full_allowed
        assert term not in full_forbidden
    assert "production_reference_decompositions.json" in bad_allowed
    assert "aesthetic_memory_v2.json" not in bad_allowed
    assert "visual_production_modules.json" not in bad_allowed
    assert "aesthetic_memory_v2.json" in bad_forbidden
    assert "visual_production_modules.json" in bad_forbidden
    assert 'if (!["run2_5_full_skill", "bad_aesthetic_memory"].includes(arm.armId)) return [];' in body
    assert 'const productionEligible = ["run2_5_full_skill", "bad_aesthetic_memory"].includes(arm.armId);' in body
    assert 'const fullProduction = arm.armId === "run2_5_full_skill";' in body
    assert re.search(r"production_reference_ids:\s*productionEligible\s*\?", body)
    assert re.search(r"aesthetic_memory_v2_ids:\s*fullProduction\s*\?", body)
    assert re.search(r"visual_production_module_ids:\s*fullProduction\s*\?", body)
    assert "fallback_policy:" in body
    assert "visual_validation_probe:" in body
    assert_contains(
        body,
        [
            "module_cinematic_cover_field",
            "module_editorial_before_after_delta",
            "module_proof_route_choreography",
            "module_system_mini_preview",
            "module_climax_hero_object",
            "module_release_handoff_gate",
        ],
    )


def test_run2_6_generator_consumes_workflow_policy_and_preserves_control_boundaries() -> None:
    body = (ROOT / "scripts" / "generate_ppt_run2_6_arms.mjs").read_text(encoding="utf-8")
    arm_order = ["prompt_only", "run1_5_skill", "run2_6_full_skill", "bad_aesthetic_memory"]

    def arm_block(arm_id: str) -> str:
        start = body.index(f'armId: "{arm_id}"')
        next_starts = [body.find(f'armId: "{next_arm}"', start + 1) for next_arm in arm_order]
        next_starts = [index for index in next_starts if index > start]
        end = min(next_starts) if next_starts else body.index("function sequenceStepsForSlide", start)
        return body[start:end]

    def section(block: str, start_marker: str, end_marker: str) -> str:
        start = block.index(start_marker)
        end = block.index(end_marker, start)
        return block[start:end]

    restricted_workflow_inputs = [
        "commercial_usecase_bank.json",
        "aesthetic_benchmark_bank.json",
        "workflow_decision_policy.json",
    ]

    assert_contains(
        body,
        [
            "prompt_only",
            "run1_5_skill",
            "run2_6_full_skill",
            "bad_aesthetic_memory",
            "commercial_usecase_bank.json",
            "aesthetic_benchmark_bank.json",
            "workflow_decision_policy.json",
            "source_brand_sanitization",
            "no_cross_arm_reuse",
        ],
    )
    prompt_allowed = section(arm_block("prompt_only"), "allowed:", "forbidden:")
    prompt_forbidden = section(arm_block("prompt_only"), "forbidden:", "palette:")
    run1_allowed = section(arm_block("run1_5_skill"), "allowed:", "forbidden:")
    run1_forbidden = section(arm_block("run1_5_skill"), "forbidden:", "palette:")
    full_allowed = section(arm_block("run2_6_full_skill"), "allowed:", "forbidden:")
    full_forbidden = section(arm_block("run2_6_full_skill"), "forbidden:", "palette:")
    bad_allowed = section(arm_block("bad_aesthetic_memory"), "allowed:", "forbidden:")
    bad_forbidden = section(arm_block("bad_aesthetic_memory"), "forbidden:", "palette:")

    for term in restricted_workflow_inputs:
        assert term not in prompt_allowed
        assert term in prompt_forbidden
        assert term not in run1_allowed
        assert term in run1_forbidden
        assert term in full_allowed
        assert term not in full_forbidden
    assert "commercial_usecase_bank.json" in bad_allowed
    assert "aesthetic_benchmark_bank.json" not in bad_allowed
    assert "workflow_decision_policy.json" not in bad_allowed
    assert "aesthetic_benchmark_bank.json" in bad_forbidden
    assert "workflow_decision_policy.json" in bad_forbidden
    assert "visual_repair_policy.json" not in full_allowed
    assert "run2_6r_visual_repair_full_skill" not in body
    assert 'const workflowEligible = ["run2_6_full_skill", "bad_aesthetic_memory"].includes(arm.armId);' in body
    assert 'const fullWorkflow = arm.armId === "run2_6_full_skill";' in body
    assert re.search(r"commercial_usecase_id:\s*workflowEligible\s*\?", body)
    assert re.search(r"aesthetic_benchmark_ids:\s*fullWorkflow\s*\?", body)
    assert re.search(r"theme_policy_id:\s*fullWorkflow\s*\?", body)
    assert re.search(r"typography_system_id:\s*fullWorkflow\s*\?", body)
    assert re.search(r"spacing_token_set_id:\s*fullWorkflow\s*\?", body)
    assert re.search(r"workflow_decision_ids:\s*fullWorkflow\s*\?", body)
    assert "source_brand_sanitization:" in body
    assert "visual_repair_policy_ids: []" in body
    assert 'visual_delta_from_run2_5: "not-applicable; existing Run 2.6 generator does not apply visual repair policy"' in body
    assert (
        'visual_repair_validation_probe: "not-applicable; existing Run 2.6 generator does not select visual repair"'
        in body
    )


def test_run2_6r_generator_consumes_visual_repair_policy_and_preserves_boundaries() -> None:
    body = (ROOT / "scripts" / "generate_ppt_run2_6r_visual_repair_arms.mjs").read_text(encoding="utf-8")
    arm_order = ["prompt_only", "run1_5_skill", "run2_6r_visual_repair_full_skill", "bad_aesthetic_memory"]

    def arm_block(arm_id: str) -> str:
        start = body.index(f'armId: "{arm_id}"')
        next_starts = [body.find(f'armId: "{next_arm}"', start + 1) for next_arm in arm_order]
        next_starts = [index for index in next_starts if index > start]
        end = min(next_starts) if next_starts else body.index("function sequenceStepsForSlide", start)
        return body[start:end]

    def section(block: str, start_marker: str, end_marker: str) -> str:
        start = block.index(start_marker)
        end = block.index(end_marker, start)
        return block[start:end]

    assert_contains(
        body,
        [
            "run2_6r_visual_repair_full_skill",
            "visual_repair_policy.json",
            "visual_repair_policy_ids",
            "visual_delta_from_run2_5",
            "visual_repair_validation_probe",
            "renderFullRepair",
            "drawEditorialClimaxSpread",
            "no_cross_arm_reuse",
        ],
    )
    assert "function renderFull(" not in body
    assert re.search(
        r'if \(arm\.armId === "run2_6r_visual_repair_full_skill"\) \{\s*'
        r"renderFullRepair\(slide, spec, arm, n\);\s*"
        r"\} else \{\s*"
        r"renderControl\(slide, spec, arm\);",
        body,
    )
    repair_start = body.index("function renderFullRepair")
    repair_end = body.index("function renderSlide", repair_start)
    repair_body = body[repair_start:repair_end]
    assert re.search(
        r'if \(spec\.role === "climax"\) \{\s*'
        r"drawEditorialClimaxSpread\(slide, arm\);\s*"
        r"return;\s*"
        r"\}\s*"
        r"simpleTitle\(slide, spec, arm, true\);",
        repair_body,
    )
    prompt_allowed = section(arm_block("prompt_only"), "allowed:", "forbidden:")
    prompt_forbidden = section(arm_block("prompt_only"), "forbidden:", "palette:")
    run1_allowed = section(arm_block("run1_5_skill"), "allowed:", "forbidden:")
    run1_forbidden = section(arm_block("run1_5_skill"), "forbidden:", "palette:")
    full_allowed = section(arm_block("run2_6r_visual_repair_full_skill"), "allowed:", "forbidden:")
    full_forbidden = section(arm_block("run2_6r_visual_repair_full_skill"), "forbidden:", "palette:")
    bad_allowed = section(arm_block("bad_aesthetic_memory"), "allowed:", "forbidden:")
    bad_forbidden = section(arm_block("bad_aesthetic_memory"), "forbidden:", "palette:")

    assert "visual_repair_policy.json" not in prompt_allowed
    assert "visual_repair_policy.json" in prompt_forbidden
    assert "visual_repair_policy.json" not in run1_allowed
    assert "visual_repair_policy.json" in run1_forbidden
    assert "visual_repair_policy.json" in full_allowed
    assert "visual_repair_policy.json" not in full_forbidden
    assert "commercial_usecase_bank.json" in bad_allowed
    assert "visual_repair_policy.json" not in bad_allowed
    assert "visual_repair_policy.json" in bad_forbidden
    assert 'const repairEligible = arm.armId === "run2_6r_visual_repair_full_skill";' in body
    assert re.search(r"visual_repair_policy_ids:\s*repairEligible\s*\?", body)
    assert re.search(r"visual_delta_from_run2_5:\s*repairEligible\s*\?", body)
    assert re.search(r"visual_repair_validation_probe:\s*repairEligible\s*\?", body)


def test_run2_7_generator_consumes_design_memory_and_preserves_control_boundaries() -> None:
    body = (ROOT / "scripts" / "generate_ppt_run2_7_data_workflow_arms.mjs").read_text(encoding="utf-8")
    arm_order = ["prompt_only", "run1_5_skill", "run2_7_full_skill", "bad_workflow_memory"]

    def arm_block(arm_id: str) -> str:
        start = body.index(f'armId: "{arm_id}"')
        next_starts = [body.find(f'armId: "{next_arm}"', start + 1) for next_arm in arm_order]
        next_starts = [index for index in next_starts if index > start]
        end = min(next_starts) if next_starts else body.index("function sequenceStepsForSlide", start)
        return body[start:end]

    def section(block: str, start_marker: str, end_marker: str) -> str:
        start = block.index(start_marker)
        end = block.index(end_marker, start)
        return block[start:end]

    restricted_run2_7_inputs = [
        "run2_7_commercial_usecase.json",
        "run2_7_multimodal_source_records.json",
        "run2_7_design_memory.json",
        "run2_7_workflow_policy.json",
    ]

    assert_contains(
        body,
        [
            "prompt_only",
            "run1_5_skill",
            "run2_7_full_skill",
            "bad_workflow_memory",
            "run2_7_commercial_usecase.json",
            "run2_7_multimodal_source_records.json",
            "run2_7_design_memory.json",
            "run2_7_workflow_policy.json",
            "renderRun27Full",
            "drawRun27Climax",
            "assertRun27MemorySelfCheck",
            "run27VisualSelfCheck",
            "run27Eligible",
            "run2_7_design_memory_ids",
            "run2_7_delta_from_run2_6r",
            "no_cross_arm_reuse",
        ],
    )
    prompt_allowed = section(arm_block("prompt_only"), "allowed:", "forbidden:")
    prompt_forbidden = section(arm_block("prompt_only"), "forbidden:", "palette:")
    run1_allowed = section(arm_block("run1_5_skill"), "allowed:", "forbidden:")
    run1_forbidden = section(arm_block("run1_5_skill"), "forbidden:", "palette:")
    full_allowed = section(arm_block("run2_7_full_skill"), "allowed:", "forbidden:")
    full_forbidden = section(arm_block("run2_7_full_skill"), "forbidden:", "palette:")
    bad_allowed = section(arm_block("bad_workflow_memory"), "allowed:", "forbidden:")
    bad_forbidden = section(arm_block("bad_workflow_memory"), "forbidden:", "palette:")

    for term in restricted_run2_7_inputs:
        assert term not in prompt_allowed
        assert term in prompt_forbidden
        assert term not in run1_allowed
        assert term in run1_forbidden
        assert term in full_allowed
        assert term not in full_forbidden

    assert "run2_7_commercial_usecase.json" in bad_allowed
    assert "run2_7_multimodal_source_records.json" in bad_allowed
    assert "run2_7_design_memory.json" not in bad_allowed
    assert "run2_7_workflow_policy.json" not in bad_allowed
    assert "run2_7_design_memory.json" in bad_forbidden
    assert "run2_7_workflow_policy.json" in bad_forbidden
    assert re.search(r"run2_7_usecase_id:\s*run27Eligible\s*\?", body)
    assert re.search(r"run2_7_source_record_ids:\s*run27Eligible\s*\?", body)
    assert re.search(r"run2_7_design_memory_ids:\s*fullRun27\s*\?", body)
    assert re.search(r"run2_7_workflow_decision_ids:\s*fullRun27\s*\?", body)
    assert re.search(r"run2_7_delta_from_run2_6r:\s*fullRun27\s*\?", body)
    assert re.search(r"run2_7_quality_gate:\s*fullRun27\s*\?", body)
    assert (PACK / "generation_briefs" / "bad_aesthetic_memory.md").exists()
    assert "bad_workflow_memory.md" not in body
    assert 'if (armId === "bad_workflow_memory") return "bad_aesthetic_memory";' in body
    assert re.search(r"generation_brief:\s*`\$\{pack\}/generation_briefs/\$\{generationBriefArmId\(arm\.armId\)\}\.md`", body)


def test_run2_8_generator_uses_executable_memory_and_preserves_boundaries() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_8_executable_memory_arms.mjs"
    assert script_path.exists(), "missing Run 2.8 executable-memory generator"
    body = script_path.read_text(encoding="utf-8")
    arm_order = ["prompt_only", "run1_5_skill", "run2_8_full_skill", "bad_memory_schema"]

    def arm_block(arm_id: str) -> str:
        start = body.index(f'armId: "{arm_id}"')
        next_starts = [body.find(f'armId: "{next_arm}"', start + 1) for next_arm in arm_order]
        next_starts = [index for index in next_starts if index > start]
        end = min(next_starts) if next_starts else len(body)
        return body[start:end]

    def section(block: str, start_marker: str, end_marker: str) -> str:
        start = block.index(start_marker)
        end = block.index(end_marker, start)
        return block[start:end]

    restricted_run2_8_inputs = [
        "run2_8_tutorial_decomposition.json",
        "run2_8_executable_design_memory.json",
        "run2_8_workflow_gate_matrix.json",
    ]

    assert_contains(
        body,
        [
            "prompt_only",
            "run1_5_skill",
            "run2_8_full_skill",
            "bad_memory_schema",
            "run2_8_tutorial_decomposition.json",
            "run2_8_executable_design_memory.json",
            "run2_8_workflow_gate_matrix.json",
            "renderRun28Full",
            "drawRun28Climax",
            "run28MemoryBindingsByRole",
            "assertRun28GateMatrixSelfCheck",
            "run2_8_memory_binding_ids",
        ],
    )
    prompt_allowed = section(arm_block("prompt_only"), "allowed:", "forbidden:")
    prompt_forbidden = section(arm_block("prompt_only"), "forbidden:", "palette:")
    run1_allowed = section(arm_block("run1_5_skill"), "allowed:", "forbidden:")
    run1_forbidden = section(arm_block("run1_5_skill"), "forbidden:", "palette:")
    full_allowed = section(arm_block("run2_8_full_skill"), "allowed:", "forbidden:")
    full_forbidden = section(arm_block("run2_8_full_skill"), "forbidden:", "palette:")
    bad_allowed = section(arm_block("bad_memory_schema"), "allowed:", "forbidden:")
    bad_forbidden = section(arm_block("bad_memory_schema"), "forbidden:", "palette:")

    for term in restricted_run2_8_inputs:
        assert term not in prompt_allowed
        assert term in prompt_forbidden
        assert term not in run1_allowed
        assert term in run1_forbidden
        assert term in full_allowed
        assert term not in full_forbidden

    assert "run2_8_tutorial_decomposition.json" in bad_allowed
    assert "run2_8_executable_design_memory.json" not in bad_allowed
    assert "run2_8_workflow_gate_matrix.json" not in bad_allowed
    assert "run2_8_executable_design_memory.json" in bad_forbidden
    assert "run2_8_workflow_gate_matrix.json" in bad_forbidden
    assert 'const fullRun28 = arm.armId === "run2_8_full_skill";' in body
    for field in EXPECTED_RUN2_8_TRACE_FIELDS:
        assert re.search(fr"{field}:\s*fullRun28\s*\?", body), field
    assert not re.search(r"run2_8_memory_binding_ids:\s*bad", body)
    for function_name in [
        "drawRun28TypeScale",
        "drawRun28SpacingZones",
        "drawRun28BeforeAfterDelta",
        "drawRun28WorkflowGate",
        "drawRun28Climax",
    ]:
        assert f'registerNativeModule(metrics, "{function_name}")' in body
    assert "const actualCodeBindingIds = Array.from(roleMetrics.codeBindingIds);" in body
    assert "run2_8_code_binding_ids: fullRun28" in body
    assert "actualCodeBindingIds" in body
    assert "gate objects ${metrics.gateObjects}" in body
    assert "metrics.workflowObjects > layoutBudget.max_gate_objects" not in body
    assert "full_arm_contract_preview_not_rendered" in body
    assert "full_arm_native_generator_rendered" in body
    render_start = body.index("function renderRun28FullSlide")
    contrast_start = body.index('} else if (spec.role === "contrast")', render_start)
    proof_start = body.index('} else if (spec.role === "proof")', render_start)
    climax_start = body.index('} else if (spec.role === "climax")', render_start)
    close_start = body.index('} else {', climax_start)
    contrast_block = body[contrast_start:proof_start]
    proof_block = body[proof_start:climax_start]
    climax_block = body[climax_start:close_start]
    assert "drawRun28SpacingZones" in contrast_block
    assert "drawRun28BeforeAfterDelta" in contrast_block
    assert "drawRun28BeforeAfterDelta" in proof_block
    assert "drawRun28SpacingZones" in proof_block
    assert "drawRun28WorkflowGate" in proof_block
    assert "drawRun28Climax(slide, heroBinding, arm, spec, metrics, { typeBinding })" in climax_block


def test_run2_9_generator_uses_visual_primitive_modules_and_preserves_boundaries() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_9_visual_primitive_arms.mjs"
    assert script_path.exists(), "missing Run 2.9 visual-primitive generator"
    body = script_path.read_text(encoding="utf-8")
    arm_order = ["prompt_only", "run1_5_skill", "run2_9_full_skill", "bad_visual_primitive_memory"]

    def arm_block(arm_id: str) -> str:
        start = body.index(f'armId: "{arm_id}"')
        next_starts = [body.find(f'armId: "{next_arm}"', start + 1) for next_arm in arm_order]
        next_starts = [index for index in next_starts if index > start]
        end = min(next_starts) if next_starts else len(body)
        return body[start:end]

    def section(block: str, start_marker: str, end_marker: str) -> str:
        start = block.index(start_marker)
        end = block.index(end_marker, start)
        return block[start:end]

    restricted_run2_9_inputs = [
        "run2_9_visual_primitive_repair.json",
        "run2_9_executable_visual_modules.json",
        "run2_9_visual_gate_matrix.json",
    ]

    assert_contains(
        body,
        [
            "prompt_only",
            "run1_5_skill",
            "run2_9_full_skill",
            "bad_visual_primitive_memory",
            "run2_9_visual_primitive_repair.json",
            "run2_9_executable_visual_modules.json",
            "run2_9_visual_gate_matrix.json",
            "drawRun29EditorialSpread",
            "drawRun29LayeredProductSurface",
            "drawRun29MotionStoryboard",
            "drawRun29ClimaxStage",
            "drawRun29TypographicField",
            "run29VisualModulesByRole",
            "assertRun29VisualGateSelfCheck",
            "registerVisualModule",
            "run2_9_code_module_ids",
        ],
    )
    prompt_allowed = section(arm_block("prompt_only"), "allowed:", "forbidden:")
    prompt_forbidden = section(arm_block("prompt_only"), "forbidden:", "palette:")
    run1_allowed = section(arm_block("run1_5_skill"), "allowed:", "forbidden:")
    run1_forbidden = section(arm_block("run1_5_skill"), "forbidden:", "palette:")
    full_allowed = section(arm_block("run2_9_full_skill"), "allowed:", "forbidden:")
    full_forbidden = section(arm_block("run2_9_full_skill"), "forbidden:", "palette:")
    bad_allowed = section(arm_block("bad_visual_primitive_memory"), "allowed:", "forbidden:")
    bad_forbidden = section(arm_block("bad_visual_primitive_memory"), "forbidden:", "palette:")

    for term in restricted_run2_9_inputs:
        assert term not in prompt_allowed
        assert term in prompt_forbidden
        assert term not in run1_allowed
        assert term in run1_forbidden
        assert term in full_allowed
        assert term not in full_forbidden

    assert "run2_9_visual_primitive_repair.json" in bad_allowed
    assert "run2_9_executable_visual_modules.json" not in bad_allowed
    assert "run2_9_visual_gate_matrix.json" not in bad_allowed
    assert "run2_9_executable_visual_modules.json" in bad_forbidden
    assert "run2_9_visual_gate_matrix.json" in bad_forbidden
    assert 'const fullRun29 = arm.armId === "run2_9_full_skill";' in body
    for field in EXPECTED_RUN2_9_TRACE_FIELDS:
        assert re.search(fr"{field}:\s*fullRun29\s*\?", body), field
    assert "const actualCodeModuleIds = Array.from(roleMetrics.visualModuleIds);" in body
    assert "run2_9_code_module_ids: fullRun29" in body
    for function_name in [
        "drawRun29EditorialSpread",
        "drawRun29LayeredProductSurface",
        "drawRun29MotionStoryboard",
        "drawRun29ClimaxStage",
        "drawRun29TypographicField",
    ]:
        assert f'registerVisualModule(metrics, "{function_name}")' in body
    render_start = body.index("function renderRun29FullSlide")
    setup_start = body.index('} else if (spec.role === "setup")', render_start)
    contrast_start = body.index('} else if (spec.role === "contrast")', render_start)
    proof_start = body.index('} else if (spec.role === "proof")', render_start)
    climax_start = body.index('} else if (spec.role === "climax")', render_start)
    close_start = body.index('} else {', climax_start)
    cover_block = body[render_start:setup_start]
    setup_block = body[setup_start:contrast_start]
    proof_block = body[proof_start:climax_start]
    climax_block = body[climax_start:close_start]
    assert "drawRun29TypographicField" in cover_block
    assert "drawRun29EditorialSpread" in cover_block
    assert "drawRun29LayeredProductSurface" in setup_block
    assert "drawRun29MotionStoryboard" in proof_block
    assert "drawRun29ClimaxStage" in climax_block


def test_run2_10_generator_uses_visual_system_modules_and_preserves_boundaries() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_10_visual_system_arms.mjs"
    assert script_path.exists(), "missing Run 2.10 visual-system generator"
    body = script_path.read_text(encoding="utf-8")
    arm_order = ["prompt_only", "run1_5_skill", "run2_10_full_skill", "bad_visual_system_memory"]

    def arm_block(arm_id: str) -> str:
        start = body.index(f'armId: "{arm_id}"')
        next_starts = [body.find(f'armId: "{next_arm}"', start + 1) for next_arm in arm_order]
        next_starts = [index for index in next_starts if index > start]
        end = min(next_starts) if next_starts else len(body)
        return body[start:end]

    def section(block: str, start_marker: str, end_marker: str) -> str:
        start = block.index(start_marker)
        end = block.index(end_marker, start)
        return block[start:end]

    restricted_run2_10_inputs = [
        "run2_10_visual_system_sources.json",
        "run2_10_visual_system_memory.json",
        "run2_10_visual_system_gate_matrix.json",
    ]

    assert_contains(
        body,
        [
            "prompt_only",
            "run1_5_skill",
            "run2_10_full_skill",
            "bad_visual_system_memory",
            "drawRun210FullBleedVisualField",
            "drawRun210ProductTheater",
            "drawRun210KineticSequence",
            "drawRun210EditorialTypeSystem",
            "drawRun210NonRectangularProofPath",
            "drawRun210CinematicClimax",
            "run210VisualSystemsByRole",
            "assertRun210VisualSystemGateSelfCheck",
            "registerRun210VisualModule",
            "run2_10_code_module_ids",
            "run2_10_shape_count_budget",
            "run2_10_asymmetry_whitespace_rule",
        ],
    )

    prompt_allowed = section(arm_block("prompt_only"), "allowed:", "forbidden:")
    prompt_forbidden = section(arm_block("prompt_only"), "forbidden:", "palette:")
    run1_allowed = section(arm_block("run1_5_skill"), "allowed:", "forbidden:")
    run1_forbidden = section(arm_block("run1_5_skill"), "forbidden:", "palette:")
    full_allowed = section(arm_block("run2_10_full_skill"), "allowed:", "forbidden:")
    full_forbidden = section(arm_block("run2_10_full_skill"), "forbidden:", "palette:")
    bad_allowed = section(arm_block("bad_visual_system_memory"), "allowed:", "forbidden:")
    bad_forbidden = section(arm_block("bad_visual_system_memory"), "forbidden:", "palette:")

    for term in restricted_run2_10_inputs:
        assert term not in prompt_allowed
        assert term in prompt_forbidden
        assert term not in run1_allowed
        assert term in run1_forbidden
        assert term in full_allowed
        assert term not in full_forbidden

    assert "run2_10_visual_system_sources.json" in bad_allowed
    assert "run2_10_visual_system_memory.json" not in bad_allowed
    assert "run2_10_visual_system_gate_matrix.json" not in bad_allowed
    assert "run2_10_visual_system_memory.json" in bad_forbidden
    assert "run2_10_visual_system_gate_matrix.json" in bad_forbidden
    assert 'const fullRun210 = arm.armId === "run2_10_full_skill";' in body
    for field in EXPECTED_RUN2_10_TRACE_FIELDS:
        assert re.search(fr"{field}:\s*fullRun210\s*\?", body), field
    assert "const actualCodeModuleIds = Array.from(roleMetrics.visualModuleIds);" in body
    assert "run2_10_code_module_ids: fullRun210" in body
    for function_name in [
        "drawRun210FullBleedVisualField",
        "drawRun210ProductTheater",
        "drawRun210KineticSequence",
        "drawRun210EditorialTypeSystem",
        "drawRun210NonRectangularProofPath",
        "drawRun210CinematicClimax",
    ]:
        assert f'registerRun210VisualModule(metrics, "{function_name}")' in body
    render_start = body.index("function renderRun210FullSlide")
    setup_start = body.index('} else if (spec.role === "setup")', render_start)
    contrast_start = body.index('} else if (spec.role === "contrast")', render_start)
    proof_start = body.index('} else if (spec.role === "proof")', render_start)
    climax_start = body.index('} else if (spec.role === "climax")', render_start)
    close_start = body.index('} else {', climax_start)
    cover_block = body[render_start:setup_start]
    setup_block = body[setup_start:contrast_start]
    contrast_block = body[contrast_start:proof_start]
    proof_block = body[proof_start:climax_start]
    climax_block = body[climax_start:close_start]
    assert "drawRun210EditorialTypeSystem" in cover_block
    assert "drawRun210FullBleedVisualField" in cover_block
    assert "drawRun210ProductTheater" in setup_block
    assert "drawRun210NonRectangularProofPath" in setup_block
    assert "drawRun210FullBleedVisualField" in contrast_block
    assert "drawRun210NonRectangularProofPath" in contrast_block
    assert "drawRun210ProductTheater" in proof_block
    assert "drawRun210KineticSequence" in proof_block
    assert "drawRun210CinematicClimax" in climax_block
    assert "drawRun210KineticSequence" in climax_block


def test_run2_13_generator_uses_thick_data_seeds_and_preserves_boundaries() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_13_thick_data_arms.mjs"
    assert script_path.exists(), "missing Run 2.13 thick-data generator"
    body = script_path.read_text(encoding="utf-8")
    arm_order = ["prompt_only", "run1_5_skill", "run2_13_full_skill", "bad_thick_data_memory"]

    def arm_block(arm_id: str) -> str:
        start = body.index(f'armId: "{arm_id}"')
        next_starts = [body.find(f'armId: "{next_arm}"', start + 1) for next_arm in arm_order]
        next_starts = [index for index in next_starts if index > start]
        end = min(next_starts) if next_starts else len(body)
        return body[start:end]

    def section(block: str, start_marker: str, end_marker: str) -> str:
        start = block.index(start_marker)
        end = block.index(end_marker, start)
        return block[start:end]

    restricted_run2_12_inputs = [
        "run2_12_thick_multimodal_evidence.json",
        "run2_12_design_memory_seed.json",
        "run2_12_workflow_gate_seed.json",
    ]

    assert_contains(
        body,
        [
            "prompt_only",
            "run1_5_skill",
            "run2_13_full_skill",
            "bad_thick_data_memory",
            "drawRun213LaunchArcRoute",
            "drawRun213TypeWhitespaceSystem",
            "drawRun213ProductDemoSequence",
            "drawRun213MetricClimaxObject",
            "drawRun213WorkflowGateRail",
            "run213MemorySeedsByRole",
            "assertRun213ThickDataGateSelfCheck",
            "registerRun213Module",
            "run2_12_code_module_ids",
            "run2_12_density_whitespace_gate",
            "run2_12_demo_sequence_gate",
            "run2_12_climax_object_gate",
        ],
    )

    prompt_allowed = section(arm_block("prompt_only"), "allowed:", "forbidden:")
    prompt_forbidden = section(arm_block("prompt_only"), "forbidden:", "palette:")
    run1_allowed = section(arm_block("run1_5_skill"), "allowed:", "forbidden:")
    run1_forbidden = section(arm_block("run1_5_skill"), "forbidden:", "palette:")
    full_allowed = section(arm_block("run2_13_full_skill"), "allowed:", "forbidden:")
    full_forbidden = section(arm_block("run2_13_full_skill"), "forbidden:", "palette:")
    bad_allowed = section(arm_block("bad_thick_data_memory"), "allowed:", "forbidden:")
    bad_forbidden = section(arm_block("bad_thick_data_memory"), "forbidden:", "palette:")

    for term in restricted_run2_12_inputs:
        assert term not in prompt_allowed
        assert term in prompt_forbidden
        assert term not in run1_allowed
        assert term in run1_forbidden
        assert term in full_allowed
        assert term not in full_forbidden

    assert "run2_12_thick_multimodal_evidence.json" in bad_allowed
    assert "run2_12_design_memory_seed.json" not in bad_allowed
    assert "run2_12_workflow_gate_seed.json" not in bad_allowed
    assert "skill_workflow.json" not in bad_allowed
    assert "run2_12_design_memory_seed.json" in bad_forbidden
    assert "run2_12_workflow_gate_seed.json" in bad_forbidden
    assert "skill_workflow.json" in bad_forbidden
    assert 'const fullRun213 = arm.armId === "run2_13_full_skill";' in body
    for field in EXPECTED_RUN2_13_TRACE_FIELDS:
        assert re.search(fr"{field}:\s*fullRun213\s*\?", body), field
    assert "const actualCodeModuleIds = Array.from(roleMetrics.visualModuleIds);" in body
    assert "run2_12_code_module_ids: fullRun213" in body
    for function_name in [
        "drawRun213LaunchArcRoute",
        "drawRun213TypeWhitespaceSystem",
        "drawRun213ProductDemoSequence",
        "drawRun213MetricClimaxObject",
        "drawRun213WorkflowGateRail",
    ]:
        assert f'registerRun213Module(metrics, "{function_name}")' in body
    render_start = body.index("function renderRun213FullSlide")
    setup_start = body.index('} else if (spec.role === "setup")', render_start)
    contrast_start = body.index('} else if (spec.role === "contrast")', render_start)
    proof_start = body.index('} else if (spec.role === "proof")', render_start)
    climax_start = body.index('} else if (spec.role === "climax")', render_start)
    close_start = body.index('} else {', climax_start)
    cover_block = body[render_start:setup_start]
    setup_block = body[setup_start:contrast_start]
    proof_block = body[proof_start:climax_start]
    climax_block = body[climax_start:close_start]
    close_block = body[close_start:]
    assert "drawRun213LaunchArcRoute" in cover_block
    assert "drawRun213TypeWhitespaceSystem" in cover_block
    assert "drawRun213LaunchArcRoute" in setup_block
    assert "drawRun213WorkflowGateRail" in setup_block
    assert "drawRun213ProductDemoSequence" in proof_block
    assert "drawRun213WorkflowGateRail" in proof_block
    assert "drawRun213MetricClimaxObject" in climax_block
    assert "drawRun213ProductDemoSequence" in climax_block
    assert "drawRun213WorkflowGateRail" in close_block


def test_run2_14_generator_hides_trace_inside_210_aesthetic_shell() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_14_aesthetic_trace_arms.mjs"
    assert script_path.exists(), "missing Run 2.14 aesthetic-trace generator"
    body = script_path.read_text(encoding="utf-8")
    arm_order = ["prompt_only", "run1_5_skill", "run2_14_full_skill", "bad_visible_workflow_memory"]

    def arm_block(arm_id: str) -> str:
        start = body.index(f'armId: "{arm_id}"')
        next_starts = [body.find(f'armId: "{next_arm}"', start + 1) for next_arm in arm_order]
        next_starts = [index for index in next_starts if index > start]
        end = min(next_starts) if next_starts else len(body)
        return body[start:end]

    def section(block: str, start_marker: str, end_marker: str) -> str:
        start = block.index(start_marker)
        end = block.index(end_marker, start)
        return block[start:end]

    run2_12_inputs = [
        "run2_12_thick_multimodal_evidence.json",
        "run2_12_design_memory_seed.json",
        "run2_12_workflow_gate_seed.json",
    ]
    run2_10_aesthetic_inputs = [
        "run2_10_visual_system_sources.json",
        "run2_10_visual_system_memory.json",
        "run2_10_visual_system_gate_matrix.json",
    ]

    assert_contains(
        body,
        [
            "prompt_only",
            "run1_5_skill",
            "run2_14_full_skill",
            "bad_visible_workflow_memory",
            "drawRun214PresentationShell",
            "drawRun214CinematicProductTheater",
            "drawRun214HiddenWorkflowRoute",
            "drawRun214MetricRevealStage",
            "drawRun214QuietReleaseHandoff",
            "assertRun214AestheticTraceGateSelfCheck",
            "registerRun214Module",
            "manifest_only_trace_public_surface",
            "visible_workflow_suppression",
            "run2_14_code_module_ids",
        ],
    )

    prompt_allowed = section(arm_block("prompt_only"), "allowed:", "forbidden:")
    prompt_forbidden = section(arm_block("prompt_only"), "forbidden:", "palette:")
    run1_allowed = section(arm_block("run1_5_skill"), "allowed:", "forbidden:")
    run1_forbidden = section(arm_block("run1_5_skill"), "forbidden:", "palette:")
    full_allowed = section(arm_block("run2_14_full_skill"), "allowed:", "forbidden:")
    full_forbidden = section(arm_block("run2_14_full_skill"), "forbidden:", "palette:")
    bad_allowed = section(arm_block("bad_visible_workflow_memory"), "allowed:", "forbidden:")
    bad_forbidden = section(arm_block("bad_visible_workflow_memory"), "forbidden:", "palette:")

    for term in [*run2_12_inputs, *run2_10_aesthetic_inputs]:
        assert term not in prompt_allowed
        assert term in prompt_forbidden
        assert term not in run1_allowed
        assert term in run1_forbidden

    for term in [*run2_12_inputs, *run2_10_aesthetic_inputs]:
        assert term in full_allowed
        assert term not in full_forbidden

    for term in run2_12_inputs:
        assert term in bad_allowed
        assert term not in bad_forbidden
    for term in run2_10_aesthetic_inputs:
        assert term not in bad_allowed
        assert term in bad_forbidden

    assert 'const fullRun214 = arm.armId === "run2_14_full_skill";' in body
    for field in EXPECTED_RUN2_14_TRACE_FIELDS:
        assert re.search(fr"{field}:\s*fullRun214\s*\?", body), field
    assert "const actualCodeModuleIds = Array.from(roleMetrics.visualModuleIds);" in body
    assert "run2_14_code_module_ids: fullRun214" in body
    for function_name in [
        "drawRun214PresentationShell",
        "drawRun214CinematicProductTheater",
        "drawRun214HiddenWorkflowRoute",
        "drawRun214MetricRevealStage",
        "drawRun214QuietReleaseHandoff",
    ]:
        assert f'registerRun214Module(metrics, "{function_name}")' in body
    render_start = body.index("function renderRun214FullSlide")
    setup_start = body.index('} else if (spec.role === "setup")', render_start)
    contrast_start = body.index('} else if (spec.role === "contrast")', render_start)
    proof_start = body.index('} else if (spec.role === "proof")', render_start)
    climax_start = body.index('} else if (spec.role === "climax")', render_start)
    close_start = body.index('} else {', climax_start)
    cover_block = body[render_start:setup_start]
    setup_block = body[setup_start:contrast_start]
    proof_block = body[proof_start:climax_start]
    climax_block = body[climax_start:close_start]
    close_block = body[close_start:]
    assert "drawRun214PresentationShell" in cover_block
    assert "drawRun214HiddenWorkflowRoute" in cover_block
    assert "drawRun214CinematicProductTheater" in setup_block
    assert "drawRun214HiddenWorkflowRoute" in setup_block
    assert "drawRun214CinematicProductTheater" in proof_block
    assert "drawRun214HiddenWorkflowRoute" in proof_block
    assert "drawRun214MetricRevealStage" in climax_block
    assert "drawRun214QuietReleaseHandoff" in close_block


def test_run2_15_layout_selector_artifacts_define_reusable_workflow() -> None:
    sources = load_json(PACK / "run2_15_layout_selector_sources.json")
    memory = load_json(PACK / "run2_15_layout_module_memory.json")
    gates = load_json(PACK / "run2_15_layout_selector_gate_matrix.json")
    trace_contract = load_json(PACK / "results" / "trace_manifest_contract.json")

    assert sources["status"] == "run2_15_layout_selector_sources_ready"
    assert memory["status"] == "run2_15_layout_module_memory_ready"
    assert gates["status"] == "run2_15_layout_selector_gate_matrix_ready"

    source_records = sources["records"]
    memory_records = memory["modules"]
    gate_records = gates["gates"]

    assert len(source_records) >= 6
    assert len(memory_records) >= 6
    assert len(gate_records) >= 5

    required_source_fields = {
        "record_id",
        "source_family",
        "derived_from_run_ids",
        "modality_mix",
        "commercial_need",
        "design_observation",
        "layout_selector_obligation",
        "typography_obligation",
        "spacing_obligation",
        "product_theater_obligation",
        "motion_beat_obligation",
        "trace_visibility_obligation",
        "anti_copy_boundary",
    }
    for record in source_records:
        assert required_source_fields <= record.keys()
        assert record["anti_copy_boundary"]
        assert record["layout_selector_obligation"]
        assert not record.get("raw_media_path")
        assert not record.get("copied_source_layout")

    source_ids = {record["record_id"] for record in source_records}
    required_module_families = {
        "editorial_cover_field",
        "product_theater_stage",
        "before_after_route",
        "metric_reveal_stage",
        "quiet_release_handoff",
        "dense_evidence_compression",
    }
    module_families = {module["module_family"] for module in memory_records}
    assert required_module_families <= module_families

    required_memory_fields = {
        "module_id",
        "module_family",
        "slide_roles",
        "source_record_ids",
        "selection_trigger",
        "composition_contract",
        "typography_contract",
        "spacing_contract",
        "asset_contract",
        "trace_visibility_contract",
        "fallback_contract",
        "native_ppt_obligations",
        "forbidden_patterns",
    }
    for module in memory_records:
        assert required_memory_fields <= module.keys()
        assert set(module["source_record_ids"]) <= source_ids
        assert module["slide_roles"]
        assert module["trace_visibility_contract"] == "manifest_viewer_qa_only_for_public_surface"
        assert module["native_ppt_obligations"]

    module_ids = {module["module_id"] for module in memory_records}
    required_gate_ids = {
        "gate_2_15_role_to_module_selection",
        "gate_2_15_text_resilience",
        "gate_2_15_trace_hidden_from_surface",
        "gate_2_15_product_theater_realism",
        "gate_2_15_metric_reveal_climax",
        "gate_2_15_bad_selector_control_boundary",
    }
    gate_ids = {gate["gate_id"] for gate in gate_records}
    assert required_gate_ids <= gate_ids
    trace_fields = set(trace_contract["per_slide_required_fields"])
    assert EXPECTED_RUN2_15_TRACE_FIELDS <= trace_fields

    for gate in gate_records:
        assert set(gate["candidate_module_ids"]) <= module_ids
        assert gate["required_selector_inputs"]
        assert gate["selection_rules"]
        assert gate["rejection_rules"]
        assert gate["trace_fields"]
        assert set(gate["trace_fields"]) <= trace_fields
        assert "run2_15_selected_layout_module_ids" in gate["trace_fields"]
        assert gate["text_resilience_probe"]
        assert gate["bad_control_probe"]


def test_run2_16_generator_uses_run2_15_selector_before_native_ppt_code() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_16_selector_rerun_arms.mjs"
    assert script_path.exists(), "missing Run 2.16 selector-driven rerun generator"
    body = script_path.read_text(encoding="utf-8")
    arm_order = ["prompt_only", "run1_5_skill", "run2_16_full_skill", "bad_selector_memory"]

    def arm_block(arm_id: str) -> str:
        start = body.index(f'armId: "{arm_id}"')
        next_starts = [body.find(f'armId: "{next_arm}"', start + 1) for next_arm in arm_order]
        next_starts = [index for index in next_starts if index > start]
        end = min(next_starts) if next_starts else len(body)
        return body[start:end]

    def section(block: str, start_marker: str, end_marker: str) -> str:
        start = block.index(start_marker)
        end = block.index(end_marker, start)
        return block[start:end]

    run2_15_selector_inputs = [
        "run2_15_layout_selector_sources.json",
        "run2_15_layout_module_memory.json",
        "run2_15_layout_selector_gate_matrix.json",
    ]
    run2_14_aesthetic_inputs = [
        "run2_10_visual_system_sources.json",
        "run2_10_visual_system_memory.json",
        "run2_10_visual_system_gate_matrix.json",
    ]

    assert_contains(
        body,
        [
            "prompt_only",
            "run1_5_skill",
            "run2_16_full_skill",
            "bad_selector_memory",
            "validateRun215SelectorSchemas",
            "assertArmInputBoundaries",
            "readRun216SelectorJsonForArm",
            "selectRun215LayoutModulesForSlide",
            "assertRun216SelectorGateSelfCheck",
            "drawRun216EditorialCoverField",
            "drawRun216ProductTheaterStage",
            "drawRun216BeforeAfterRoute",
            "drawRun216MetricRevealStage",
            "drawRun216QuietReleaseHandoff",
            "drawRun216DenseEvidenceCompression",
            "manifest_viewer_qa_only_for_public_surface",
            "selector_gate_matrix_executed_before_native_ppt_generation",
        ],
    )

    prompt_allowed = section(arm_block("prompt_only"), "allowed:", "forbidden:")
    prompt_forbidden = section(arm_block("prompt_only"), "forbidden:", "palette:")
    run1_allowed = section(arm_block("run1_5_skill"), "allowed:", "forbidden:")
    run1_forbidden = section(arm_block("run1_5_skill"), "forbidden:", "palette:")
    full_allowed = section(arm_block("run2_16_full_skill"), "allowed:", "forbidden:")
    full_forbidden = section(arm_block("run2_16_full_skill"), "forbidden:", "palette:")
    bad_allowed = section(arm_block("bad_selector_memory"), "allowed:", "forbidden:")
    bad_forbidden = section(arm_block("bad_selector_memory"), "forbidden:", "palette:")

    for term in [*run2_15_selector_inputs, *run2_14_aesthetic_inputs]:
        assert term not in prompt_allowed
        assert term in prompt_forbidden
        assert term not in run1_allowed
        assert term in run1_forbidden
        assert term in full_allowed
        assert term not in full_forbidden

    assert "run2_15_layout_selector_sources.json" in bad_allowed
    assert "run2_15_layout_module_memory.json" not in bad_allowed
    assert "run2_15_layout_selector_gate_matrix.json" not in bad_allowed
    assert "run2_15_layout_module_memory.json" in bad_forbidden
    assert "run2_15_layout_selector_gate_matrix.json" in bad_forbidden

    assert 'const fullRun216 = arm.armId === "run2_16_full_skill";' in body
    for field in EXPECTED_RUN2_15_TRACE_FIELDS:
        assert re.search(fr"{field}:\s*fullRun216\s*\?", body), field
    for function_name in [
        "drawRun216EditorialCoverField",
        "drawRun216ProductTheaterStage",
        "drawRun216BeforeAfterRoute",
        "drawRun216MetricRevealStage",
        "drawRun216QuietReleaseHandoff",
        "drawRun216DenseEvidenceCompression",
    ]:
        assert f'registerRun216Module(metrics, "{function_name}")' in body
    assert "throw new Error(`${arm.armId} input boundary does not permit reading ${relPath}`)" in body
    assert "run2_15_layout_module_memory.json" in arm_block("bad_selector_memory")
    assert "run2_15_layout_selector_gate_matrix.json" in arm_block("bad_selector_memory")


def test_run2_16_selector_runtime_guards_block_bad_arm_and_select_gates() -> None:
    node_script = """
import fs from "node:fs";
import path from "node:path";
const mod = await import("./scripts/generate_ppt_run2_16_selector_rerun_arms.mjs");
const badArm = mod.armSpecs.find((arm) => arm.armId === "bad_selector_memory");
const fullArm = mod.armSpecs.find((arm) => arm.armId === "run2_16_full_skill");
let blocked = false;
try {
  mod.readRun216SelectorJsonForArm(badArm, mod.RUN2_16_INPUTS.selectorGateMatrix);
} catch (error) {
  blocked = String(error.message).includes("input boundary does not permit reading");
}
if (!blocked) throw new Error("bad selector arm was able to read selector gate matrix");
const root = process.cwd();
const sources = JSON.parse(fs.readFileSync(path.join(root, mod.RUN2_16_INPUTS.selectorSources), "utf8"));
const memory = JSON.parse(fs.readFileSync(path.join(root, mod.RUN2_16_INPUTS.selectorMemory), "utf8"));
const gates = JSON.parse(fs.readFileSync(path.join(root, mod.RUN2_16_INPUTS.selectorGateMatrix), "utf8"));
mod.validateRun215SelectorSchemas(sources, memory, gates);
const climax = mod.selectRun215LayoutModulesForSlide("climax", memory, gates);
const moduleIds = climax.modules.map((item) => item.module_id);
const gateIds = climax.gates.map((item) => item.gate_id);
if (!moduleIds.includes("module_2_15_metric_reveal_stage")) throw new Error("climax did not select metric reveal module");
if (!gateIds.includes("gate_2_15_metric_reveal_climax")) throw new Error("climax did not consume metric reveal gate");
if (!fullArm.allowed.includes(mod.RUN2_16_INPUTS.selectorGateMatrix)) throw new Error("full arm does not allow selector gate matrix");
"""
    completed = subprocess.run(
        ["node", "--input-type=module", "-e", node_script],
        cwd=ROOT,
        capture_output=True,
        text=True,
        timeout=20,
    )

    assert completed.returncode == 0, completed.stderr


def test_run2_19_generator_consumes_run2_18_thickness_pack_before_native_ppt_code() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_19_thickness_rerun_arms.mjs"
    assert script_path.exists(), "missing Run 2.19 thickness rerun generator"
    body = script_path.read_text(encoding="utf-8")
    arm_order = ["prompt_only", "run1_5_skill", "run2_19_full_skill", "bad_thickness_memory"]

    def arm_block(arm_id: str) -> str:
        start = body.index(f'armId: "{arm_id}"')
        next_starts = [body.find(f'armId: "{next_arm}"', start + 1) for next_arm in arm_order]
        next_starts = [index for index in next_starts if index > start]
        end = min(next_starts) if next_starts else len(body)
        return body[start:end]

    def section(block: str, start_marker: str, end_marker: str) -> str:
        start = block.index(start_marker)
        end = block.index(end_marker, start)
        return block[start:end]

    run2_18_inputs = [
        "run2_18_multimodal_evidence_expansion.json",
        "run2_18_design_memory_expansion.json",
        "run2_18_workflow_gate_expansion.json",
    ]

    assert_contains(
        body,
        [
            "prompt_only",
            "run1_5_skill",
            "run2_19_full_skill",
            "bad_thickness_memory",
            "validateRun218ThicknessSchemas",
            "assertArmInputBoundaries",
            "readRun219ThicknessJsonForArm",
            "selectRun218ThicknessForSlide",
            "assertRun219ThicknessGateSelfCheck",
            "drawRun219LaunchIdentityField",
            "drawRun219ProductSurfaceTheater",
            "drawRun219DemoSequenceRoute",
            "drawRun219MetricClimaxObject",
            "drawRun219DensityGate",
            "thickness_pack_executed_before_native_ppt_generation",
        ],
    )

    prompt_allowed = section(arm_block("prompt_only"), "allowed:", "forbidden:")
    prompt_forbidden = section(arm_block("prompt_only"), "forbidden:", "palette:")
    run1_allowed = section(arm_block("run1_5_skill"), "allowed:", "forbidden:")
    run1_forbidden = section(arm_block("run1_5_skill"), "forbidden:", "palette:")
    full_allowed = section(arm_block("run2_19_full_skill"), "allowed:", "forbidden:")
    full_forbidden = section(arm_block("run2_19_full_skill"), "forbidden:", "palette:")
    bad_allowed = section(arm_block("bad_thickness_memory"), "allowed:", "forbidden:")
    bad_forbidden = section(arm_block("bad_thickness_memory"), "forbidden:", "palette:")

    for term in run2_18_inputs:
        assert term not in prompt_allowed
        assert term in prompt_forbidden
        assert term not in run1_allowed
        assert term in run1_forbidden
        assert term in full_allowed
        assert term not in full_forbidden

    assert "run2_18_multimodal_evidence_expansion.json" in bad_allowed
    assert "run2_18_design_memory_expansion.json" not in bad_allowed
    assert "run2_18_workflow_gate_expansion.json" not in bad_allowed
    assert "run2_18_design_memory_expansion.json" in bad_forbidden
    assert "run2_18_workflow_gate_expansion.json" in bad_forbidden

    assert 'const fullRun219 = arm.armId === "run2_19_full_skill";' in body
    for field in EXPECTED_RUN2_19_TRACE_FIELDS:
        assert re.search(fr"{field}:\s*fullRun219\s*\?", body), field
    for function_name in [
        "drawRun219LaunchIdentityField",
        "drawRun219ProductSurfaceTheater",
        "drawRun219DemoSequenceRoute",
        "drawRun219MetricClimaxObject",
        "drawRun219DensityGate",
    ]:
        assert f'registerRun219Module(metrics, "{function_name}")' in body


def test_run2_19_thickness_runtime_guards_block_bad_arm_and_select_gates() -> None:
    node_script = """
import fs from "node:fs";
import path from "node:path";
const mod = await import("./scripts/generate_ppt_run2_19_thickness_rerun_arms.mjs");
const badArm = mod.armSpecs.find((arm) => arm.armId === "bad_thickness_memory");
const fullArm = mod.armSpecs.find((arm) => arm.armId === "run2_19_full_skill");
let blocked = false;
try {
  mod.readRun219ThicknessJsonForArm(badArm, mod.RUN2_19_INPUTS.workflowGates);
} catch (error) {
  blocked = String(error.message).includes("input boundary does not permit reading");
}
if (!blocked) throw new Error("bad thickness arm was able to read workflow gates");
const root = process.cwd();
const evidence = JSON.parse(fs.readFileSync(path.join(root, mod.RUN2_19_INPUTS.evidence), "utf8"));
const memory = JSON.parse(fs.readFileSync(path.join(root, mod.RUN2_19_INPUTS.memory), "utf8"));
const gates = JSON.parse(fs.readFileSync(path.join(root, mod.RUN2_19_INPUTS.workflowGates), "utf8"));
mod.validateRun218ThicknessSchemas(evidence, memory, gates);
const climax = mod.selectRun218ThicknessForSlide("climax", evidence, memory, gates);
const memoryIds = climax.memories.map((item) => item.memory_id);
const gateIds = climax.gates.map((item) => item.gate_id);
if (!memoryIds.includes("memory_2_18_metric_climax_object")) throw new Error("climax did not select metric climax memory");
if (!gateIds.includes("gate_2_18_metric_climax_selection")) throw new Error("climax did not consume metric climax gate");
if (!fullArm.allowed.includes(mod.RUN2_19_INPUTS.workflowGates)) throw new Error("full arm does not allow Run 2.18 workflow gates");
"""
    completed = subprocess.run(
        ["node", "--input-type=module", "-e", node_script],
        cwd=ROOT,
        capture_output=True,
        text=True,
        timeout=20,
    )

    assert completed.returncode == 0, completed.stderr


def test_run2_19_records_thickness_rerun_result() -> None:
    result = (PACK / "results" / "run2_19_thickness_rerun_result.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_19_thickness_rerun_result.json")

    assert result_json["status"] == "rerun_completed_public_blocked"
    assert result_json["public_ready"] is False
    assert result_json["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert result_json["run_id"] == "2.19"
    assert result_json["rerun"]["best_internal_arm"] == "run2_19_full_skill"
    assert result_json["rerun"]["best_internal_arm_verdict"] == "thickness_pack_executed_before_native_ppt_generation"
    assert result_json["input_chain"]["evidence"] == "docs/product/ppt-run2-data-skill-quality/run2_18_multimodal_evidence_expansion.json"
    assert result_json["input_chain"]["memory"] == "docs/product/ppt-run2-data-skill-quality/run2_18_design_memory_expansion.json"
    assert result_json["input_chain"]["workflow_gates"] == "docs/product/ppt-run2-data-skill-quality/run2_18_workflow_gate_expansion.json"
    assert result_json["control_boundary"]["bad_thickness_memory"].startswith("evidence_only")
    assert result_json["rerun"]["combined_contact_sheet"].endswith("run2-19-four-arm-contact-sheet.png")
    assert result_json["rerun"]["full_skill_series_sheet"].endswith("run2-full-skill-series-horizontal.png")
    assert result_json["rerun"]["generated_outputs_committed"] is False
    assert_contains(
        result,
        [
            "Run 2.19",
            "Run 2.18 thickness pack",
            "four-arm rerun",
            "bad_thickness_memory",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_19_generated_rerun() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "Run 2.19",
            "ppt-run2-19-prompt-only",
            "ppt-run2-19-run1-5-skill",
            "ppt-run2-19-full-vulca",
            "ppt-run2-19-bad-thickness-memory",
            "run2_19_thickness_rerun_result.json",
        ],
    )


def test_run2_20_trace_effectiveness_audit_script_computes_2_19_data_use(tmp_path: Path) -> None:
    script_path = ROOT / "scripts" / "audit_ppt_run2_20_trace_effectiveness.py"
    assert script_path.exists(), "missing Run 2.20 trace effectiveness audit script"

    result_json = tmp_path / "run2_20_trace_effectiveness_audit.json"
    result_md = tmp_path / "run2_20_trace_effectiveness_audit.md"
    completed = subprocess.run(
        [
            sys.executable,
            str(script_path),
            "--result-json",
            str(result_json),
            "--result-md",
            str(result_md),
        ],
        cwd=ROOT,
        capture_output=True,
        text=True,
        timeout=20,
    )

    assert completed.returncode == 0, completed.stderr
    audit = load_json(result_json)
    report = result_md.read_text(encoding="utf-8")

    assert audit["status"] == "run2_20_trace_effectiveness_audit_public_blocked"
    assert audit["run_id"] == "2.20"
    assert audit["source_generated_run"] == "2.19"
    assert audit["creates_new_ppt_deck"] is False
    assert audit["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert audit["public_ready"] is False

    coverage = audit["trace_effectiveness"]
    assert coverage["full_arm"]["slide_count"] == 6
    assert coverage["full_arm"]["evidence_records_selected"] == coverage["inventory"]["run2_18_evidence_records"] == 8
    assert coverage["full_arm"]["memory_records_selected"] == coverage["inventory"]["run2_18_memory_records"] == 6
    assert coverage["full_arm"]["workflow_gates_selected"] == coverage["inventory"]["run2_18_workflow_gates"] == 6
    assert coverage["full_arm"]["all_slides_have_memory_gate_and_code"] is True
    assert coverage["full_arm"]["layout_budget_passed_all_slides"] is True
    assert coverage["full_arm"]["selected_code_module_count"] >= 5
    assert not coverage["full_arm"]["unused_evidence_record_ids"]
    assert not coverage["full_arm"]["unused_memory_ids"]
    assert not coverage["full_arm"]["unused_gate_ids"]

    bad = audit["control_boundary"]["bad_thickness_memory"]
    assert bad["uses_evidence_only"] is True
    assert bad["selected_evidence_slide_count"] == 6
    assert bad["selected_memory_ids"] == []
    assert bad["selected_workflow_gate_ids"] == []
    assert bad["selected_code_module_ids"] == []
    assert audit["gate_summary"]["data_workflow_trace_effectiveness_gate"] == "pass_internal_only"
    assert audit["gate_summary"]["public_release_gate"] == "blocked"
    assert "Run 2.20" in report
    assert "trace effectiveness" in report
    assert "Run 2.19" in report
    assert "public blocked" in report


def test_run2_20_records_trace_effectiveness_audit_result() -> None:
    result = (PACK / "results" / "run2_20_trace_effectiveness_audit.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_20_trace_effectiveness_audit.json")

    assert result_json["status"] == "run2_20_trace_effectiveness_audit_public_blocked"
    assert result_json["source_generated_run"] == "2.19"
    assert result_json["creates_new_ppt_deck"] is False
    assert result_json["trace_effectiveness"]["full_arm"]["all_slides_have_memory_gate_and_code"] is True
    assert result_json["control_boundary"]["bad_thickness_memory"]["uses_evidence_only"] is True
    assert result_json["next_required_action"].startswith("thicken_visual_decision")
    assert_contains(
        result,
        [
            "Run 2.20",
            "trace effectiveness",
            "Run 2.19",
            "Run 2.18 thickness pack",
            "bad_thickness_memory",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_embeds_run2_20_trace_effectiveness_audit() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "run2_20_trace_effectiveness_audit.json",
            "Run 2.20 trace effectiveness audit",
            "trace effectiveness",
            "bad_thickness_memory",
        ],
    )


def test_run2_21_builder_writes_visual_decision_memory_artifacts(tmp_path: Path) -> None:
    script_path = ROOT / "scripts" / "build_ppt_run2_21_visual_decision_memory.py"
    assert script_path.exists(), "missing Run 2.21 visual-decision memory builder"

    result_json = tmp_path / "run2_21_visual_decision_memory_result.json"
    result_md = tmp_path / "run2_21_visual_decision_memory_result.md"
    completed = subprocess.run(
        [
            sys.executable,
            str(script_path),
            "--out-dir",
            str(tmp_path),
            "--result-json",
            str(result_json),
            "--result-md",
            str(result_md),
        ],
        cwd=ROOT,
        capture_output=True,
        text=True,
        timeout=20,
    )

    assert completed.returncode == 0, completed.stderr
    decision_memory = load_json(tmp_path / "run2_21_visual_decision_memory.json")
    selector_gates = load_json(tmp_path / "run2_21_per_role_selector_gates.json")
    rejection_matrix = load_json(tmp_path / "run2_21_evidence_rejection_matrix.json")
    result = load_json(result_json)

    assert decision_memory["status"] == "run2_21_visual_decision_memory_ready_public_blocked"
    assert decision_memory["creates_new_ppt_deck"] is False
    assert decision_memory["derived_from"] == [
        "run2_18_multimodal_evidence_expansion.json",
        "run2_18_design_memory_expansion.json",
        "run2_18_workflow_gate_expansion.json",
        "results/run2_20_trace_effectiveness_audit.json",
    ]
    assert len(decision_memory["visual_decision_memory"]) == 6
    all_evidence_ids = {record["record_id"] for record in load_json(PACK / "run2_18_multimodal_evidence_expansion.json")["records"]}
    for record in decision_memory["visual_decision_memory"]:
        assert record["role"] in {"cover", "setup", "contrast", "proof", "climax", "close"}
        assert record["primary_evidence_id"] in all_evidence_ids
        assert 1 <= len(record["secondary_evidence_ids"]) <= 2
        assert set(record["secondary_evidence_ids"]).issubset(all_evidence_ids)
        rejected_ids = {item["evidence_id"] for item in record["rejected_evidence"]}
        assert rejected_ids
        assert rejected_ids.isdisjoint({record["primary_evidence_id"], *record["secondary_evidence_ids"]})
        assert {record["primary_evidence_id"], *record["secondary_evidence_ids"], *rejected_ids} == all_evidence_ids
        assert all(item["reason"] for item in record["rejected_evidence"])
        assert record["selected_memory_ids"]
        assert record["selected_gate_ids"]
        for field in [
            "typography_decision",
            "spacing_decision",
            "composition_decision",
            "proof_object_decision",
            "code_generation_obligation",
            "visual_quality_risk",
            "source_boundary",
        ]:
            assert record[field], f"{record['decision_id']} missing {field}"

    assert selector_gates["status"] == "run2_21_per_role_selector_gates_ready_public_blocked"
    assert len(selector_gates["gates"]) == 6
    for gate in selector_gates["gates"]:
        assert gate["required_primary_evidence_count"] == 1
        assert gate["max_secondary_evidence_count"] == 2
        assert gate["required_visual_decision_memory_id"].startswith("vdm_2_21_")
        assert "run2_21_primary_evidence_id" in gate["required_trace_fields"]
        assert "run2_21_rejected_evidence_reasons" in gate["required_trace_fields"]
        assert gate["public_surface_policy"] == "trace_suppressed_from_public_slide_surface"
        assert gate["release_boundary"].startswith("public_blocked")

    assert rejection_matrix["status"] == "run2_21_evidence_rejection_matrix_ready_public_blocked"
    assert len(rejection_matrix["role_records"]) == 6
    for role_record in rejection_matrix["role_records"]:
        assert role_record["all_evidence_accounted_for"] is True
        assert role_record["primary_evidence_id"] in all_evidence_ids
        assert role_record["rejected_evidence"]
        assert all(item["reason"] for item in role_record["rejected_evidence"])

    assert result["status"] == "run2_21_visual_decision_memory_ready_public_blocked"
    assert result["creates_new_ppt_deck"] is False
    assert result["delivery_artifacts"] == {
        "pptx_paths": [],
        "rendered_slide_paths": [],
        "contact_sheet_paths": [],
        "html_motion_renderer_paths": [],
    }
    assert result["public_ready"] is False
    assert result["next_required_action"].startswith("consume_run2_21_visual_decision_memory")
    assert not list(tmp_path.glob("*.pptx"))


def test_run2_21_records_visual_decision_memory_result() -> None:
    result = (PACK / "results" / "run2_21_visual_decision_memory_result.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_21_visual_decision_memory_result.json")
    decision_memory = load_json(PACK / "run2_21_visual_decision_memory.json")
    selector_gates = load_json(PACK / "run2_21_per_role_selector_gates.json")
    rejection_matrix = load_json(PACK / "run2_21_evidence_rejection_matrix.json")

    assert result_json["status"] == "run2_21_visual_decision_memory_ready_public_blocked"
    assert result_json["source_audit_run"] == "2.20"
    assert result_json["creates_new_ppt_deck"] is False
    assert result_json["delivery_artifacts"]["pptx_paths"] == []
    assert result_json["delivery_artifacts"]["rendered_slide_paths"] == []
    assert result_json["artifact_counts"] == {
        "visual_decision_memory": 6,
        "per_role_selector_gates": 6,
        "evidence_rejection_records": 6,
    }
    assert len(decision_memory["visual_decision_memory"]) == 6
    assert len(selector_gates["gates"]) == 6
    assert len(rejection_matrix["role_records"]) == 6
    assert_contains(
        result,
        [
            "Run 2.21",
            "visual-decision memory",
            "per-role selector",
            "evidence rejection matrix",
            "Run 2.20",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_embeds_run2_21_visual_decision_memory() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "run2_21_visual_decision_memory.json",
            "run2_21_per_role_selector_gates.json",
            "run2_21_evidence_rejection_matrix.json",
            "run2_21_visual_decision_memory_result.json",
            "Run 2.21 visual-decision memory",
        ],
    )


def test_run2_22_generator_consumes_run2_21_selector_memory_before_native_ppt_code() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_22_selector_memory_arms.mjs"
    assert script_path.exists(), "missing Run 2.22 selector-memory rerun generator"
    body = script_path.read_text(encoding="utf-8")
    arm_order = ["prompt_only", "run1_5_skill", "run2_22_full_selector_memory", "bad_selector_memory"]

    def arm_block(arm_id: str) -> str:
        start = body.index(f'armId: "{arm_id}"')
        next_starts = [body.find(f'armId: "{next_arm}"', start + 1) for next_arm in arm_order]
        next_starts = [index for index in next_starts if index > start]
        end = min(next_starts) if next_starts else len(body)
        return body[start:end]

    def section(block: str, start_marker: str, end_marker: str) -> str:
        start = block.index(start_marker)
        end = block.index(end_marker, start)
        return block[start:end]

    run2_21_inputs = [
        "run2_21_visual_decision_memory.json",
        "run2_21_per_role_selector_gates.json",
        "run2_21_evidence_rejection_matrix.json",
    ]

    assert_contains(
        body,
        [
            "prompt_only",
            "run1_5_skill",
            "run2_22_full_selector_memory",
            "bad_selector_memory",
            "validateRun221SelectorSchemas",
            "assertArmInputBoundaries",
            "readRun222SelectorJsonForArm",
            "selectRun221ForSlide",
            "assertRun222SelectorGateSelfCheck",
            "drawRun222CinematicSelectorField",
            "drawRun222ProductTheaterSurface",
            "drawRun222EvidenceRoute",
            "drawRun222ClimaxEditorialStage",
            "selector_memory_executed_before_native_ppt_generation",
        ],
    )

    prompt_allowed = section(arm_block("prompt_only"), "allowed:", "forbidden:")
    prompt_forbidden = section(arm_block("prompt_only"), "forbidden:", "palette:")
    run1_allowed = section(arm_block("run1_5_skill"), "allowed:", "forbidden:")
    run1_forbidden = section(arm_block("run1_5_skill"), "forbidden:", "palette:")
    full_allowed = section(arm_block("run2_22_full_selector_memory"), "allowed:", "forbidden:")
    full_forbidden = section(arm_block("run2_22_full_selector_memory"), "forbidden:", "palette:")
    bad_allowed = section(arm_block("bad_selector_memory"), "allowed:", "forbidden:")
    bad_forbidden = section(arm_block("bad_selector_memory"), "forbidden:", "palette:")

    for term in run2_21_inputs:
        assert term not in prompt_allowed
        assert term in prompt_forbidden
        assert term not in run1_allowed
        assert term in run1_forbidden
        assert term in full_allowed
        assert term not in full_forbidden

    assert "run2_21_visual_decision_memory.json" in bad_allowed
    assert "run2_21_per_role_selector_gates.json" not in bad_allowed
    assert "run2_21_evidence_rejection_matrix.json" not in bad_allowed
    assert "run2_21_per_role_selector_gates.json" in bad_forbidden
    assert "run2_21_evidence_rejection_matrix.json" in bad_forbidden

    assert 'const fullRun222 = arm.armId === "run2_22_full_selector_memory";' in body
    for field in EXPECTED_RUN2_22_TRACE_FIELDS:
        assert re.search(fr"{field}:\s*fullRun222\s*\?", body), field
    for function_name in [
        "drawRun222CinematicSelectorField",
        "drawRun222ProductTheaterSurface",
        "drawRun222EvidenceRoute",
        "drawRun222ClimaxEditorialStage",
    ]:
        assert f'registerRun222Module(metrics, "{function_name}")' in body


def test_run2_22_selector_runtime_guards_block_bad_arm_and_select_role_memory() -> None:
    node_script = """
import fs from "node:fs";
import path from "node:path";
const mod = await import("./scripts/generate_ppt_run2_22_selector_memory_arms.mjs");
const badArm = mod.armSpecs.find((arm) => arm.armId === "bad_selector_memory");
const fullArm = mod.armSpecs.find((arm) => arm.armId === "run2_22_full_selector_memory");
let blocked = false;
try {
  mod.readRun222SelectorJsonForArm(badArm, mod.RUN2_22_INPUTS.selectorGates);
} catch (error) {
  blocked = String(error.message).includes("input boundary does not permit reading");
}
if (!blocked) throw new Error("bad selector arm was able to read Run 2.21 selector gates");
let rejectionBlocked = false;
try {
  mod.readRun222SelectorJsonForArm(badArm, mod.RUN2_22_INPUTS.rejectionMatrix);
} catch (error) {
  rejectionBlocked = String(error.message).includes("input boundary does not permit reading");
}
if (!rejectionBlocked) throw new Error("bad selector arm was able to read Run 2.21 rejection matrix");
const badTrace = mod.traceFor(badArm);
for (const slide of badTrace.slides) {
  if (slide.run2_21_selector_gate_id) throw new Error("bad selector trace leaked selector gate id");
  if ((slide.run2_21_rejected_evidence_reasons || []).length) throw new Error("bad selector trace leaked rejection reasons");
  if ((slide.run2_22_code_module_ids || []).length) throw new Error("bad selector trace leaked Run 2.22 code modules");
}
const root = process.cwd();
const decisionMemory = JSON.parse(fs.readFileSync(path.join(root, mod.RUN2_22_INPUTS.decisionMemory), "utf8"));
const selectorGates = JSON.parse(fs.readFileSync(path.join(root, mod.RUN2_22_INPUTS.selectorGates), "utf8"));
const rejectionMatrix = JSON.parse(fs.readFileSync(path.join(root, mod.RUN2_22_INPUTS.rejectionMatrix), "utf8"));
mod.validateRun221SelectorSchemas(decisionMemory, selectorGates, rejectionMatrix);
const cover = mod.selectRun221ForSlide("cover", decisionMemory, selectorGates, rejectionMatrix);
if (cover.decision.decision_id !== "vdm_2_21_cover") throw new Error("cover did not select the Run 2.21 cover memory");
if (cover.gate.gate_id !== "gate_2_21_cover_selector") throw new Error("cover did not select the Run 2.21 cover selector gate");
if (cover.decision.primary_evidence_id !== cover.rejection.primary_evidence_id) throw new Error("cover primary evidence mismatch");
if (!fullArm.allowed.includes(mod.RUN2_22_INPUTS.rejectionMatrix)) throw new Error("full arm does not allow Run 2.21 rejection matrix");
"""
    completed = subprocess.run(
        ["node", "--input-type=module", "-e", node_script],
        cwd=ROOT,
        capture_output=True,
        text=True,
        timeout=20,
    )

    assert completed.returncode == 0, completed.stderr


def test_run2_22_records_selector_rerun_result() -> None:
    result = (PACK / "results" / "run2_22_selector_rerun_result.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_22_selector_rerun_result.json")

    assert result_json["status"] == "rerun_completed_public_blocked"
    assert result_json["public_ready"] is False
    assert result_json["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert result_json["run_id"] == "2.22"
    assert result_json["rerun"]["best_internal_arm"] == "run2_22_full_selector_memory"
    assert result_json["rerun"]["best_internal_arm_verdict"] == "selector_memory_executed_before_native_ppt_generation"
    assert result_json["input_chain"]["visual_decision_memory"] == "docs/product/ppt-run2-data-skill-quality/run2_21_visual_decision_memory.json"
    assert result_json["input_chain"]["selector_gates"] == "docs/product/ppt-run2-data-skill-quality/run2_21_per_role_selector_gates.json"
    assert result_json["input_chain"]["rejection_matrix"] == "docs/product/ppt-run2-data-skill-quality/run2_21_evidence_rejection_matrix.json"
    assert result_json["control_boundary"]["bad_selector_memory"].startswith("decision_memory_only")
    assert result_json["remaining_public_release_gates"] == [
        "human_visual_review",
        "native_or_cross_platform_render_inspection",
        "motion_or_video_review",
        "source_boundary_review",
        "human_release_approval",
    ]
    assert result_json["rerun"]["combined_contact_sheet"].endswith("run2-22-four-arm-contact-sheet.png")
    assert result_json["rerun"]["full_skill_series_sheet"].endswith("run2-full-skill-series-horizontal.png")
    assert result_json["rerun"]["generated_outputs_committed"] is False
    assert_contains(
        result,
        [
            "Run 2.22",
            "Run 2.21 visual-decision memory",
            "four-arm rerun",
            "bad_selector_memory",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_22_generated_rerun() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "Run 2.22",
            "ppt-run2-22-prompt-only",
            "ppt-run2-22-run1-5-skill",
            "ppt-run2-22-full-vulca",
            "ppt-run2-22-bad-selector-memory",
            "run2_22_selector_rerun_result.json",
        ],
    )


def test_run2_23_selector_effectiveness_audit_compares_2_19_and_2_22(tmp_path: Path) -> None:
    script_path = ROOT / "scripts" / "audit_ppt_run2_23_selector_effectiveness.py"
    assert script_path.exists(), "missing Run 2.23 selector effectiveness audit script"

    result_json = tmp_path / "run2_23_selector_effectiveness_audit.json"
    result_md = tmp_path / "run2_23_selector_effectiveness_audit.md"
    completed = subprocess.run(
        [
            sys.executable,
            str(script_path),
            "--result-json",
            str(result_json),
            "--result-md",
            str(result_md),
        ],
        cwd=ROOT,
        capture_output=True,
        text=True,
        timeout=20,
    )

    assert completed.returncode == 0, completed.stderr
    audit = load_json(result_json)
    report = result_md.read_text(encoding="utf-8")

    assert audit["status"] == "run2_23_selector_effectiveness_audit_public_blocked"
    assert audit["run_id"] == "2.23"
    assert audit["source_generated_run"] == "2.22"
    assert audit["comparison_baseline_run"] == "2.19"
    assert audit["source_selector_layer"] == "2.21"
    assert audit["creates_new_ppt_deck"] is False
    assert audit["public_ready"] is False
    assert audit["stage_policy"] == "repeat_same_five_layers_not_run3"

    selector = audit["selector_effectiveness"]
    assert selector["full_arm"]["arm_id"] == "run2_22_full_selector_memory"
    assert selector["full_arm"]["slide_count"] == 6
    assert selector["full_arm"]["visual_decision_memory_records_selected"] == 6
    assert selector["full_arm"]["selector_gates_selected"] == 6
    assert selector["full_arm"]["all_slides_have_selector_gate_and_code"] is True
    assert selector["full_arm"]["primary_evidence_per_slide_all"] is True
    assert selector["full_arm"]["secondary_evidence_within_cap_all"] is True
    assert selector["full_arm"]["rejection_reasons_present_all"] is True
    assert selector["full_arm"]["public_surface_policy_suppressed_all"] is True
    assert selector["comparison_to_run2_19"]["roles_with_code_module_delta"] >= 5
    assert len(selector["full_arm"]["role_records"]) == 6

    bad = audit["control_boundary"]["bad_selector_memory"]
    assert bad["decision_memory_only"] is True
    assert bad["blocks_selector_gates"] is True
    assert bad["blocks_rejection_matrix"] is True
    assert bad["selected_selector_gate_ids"] == []
    assert bad["selected_rejected_evidence_reasons"] == []
    assert bad["selected_code_module_ids"] == []

    assert len(audit["chain_records"]) == 6
    assert audit["gate_summary"]["selector_memory_gate"] == "pass_internal_only"
    assert audit["gate_summary"]["public_release_gate"] == "blocked"
    assert_contains(
        report,
        [
            "Run 2.23",
            "selector effectiveness",
            "Run 2.22",
            "Run 2.19",
            "bad_selector_memory",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_run2_23_records_selector_effectiveness_audit_result() -> None:
    result = (PACK / "results" / "run2_23_selector_effectiveness_audit.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_23_selector_effectiveness_audit.json")

    assert result_json["status"] == "run2_23_selector_effectiveness_audit_public_blocked"
    assert result_json["source_generated_run"] == "2.22"
    assert result_json["comparison_baseline_run"] == "2.19"
    assert result_json["creates_new_ppt_deck"] is False
    assert result_json["selector_effectiveness"]["full_arm"]["all_slides_have_selector_gate_and_code"] is True
    assert result_json["control_boundary"]["bad_selector_memory"]["decision_memory_only"] is True
    assert result_json["next_required_action"].startswith("human_review_run2_22_visual_delta")
    assert_contains(
        result,
        [
            "Run 2.23",
            "selector effectiveness",
            "Run 2.22",
            "Run 2.19",
            "bad_selector_memory",
            "public blocked",
        ],
    )


def test_ppt_run_html_viewer_embeds_run2_23_selector_effectiveness_audit() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "run2_23_selector_effectiveness_audit.json",
            "Run 2.23 selector effectiveness audit",
            "selector effectiveness",
            "bad_selector_memory",
        ],
    )


def test_run2_24_builder_writes_single_usecase_content_visual_evidence_pack(tmp_path: Path) -> None:
    script_path = ROOT / "scripts" / "build_ppt_run2_24_single_usecase_content_evidence.py"
    assert script_path.exists(), "missing Run 2.24 single-usecase content/evidence builder"

    result_json = tmp_path / "run2_24_single_usecase_thickening_result.json"
    result_md = tmp_path / "run2_24_single_usecase_thickening_result.md"
    completed = subprocess.run(
        [
            sys.executable,
            str(script_path),
            "--out-dir",
            str(tmp_path),
            "--result-json",
            str(result_json),
            "--result-md",
            str(result_md),
        ],
        cwd=ROOT,
        capture_output=True,
        text=True,
        timeout=20,
    )

    assert completed.returncode == 0, completed.stderr
    content_memory = load_json(tmp_path / "run2_24_single_usecase_content_memory.json")
    visual_assets = load_json(tmp_path / "run2_24_visual_evidence_asset_memory.json")
    workflow_gates = load_json(tmp_path / "run2_24_content_visual_workflow_gates.json")
    result = load_json(result_json)
    report = result_md.read_text(encoding="utf-8")

    assert result["status"] == "run2_24_single_usecase_content_visual_evidence_ready_public_blocked"
    assert result["run_id"] == "2.24"
    assert result["source_audit_run"] == "2.23"
    assert result["selected_usecase_id"] == "usecase_design_to_production_platform_launch"
    assert result["creates_new_ppt_deck"] is False
    assert result["public_ready"] is False
    assert result["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert result["next_required_action"].startswith("consume_run2_24_single_usecase_pack")
    assert result["delivery_artifacts"] == {
        "pptx_paths": [],
        "rendered_slide_paths": [],
        "contact_sheet_paths": [],
        "html_motion_renderer_paths": [],
    }

    assert content_memory["status"] == "run2_24_single_usecase_content_memory_ready_public_blocked"
    assert content_memory["selected_usecase"]["id"] == "usecase_design_to_production_platform_launch"
    assert content_memory["story_policy"]["single_primary_usecase"] is True
    assert len(content_memory["slide_content_memory"]) == 6
    for record in content_memory["slide_content_memory"]:
        assert record["role"] in {"cover", "setup", "contrast", "proof", "climax", "close"}
        assert record["content_memory_id"].startswith("content_2_24_")
        assert record["headline"]
        assert record["support_line"]
        assert len(record["business_proof_points"]) >= 2
        assert len(record["visual_evidence_slot_ids"]) >= 2
        assert record["trace_fields_required"] == [
            "run2_24_selected_usecase_id",
            "run2_24_content_memory_id",
            "run2_24_visual_evidence_slot_ids",
            "run2_24_content_density_gate_id",
        ]
        assert "screenshots" in record["forbidden_source_materials"]
        assert "brand marks" in record["forbidden_source_materials"]

    assert visual_assets["status"] == "run2_24_visual_evidence_asset_memory_ready_public_blocked"
    assert visual_assets["storage_policy"]["raw_media"] == "forbidden"
    assert len(visual_assets["visual_evidence_assets"]) == 12
    for asset in visual_assets["visual_evidence_assets"]:
        assert asset["asset_id"].startswith("asset_2_24_")
        assert asset["role"] in {"cover", "setup", "contrast", "proof", "climax", "close"}
        assert asset["native_ppt_strategy"]
        assert asset["content_payload"]
        assert asset["source_boundary"].startswith("Derived")
        assert asset["public_surface_allowed"] is True

    assert workflow_gates["status"] == "run2_24_content_visual_workflow_gates_ready_public_blocked"
    assert len(workflow_gates["gates"]) == 6
    for gate in workflow_gates["gates"]:
        assert gate["selected_usecase_id"] == "usecase_design_to_production_platform_launch"
        assert gate["min_business_proof_points"] >= 2
        assert gate["min_visual_evidence_slots"] >= 2
        assert gate["forbid_cross_case_primary_story"] is True
        assert "run2_24_selected_usecase_id" in gate["required_trace_fields"]
        assert gate["public_release_gate"] == "blocked"

    assert_contains(
        report,
        [
            "Run 2.24",
            "single usecase",
            "content memory",
            "visual evidence",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )
    assert not list(tmp_path.glob("*.pptx"))


def test_run2_24_records_single_usecase_content_visual_evidence_pack() -> None:
    result = (PACK / "results" / "run2_24_single_usecase_thickening_result.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_24_single_usecase_thickening_result.json")
    content_memory = load_json(PACK / "run2_24_single_usecase_content_memory.json")
    visual_assets = load_json(PACK / "run2_24_visual_evidence_asset_memory.json")
    workflow_gates = load_json(PACK / "run2_24_content_visual_workflow_gates.json")
    workflow = load_json(PACK / "skill_workflow.json")

    assert result_json["status"] == "run2_24_single_usecase_content_visual_evidence_ready_public_blocked"
    assert result_json["selected_usecase_id"] == "usecase_design_to_production_platform_launch"
    assert result_json["creates_new_ppt_deck"] is False
    assert result_json["public_ready"] is False
    assert result_json["artifact_counts"] == {
        "slide_content_memory": 6,
        "visual_evidence_assets": 12,
        "content_visual_workflow_gates": 6,
    }
    assert content_memory["story_policy"]["single_primary_usecase"] is True
    assert len(content_memory["slide_content_memory"]) == 6
    assert len(visual_assets["visual_evidence_assets"]) == 12
    assert len(workflow_gates["gates"]) == 6
    assert workflow["status"] == "run2_53_product_surface_scene_workflow_directed_public_blocked"
    assert {stage["id"] for stage in workflow["stages"]} >= {
        "lock_run2_24_single_usecase_content_memory",
        "compile_run2_24_visual_evidence_asset_memory",
        "apply_run2_24_content_visual_workflow_gates",
        "compile_run2_35_visual_evidence_asset_realism_memory",
        "compile_run2_35_editorial_composition_memory",
        "apply_run2_35_visual_evidence_workflow_gates",
        "compile_run2_38_public_video_slide_direction_memory",
        "compile_run2_38_per_slide_visual_recipe_memory",
        "apply_run2_38_public_video_workflow_gates",
    }
    assert any(
        trigger["id"] == "run2_24_single_usecase_pack_required_before_next_rerun"
        for trigger in workflow["repair_triggers"]
    )
    assert any(
        trigger["id"] == "run2_35_visual_evidence_realism_required_before_run2_36_rerun"
        for trigger in workflow["repair_triggers"]
    )
    assert_contains(
        result,
        [
            "Run 2.24",
            "usecase_design_to_production_platform_launch",
            "content memory",
            "visual evidence",
            "public blocked",
        ],
    )


def test_ppt_run_html_viewer_embeds_run2_24_single_usecase_content_visual_evidence_pack() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "run2_24_single_usecase_content_memory.json",
            "run2_24_visual_evidence_asset_memory.json",
            "run2_24_content_visual_workflow_gates.json",
            "run2_24_single_usecase_thickening_result.json",
            "Run 2.24 single-usecase content + visual evidence",
        ],
    )


def test_run2_25_generator_consumes_run2_24_pack_before_native_ppt_code() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_25_single_usecase_arms.mjs"
    assert script_path.exists(), "missing Run 2.25 single-usecase rerun generator"
    body = script_path.read_text(encoding="utf-8")
    arm_order = [
        "prompt_only",
        "run1_5_skill",
        "run2_25_full_single_usecase_content_visual",
        "bad_content_visual_memory",
    ]

    def arm_block(arm_id: str) -> str:
        start = body.index(f'armId: "{arm_id}"')
        next_starts = [body.find(f'armId: "{next_arm}"', start + 1) for next_arm in arm_order]
        next_starts = [index for index in next_starts if index > start]
        end = min(next_starts) if next_starts else len(body)
        return body[start:end]

    def section(block: str, start_marker: str, end_marker: str) -> str:
        start = block.index(start_marker)
        end = block.index(end_marker, start)
        return block[start:end]

    run2_24_inputs = [
        "run2_24_single_usecase_content_memory.json",
        "run2_24_visual_evidence_asset_memory.json",
        "run2_24_content_visual_workflow_gates.json",
    ]

    assert_contains(
        body,
        [
            "validateRun224SingleUsecaseSchemas",
            "assertRun225ArmInputBoundaries",
            "readRun225PackJsonForArm",
            "selectRun224ForSlide",
            "assertRun225ContentVisualGateSelfCheck",
            "drawRun225LaunchField",
            "drawRun225SelectedRouteMap",
            "drawRun225ContentEvidenceSurface",
            "drawRun225ClimaxProofObject",
            "run2_24_pack_executed_before_native_ppt_generation",
        ],
    )

    prompt_allowed = section(arm_block("prompt_only"), "allowed:", "forbidden:")
    prompt_forbidden = section(arm_block("prompt_only"), "forbidden:", "palette:")
    run1_allowed = section(arm_block("run1_5_skill"), "allowed:", "forbidden:")
    run1_forbidden = section(arm_block("run1_5_skill"), "forbidden:", "palette:")
    full_allowed = section(arm_block("run2_25_full_single_usecase_content_visual"), "allowed:", "forbidden:")
    full_forbidden = section(arm_block("run2_25_full_single_usecase_content_visual"), "forbidden:", "palette:")
    bad_allowed = section(arm_block("bad_content_visual_memory"), "allowed:", "forbidden:")
    bad_forbidden = section(arm_block("bad_content_visual_memory"), "forbidden:", "palette:")

    for term in run2_24_inputs:
        assert term not in prompt_allowed
        assert term in prompt_forbidden
        assert term not in run1_allowed
        assert term in run1_forbidden
        assert term in full_allowed
        assert term not in full_forbidden
        assert term not in bad_allowed
        assert term in bad_forbidden

    assert 'const fullRun225 = arm.armId === "run2_25_full_single_usecase_content_visual";' in body
    for field in [
        "run2_24_selected_usecase_id",
        "run2_24_content_memory_id",
        "run2_24_visual_evidence_slot_ids",
        "run2_24_content_density_gate_id",
        "run2_25_content_visual_execution_status",
        "run2_25_code_module_ids",
    ]:
        assert re.search(fr"{field}:\s*fullRun225\s*\?", body), field
    for function_name in [
        "drawRun225LaunchField",
        "drawRun225SelectedRouteMap",
        "drawRun225ContentEvidenceSurface",
        "drawRun225ClimaxProofObject",
    ]:
        assert f'registerRun225Module(metrics, "{function_name}")' in body


def test_run2_25_runtime_guards_block_bad_arm_and_select_single_usecase_pack() -> None:
    node_script = """
const mod = await import("./scripts/generate_ppt_run2_25_single_usecase_arms.mjs");
const badArm = mod.armSpecs.find((arm) => arm.armId === "bad_content_visual_memory");
const fullArm = mod.armSpecs.find((arm) => arm.armId === "run2_25_full_single_usecase_content_visual");
for (const input of mod.RUN2_25_DATA_INPUTS) {
  let blocked = false;
  try {
    mod.readRun225PackJsonForArm(badArm, input);
  } catch (error) {
    blocked = String(error.message).includes("input boundary does not permit reading");
  }
  if (!blocked) throw new Error(`bad content/visual arm was able to read ${input}`);
}
const badTrace = mod.traceFor(badArm);
for (const slide of badTrace.slides) {
  if (slide.run2_24_content_memory_id) throw new Error("bad control trace leaked Run 2.24 content memory id");
  if ((slide.run2_24_visual_evidence_slot_ids || []).length) throw new Error("bad control trace leaked Run 2.24 visual evidence slots");
  if (slide.run2_24_content_density_gate_id) throw new Error("bad control trace leaked Run 2.24 density gate");
  if ((slide.run2_25_code_module_ids || []).length) throw new Error("bad control trace leaked Run 2.25 code modules");
}
const fullTrace = mod.traceFor(fullArm);
if (fullTrace.selected_usecase_id !== "usecase_design_to_production_platform_launch") throw new Error("full trace did not lock the selected usecase");
for (const slide of fullTrace.slides) {
  if (slide.run2_24_selected_usecase_id !== "usecase_design_to_production_platform_launch") throw new Error(`slide ${slide.slide_id} selected the wrong usecase`);
  if (!slide.run2_24_content_memory_id.startsWith("content_2_24_")) throw new Error(`slide ${slide.slide_id} missing content memory`);
  if ((slide.run2_24_visual_evidence_slot_ids || []).length < 2) throw new Error(`slide ${slide.slide_id} missing visual evidence slots`);
  if (!slide.run2_24_content_density_gate_id.startsWith("gate_2_24_")) throw new Error(`slide ${slide.slide_id} missing density gate`);
  if (slide.run2_25_content_visual_execution_status !== "run2_24_pack_executed_before_native_ppt_generation") throw new Error(`slide ${slide.slide_id} missing execution status`);
  if ((slide.run2_25_code_module_ids || []).length === 0) throw new Error(`slide ${slide.slide_id} missing code module ids`);
}
const selected = mod.selectRun224ForSlide("climax", mod.loadRun225ContractData(fullArm));
if (selected.content.content_memory_id !== "content_2_24_climax") throw new Error("climax did not select Run 2.24 content memory");
if (selected.assets.length < 2) throw new Error("climax did not select visual evidence assets");
if (!fullArm.allowed.includes(mod.RUN2_25_INPUTS.visualAssets)) throw new Error("full arm does not allow Run 2.24 visual assets");
"""
    completed = subprocess.run(
        ["node", "--input-type=module", "-e", node_script],
        cwd=ROOT,
        capture_output=True,
        text=True,
        timeout=20,
    )

    assert completed.returncode == 0, completed.stderr


def test_run2_25_content_evidence_surface_geometry_prevents_crushed_text() -> None:
    node_script = """
const mod = await import("./scripts/generate_ppt_run2_25_single_usecase_arms.mjs");
for (const width of [458, 552, 760, 870]) {
  const geom = mod.run225ContentEvidenceSurfaceGeometry({ x: 0, y: 0, w: width, h: 242 });
  if (geom.headline.w < 180) throw new Error(`headline width crushed for ${width}: ${geom.headline.w}`);
  if (geom.proofPoint.w < 150) throw new Error(`proof point width crushed for ${width}: ${geom.proofPoint.w}`);
  if (geom.rail.w < 176) throw new Error(`visual rail width crushed for ${width}: ${geom.rail.w}`);
  if (geom.assetCard.w < 140) throw new Error(`asset card width crushed for ${width}: ${geom.assetCard.w}`);
}
const medium = mod.run225ContentEvidenceSurfaceGeometry({ x: 0, y: 0, w: 640, h: 346 });
if (medium.mode !== "medium") throw new Error(`640px surface should use medium geometry: ${medium.mode}`);
if (medium.headline.h < 92) throw new Error(`medium headline height too short: ${medium.headline.h}`);
if (medium.proofPoint.y < medium.headline.y + medium.headline.h + 16) throw new Error("medium proof points collide with headline");
if (medium.proofPoint.compress !== true) throw new Error("medium proof points should use compressed copy");
if (medium.proofPoint.count !== 2) throw new Error(`medium surface should show two proof points, got ${medium.proofPoint.count}`);
const compact = mod.run225ContentEvidenceSurfaceGeometry({ x: 0, y: 0, w: 458, h: 242 });
if (compact.mode !== "compact") throw new Error("458px surface should use compact geometry");
const wide = mod.run225ContentEvidenceSurfaceGeometry({ x: 0, y: 0, w: 870, h: 368 });
if (wide.mode !== "wide") throw new Error("870px surface should use wide geometry");
"""
    completed = subprocess.run(
        ["node", "--input-type=module", "-e", node_script],
        cwd=ROOT,
        capture_output=True,
        text=True,
        timeout=20,
    )

    assert completed.returncode == 0, completed.stderr


def test_run2_25_records_single_usecase_generated_rerun_result() -> None:
    result = (PACK / "results" / "run2_25_single_usecase_rerun_result.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_25_single_usecase_rerun_result.json")

    assert result_json["status"] == "run2_25_single_usecase_rerun_public_blocked"
    assert result_json["public_ready"] is False
    assert result_json["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert result_json["run_id"] == "2.25"
    assert result_json["selected_usecase_id"] == "usecase_design_to_production_platform_launch"
    assert result_json["rerun"]["best_internal_arm"] == "run2_25_full_single_usecase_content_visual"
    assert result_json["rerun"]["best_internal_arm_verdict"] == "run2_24_pack_executed_before_native_ppt_generation"
    assert result_json["input_chain"]["content_memory"] == "docs/product/ppt-run2-data-skill-quality/run2_24_single_usecase_content_memory.json"
    assert result_json["input_chain"]["visual_evidence_asset_memory"] == "docs/product/ppt-run2-data-skill-quality/run2_24_visual_evidence_asset_memory.json"
    assert result_json["input_chain"]["content_visual_workflow_gates"] == "docs/product/ppt-run2-data-skill-quality/run2_24_content_visual_workflow_gates.json"
    assert result_json["control_boundary"]["bad_content_visual_memory"].startswith("selected_usecase_label_only")
    assert result_json["rerun"]["combined_contact_sheet"].endswith("run2-25-four-arm-contact-sheet.png")
    assert result_json["rerun"]["full_skill_series_sheet"].endswith("run2-full-skill-series-horizontal.png")
    assert result_json["rerun"]["generated_outputs_committed"] is False
    assert_contains(
        result,
        [
            "Run 2.25",
            "Run 2.24 single-usecase content memory",
            "four-arm rerun",
            "bad_content_visual_memory",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_25_generated_rerun() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "Run 2.25",
            "ppt-run2-25-prompt-only",
            "ppt-run2-25-run1-5-skill",
            "ppt-run2-25-full-vulca",
            "ppt-run2-25-bad-content-visual-memory",
            "run2_25_single_usecase_rerun_result.json",
        ],
    )


def test_run2_26_visual_module_quality_audit_scores_run2_25_outputs(tmp_path: Path) -> None:
    script_path = ROOT / "scripts" / "audit_ppt_run2_26_visual_module_quality.py"
    assert script_path.exists(), "missing Run 2.26 visual module quality audit script"

    result_json = tmp_path / "run2_26_visual_module_quality_audit.json"
    result_md = tmp_path / "run2_26_visual_module_quality_audit.md"
    completed = subprocess.run(
        [
            sys.executable,
            str(script_path),
            "--result-json",
            str(result_json),
            "--result-md",
            str(result_md),
        ],
        cwd=ROOT,
        capture_output=True,
        text=True,
        timeout=20,
    )

    assert completed.returncode == 0, completed.stderr
    audit = load_json(result_json)
    report = result_md.read_text(encoding="utf-8")

    assert audit["status"] == "run2_26_visual_module_quality_audit_public_blocked"
    assert audit["run_id"] == "2.26"
    assert audit["source_generated_run"] == "2.25"
    assert audit["creates_new_ppt_deck"] is False
    assert audit["public_ready"] is False
    assert audit["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert audit["input_chain"]["run2_25_full_trace_manifest"].endswith("ppt-run2-25-full-vulca/trace_manifest.json")
    assert audit["input_chain"]["run2_25_full_layout_dir"].endswith("ppt-run2-25-full-vulca/layout/final")
    assert audit["no_new_deck_proof"]["status"] == "pass"
    assert audit["no_new_deck_proof"]["matched_run2_26_outputs"] == []
    assert audit["no_new_deck_proof"]["new_pptx_created"] is False
    assert audit["quality_summary"]["module_quality_gate"] == "pass_internal_only"
    assert audit["quality_summary"]["public_release_gate"] == "blocked"
    assert audit["quality_summary"]["top_next_module_to_thicken"] == "drawRun225ContentEvidenceSurface"
    assert audit["quality_summary"]["roles_with_visible_layout_defects"] == []
    assert audit["quality_summary"]["roles_with_compressed_proof_surface"]
    assert audit["issue_categories"] == [
        "layout_geometry",
        "content_density",
        "visual_evidence_visibility",
        "composition_hierarchy",
        "climax_impact",
    ]
    assert len(audit["role_records"]) == 6
    for record in audit["role_records"]:
        assert record["role"] in {"cover", "setup", "contrast", "proof", "climax", "close"}
        assert record["geometry"]["crushed_text_box_count"] == 0
        assert record["geometry"]["max_line_count"] <= 14
        assert record["content_density"]["visual_evidence_slot_count"] >= 2
        assert record["visual_evidence_visibility"]["trace_slots_present"] is True
        assert record["status"] in {"pass_internal_only", "needs_next_module_thickening"}
        assert record["recommended_next_action"]

    assert_contains(
        report,
        [
            "Run 2.26",
            "visual module quality",
            "Run 2.25",
            "drawRun225ContentEvidenceSurface",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_run2_26_records_visual_module_quality_audit_result() -> None:
    result = (PACK / "results" / "run2_26_visual_module_quality_audit.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_26_visual_module_quality_audit.json")

    assert result_json["status"] == "run2_26_visual_module_quality_audit_public_blocked"
    assert result_json["source_generated_run"] == "2.25"
    assert result_json["creates_new_ppt_deck"] is False
    assert result_json["public_ready"] is False
    assert result_json["no_new_deck_proof"]["status"] == "pass"
    assert result_json["no_new_deck_proof"]["matched_run2_26_outputs"] == []
    assert result_json["no_new_deck_proof"]["new_pptx_created"] is False
    assert result_json["quality_summary"]["top_next_module_to_thicken"] == "drawRun225ContentEvidenceSurface"
    assert result_json["next_required_action"].startswith("thicken_drawRun225ContentEvidenceSurface")
    assert_contains(
        result,
        [
            "Run 2.26",
            "visual module quality",
            "Run 2.25",
            "public blocked",
        ],
    )


def test_ppt_run_html_viewer_embeds_run2_26_visual_module_quality_audit() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "run2_26_visual_module_quality_audit.json",
            "Run 2.26 visual module quality audit",
            "drawRun225ContentEvidenceSurface",
            "roles_with_compressed_proof_surface",
            "no_new_deck_proof",
        ],
    )


def test_run2_27_generator_consumes_run2_26_audit_before_native_ppt_code() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_27_content_surface_thickening_arms.mjs"
    assert script_path.exists(), "missing Run 2.27 content surface thickening generator"
    body = script_path.read_text(encoding="utf-8")
    arm_order = [
        "prompt_only",
        "run1_5_skill",
        "run2_27_full_content_surface_thickening",
        "bad_surface_thickening_memory",
    ]

    def arm_block(arm_id: str) -> str:
        start = body.index(f'armId: "{arm_id}"')
        next_starts = [body.find(f'armId: "{next_arm}"', start + 1) for next_arm in arm_order]
        next_starts = [index for index in next_starts if index > start]
        end = min(next_starts) if next_starts else len(body)
        return body[start:end]

    def section(block: str, start_marker: str, end_marker: str) -> str:
        start = block.index(start_marker)
        end = block.index(end_marker, start)
        return block[start:end]

    required_inputs = [
        "run2_24_single_usecase_content_memory.json",
        "run2_24_visual_evidence_asset_memory.json",
        "run2_24_content_visual_workflow_gates.json",
        "run2_26_visual_module_quality_audit.json",
    ]

    assert_contains(
        body,
        [
            "validateRun226VisualModuleAudit",
            "loadRun227ContractData",
            "drawRun227ContentEvidenceSurface",
            "run227ContentEvidenceSurfaceGeometry",
            "run2_27_content_surface_thickening_execution_status",
            "run2_27_code_module_ids",
            "thicken_drawRun225ContentEvidenceSurface_before_run2_27_rerun",
            "run2_26_visual_module_quality_audit_public_blocked",
        ],
    )

    prompt_allowed = section(arm_block("prompt_only"), "allowed:", "forbidden:")
    prompt_forbidden = section(arm_block("prompt_only"), "forbidden:", "palette:")
    run1_allowed = section(arm_block("run1_5_skill"), "allowed:", "forbidden:")
    run1_forbidden = section(arm_block("run1_5_skill"), "forbidden:", "palette:")
    full_allowed = section(arm_block("run2_27_full_content_surface_thickening"), "allowed:", "forbidden:")
    full_forbidden = section(arm_block("run2_27_full_content_surface_thickening"), "forbidden:", "palette:")
    bad_allowed = section(arm_block("bad_surface_thickening_memory"), "allowed:", "forbidden:")
    bad_forbidden = section(arm_block("bad_surface_thickening_memory"), "forbidden:", "palette:")

    for term in required_inputs:
        assert term not in prompt_allowed
        assert term in prompt_forbidden
        assert term not in run1_allowed
        assert term in run1_forbidden
        assert term in full_allowed
        assert term not in full_forbidden
        assert term not in bad_allowed
        assert term in bad_forbidden

    assert 'const fullRun227 = arm.armId === "run2_27_full_content_surface_thickening";' in body
    assert 'registerRun227Module(metrics, "drawRun227ContentEvidenceSurface")' in body
    for field in [
        "run2_26_source_audit_status",
        "run2_26_roles_with_compressed_proof_surface",
        "run2_27_content_surface_thickening_execution_status",
        "run2_27_required_code_module_ids",
        "run2_27_code_module_ids",
    ]:
        assert re.search(fr"{field}:\s*fullRun227\s*\?", body), field


def test_run2_27_surface_geometry_shows_all_proof_points_without_compression() -> None:
    node_script = """
const mod = await import("./scripts/generate_ppt_run2_27_content_surface_thickening_arms.mjs");
for (const width of [540, 604, 760, 870]) {
  const geom = mod.run227ContentEvidenceSurfaceGeometry({ x: 0, y: 0, w: width, h: 360 });
  if (geom.proofPoint.count !== 3) throw new Error(`Run 2.27 must show all three proof points for ${width}, got ${geom.proofPoint.count}`);
  if (geom.proofPoint.compress !== false) throw new Error(`Run 2.27 must not compress proof copy for ${width}`);
  if (geom.proofPoint.w < 280) throw new Error(`Run 2.27 proof row width too small for ${width}: ${geom.proofPoint.w}`);
  if (geom.proofPoint.h < 32) throw new Error(`Run 2.27 proof row height too short for ${width}: ${geom.proofPoint.h}`);
  if (geom.assetCard.w < 132) throw new Error(`Run 2.27 asset card width too small for ${width}: ${geom.assetCard.w}`);
}
const compact = mod.run227ContentEvidenceSurfaceGeometry({ x: 0, y: 0, w: 540, h: 342 });
if (compact.mode !== "compact-thick") throw new Error(`540px surface should use compact-thick geometry: ${compact.mode}`);
if (compact.proofPoint.y < compact.headline.y + compact.headline.h + 12) throw new Error("proof rows collide with headline in compact-thick geometry");
"""
    completed = subprocess.run(
        ["node", "--input-type=module", "-e", node_script],
        cwd=ROOT,
        capture_output=True,
        text=True,
        timeout=20,
    )

    assert completed.returncode == 0, completed.stderr


def test_run2_27_runtime_guards_block_bad_arm_and_select_audit_targets() -> None:
    node_script = """
const mod = await import("./scripts/generate_ppt_run2_27_content_surface_thickening_arms.mjs");
const badArm = mod.armSpecs.find((arm) => arm.armId === "bad_surface_thickening_memory");
const fullArm = mod.armSpecs.find((arm) => arm.armId === "run2_27_full_content_surface_thickening");
for (const input of mod.RUN2_27_DATA_INPUTS) {
  let blocked = false;
  try {
    mod.readRun227PackJsonForArm(badArm, input);
  } catch (error) {
    blocked = String(error.message).includes("input boundary does not permit reading");
  }
  if (!blocked) throw new Error(`bad surface-thickening arm was able to read ${input}`);
}
const audit = mod.loadRun226VisualModuleAudit(fullArm);
if (audit.quality_summary.top_next_module_to_thicken !== "drawRun225ContentEvidenceSurface") throw new Error("Run 2.27 loaded wrong audit target");
if (!audit.quality_summary.roles_with_compressed_proof_surface.includes("setup")) throw new Error("Run 2.27 audit target missing setup");
const badTrace = mod.traceFor(badArm);
for (const slide of badTrace.slides) {
  if (slide.run2_26_source_audit_status) throw new Error("bad control trace leaked Run 2.26 audit");
  if ((slide.run2_27_code_module_ids || []).length) throw new Error("bad control trace leaked Run 2.27 code modules");
}
const fullTrace = mod.traceFor(fullArm);
if (fullTrace.source_audit_run_id !== "2.26") throw new Error("full trace did not bind Run 2.26 audit");
for (const slide of fullTrace.slides) {
  if (slide.role === "setup" || slide.role === "contrast" || slide.role === "close") {
    if (!slide.run2_27_required_code_module_ids.includes("drawRun227ContentEvidenceSurface")) throw new Error(`${slide.role} missing thickened surface requirement`);
    if (!slide.run2_27_code_module_ids.includes("drawRun227ContentEvidenceSurface")) throw new Error(`${slide.role} missing thickened surface trace`);
  }
}
"""
    completed = subprocess.run(
        ["node", "--input-type=module", "-e", node_script],
        cwd=ROOT,
        capture_output=True,
        text=True,
        timeout=20,
    )

    assert completed.returncode == 0, completed.stderr


def test_run2_27_records_content_surface_thickening_rerun_result() -> None:
    result = (PACK / "results" / "run2_27_content_surface_thickening_rerun_result.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_27_content_surface_thickening_rerun_result.json")

    assert result_json["status"] == "run2_27_content_surface_thickening_rerun_public_blocked"
    assert result_json["public_ready"] is False
    assert result_json["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert result_json["source_audit_run_id"] == "2.26"
    assert result_json["rerun"]["best_internal_arm"] == "run2_27_full_content_surface_thickening"
    assert result_json["rerun"]["best_internal_arm_verdict"] == "run2_26_audit_target_executed_before_native_ppt_generation"
    assert result_json["quality_delta"]["target_module"] == "drawRun225ContentEvidenceSurface"
    assert result_json["quality_delta"]["replacement_module"] == "drawRun227ContentEvidenceSurface"
    assert result_json["quality_delta"]["roles_with_compressed_proof_surface_after"] == []
    assert result_json["visual_quality_boundary"] == (
        "content_surface_thickening_proof_only_not_public_video_grade_aesthetic_or_human_release_approval"
    )
    assert result_json["rerun"]["combined_contact_sheet"].endswith("run2-27-four-arm-contact-sheet.png")
    assert result_json["rerun"]["full_skill_series_sheet"].endswith("run2-full-skill-series-horizontal.png")
    assert_contains(
        result,
        [
            "Run 2.27",
            "Run 2.26 visual module quality audit",
            "drawRun227ContentEvidenceSurface",
            "four-arm rerun",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_27_content_surface_thickening_rerun() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "Run 2.27",
            "ppt-run2-27-prompt-only",
            "ppt-run2-27-run1-5-skill",
            "ppt-run2-27-full-vulca",
            "ppt-run2-27-bad-surface-thickening-memory",
            "run2_27_content_surface_thickening_rerun_result.json",
            "drawRun227ContentEvidenceSurface",
        ],
    )


def test_run2_28_evidence_chain_view_model_binds_source_data_workflow_and_slide_surface() -> None:
    model = load_json(PACK / "run2_28_evidence_chain_view_model.json")
    source_records = load_json(PACK / "run2_7_multimodal_source_records.json")
    content_memory = load_json(PACK / "run2_24_single_usecase_content_memory.json")
    workflow_gates = load_json(PACK / "run2_24_content_visual_workflow_gates.json")

    assert model["status"] == "run2_28_evidence_chain_view_model_ready_public_blocked"
    assert model["selected_usecase_id"] == "usecase_design_to_production_platform_launch"
    assert model["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert model["source_layer"] == "run2_7_multimodal_source_records.json"
    assert model["content_layer"] == "run2_24_single_usecase_content_memory.json"
    assert model["workflow_layer"] == "run2_24_content_visual_workflow_gates.json"

    source_ids = {record["id"] for record in source_records["records"]}
    content_by_role = {record["role"]: record for record in content_memory["slide_content_memory"]}
    gate_by_role = {gate["role"]: gate for gate in workflow_gates["gates"]}
    chains = model["slide_evidence_chains"]

    assert [chain["role"] for chain in chains] == ["cover", "setup", "contrast", "proof", "climax", "close"]
    for chain in chains:
        role = chain["role"]
        content = content_by_role[role]
        gate = gate_by_role[role]

        assert chain["content_memory_id"] == content["content_memory_id"]
        assert chain["workflow_gate_id"] == gate["gate_id"]
        assert set(chain["visual_evidence_slot_ids"]) == set(content["visual_evidence_slot_ids"])
        assert set(chain["multimodal_source_record_ids"]) <= source_ids
        assert chain["extracted_design_rule"]
        assert chain["workflow_decision"]
        assert chain["native_ppt_surface_instruction"]
        assert chain["viewer_inspection_prompt"]
        assert chain["negative_control_failure_mode"]
        assert set(chain["required_trace_fields"]) >= {
            "run2_28_evidence_chain_id",
            "run2_28_multimodal_source_record_ids",
            "run2_28_extracted_design_rule",
            "run2_28_workflow_decision",
            "run2_28_native_surface_module_id",
        }

    assert model["viewer_contract"]["show_chain_in_html_viewer"] is True
    assert model["viewer_contract"]["chain_order"] == [
        "source evidence",
        "extracted design rule",
        "workflow decision",
        "generated slide surface",
    ]
    assert model["public_release_gate"] == "blocked"


def test_run2_28_generator_consumes_evidence_chain_before_native_ppt_code() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_28_evidence_chain_arms.mjs"
    assert script_path.exists(), "missing Run 2.28 evidence-chain generator"
    body = script_path.read_text(encoding="utf-8")
    arm_order = [
        "prompt_only",
        "run1_5_skill",
        "run2_28_full_evidence_chain",
        "bad_evidence_chain_memory",
    ]

    def arm_block(arm_id: str) -> str:
        start = body.index(f'armId: "{arm_id}"')
        next_starts = [body.find(f'armId: "{next_arm}"', start + 1) for next_arm in arm_order]
        next_starts = [index for index in next_starts if index > start]
        end = min(next_starts) if next_starts else len(body)
        return body[start:end]

    def section(block: str, start_marker: str, end_marker: str) -> str:
        start = block.index(start_marker)
        end = block.index(end_marker, start)
        return block[start:end]

    required_inputs = [
        "run2_28_evidence_chain_view_model.json",
        "run2_24_single_usecase_content_memory.json",
        "run2_24_visual_evidence_asset_memory.json",
        "run2_24_content_visual_workflow_gates.json",
        "run2_27_content_surface_thickening_rerun_result.json",
    ]

    assert_contains(
        body,
        [
            "validateRun228EvidenceChainViewModel",
            "loadRun228ContractData",
            "drawRun228EvidenceChainSurface",
            "drawRun228WorkflowTracePanel",
            "run228EvidenceChainSurfaceGeometry",
            "run2_28_evidence_chain_execution_status",
            "run2_28_native_surface_module_id",
            "source evidence",
            "extracted design rule",
            "workflow decision",
            "generated slide surface",
        ],
    )

    prompt_allowed = section(arm_block("prompt_only"), "allowed:", "forbidden:")
    prompt_forbidden = section(arm_block("prompt_only"), "forbidden:", "palette:")
    run1_allowed = section(arm_block("run1_5_skill"), "allowed:", "forbidden:")
    run1_forbidden = section(arm_block("run1_5_skill"), "forbidden:", "palette:")
    full_allowed = section(arm_block("run2_28_full_evidence_chain"), "allowed:", "forbidden:")
    full_forbidden = section(arm_block("run2_28_full_evidence_chain"), "forbidden:", "palette:")
    bad_allowed = section(arm_block("bad_evidence_chain_memory"), "allowed:", "forbidden:")
    bad_forbidden = section(arm_block("bad_evidence_chain_memory"), "forbidden:", "palette:")

    for term in required_inputs:
        assert term not in prompt_allowed
        assert term in prompt_forbidden
        assert term not in run1_allowed
        assert term in run1_forbidden
        assert term in full_allowed
        assert term not in full_forbidden
        assert term not in bad_allowed
        assert term in bad_forbidden

    assert 'const fullRun228 = arm.armId === "run2_28_full_evidence_chain";' in body
    assert 'registerRun228Module(metrics, "drawRun228EvidenceChainSurface")' in body
    for field in [
        "run2_28_evidence_chain_id",
        "run2_28_multimodal_source_record_ids",
        "run2_28_extracted_design_rule",
        "run2_28_workflow_decision",
        "run2_28_native_surface_module_id",
        "run2_28_evidence_chain_execution_status",
    ]:
        assert re.search(fr"{field}:\s*fullRun228\s*\?", body), field


def test_ppt_run_html_viewer_mentions_run2_28_evidence_chain_rerun() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "Run 2.28",
            "ppt-run2-28-prompt-only",
            "ppt-run2-28-run1-5-skill",
            "ppt-run2-28-full-vulca",
            "ppt-run2-28-bad-evidence-chain-memory",
            "run2_28_evidence_chain_view_model.json",
            "run2_28_evidence_chain_rerun_result.json",
            "source evidence",
            "extracted design rule",
            "workflow decision",
            "generated slide surface",
        ],
    )


def test_run2_29_presentation_synthesis_memory_compresses_evidence_chain_without_dropping_trace() -> None:
    memory = load_json(PACK / "run2_29_presentation_synthesis_memory.json")
    evidence_chain = load_json(PACK / "run2_28_evidence_chain_view_model.json")
    prior_result = load_json(PACK / "results" / "run2_28_evidence_chain_rerun_result.json")

    assert memory["status"] == "run2_29_presentation_synthesis_memory_ready_public_blocked"
    assert memory["selected_usecase_id"] == "usecase_design_to_production_platform_launch"
    assert memory["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert memory["source_evidence_chain_layer"] == "run2_28_evidence_chain_view_model.json"
    assert memory["prior_rerun_result"] == "run2_28_evidence_chain_rerun_result.json"
    assert prior_result["status"] == "run2_28_evidence_chain_rerun_public_blocked"
    assert memory["surface_policy"]["primary_reader_experience"] == "presentation_first_surface"
    assert memory["surface_policy"]["secondary_reviewer_experience"] == "compressed_evidence_spine"
    assert memory["surface_policy"]["forbidden_primary_surface"] == "four_column_audit_table"

    chains_by_role = {chain["role"]: chain for chain in evidence_chain["slide_evidence_chains"]}
    records = memory["slide_synthesis_records"]
    assert [record["role"] for record in records] == ["cover", "setup", "contrast", "proof", "climax", "close"]
    assert {record["evidence_chain_id"] for record in records} == {
        chain["evidence_chain_id"] for chain in evidence_chain["slide_evidence_chains"]
    }

    required_trace_fields = {
        "run2_28_evidence_chain_id",
        "run2_28_multimodal_source_record_ids",
        "run2_28_extracted_design_rule",
        "run2_28_workflow_decision",
        "run2_28_native_surface_module_id",
        "run2_29_synthesis_record_id",
        "run2_29_public_surface_mode",
        "run2_29_trace_surface_mode",
        "run2_29_presentation_module_id",
        "run2_29_chain_compression_policy",
    }
    for record in records:
        chain = chains_by_role[record["role"]]
        assert record["evidence_chain_id"] == chain["evidence_chain_id"]
        assert record["source_native_surface_module_id"] == chain["native_surface_module_id"]
        assert record["trace_surface_mode"] == "manifest_and_html_viewer_full_chain_visible"
        assert record["presentation_module_id"].startswith("drawRun229")
        assert record["visible_on_slide_evidence_spine_steps"] == [
            "source evidence",
            "extracted design rule",
            "workflow decision",
            "generated slide surface",
        ]
        assert "do_not_render_four_column_audit_table_as_primary_surface" in record["chain_compression_policy"]
        assert "make_evidence_chain_secondary_to_editorial_claim_and_proof_object" in record["chain_compression_policy"]
        assert set(record["preserved_trace_fields"]) >= required_trace_fields
        assert record["primary_slide_surface_contract"]
        assert record["reviewer_inspection_prompt"]

    assert memory["public_release_gate"] == "blocked"


def test_run2_29_generator_consumes_presentation_synthesis_before_native_ppt_code() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_29_presentation_synthesis_arms.mjs"
    assert script_path.exists(), "missing Run 2.29 presentation-synthesis generator"
    body = script_path.read_text(encoding="utf-8")
    arm_order = [
        "prompt_only",
        "run1_5_skill",
        "run2_29_full_presentation_synthesis",
        "bad_presentation_synthesis_memory",
    ]

    def arm_block(arm_id: str) -> str:
        start = body.index(f'armId: "{arm_id}"')
        next_starts = [body.find(f'armId: "{next_arm}"', start + 1) for next_arm in arm_order]
        next_starts = [index for index in next_starts if index > start]
        end = min(next_starts) if next_starts else len(body)
        return body[start:end]

    def section(block: str, start_marker: str, end_marker: str) -> str:
        start = block.index(start_marker)
        end = block.index(end_marker, start)
        return block[start:end]

    required_inputs = [
        "run2_29_presentation_synthesis_memory.json",
        "run2_28_evidence_chain_view_model.json",
        "run2_24_single_usecase_content_memory.json",
        "run2_24_visual_evidence_asset_memory.json",
        "run2_24_content_visual_workflow_gates.json",
        "run2_28_evidence_chain_rerun_result.json",
    ]

    assert_contains(
        body,
        [
            "validateRun229PresentationSynthesisMemory",
            "loadRun229ContractData",
            "drawRun229CompressedEvidenceSpine",
            "drawRun229EditorialLaunchFrame",
            "drawRun229HeroProofScene",
            "drawRun229DecisionHandoff",
            "run2_29_presentation_synthesis_execution_status",
            "run2_29_public_surface_mode",
            "run2_29_trace_surface_mode",
            "run2_29_presentation_module_id",
            "presentation_first_surface_rendered_with_secondary_evidence_spine",
            "do_not_render_four_column_audit_table_as_primary_surface",
            "compressed evidence spine",
            "presentation-first surface",
        ],
    )

    prompt_allowed = section(arm_block("prompt_only"), "allowed:", "forbidden:")
    prompt_forbidden = section(arm_block("prompt_only"), "forbidden:", "palette:")
    run1_allowed = section(arm_block("run1_5_skill"), "allowed:", "forbidden:")
    run1_forbidden = section(arm_block("run1_5_skill"), "forbidden:", "palette:")
    full_allowed = section(arm_block("run2_29_full_presentation_synthesis"), "allowed:", "forbidden:")
    full_forbidden = section(arm_block("run2_29_full_presentation_synthesis"), "forbidden:", "palette:")
    bad_allowed = section(arm_block("bad_presentation_synthesis_memory"), "allowed:", "forbidden:")
    bad_forbidden = section(arm_block("bad_presentation_synthesis_memory"), "forbidden:", "palette:")

    for term in required_inputs:
        assert term not in prompt_allowed
        assert term in prompt_forbidden
        assert term not in run1_allowed
        assert term in run1_forbidden
        assert term in full_allowed
        assert term not in full_forbidden
        assert term not in bad_allowed
        assert term in bad_forbidden

    assert 'const fullRun229 = arm.armId === "run2_29_full_presentation_synthesis";' in body
    assert 'registerRun229Module(metrics, "drawRun229CompressedEvidenceSpine")' in body
    for field in [
        "run2_29_synthesis_record_id",
        "run2_29_public_surface_mode",
        "run2_29_trace_surface_mode",
        "run2_29_presentation_module_id",
        "run2_29_chain_compression_policy",
        "run2_29_presentation_synthesis_execution_status",
    ]:
        assert re.search(fr"{field}:\s*fullRun229\s*\?", body), field


def test_ppt_run_html_viewer_mentions_run2_29_presentation_synthesis_rerun() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "Run 2.29",
            "ppt-run2-29-prompt-only",
            "ppt-run2-29-run1-5-skill",
            "ppt-run2-29-full-vulca",
            "ppt-run2-29-bad-presentation-synthesis-memory",
            "run2_29_presentation_synthesis_memory.json",
            "run2_29_presentation_synthesis_rerun_result.json",
            "presentation-first surface",
            "compressed evidence spine",
            "four-column audit table",
            "manifest_and_html_viewer_full_chain_visible",
        ],
    )


def test_run2_30_presentation_synthesis_audit_scores_run2_29_outputs(tmp_path: Path) -> None:
    script_path = ROOT / "scripts" / "audit_ppt_run2_30_presentation_synthesis.py"
    assert script_path.exists(), "missing Run 2.30 presentation-synthesis audit script"

    result_json = tmp_path / "run2_30_presentation_synthesis_audit.json"
    result_md = tmp_path / "run2_30_presentation_synthesis_audit.md"
    subprocess.run(
        [
            sys.executable,
            str(script_path),
            "--result-json",
            str(result_json),
            "--result-md",
            str(result_md),
        ],
        cwd=ROOT,
        check=True,
    )

    audit = load_json(result_json)
    assert audit["schema_version"] == "ppt_run2_presentation_synthesis_audit.v1"
    assert audit["run_id"] == "2.30"
    assert audit["status"] == "run2_30_presentation_synthesis_audit_public_blocked"
    assert audit["source_generated_run"] == "2.29"
    assert audit["comparison_baseline_run"] == "2.28"
    assert audit["source_synthesis_layer"] == "2.29"
    assert audit["creates_new_ppt_deck"] is False
    assert audit["public_ready"] is False
    assert audit["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert audit["input_chain"]["run2_29_full_trace_manifest"].endswith("ppt-run2-29-full-vulca/trace_manifest.json")
    assert audit["input_chain"]["run2_29_bad_trace_manifest"].endswith(
        "ppt-run2-29-bad-presentation-synthesis-memory/trace_manifest.json"
    )
    assert audit["input_chain"]["run2_28_full_trace_manifest"].endswith("ppt-run2-28-full-vulca/trace_manifest.json")
    assert audit["no_new_deck_proof"]["new_pptx_created"] is False
    assert audit["no_new_deck_proof"]["status"] == "pass"

    trace = audit["trace_closure"]
    assert trace["full_arm"]["arm_id"] == "run2_29_full_presentation_synthesis"
    assert trace["full_arm"]["slide_count"] == 6
    assert trace["full_arm"]["presentation_synthesis_records_selected"] == 6
    assert trace["full_arm"]["compressed_evidence_spine_modules_called"] == 6
    assert trace["full_arm"]["presentation_first_execution_status_slides"] == 6
    assert trace["full_arm"]["run2_28_chain_fields_preserved"] == 6
    assert trace["bad_control"]["arm_id"] == "bad_presentation_synthesis_memory"
    assert trace["bad_control"]["presentation_synthesis_fields_leaked"] == 0

    comparison = audit["comparison_to_run2_28"]
    assert comparison["audit_table_demoted_to_secondary_spine"] is True
    assert comparison["full_chain_preserved_in_trace"] is True
    assert comparison["primary_surface_delta"] == "four_column_audit_table_to_presentation_first_surface"

    summary = audit["quality_summary"]
    assert summary["presentation_synthesis_gate"] == "pass_internal_only"
    assert summary["public_release_gate"] == "blocked"
    assert summary["top_next_layer_to_thicken"] == "spine_readability_and_climax_consistency"
    assert set(summary["roles_with_dense_spine_text"]) == {"cover", "setup", "contrast", "proof", "climax", "close"}
    assert summary["roles_with_climax_style_shift"] == ["climax"]

    assert len(audit["role_records"]) == 6
    for record in audit["role_records"]:
        assert record["presentation_first_surface"]["status"] == "pass"
        assert record["trace_closure"]["run2_28_chain_preserved"] is True
        assert record["trace_closure"]["synthesis_record_selected"] is True
        assert record["evidence_spine"]["compressed_spine_module_called"] is True
        assert record["evidence_spine"]["min_spine_font_size"] < 8
        assert "spine_readability" in record["issue_categories"]


def test_run2_30_records_presentation_synthesis_audit_result() -> None:
    result = (PACK / "results" / "run2_30_presentation_synthesis_audit.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_30_presentation_synthesis_audit.json")

    assert result_json["status"] == "run2_30_presentation_synthesis_audit_public_blocked"
    assert result_json["source_generated_run"] == "2.29"
    assert result_json["comparison_baseline_run"] == "2.28"
    assert result_json["quality_summary"]["top_next_layer_to_thicken"] == "spine_readability_and_climax_consistency"
    assert result_json["next_required_action"] == "thicken_run2_29_spine_readability_and_climax_consistency_before_run2_31_rerun"
    assert_contains(
        result,
        [
            "Run 2.30 Presentation Synthesis Audit",
            "audit-only",
            "2.29",
            "2.28",
            "compressed evidence spine",
            "presentation-first surface",
            "spine_readability_and_climax_consistency",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_embeds_run2_30_presentation_synthesis_audit() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "run2_30_presentation_synthesis_audit.json",
            "Run 2.30 presentation synthesis audit",
            "audit_table_demoted_to_secondary_spine",
            "spine_readability_and_climax_consistency",
            "presentation_synthesis_gate",
            "full_chain_preserved_in_trace",
            "creates no new PPT deck",
        ],
    )


def test_run2_31_generator_consumes_run2_30_audit_before_native_ppt_code() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_31_spine_climax_repair_arms.mjs"
    assert script_path.exists(), "missing Run 2.31 spine/climax repair generator"
    body = script_path.read_text(encoding="utf-8")
    arm_order = [
        "prompt_only",
        "run1_5_skill",
        "run2_31_full_spine_climax_repair",
        "bad_spine_climax_repair_memory",
    ]

    def arm_block(arm_id: str) -> str:
        start = body.index(f'armId: "{arm_id}"')
        next_starts = [body.find(f'armId: "{next_arm}"', start + 1) for next_arm in arm_order]
        next_starts = [index for index in next_starts if index > start]
        end = min(next_starts) if next_starts else len(body)
        return body[start:end]

    def section(block: str, start_marker: str, end_marker: str) -> str:
        start = block.index(start_marker)
        end = block.index(end_marker, start)
        return block[start:end]

    required_inputs = [
        "run2_30_presentation_synthesis_audit.json",
        "run2_29_presentation_synthesis_memory.json",
        "run2_28_evidence_chain_view_model.json",
        "run2_24_single_usecase_content_memory.json",
        "run2_24_visual_evidence_asset_memory.json",
        "run2_24_content_visual_workflow_gates.json",
        "run2_29_presentation_synthesis_rerun_result.json",
    ]

    assert_contains(
        body,
        [
            "validateRun230PresentationSynthesisAudit",
            "loadRun231ContractData",
            "drawRun231ReadableEvidenceSpine",
            "drawRun231HeroProofScene",
            "spine_readability_and_climax_consistency",
            "run2_30_source_audit_status",
            "run2_30_top_next_layer_to_thicken",
            "run2_31_spine_min_font_size_target",
            "run2_31_climax_style_policy",
            "run2_31_spine_climax_repair_execution_status",
            "spine_readability_and_climax_consistency_repaired_before_native_ppt_generation",
        ],
    )

    prompt_allowed = section(arm_block("prompt_only"), "allowed:", "forbidden:")
    prompt_forbidden = section(arm_block("prompt_only"), "forbidden:", "palette:")
    run1_allowed = section(arm_block("run1_5_skill"), "allowed:", "forbidden:")
    run1_forbidden = section(arm_block("run1_5_skill"), "forbidden:", "palette:")
    full_allowed = section(arm_block("run2_31_full_spine_climax_repair"), "allowed:", "forbidden:")
    full_forbidden = section(arm_block("run2_31_full_spine_climax_repair"), "forbidden:", "palette:")
    bad_allowed = section(arm_block("bad_spine_climax_repair_memory"), "allowed:", "forbidden:")
    bad_forbidden = section(arm_block("bad_spine_climax_repair_memory"), "forbidden:", "palette:")

    for term in required_inputs:
        assert term not in prompt_allowed
        assert term in prompt_forbidden
        assert term not in run1_allowed
        assert term in run1_forbidden
        assert term in full_allowed
        assert term not in full_forbidden
        assert term not in bad_allowed
        assert term in bad_forbidden

    assert 'const fullRun231 = arm.armId === "run2_31_full_spine_climax_repair";' in body
    assert 'registerRun231Module(metrics, "drawRun231ReadableEvidenceSpine")' in body
    for field in [
        "run2_30_source_audit_status",
        "run2_30_top_next_layer_to_thicken",
        "run2_31_spine_min_font_size_target",
        "run2_31_climax_style_policy",
        "run2_31_spine_climax_repair_execution_status",
    ]:
        assert re.search(fr"{field}:\s*fullRun231\s*\?", body), field


def test_run2_31_records_spine_climax_repair_rerun_result() -> None:
    result = (PACK / "results" / "run2_31_spine_climax_repair_rerun_result.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_31_spine_climax_repair_rerun_result.json")

    assert result_json["status"] == "run2_31_spine_climax_repair_rerun_public_blocked"
    assert result_json["source_audit_run_id"] == "2.30"
    assert result_json["input_chain"]["presentation_synthesis_audit"].endswith(
        "run2_30_presentation_synthesis_audit.json"
    )
    assert result_json["rerun"]["best_internal_arm"] == "run2_31_full_spine_climax_repair"
    assert result_json["rerun"]["best_internal_arm_verdict"] == (
        "spine_readability_and_climax_consistency_repaired_before_native_ppt_generation"
    )
    assert result_json["quality_delta"]["target_layer"] == "spine_readability_and_climax_consistency"
    assert result_json["quality_delta"]["spine_min_font_size_target"] >= 8
    assert result_json["quality_delta"]["climax_style_policy"] == "high_contrast_climax_with_shared_light_editorial_frame"
    assert result_json["rerun"]["combined_contact_sheet"].endswith("run2-31-four-arm-contact-sheet.png")
    assert result_json["rerun"]["full_skill_series_sheet"].endswith("run2-full-skill-series-horizontal.png")
    assert_contains(
        result,
        [
            "Run 2.31",
            "Run 2.30 presentation synthesis audit",
            "drawRun231ReadableEvidenceSpine",
            "drawRun231HeroProofScene",
            "spine_readability_and_climax_consistency",
            "four-arm rerun",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_31_spine_climax_repair_rerun() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "Run 2.31",
            "ppt-run2-31-prompt-only",
            "ppt-run2-31-run1-5-skill",
            "ppt-run2-31-full-vulca",
            "ppt-run2-31-bad-spine-climax-repair-memory",
            "run2_31_spine_climax_repair_rerun_result.json",
            "drawRun231ReadableEvidenceSpine",
            "drawRun231HeroProofScene",
            "spine_readability_and_climax_consistency",
        ],
    )


def test_run2_32_spine_climax_repair_audit_scores_run2_31_outputs(tmp_path: Path) -> None:
    script_path = ROOT / "scripts" / "audit_ppt_run2_32_spine_climax_repair.py"
    assert script_path.exists(), "missing Run 2.32 spine/climax repair audit script"

    result_json = tmp_path / "run2_32_spine_climax_repair_audit.json"
    result_md = tmp_path / "run2_32_spine_climax_repair_audit.md"
    subprocess.run(
        [
            sys.executable,
            str(script_path),
            "--result-json",
            str(result_json),
            "--result-md",
            str(result_md),
        ],
        cwd=ROOT,
        check=True,
    )

    audit = load_json(result_json)
    assert audit["schema_version"] == "ppt_run2_spine_climax_repair_audit.v1"
    assert audit["run_id"] == "2.32"
    assert audit["status"] == "run2_32_spine_climax_repair_audit_public_blocked"
    assert audit["source_generated_run"] == "2.31"
    assert audit["source_audit_run"] == "2.30"
    assert audit["source_repair_run"] == "2.31"
    assert audit["creates_new_ppt_deck"] is False
    assert audit["public_ready"] is False
    assert audit["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert audit["input_chain"]["run2_31_full_trace_manifest"].endswith("ppt-run2-31-full-vulca/trace_manifest.json")
    assert audit["input_chain"]["run2_31_bad_trace_manifest"].endswith(
        "ppt-run2-31-bad-spine-climax-repair-memory/trace_manifest.json"
    )
    assert audit["input_chain"]["run2_30_presentation_synthesis_audit"].endswith(
        "run2_30_presentation_synthesis_audit.json"
    )
    assert audit["input_chain"]["run2_31_spine_climax_repair_rerun_result"].endswith(
        "run2_31_spine_climax_repair_rerun_result.json"
    )
    assert audit["no_new_deck_proof"]["new_pptx_created"] is False
    assert audit["no_new_deck_proof"]["status"] == "pass"

    trace = audit["trace_closure"]
    assert trace["full_arm"]["arm_id"] == "run2_31_full_spine_climax_repair"
    assert trace["full_arm"]["slide_count"] == 6
    assert trace["full_arm"]["run2_30_audit_consumed_slides"] == 6
    assert trace["full_arm"]["repair_execution_status_slides"] == 6
    assert trace["full_arm"]["readable_evidence_spine_modules_called"] == 6
    assert trace["full_arm"]["climax_style_policy_slides"] == 6
    assert trace["bad_control"]["arm_id"] == "bad_spine_climax_repair_memory"
    assert trace["bad_control"]["repair_fields_leaked"] == 0

    repair = audit["repair_verification"]
    assert repair["spine_font_target_met_by_contract"] is True
    assert repair["spine_min_font_size_target"] >= 8
    assert repair["climax_style_policy_enforced"] is True
    assert repair["climax_hero_object_canvas_share"] > 0.5
    assert repair["repair_target_closed"] is True

    summary = audit["quality_summary"]
    assert summary["repair_gate"] == "pass_internal_only"
    assert summary["public_release_gate"] == "blocked"
    assert summary["top_next_layer_to_thicken"] == "main_surface_information_density_and_visual_evidence_realism"
    assert set(summary["roles_with_low_visual_evidence_density"]) == {"cover", "setup", "contrast", "close"}
    assert "spine_readability_and_climax_consistency" in summary["closed_target_layers"]

    assert len(audit["role_records"]) == 6
    climax_records = [record for record in audit["role_records"] if record["role"] == "climax"]
    assert len(climax_records) == 1
    for record in audit["role_records"]:
        assert record["repair_trace_closure"]["run2_30_audit_status_present"] is True
        assert record["spine_repair"]["readable_spine_module_called"] is True
        assert record["spine_repair"]["spine_min_font_size_target"] >= 8
        assert record["climax_repair"]["climax_style_policy"] == "high_contrast_climax_with_shared_light_editorial_frame"
    assert climax_records[0]["climax_repair"]["hero_object_canvas_share"] > 0.5


def test_run2_32_records_spine_climax_repair_audit_result() -> None:
    result = (PACK / "results" / "run2_32_spine_climax_repair_audit.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_32_spine_climax_repair_audit.json")

    assert result_json["status"] == "run2_32_spine_climax_repair_audit_public_blocked"
    assert result_json["source_generated_run"] == "2.31"
    assert result_json["quality_summary"]["top_next_layer_to_thicken"] == (
        "main_surface_information_density_and_visual_evidence_realism"
    )
    assert result_json["next_required_action"] == (
        "thicken_run2_31_main_surface_information_density_and_visual_evidence_realism_before_run2_33_rerun"
    )
    assert_contains(
        result,
        [
            "Run 2.32 Spine/Climax Repair Audit",
            "audit-only",
            "2.31",
            "2.30",
            "drawRun231ReadableEvidenceSpine",
            "drawRun231HeroProofScene",
            "main_surface_information_density_and_visual_evidence_realism",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_embeds_run2_32_spine_climax_repair_audit() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "run2_32_spine_climax_repair_audit.json",
            "Run 2.32 spine/climax repair audit",
            "repair_target_closed",
            "main_surface_information_density_and_visual_evidence_realism",
            "roles_with_low_visual_evidence_density",
            "creates no new PPT deck",
        ],
    )


def test_run2_33_generator_consumes_run2_32_audit_before_native_ppt_code() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_33_main_surface_visual_evidence_arms.mjs"
    assert script_path.exists(), "missing Run 2.33 main-surface visual-evidence generator"
    body = script_path.read_text(encoding="utf-8")
    arm_order = [
        "prompt_only",
        "run1_5_skill",
        "run2_33_full_main_surface_visual_evidence",
        "bad_main_surface_visual_evidence_memory",
    ]

    def arm_block(arm_id: str) -> str:
        start = body.index(f'armId: "{arm_id}"')
        next_starts = [body.find(f'armId: "{next_arm}"', start + 1) for next_arm in arm_order]
        next_starts = [index for index in next_starts if index > start]
        end = min(next_starts) if next_starts else len(body)
        return body[start:end]

    def section(block: str, start_marker: str, end_marker: str) -> str:
        start = block.index(start_marker)
        end = block.index(end_marker, start)
        return block[start:end]

    required_inputs = [
        "run2_32_spine_climax_repair_audit.json",
        "run2_31_spine_climax_repair_rerun_result.json",
        "run2_30_presentation_synthesis_audit.json",
        "run2_29_presentation_synthesis_memory.json",
        "run2_28_evidence_chain_view_model.json",
        "run2_24_single_usecase_content_memory.json",
        "run2_24_visual_evidence_asset_memory.json",
        "run2_24_content_visual_workflow_gates.json",
    ]

    assert_contains(
        body,
        [
            "validateRun232SpineClimaxRepairAudit",
            "loadRun233ContractData",
            "drawRun233MainSurfaceEvidenceLayer",
            "drawRun233VisualEvidenceStoryboard",
            "drawRun233ReadableEvidenceSpine",
            "main_surface_information_density_and_visual_evidence_realism",
            "run2_32_source_audit_status",
            "run2_32_top_next_layer_to_thicken",
            "run2_33_visual_evidence_object_min_target",
            "run2_33_main_surface_visual_evidence_execution_status",
            "main_surface_information_density_and_visual_evidence_realism_thickened_before_native_ppt_generation",
        ],
    )

    prompt_allowed = section(arm_block("prompt_only"), "allowed:", "forbidden:")
    prompt_forbidden = section(arm_block("prompt_only"), "forbidden:", "palette:")
    run1_allowed = section(arm_block("run1_5_skill"), "allowed:", "forbidden:")
    run1_forbidden = section(arm_block("run1_5_skill"), "forbidden:", "palette:")
    full_allowed = section(arm_block("run2_33_full_main_surface_visual_evidence"), "allowed:", "forbidden:")
    full_forbidden = section(arm_block("run2_33_full_main_surface_visual_evidence"), "forbidden:", "palette:")
    bad_allowed = section(arm_block("bad_main_surface_visual_evidence_memory"), "allowed:", "forbidden:")
    bad_forbidden = section(arm_block("bad_main_surface_visual_evidence_memory"), "forbidden:", "palette:")

    for term in required_inputs:
        assert term not in prompt_allowed
        assert term in prompt_forbidden
        assert term not in run1_allowed
        assert term in run1_forbidden
        assert term in full_allowed
        assert term not in full_forbidden
        assert term not in bad_allowed
        assert term in bad_forbidden

    assert 'const fullRun233 = arm.armId === "run2_33_full_main_surface_visual_evidence";' in body
    assert 'registerRun233Module(metrics, "drawRun233MainSurfaceEvidenceLayer")' in body
    for field in [
        "run2_32_source_audit_status",
        "run2_32_top_next_layer_to_thicken",
        "run2_33_visual_evidence_object_min_target",
        "run2_33_main_surface_visual_evidence_execution_status",
    ]:
        assert re.search(fr"{field}:\s*fullRun233\s*\?", body), field


def test_run2_33_records_main_surface_visual_evidence_rerun_result() -> None:
    result = (PACK / "results" / "run2_33_main_surface_visual_evidence_rerun_result.md").read_text(
        encoding="utf-8"
    )
    result_json = load_json(PACK / "results" / "run2_33_main_surface_visual_evidence_rerun_result.json")
    presentations = ROOT / "outputs" / "019e7d9c-532a-70b3-8892-fa3ae42baef2" / "presentations"
    full_trace = load_json(presentations / "ppt-run2-33-full-vulca" / "trace_manifest.json")

    assert result_json["status"] == "run2_33_main_surface_visual_evidence_rerun_public_blocked"
    assert result_json["source_audit_run_id"] == "2.32"
    assert result_json["input_chain"]["spine_climax_repair_audit"].endswith(
        "run2_32_spine_climax_repair_audit.json"
    )
    assert result_json["rerun"]["best_internal_arm"] == "run2_33_full_main_surface_visual_evidence"
    assert result_json["rerun"]["best_internal_arm_verdict"] == (
        "main_surface_information_density_and_visual_evidence_realism_thickened_before_native_ppt_generation"
    )
    assert result_json["quality_delta"]["target_layer"] == "main_surface_information_density_and_visual_evidence_realism"
    assert result_json["quality_delta"]["visual_evidence_object_min_target"] >= 2
    assert result_json["quality_delta"]["source_repair_target_closed"] is True
    assert result_json["rerun"]["combined_contact_sheet"].endswith("run2-33-four-arm-contact-sheet.png")
    assert result_json["rerun"]["full_skill_series_sheet"].endswith("run2-full-skill-series-horizontal.png")

    assert full_trace["arm_id"] == "run2_33_full_main_surface_visual_evidence"
    assert len(full_trace["slides"]) == 6
    for slide in full_trace["slides"]:
        assert slide["run2_32_top_next_layer_to_thicken"] == "main_surface_information_density_and_visual_evidence_realism"
        assert slide["run2_33_visual_evidence_object_min_target"] >= 2
        assert slide["run2_33_main_surface_visual_evidence_execution_status"] == (
            "main_surface_information_density_and_visual_evidence_realism_thickened_before_native_ppt_generation"
        )
        assert slide["layout_metrics"]["visual_evidence_objects"] >= 2
        assert "drawRun233MainSurfaceEvidenceLayer" in slide["run2_33_code_module_ids"]

    assert_contains(
        result,
        [
            "Run 2.33",
            "Run 2.32 spine/climax repair audit",
            "drawRun233MainSurfaceEvidenceLayer",
            "drawRun233VisualEvidenceStoryboard",
            "main_surface_information_density_and_visual_evidence_realism",
            "four-arm rerun",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_33_main_surface_visual_evidence_rerun() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "Run 2.33",
            "ppt-run2-33-prompt-only",
            "ppt-run2-33-run1-5-skill",
            "ppt-run2-33-full-vulca",
            "ppt-run2-33-bad-main-surface-visual-evidence-memory",
            "run2_33_main_surface_visual_evidence_rerun_result.json",
            "drawRun233MainSurfaceEvidenceLayer",
            "drawRun233VisualEvidenceStoryboard",
            "main_surface_information_density_and_visual_evidence_realism",
        ],
    )


def test_run2_34_main_surface_visual_evidence_audit_scores_run2_33_outputs(tmp_path: Path) -> None:
    script_path = ROOT / "scripts" / "audit_ppt_run2_34_main_surface_visual_evidence.py"
    assert script_path.exists(), "missing Run 2.34 main-surface visual-evidence audit script"

    result_json = tmp_path / "run2_34_main_surface_visual_evidence_audit.json"
    result_md = tmp_path / "run2_34_main_surface_visual_evidence_audit.md"
    subprocess.run(
        [
            sys.executable,
            str(script_path),
            "--result-json",
            str(result_json),
            "--result-md",
            str(result_md),
        ],
        cwd=ROOT,
        check=True,
    )

    audit = load_json(result_json)
    assert audit["schema_version"] == "ppt_run2_main_surface_visual_evidence_audit.v1"
    assert audit["run_id"] == "2.34"
    assert audit["status"] == "run2_34_main_surface_visual_evidence_audit_public_blocked"
    assert audit["source_generated_run"] == "2.33"
    assert audit["source_audit_run"] == "2.32"
    assert audit["source_rerun_run"] == "2.33"
    assert audit["creates_new_ppt_deck"] is False
    assert audit["public_ready"] is False
    assert audit["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert audit["input_chain"]["run2_33_full_trace_manifest"].endswith("ppt-run2-33-full-vulca/trace_manifest.json")
    assert audit["input_chain"]["run2_33_bad_trace_manifest"].endswith(
        "ppt-run2-33-bad-main-surface-visual-evidence-memory/trace_manifest.json"
    )
    assert audit["input_chain"]["run2_33_rerun_result"].endswith(
        "run2_33_main_surface_visual_evidence_rerun_result.json"
    )
    assert audit["input_chain"]["run2_32_spine_climax_repair_audit"].endswith(
        "run2_32_spine_climax_repair_audit.json"
    )
    assert audit["no_new_deck_proof"]["new_pptx_created"] is False
    assert audit["no_new_deck_proof"]["status"] == "pass"

    trace = audit["trace_closure"]
    assert trace["full_arm"]["arm_id"] == "run2_33_full_main_surface_visual_evidence"
    assert trace["full_arm"]["slide_count"] == 6
    assert trace["full_arm"]["run2_32_audit_consumed_slides"] == 6
    assert trace["full_arm"]["main_surface_execution_status_slides"] == 6
    assert trace["full_arm"]["main_surface_evidence_layer_modules_called"] == 6
    assert trace["full_arm"]["visual_evidence_storyboard_modules_called"] == 6
    assert trace["full_arm"]["visual_evidence_object_target_slides"] == 6
    assert trace["bad_control"]["arm_id"] == "bad_main_surface_visual_evidence_memory"
    assert trace["bad_control"]["main_surface_fields_leaked"] == 0

    verification = audit["main_surface_verification"]
    assert verification["source_audit_target"] == "main_surface_information_density_and_visual_evidence_realism"
    assert verification["source_rerun_verdict"] == (
        "main_surface_information_density_and_visual_evidence_realism_thickened_before_native_ppt_generation"
    )
    assert verification["visual_evidence_object_target_met"] is True
    assert verification["main_surface_modules_called_by_contract"] is True
    assert verification["delivery_gate"] == "internal-demo-ok-public-blocked"
    assert verification["static_no_animation"] is True
    assert verification["human_review_required"] is True

    summary = audit["quality_summary"]
    assert summary["main_surface_visual_evidence_gate"] == "pass_internal_only"
    assert summary["public_release_gate"] == "blocked"
    assert "main_surface_information_density_and_visual_evidence_realism" in summary["closed_target_layers"]
    assert summary["top_next_layer_to_thicken"] == (
        "usecase_specific_visual_evidence_asset_realism_and_editorial_composition"
    )
    assert set(summary["roles_with_weak_editorial_anchor"]) == {"cover", "setup", "contrast", "close"}
    assert set(summary["roles_with_schematic_visual_evidence"]) == {
        "cover",
        "setup",
        "contrast",
        "proof",
        "climax",
        "close",
    }

    assert len(audit["role_records"]) == 6
    for record in audit["role_records"]:
        assert record["main_surface_trace_closure"]["run2_32_audit_status_present"] is True
        assert record["main_surface_trace_closure"]["execution_status_present"] is True
        assert record["main_surface_modules"]["main_surface_evidence_layer_called"] is True
        assert record["main_surface_modules"]["visual_evidence_storyboard_called"] is True
        assert record["main_surface_density"]["visual_evidence_objects"] >= 2
        assert record["main_surface_density"]["visual_evidence_object_target_met"] is True
        assert record["visual_evidence_realism"]["schematic_visual_evidence"] is True
        assert record["bad_control_leaked_field_count"] == 0


def test_run2_34_records_main_surface_visual_evidence_audit_result() -> None:
    result = (PACK / "results" / "run2_34_main_surface_visual_evidence_audit.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_34_main_surface_visual_evidence_audit.json")

    assert result_json["status"] == "run2_34_main_surface_visual_evidence_audit_public_blocked"
    assert result_json["source_generated_run"] == "2.33"
    assert result_json["quality_summary"]["top_next_layer_to_thicken"] == (
        "usecase_specific_visual_evidence_asset_realism_and_editorial_composition"
    )
    assert result_json["next_required_action"] == (
        "thicken_run2_33_visual_evidence_asset_realism_and_editorial_composition_before_run2_35_data_workflow"
    )
    assert_contains(
        result,
        [
            "Run 2.34 Main Surface Visual Evidence Audit",
            "audit-only",
            "2.33",
            "2.32",
            "drawRun233MainSurfaceEvidenceLayer",
            "drawRun233VisualEvidenceStoryboard",
            "usecase_specific_visual_evidence_asset_realism_and_editorial_composition",
            "Run 2.35 data/workflow",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_embeds_run2_34_main_surface_visual_evidence_audit() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "run2_34_main_surface_visual_evidence_audit.json",
            "Run 2.34 main-surface visual-evidence audit",
            "creates no new PPT deck",
            "main_surface_visual_evidence_gate",
            "roles_with_schematic_visual_evidence",
            "usecase_specific_visual_evidence_asset_realism_and_editorial_composition",
            "Run 2.35 data/workflow",
        ],
    )


def test_run2_35_builder_writes_visual_evidence_realism_workflow_pack(tmp_path: Path) -> None:
    script_path = ROOT / "scripts" / "build_ppt_run2_35_visual_evidence_realism_workflow.py"
    assert script_path.exists(), "missing Run 2.35 visual-evidence realism workflow builder"

    result_json = tmp_path / "run2_35_visual_evidence_realism_workflow_result.json"
    result_md = tmp_path / "run2_35_visual_evidence_realism_workflow_result.md"
    completed = subprocess.run(
        [
            sys.executable,
            str(script_path),
            "--out-dir",
            str(tmp_path),
            "--result-json",
            str(result_json),
            "--result-md",
            str(result_md),
        ],
        cwd=ROOT,
        capture_output=True,
        text=True,
        timeout=20,
    )

    assert completed.returncode == 0, completed.stderr
    realism_memory = load_json(tmp_path / "run2_35_visual_evidence_asset_realism_memory.json")
    composition_memory = load_json(tmp_path / "run2_35_editorial_composition_memory.json")
    workflow_gates = load_json(tmp_path / "run2_35_visual_evidence_workflow_gates.json")
    result = load_json(result_json)
    report = result_md.read_text(encoding="utf-8")

    assert result["status"] == "run2_35_visual_evidence_realism_workflow_ready_public_blocked"
    assert result["run_id"] == "2.35"
    assert result["source_audit_run"] == "2.34"
    assert result["selected_usecase_id"] == "usecase_design_to_production_platform_launch"
    assert result["creates_new_ppt_deck"] is False
    assert result["public_ready"] is False
    assert result["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert result["next_required_action"] == "consume_run2_35_visual_evidence_realism_workflow_before_run2_36_rerun"
    assert result["delivery_artifacts"] == {
        "pptx_paths": [],
        "rendered_slide_paths": [],
        "contact_sheet_paths": [],
        "html_motion_renderer_paths": [],
    }

    assert realism_memory["status"] == "run2_35_visual_evidence_asset_realism_memory_ready_public_blocked"
    assert realism_memory["selected_usecase_id"] == "usecase_design_to_production_platform_launch"
    assert realism_memory["source_audit_target"] == "usecase_specific_visual_evidence_asset_realism_and_editorial_composition"
    assert realism_memory["storage_policy"]["raw_media"] == "forbidden"
    assert realism_memory["storage_policy"]["copied_screenshots"] == "forbidden"
    assert len(realism_memory["visual_evidence_asset_realism_records"]) == 12
    for record in realism_memory["visual_evidence_asset_realism_records"]:
        assert record["realism_memory_id"].startswith("realism_2_35_")
        assert record["role"] in {"cover", "setup", "contrast", "proof", "climax", "close"}
        assert record["parent_run2_24_asset_id"].startswith("asset_2_24_")
        assert record["source_ids"]
        assert record["observable_product_state"]
        assert record["business_context_caption"]
        assert record["audience_question_answered"]
        assert record["native_ppt_realism_strategy"]
        assert record["anti_schematic_constraints"]
        assert "generic block diagram" in " ".join(record["anti_schematic_constraints"])
        assert "run2_35_visual_evidence_asset_realism_ids" in record["required_trace_fields"]
        assert record["public_surface_allowed"] is True
        assert record["source_boundary"].startswith("Derived")

    assert composition_memory["status"] == "run2_35_editorial_composition_memory_ready_public_blocked"
    assert len(composition_memory["editorial_composition_records"]) == 6
    weak_roles = {"cover", "setup", "contrast", "close"}
    for record in composition_memory["editorial_composition_records"]:
        assert record["composition_memory_id"].startswith("composition_2_35_")
        assert record["role"] in {"cover", "setup", "contrast", "proof", "climax", "close"}
        assert record["editorial_anchor_object"]
        assert record["hero_canvas_share_target"] >= 0.35
        assert record["composition_obligations"]
        assert "equal-weight schematic panels" in record["forbidden_patterns"]
        assert "run2_35_editorial_composition_memory_id" in record["required_trace_fields"]
        if record["role"] in weak_roles:
            assert record["repairs_run2_34_weak_editorial_anchor"] is True

    assert workflow_gates["status"] == "run2_35_visual_evidence_workflow_gates_ready_public_blocked"
    assert len(workflow_gates["gates"]) == 6
    for gate in workflow_gates["gates"]:
        assert gate["selected_usecase_id"] == "usecase_design_to_production_platform_launch"
        assert gate["required_realism_memory_ids"]
        assert gate["required_editorial_composition_memory_id"].startswith("composition_2_35_")
        assert gate["min_realistic_visual_evidence_objects"] >= 2
        assert gate["forbid_generic_block_diagrams"] is True
        assert "run2_35_visual_evidence_asset_realism_ids" in gate["required_trace_fields"]
        assert "run2_35_realism_gate_id" in gate["required_trace_fields"]
        assert gate["next_rerun_contract"] == "must_be_consumed_before_run2_36_four_arm_rerun"
        assert gate["public_release_gate"] == "blocked"

    assert_contains(
        report,
        [
            "Run 2.35",
            "visual evidence asset realism",
            "editorial composition",
            "workflow gates",
            "Run 2.35 data/workflow",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )
    assert not list(tmp_path.glob("*.pptx"))


def test_run2_35_records_visual_evidence_realism_workflow_result() -> None:
    result = (PACK / "results" / "run2_35_visual_evidence_realism_workflow_result.md").read_text(
        encoding="utf-8"
    )
    result_json = load_json(PACK / "results" / "run2_35_visual_evidence_realism_workflow_result.json")
    realism_memory = load_json(PACK / "run2_35_visual_evidence_asset_realism_memory.json")
    composition_memory = load_json(PACK / "run2_35_editorial_composition_memory.json")
    workflow_gates = load_json(PACK / "run2_35_visual_evidence_workflow_gates.json")
    workflow = load_json(PACK / "skill_workflow.json")

    assert result_json["status"] == "run2_35_visual_evidence_realism_workflow_ready_public_blocked"
    assert result_json["source_audit_run"] == "2.34"
    assert result_json["selected_usecase_id"] == "usecase_design_to_production_platform_launch"
    assert result_json["creates_new_ppt_deck"] is False
    assert result_json["public_ready"] is False
    assert result_json["artifact_counts"] == {
        "visual_evidence_asset_realism_records": 12,
        "editorial_composition_records": 6,
        "visual_evidence_workflow_gates": 6,
    }
    assert len(realism_memory["visual_evidence_asset_realism_records"]) == 12
    assert len(composition_memory["editorial_composition_records"]) == 6
    assert len(workflow_gates["gates"]) == 6
    assert workflow["status"] == "run2_53_product_surface_scene_workflow_directed_public_blocked"
    assert {stage["id"] for stage in workflow["stages"]} >= {
        "compile_run2_35_visual_evidence_asset_realism_memory",
        "compile_run2_35_editorial_composition_memory",
        "apply_run2_35_visual_evidence_workflow_gates",
        "compile_run2_38_public_video_slide_direction_memory",
        "compile_run2_38_per_slide_visual_recipe_memory",
        "apply_run2_38_public_video_workflow_gates",
    }
    assert any(
        trigger["id"] == "run2_35_visual_evidence_realism_required_before_run2_36_rerun"
        for trigger in workflow["repair_triggers"]
    )
    assert_contains(
        result,
        [
            "Run 2.35",
            "usecase_specific_visual_evidence_asset_realism_and_editorial_composition",
            "visual evidence asset realism",
            "editorial composition",
            "Run 2.36",
            "public blocked",
        ],
    )


def test_ppt_run_html_viewer_embeds_run2_35_visual_evidence_realism_workflow() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "run2_35_visual_evidence_asset_realism_memory.json",
            "run2_35_editorial_composition_memory.json",
            "run2_35_visual_evidence_workflow_gates.json",
            "run2_35_visual_evidence_realism_workflow_result.json",
            "Run 2.35 visual evidence realism workflow",
            "usecase_specific_visual_evidence_asset_realism_and_editorial_composition",
            "Run 2.36",
        ],
    )


def test_ppt_run_html_viewer_loads_run2_35_reference_data(tmp_path: Path) -> None:
    presentations = ROOT / "outputs" / "019e7d9c-532a-70b3-8892-fa3ae42baef2" / "presentations"
    data = build_data(presentations, tmp_path / "ppt-run-viewer.html")
    refs = data["references"]

    assert refs["run235ResultStatus"] == "run2_35_visual_evidence_realism_workflow_ready_public_blocked"
    assert refs["run235Result"]["target_layer"] == (
        "usecase_specific_visual_evidence_asset_realism_and_editorial_composition"
    )
    assert refs["run235Result"]["creates_new_ppt_deck"] is False
    assert refs["run235Result"]["next_required_action"] == (
        "consume_run2_35_visual_evidence_realism_workflow_before_run2_36_rerun"
    )
    assert len(refs["run235RealismMemory"]) == 12
    assert len(refs["run235CompositionMemory"]) == 6
    assert len(refs["run235WorkflowGates"]) == 6
    assert {
        "run2_35_visual_evidence_asset_realism_ids",
        "run2_35_editorial_composition_memory_id",
        "run2_35_realism_gate_id",
    } <= set(refs["run235WorkflowGates"][0]["required_trace_fields"])


def test_run2_36_generator_consumes_run2_35_realism_workflow_before_native_ppt_code() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_36_visual_evidence_realism_arms.mjs"
    assert script_path.exists(), "missing Run 2.36 visual-evidence realism generator"
    body = script_path.read_text(encoding="utf-8")
    arm_order = [
        "prompt_only",
        "run1_5_skill",
        "run2_36_full_visual_evidence_realism",
        "bad_visual_evidence_realism_memory",
    ]

    def arm_block(arm_id: str) -> str:
        start = body.index(f'armId: "{arm_id}"')
        next_starts = [body.find(f'armId: "{next_arm}"', start + 1) for next_arm in arm_order]
        next_starts = [index for index in next_starts if index > start]
        end = min(next_starts) if next_starts else len(body)
        return body[start:end]

    def section(block: str, start_marker: str, end_marker: str) -> str:
        start = block.index(start_marker)
        end = block.index(end_marker, start)
        return block[start:end]

    required_inputs = [
        "run2_35_visual_evidence_asset_realism_memory.json",
        "run2_35_editorial_composition_memory.json",
        "run2_35_visual_evidence_workflow_gates.json",
        "run2_35_visual_evidence_realism_workflow_result.json",
        "run2_34_main_surface_visual_evidence_audit.json",
        "run2_33_main_surface_visual_evidence_rerun_result.json",
    ]

    assert_contains(
        body,
        [
            "validateRun235VisualEvidenceRealismWorkflow",
            "loadRun236ContractData",
            "drawRun236RealisticProductState",
            "drawRun236EditorialAnchorObject",
            "drawRun236RealismGateRibbon",
            "usecase_specific_visual_evidence_asset_realism_and_editorial_composition",
            "run2_35_visual_evidence_asset_realism_ids",
            "run2_35_editorial_composition_memory_id",
            "run2_35_realism_gate_id",
            "run2_35_observable_product_state",
            "run2_35_hero_canvas_share_target",
            "run2_35_visual_evidence_realism_execution_status",
            "usecase_specific_visual_evidence_asset_realism_and_editorial_composition_consumed_before_native_ppt_generation",
        ],
    )

    prompt_allowed = section(arm_block("prompt_only"), "allowed:", "forbidden:")
    prompt_forbidden = section(arm_block("prompt_only"), "forbidden:", "palette:")
    run1_allowed = section(arm_block("run1_5_skill"), "allowed:", "forbidden:")
    run1_forbidden = section(arm_block("run1_5_skill"), "forbidden:", "palette:")
    full_allowed = section(arm_block("run2_36_full_visual_evidence_realism"), "allowed:", "forbidden:")
    full_forbidden = section(arm_block("run2_36_full_visual_evidence_realism"), "forbidden:", "palette:")
    bad_allowed = section(arm_block("bad_visual_evidence_realism_memory"), "allowed:", "forbidden:")
    bad_forbidden = section(arm_block("bad_visual_evidence_realism_memory"), "forbidden:", "palette:")

    for term in required_inputs:
        assert term not in prompt_allowed
        assert term in prompt_forbidden
        assert term not in run1_allowed
        assert term in run1_forbidden
        assert term in full_allowed
        assert term not in full_forbidden
        assert term not in bad_allowed
        assert term in bad_forbidden

    assert 'const fullRun236 = arm.armId === "run2_36_full_visual_evidence_realism";' in body
    assert 'registerRun236Module(metrics, "drawRun236RealisticProductState")' in body
    for field in [
        "run2_35_visual_evidence_asset_realism_ids",
        "run2_35_editorial_composition_memory_id",
        "run2_35_realism_gate_id",
        "run2_35_observable_product_state",
        "run2_35_hero_canvas_share_target",
        "run2_35_visual_evidence_realism_execution_status",
    ]:
        assert re.search(fr"{field}:\s*fullRun236\s*\?", body), field


def test_run2_36_records_visual_evidence_realism_rerun_result() -> None:
    result = (PACK / "results" / "run2_36_visual_evidence_realism_rerun_result.md").read_text(
        encoding="utf-8"
    )
    result_json = load_json(PACK / "results" / "run2_36_visual_evidence_realism_rerun_result.json")
    presentations = ROOT / "outputs" / "019e7d9c-532a-70b3-8892-fa3ae42baef2" / "presentations"
    full_trace = load_json(presentations / "ppt-run2-36-full-vulca" / "trace_manifest.json")
    bad_trace = load_json(
        presentations / "ppt-run2-36-bad-visual-evidence-realism-memory" / "trace_manifest.json"
    )

    assert result_json["status"] == "run2_36_visual_evidence_realism_rerun_public_blocked"
    assert result_json["source_data_workflow_run_id"] == "2.35"
    assert result_json["input_chain"]["visual_evidence_asset_realism_memory"].endswith(
        "run2_35_visual_evidence_asset_realism_memory.json"
    )
    assert result_json["input_chain"]["editorial_composition_memory"].endswith(
        "run2_35_editorial_composition_memory.json"
    )
    assert result_json["input_chain"]["visual_evidence_workflow_gates"].endswith(
        "run2_35_visual_evidence_workflow_gates.json"
    )
    assert result_json["rerun"]["best_internal_arm"] == "run2_36_full_visual_evidence_realism"
    assert result_json["rerun"]["best_internal_arm_verdict"] == (
        "usecase_specific_visual_evidence_asset_realism_and_editorial_composition_consumed_before_native_ppt_generation"
    )
    assert result_json["quality_delta"]["target_layer"] == (
        "usecase_specific_visual_evidence_asset_realism_and_editorial_composition"
    )
    assert result_json["quality_delta"]["run2_35_realism_records_consumed"] == 12
    assert result_json["quality_delta"]["run2_35_composition_records_consumed"] == 6
    assert result_json["quality_delta"]["run2_35_workflow_gates_consumed"] == 6
    assert result_json["rerun"]["combined_contact_sheet"].endswith("run2-36-four-arm-contact-sheet.png")
    assert result_json["rerun"]["full_skill_series_sheet"].endswith("run2-full-skill-series-horizontal.png")

    assert full_trace["arm_id"] == "run2_36_full_visual_evidence_realism"
    assert full_trace["run2_35_workflow_status"] == "run2_35_visual_evidence_realism_workflow_ready_public_blocked"
    assert len(full_trace["slides"]) == 6
    for slide in full_trace["slides"]:
        assert len(slide["run2_35_visual_evidence_asset_realism_ids"]) >= 2
        assert slide["run2_35_editorial_composition_memory_id"].startswith("composition_2_35_")
        assert slide["run2_35_realism_gate_id"].startswith("gate_2_35_")
        assert slide["run2_35_observable_product_state"]
        assert slide["run2_35_hero_canvas_share_target"] >= 0.35
        assert slide["run2_35_visual_evidence_realism_execution_status"] == (
            "usecase_specific_visual_evidence_asset_realism_and_editorial_composition_consumed_before_native_ppt_generation"
        )
        assert slide["layout_metrics"]["realistic_visual_evidence_objects"] >= 2
        assert slide["layout_metrics"]["hero_object_canvas_share"] >= 0.35
        assert "drawRun236RealisticProductState" in slide["run2_36_code_module_ids"]
        assert "drawRun236EditorialAnchorObject" in slide["run2_36_code_module_ids"]
        assert "drawRun236RealismGateRibbon" in slide["run2_36_code_module_ids"]

    assert bad_trace["arm_id"] == "bad_visual_evidence_realism_memory"
    for slide in bad_trace["slides"]:
        assert slide["run2_35_visual_evidence_asset_realism_ids"] == []
        assert slide["run2_35_editorial_composition_memory_id"] == ""
        assert slide["run2_35_realism_gate_id"] == ""

    assert_contains(
        result,
        [
            "Run 2.36",
            "Run 2.35 visual evidence realism workflow",
            "drawRun236RealisticProductState",
            "drawRun236EditorialAnchorObject",
            "usecase_specific_visual_evidence_asset_realism_and_editorial_composition",
            "four-arm rerun",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_36_visual_evidence_realism_rerun() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "Run 2.36",
            "ppt-run2-36-prompt-only",
            "ppt-run2-36-run1-5-skill",
            "ppt-run2-36-full-vulca",
            "ppt-run2-36-bad-visual-evidence-realism-memory",
            "run2_36_visual_evidence_realism_rerun_result.json",
            "drawRun236RealisticProductState",
            "drawRun236EditorialAnchorObject",
            "usecase_specific_visual_evidence_asset_realism_and_editorial_composition",
        ],
    )


def test_run2_37_visual_quality_audit_scores_run2_36_outputs(tmp_path: Path) -> None:
    script_path = ROOT / "scripts" / "audit_ppt_run2_37_visual_quality.py"
    assert script_path.exists(), "missing Run 2.37 visual-quality audit script"

    result_json = tmp_path / "run2_37_visual_quality_audit.json"
    result_md = tmp_path / "run2_37_visual_quality_audit.md"
    subprocess.run(
        [
            sys.executable,
            str(script_path),
            "--result-json",
            str(result_json),
            "--result-md",
            str(result_md),
        ],
        cwd=ROOT,
        check=True,
    )

    audit = load_json(result_json)
    assert audit["schema_version"] == "ppt_run2_visual_quality_audit.v1"
    assert audit["run_id"] == "2.37"
    assert audit["status"] == "run2_37_visual_quality_audit_public_blocked"
    assert audit["source_generated_run"] == "2.36"
    assert audit["source_data_workflow_run"] == "2.35"
    assert audit["creates_new_ppt_deck"] is False
    assert audit["public_ready"] is False
    assert audit["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert audit["input_chain"]["run2_36_full_trace_manifest"].endswith("ppt-run2-36-full-vulca/trace_manifest.json")
    assert audit["input_chain"]["run2_36_bad_trace_manifest"].endswith(
        "ppt-run2-36-bad-visual-evidence-realism-memory/trace_manifest.json"
    )
    assert audit["input_chain"]["run2_36_rerun_result"].endswith(
        "run2_36_visual_evidence_realism_rerun_result.json"
    )
    assert audit["input_chain"]["run2_35_visual_evidence_realism_workflow_result"].endswith(
        "run2_35_visual_evidence_realism_workflow_result.json"
    )
    assert audit["input_chain"]["run2_36_four_arm_contact_sheet"].endswith("run2-36-four-arm-contact-sheet.png")
    assert audit["no_new_deck_proof"]["new_pptx_created"] is False
    assert audit["no_new_deck_proof"]["status"] == "pass"

    trace = audit["trace_closure"]
    assert trace["full_arm"]["arm_id"] == "run2_36_full_visual_evidence_realism"
    assert trace["full_arm"]["slide_count"] == 6
    assert trace["full_arm"]["run2_35_workflow_consumed_slides"] == 6
    assert trace["full_arm"]["realism_ids_bound_slides"] == 6
    assert trace["full_arm"]["required_run236_modules_called_slides"] == 6
    assert trace["bad_control"]["arm_id"] == "bad_visual_evidence_realism_memory"
    assert trace["bad_control"]["visual_evidence_realism_fields_leaked"] == 0

    assessment = audit["visual_quality_assessment"]
    assert assessment["data_consumption_gate"] == "pass_internal_only"
    assert assessment["workflow_proof_gate"] == "pass_internal_only"
    assert assessment["design_quality_gate"] == "blocked"
    assert assessment["public_video_readiness"] == "blocked"
    assert assessment["root_cause_primary"] == "visual_module_language_too_repetitive_and_card_like"
    assert assessment["repeated_layout_signature_count"] == 6
    assert assessment["unique_layout_signature_count"] == 1
    assert assessment["top_next_layer_to_thicken"] == "public_video_grade_slide_direction_and_per_slide_visual_recipe"
    assert set(assessment["roles_with_repetitive_card_layout"]) == {
        "cover",
        "setup",
        "contrast",
        "proof",
        "climax",
        "close",
    }
    assert set(assessment["roles_with_insufficient_public_aesthetic"]) == {
        "cover",
        "setup",
        "contrast",
        "proof",
        "climax",
        "close",
    }

    assert len(audit["role_records"]) == 6
    for record in audit["role_records"]:
        assert record["run2_36_data_consumed"] is True
        assert record["workflow_gate_exposed"] is True
        assert record["public_video_grade"] is False
        assert "repetitive_card_grid_language" in record["aesthetic_failure_reasons"]
        assert "drawRun236RealisticProductState" in record["run2_36_code_module_ids"]
        assert record["layout_signature"] == "editorial_anchor_object+two_product_state_cards+gate_ribbon"


def test_run2_37_records_visual_quality_audit_result() -> None:
    result = (PACK / "results" / "run2_37_visual_quality_audit.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_37_visual_quality_audit.json")

    assert result_json["status"] == "run2_37_visual_quality_audit_public_blocked"
    assert result_json["source_generated_run"] == "2.36"
    assert result_json["user_feedback"] == "Run 2.36 effect feels visually average"
    assert result_json["visual_quality_assessment"]["design_quality_gate"] == "blocked"
    assert result_json["visual_quality_assessment"]["top_next_layer_to_thicken"] == (
        "public_video_grade_slide_direction_and_per_slide_visual_recipe"
    )
    assert result_json["next_required_action"] == (
        "build_run2_38_public_video_grade_visual_direction_memory_and_workflow_gates"
    )
    assert_contains(
        result,
        [
            "Run 2.37 Visual Quality Audit",
            "visually average",
            "audit-only",
            "2.36",
            "2.35",
            "data consumption passes",
            "design quality is blocked",
            "visual_module_language_too_repetitive_and_card_like",
            "public_video_grade_slide_direction_and_per_slide_visual_recipe",
            "Run 2.38 data/workflow",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_embeds_run2_37_visual_quality_audit() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "run2_37_visual_quality_audit.json",
            "Run 2.37 visual quality audit",
            "visual_module_language_too_repetitive_and_card_like",
            "public_video_grade_slide_direction_and_per_slide_visual_recipe",
            "roles_with_repetitive_card_layout",
            "roles_with_insufficient_public_aesthetic",
            "Run 2.38 data/workflow",
        ],
    )


def test_run2_38_builder_writes_public_video_visual_direction_workflow_pack(tmp_path: Path) -> None:
    script_path = ROOT / "scripts" / "build_ppt_run2_38_public_video_visual_direction_workflow.py"
    assert script_path.exists(), "missing Run 2.38 public-video visual-direction workflow builder"

    result_json = tmp_path / "run2_38_public_video_visual_direction_workflow_result.json"
    result_md = tmp_path / "run2_38_public_video_visual_direction_workflow_result.md"
    subprocess.run(
        [
            sys.executable,
            str(script_path),
            "--out-dir",
            str(tmp_path),
            "--result-json",
            str(result_json),
            "--result-md",
            str(result_md),
        ],
        cwd=ROOT,
        check=True,
    )

    result = load_json(result_json)
    direction_memory = load_json(tmp_path / "run2_38_public_video_slide_direction_memory.json")
    recipe_memory = load_json(tmp_path / "run2_38_per_slide_visual_recipe_memory.json")
    workflow_gates = load_json(tmp_path / "run2_38_public_video_workflow_gates.json")

    assert result["status"] == "run2_38_public_video_visual_direction_workflow_ready_public_blocked"
    assert result["run_id"] == "2.38"
    assert result["source_audit_run"] == "2.37"
    assert result["source_generated_run"] == "2.36"
    assert result["creates_new_ppt_deck"] is False
    assert result["public_ready"] is False
    assert result["target_layer"] == "public_video_grade_slide_direction_and_per_slide_visual_recipe"
    assert result["input_chain"]["visual_quality_audit"].endswith("run2_37_visual_quality_audit.json")
    assert result["input_chain"]["visual_evidence_realism_rerun_result"].endswith(
        "run2_36_visual_evidence_realism_rerun_result.json"
    )
    assert result["artifact_counts"] == {
        "public_video_slide_direction_records": 6,
        "per_slide_visual_recipe_records": 6,
        "public_video_workflow_gates": 6,
    }
    assert result["delivery_artifacts"]["pptx_paths"] == []
    assert result["next_required_action"] == (
        "consume_run2_38_public_video_visual_direction_workflow_before_run2_39_rerun"
    )

    assert direction_memory["status"] == "run2_38_public_video_slide_direction_memory_ready_public_blocked"
    assert direction_memory["source_user_feedback"] == "Run 2.36 effect feels visually average"
    direction_records = direction_memory["public_video_slide_direction_records"]
    assert [record["role"] for record in direction_records] == ["cover", "setup", "contrast", "proof", "climax", "close"]
    assert {record["visual_rhythm_id"] for record in direction_records} == {
        "poster_reveal",
        "failure_path_scene",
        "asymmetric_before_after",
        "product_workflow_surface",
        "cinematic_climax_object",
        "decision_handoff_path",
    }
    assert {record["public_video_scene_type"] for record in direction_records} == {
        "launch_poster",
        "failure_scene",
        "before_after_product_state",
        "product_workflow_surface",
        "cinematic_climax_object",
        "decision_handoff_path",
    }
    for record in direction_records:
        assert record["source_audit_finding"] == "visual_module_language_too_repetitive_and_card_like"
        assert record["first_read_object"]
        assert record["commercial_story_payload"]["specific_business_object"]
        assert record["commercial_story_payload"]["viewer_takeaway"]
        assert "same_visual_module_signature_across_all_roles" in record["forbidden_visible_patterns"]
        assert "audit ribbon as public composition" in record["forbidden_visible_patterns"]
        assert "run2_38_visual_direction_memory_id" in record["required_trace_fields"]
        assert "run2_38_public_video_execution_status" in record["required_trace_fields"]
        assert len(record["public_video_grade_acceptance_checks"]) >= 4

    assert recipe_memory["status"] == "run2_38_per_slide_visual_recipe_memory_ready_public_blocked"
    recipe_records = recipe_memory["per_slide_visual_recipe_records"]
    assert [record["role"] for record in recipe_records] == ["cover", "setup", "contrast", "proof", "climax", "close"]
    assert len({record["layout_signature_target"] for record in recipe_records}) == 6
    for record in recipe_records:
        assert record["forbid_run2_36_dominant_layout_signature"] == (
            "editorial_anchor_object+two_product_state_cards+gate_ribbon"
        )
        assert record["show_workflow_gate_as_public_ribbon"] is False
        assert record["primary_visual_weight_target"] >= 0.55
        assert record["typography_recipe"]
        assert record["spacing_recipe"]
        assert record["visual_evidence_recipe"]
        assert record["motion_beat_recipe"]
        assert "run2_38_per_slide_visual_recipe_id" in record["required_trace_fields"]

    assert workflow_gates["status"] == "run2_38_public_video_workflow_gates_ready_public_blocked"
    assert workflow_gates["visual_rhythm_diversity_contract"]["min_unique_visual_rhythms"] == 6
    assert workflow_gates["visual_rhythm_diversity_contract"]["max_repeated_layout_signature_allowed"] == 1
    for gate in workflow_gates["gates"]:
        assert gate["forbid_run2_36_dominant_layout_signature"] is True
        assert gate["required_public_video_slide_direction_memory_id"].startswith("direction_2_38_")
        assert gate["required_per_slide_visual_recipe_memory_id"].startswith("recipe_2_38_")
        assert gate["next_rerun_contract"] == "must_be_consumed_before_run2_39_four_arm_rerun"
        assert "The generated slide must use the Run 2.38 per-slide recipe before native PPT drawing." in gate["pass_fail_checks"]
        assert "The public surface must not expose workflow gates as audit ribbons." in gate["pass_fail_checks"]


def test_run2_38_records_public_video_visual_direction_workflow_result() -> None:
    result = (PACK / "results" / "run2_38_public_video_visual_direction_workflow_result.md").read_text(
        encoding="utf-8"
    )
    result_json = load_json(PACK / "results" / "run2_38_public_video_visual_direction_workflow_result.json")
    direction_memory = load_json(PACK / "run2_38_public_video_slide_direction_memory.json")
    recipe_memory = load_json(PACK / "run2_38_per_slide_visual_recipe_memory.json")
    workflow_gates = load_json(PACK / "run2_38_public_video_workflow_gates.json")
    workflow = load_json(PACK / "skill_workflow.json")

    assert result_json["status"] == "run2_38_public_video_visual_direction_workflow_ready_public_blocked"
    assert result_json["source_audit_run"] == "2.37"
    assert result_json["target_layer"] == "public_video_grade_slide_direction_and_per_slide_visual_recipe"
    assert result_json["artifact_counts"]["public_video_slide_direction_records"] == 6
    assert result_json["artifact_counts"]["per_slide_visual_recipe_records"] == 6
    assert result_json["artifact_counts"]["public_video_workflow_gates"] == 6
    assert result_json["next_required_action"] == (
        "consume_run2_38_public_video_visual_direction_workflow_before_run2_39_rerun"
    )
    assert len(direction_memory["public_video_slide_direction_records"]) == 6
    assert len(recipe_memory["per_slide_visual_recipe_records"]) == 6
    assert len(workflow_gates["gates"]) == 6
    assert "compile_run2_38_public_video_slide_direction_memory" in {stage["id"] for stage in workflow["stages"]}
    assert "apply_run2_38_public_video_workflow_gates" in {stage["id"] for stage in workflow["stages"]}
    assert_contains(
        result,
        [
            "Run 2.38 Public Video Visual Direction Workflow",
            "data/workflow-only",
            "Run 2.37",
            "visual_module_language_too_repetitive_and_card_like",
            "public_video_grade_slide_direction_and_per_slide_visual_recipe",
            "run2_38_public_video_slide_direction_memory.json",
            "run2_38_per_slide_visual_recipe_memory.json",
            "run2_38_public_video_workflow_gates.json",
            "Run 2.39 four-arm rerun",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_embeds_run2_38_public_video_visual_direction_workflow() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "run2_38_public_video_visual_direction_workflow_result.json",
            "Run 2.38 public-video visual direction workflow",
            "run2_38_public_video_slide_direction_memory.json",
            "run2_38_per_slide_visual_recipe_memory.json",
            "run2_38_public_video_workflow_gates.json",
            "public_video_grade_slide_direction_and_per_slide_visual_recipe",
            "visual_rhythm_diversity_contract",
            "must_be_consumed_before_run2_39_four_arm_rerun",
            "Run 2.39 four-arm rerun",
        ],
    )


def test_run2_39_generator_consumes_run2_38_public_video_workflow_before_native_ppt_code() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_39_public_video_visual_direction_arms.mjs"
    assert script_path.exists(), "missing Run 2.39 public-video visual-direction generator"
    body = script_path.read_text(encoding="utf-8")

    assert_sequence(
        body,
        [
            'armId: "prompt_only"',
            'armId: "run1_5_skill"',
            'armId: "run2_39_full_public_video_visual_direction"',
            'armId: "bad_public_video_visual_direction_memory"',
        ],
    )
    assert_contains(
        body,
        [
            "run2_38_public_video_slide_direction_memory.json",
            "run2_38_per_slide_visual_recipe_memory.json",
            "run2_38_public_video_workflow_gates.json",
            "run2_38_public_video_visual_direction_workflow_result.json",
            "run2_37_visual_quality_audit.json",
            "run2_36_visual_evidence_realism_rerun_result.json",
            "validateRun238PublicVideoWorkflow",
            "loadRun239ContractData",
            "drawRun239LaunchPosterStage",
            "drawRun239FailurePathScene",
            "drawRun239AsymmetricBeforeAfterState",
            "drawRun239ProductWorkflowSurface",
            "drawRun239CinematicClimaxObject",
            "drawRun239DecisionHandoffPath",
            "public_video_grade_slide_direction_and_per_slide_visual_recipe_consumed_before_native_ppt_generation",
            "visual_rhythm_diversity_contract",
            "run2_38_visual_direction_memory_id",
            "run2_38_per_slide_visual_recipe_id",
            "run2_38_visual_rhythm_id",
            "run2_38_layout_signature_target",
            "run2_38_public_video_execution_status",
        ],
    )
    assert 'const fullRun239 = arm.armId === "run2_39_full_public_video_visual_direction";' in body
    assert 'registerRun239Module(metrics, "drawRun239LaunchPosterStage")' in body
    assert 'registerRun239Module(metrics, "drawRun239CinematicClimaxObject")' in body
    allowed_section = body[body.index('armId: "run2_39_full_public_video_visual_direction"') : body.index('armId: "bad_public_video_visual_direction_memory"')]
    for filename in [
        "run2_38_public_video_slide_direction_memory.json",
        "run2_38_per_slide_visual_recipe_memory.json",
        "run2_38_public_video_workflow_gates.json",
        "run2_38_public_video_visual_direction_workflow_result.json",
    ]:
        assert filename in allowed_section
    bad_section = body[body.index('armId: "bad_public_video_visual_direction_memory"') :]
    for filename in [
        "run2_38_public_video_slide_direction_memory.json",
        "run2_38_per_slide_visual_recipe_memory.json",
        "run2_38_public_video_workflow_gates.json",
        "run2_38_public_video_visual_direction_workflow_result.json",
    ]:
        assert filename in bad_section


def test_run2_39_records_public_video_visual_direction_rerun_result() -> None:
    result = (PACK / "results" / "run2_39_public_video_visual_direction_rerun_result.md").read_text(
        encoding="utf-8"
    )
    result_json = load_json(PACK / "results" / "run2_39_public_video_visual_direction_rerun_result.json")
    presentations = ROOT / "outputs" / "019e7d9c-532a-70b3-8892-fa3ae42baef2" / "presentations"
    full_trace = load_json(presentations / "ppt-run2-39-full-vulca" / "trace_manifest.json")
    bad_trace = load_json(
        presentations / "ppt-run2-39-bad-public-video-visual-direction-memory" / "trace_manifest.json"
    )

    assert result_json["status"] == "run2_39_public_video_visual_direction_rerun_public_blocked"
    assert result_json["source_data_workflow_run_id"] == "2.38"
    assert result_json["input_chain"]["public_video_slide_direction_memory"].endswith(
        "run2_38_public_video_slide_direction_memory.json"
    )
    assert result_json["input_chain"]["per_slide_visual_recipe_memory"].endswith(
        "run2_38_per_slide_visual_recipe_memory.json"
    )
    assert result_json["input_chain"]["public_video_workflow_gates"].endswith(
        "run2_38_public_video_workflow_gates.json"
    )
    assert result_json["input_chain"]["visual_quality_audit"].endswith("run2_37_visual_quality_audit.json")
    assert result_json["rerun"]["best_internal_arm"] == "run2_39_full_public_video_visual_direction"
    assert result_json["rerun"]["best_internal_arm_verdict"] == (
        "public_video_grade_slide_direction_and_per_slide_visual_recipe_consumed_before_native_ppt_generation"
    )
    assert result_json["quality_delta"]["target_layer"] == (
        "public_video_grade_slide_direction_and_per_slide_visual_recipe"
    )
    assert result_json["quality_delta"]["run2_38_direction_records_consumed"] == 6
    assert result_json["quality_delta"]["run2_38_recipe_records_consumed"] == 6
    assert result_json["quality_delta"]["run2_38_workflow_gates_consumed"] == 6
    assert result_json["quality_delta"]["unique_visual_rhythms"] == 6
    assert result_json["rerun"]["combined_contact_sheet"].endswith("run2-39-four-arm-contact-sheet.png")
    assert result_json["rerun"]["full_skill_series_sheet"].endswith("run2-full-skill-series-horizontal.png")

    assert full_trace["arm_id"] == "run2_39_full_public_video_visual_direction"
    assert full_trace["run2_38_workflow_gate_status"] == "run2_38_public_video_workflow_gates_ready_public_blocked"
    assert len(full_trace["slides"]) == 6
    rhythms = {slide["run2_38_visual_rhythm_id"] for slide in full_trace["slides"]}
    assert rhythms == {
        "poster_reveal",
        "failure_path_scene",
        "asymmetric_before_after",
        "product_workflow_surface",
        "cinematic_climax_object",
        "decision_handoff_path",
    }
    for slide in full_trace["slides"]:
        assert slide["run2_38_visual_direction_memory_id"].startswith("direction_2_38_")
        assert slide["run2_38_per_slide_visual_recipe_id"].startswith("recipe_2_38_")
        assert slide["run2_38_layout_signature_target"]
        assert slide["run2_38_public_video_scene_type"]
        assert slide["run2_38_first_read_object"]
        assert slide["run2_38_public_video_execution_status"] == (
            "public_video_grade_slide_direction_and_per_slide_visual_recipe_consumed_before_native_ppt_generation"
        )
        assert slide["layout_metrics"]["primary_visual_weight"] >= 0.55
        assert slide["layout_metrics"]["run2_36_dominant_layout_signature_reused"] is False
        assert len(slide["run2_39_code_module_ids"]) == 1
        assert slide["run2_39_code_module_ids"][0].startswith("drawRun239")

    assert bad_trace["arm_id"] == "bad_public_video_visual_direction_memory"
    for slide in bad_trace["slides"]:
        assert slide["run2_38_visual_direction_memory_id"] == ""
        assert slide["run2_38_per_slide_visual_recipe_id"] == ""
        assert slide["run2_38_visual_rhythm_id"] == ""

    assert_contains(
        result,
        [
            "Run 2.39",
            "Run 2.38 public-video visual direction workflow",
            "drawRun239LaunchPosterStage",
            "drawRun239CinematicClimaxObject",
            "public_video_grade_slide_direction_and_per_slide_visual_recipe",
            "four-arm rerun",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_39_public_video_visual_direction_rerun() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "Run 2.39",
            "ppt-run2-39-prompt-only",
            "ppt-run2-39-run1-5-skill",
            "ppt-run2-39-full-vulca",
            "ppt-run2-39-bad-public-video-visual-direction-memory",
            "run2_39_public_video_visual_direction_rerun_result.json",
            "drawRun239LaunchPosterStage",
            "drawRun239CinematicClimaxObject",
            "public_video_grade_slide_direction_and_per_slide_visual_recipe",
        ],
    )


def test_ppt_run_html_viewer_generated_includes_run2_39() -> None:
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    assert_contains(
        viewer,
        [
            "Run 2.39",
            "run2-39-four-arm-contact-sheet.png",
            "ppt-run2-39-prompt-only",
            "ppt-run2-39-run1-5-skill",
            "ppt-run2-39-full-vulca",
            "ppt-run2-39-bad-public-video-visual-direction-memory",
            "run2_39_public_video_visual_direction_rerun_result.json",
        ],
    )


def test_ppt_run_html_viewer_generated_latest_run2_50() -> None:
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    assert_contains(
        viewer,
        [
            "Run 2.40",
            "run2-40-four-arm-contact-sheet.png",
            "ppt-run2-40-prompt-only",
            "ppt-run2-40-run1-5-skill",
            "ppt-run2-40-full-vulca",
            "ppt-run2-40-bad-trace-visible-visual-compiler",
            "run2_40_visual_compiler_rerun_result.json",
            "Run 2.41",
            "run2-41-four-arm-contact-sheet.png",
            "ppt-run2-41-prompt-only",
            "ppt-run2-41-run1-5-skill",
            "ppt-run2-41-full-vulca",
            "ppt-run2-41-bad-thin-content-visual-asset-compiler",
            "run2_41_content_visual_asset_compiler_rerun_result.json",
            "Run 2.44",
            "run2-44-four-arm-contact-sheet.png",
            "ppt-run2-44-prompt-only",
            "ppt-run2-44-run1-5-skill",
            "ppt-run2-44-full-vulca",
            "ppt-run2-44-bad-run2-43-name-only-geometry",
            "run2_44_semantic_geometry_rerun_result.json",
            '"latestRunId": "2.77"',
            "Run 2.47",
            "run2-47-four-arm-contact-sheet.png",
            "ppt-run2-47-prompt-only",
            "ppt-run2-47-run1-5-skill",
            "ppt-run2-47-full-vulca",
            "ppt-run2-47-bad-missing-composition-grammar",
            "run2_47_composition_grammar_rerun_result.json",
            "Run 2.50",
            "run2-50-four-arm-contact-sheet.png",
            "ppt-run2-50-prompt-only",
            "ppt-run2-50-run1-5-skill",
            "ppt-run2-50-full-vulca",
            "ppt-run2-50-bad-missing-run2-49-repair-pack",
            "run2_50_readability_density_renderer_rerun_result.json",
        ],
    )


def test_run2_40_generator_uses_same_data_with_hidden_trace_visual_compiler() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_40_visual_compiler_arms.mjs"
    assert script_path.exists(), "missing Run 2.40 visual compiler generator"
    body = script_path.read_text(encoding="utf-8")

    assert_sequence(
        body,
        [
            'armId: "prompt_only"',
            'armId: "run1_5_skill"',
            'armId: "run2_40_full_visual_compiler"',
            'armId: "bad_trace_visible_visual_compiler"',
        ],
    )
    assert_contains(
        body,
        [
            "run2_39_public_video_visual_direction_rerun_result.json",
            "run2_38_public_video_slide_direction_memory.json",
            "run2_38_per_slide_visual_recipe_memory.json",
            "run2_38_public_video_workflow_gates.json",
            "visualCompilerTransform",
            "validateRun240SameDataVisualCompiler",
            "same_run2_38_run2_39_data_no_database_expansion",
            "trace_to_hidden_constraints_public_surface_composition",
            "public_surface_machinery_hidden",
            "drawRun240EditorialPoster",
            "drawRun240UsecaseScene",
            "drawRun240TransformationSpread",
            "drawRun240ProductMoment",
            "drawRun240CinematicResult",
            "drawRun240DecisionScene",
            "run2_40_visual_compiler_policy",
            "run2_40_public_surface_machinery_hidden",
            "run2_40_visible_machinery_terms",
        ],
    )
    full_section = body[
        body.index('armId: "run2_40_full_visual_compiler"') : body.index('armId: "bad_trace_visible_visual_compiler"')
    ]
    bad_section = body[body.index('armId: "bad_trace_visible_visual_compiler"') :]
    for filename in [
        "run2_39_public_video_visual_direction_rerun_result.json",
        "run2_38_public_video_slide_direction_memory.json",
        "run2_38_per_slide_visual_recipe_memory.json",
        "run2_38_public_video_workflow_gates.json",
    ]:
        assert filename in full_section
        assert filename in bad_section


def test_run2_40_records_same_data_visual_compiler_rerun_result() -> None:
    result = (PACK / "results" / "run2_40_visual_compiler_rerun_result.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_40_visual_compiler_rerun_result.json")
    presentations = ROOT / "outputs" / "019e7d9c-532a-70b3-8892-fa3ae42baef2" / "presentations"
    full_trace = load_json(presentations / "ppt-run2-40-full-vulca" / "trace_manifest.json")
    bad_trace = load_json(presentations / "ppt-run2-40-bad-trace-visible-visual-compiler" / "trace_manifest.json")

    assert result_json["status"] == "run2_40_visual_compiler_rerun_public_blocked"
    assert result_json["source_data_workflow_run_id"] == "2.38"
    assert result_json["source_generated_run_id"] == "2.39"
    assert result_json["database_expansion"] is False
    assert result_json["input_chain"]["run2_39_rerun_result"].endswith(
        "run2_39_public_video_visual_direction_rerun_result.json"
    )
    assert result_json["rerun"]["best_internal_arm"] == "run2_40_full_visual_compiler"
    assert result_json["rerun"]["best_internal_arm_verdict"] == (
        "same_run2_38_run2_39_data_no_database_expansion_visual_compiler_improves_public_surface"
    )
    assert result_json["quality_delta"]["target_layer"] == "visual_compiler_hidden_trace_public_composition"
    assert result_json["quality_delta"]["same_data_control"] is True
    assert result_json["quality_delta"]["new_database_records_added"] == 0
    assert result_json["quality_delta"]["public_surface_machinery_hidden_slides"] == 6
    assert result_json["rerun"]["combined_contact_sheet"].endswith("run2-40-four-arm-contact-sheet.png")
    assert result_json["rerun"]["full_skill_series_sheet"].endswith("run2-full-skill-series-horizontal.png")

    assert full_trace["arm_id"] == "run2_40_full_visual_compiler"
    assert full_trace["run2_40_visual_compiler_status"] == "same_data_visual_compiler_applied"
    assert full_trace["run2_39_source_rerun_status"] == "run2_39_public_video_visual_direction_rerun_public_blocked"
    assert len(full_trace["slides"]) == 6
    for slide in full_trace["slides"]:
        assert slide["run2_38_visual_direction_memory_id"].startswith("direction_2_38_")
        assert slide["run2_38_per_slide_visual_recipe_id"].startswith("recipe_2_38_")
        assert slide["run2_40_visual_compiler_policy"] == "trace_to_hidden_constraints_public_surface_composition"
        assert slide["run2_40_public_surface_machinery_hidden"] is True
        assert slide["run2_40_visible_machinery_terms"] == 0
        assert slide["layout_metrics"]["presentation_surface_weight"] >= 0.7
        assert slide["layout_metrics"]["trace_panel_visible"] is False
        assert slide["layout_metrics"]["gate_ribbon_visible"] is False
        assert len(slide["run2_40_code_module_ids"]) == 1
        assert slide["run2_40_code_module_ids"][0].startswith("drawRun240")

    assert bad_trace["arm_id"] == "bad_trace_visible_visual_compiler"
    assert bad_trace["run2_40_visual_compiler_status"] == "same_data_trace_visible_control"
    for slide in bad_trace["slides"]:
        assert slide["run2_38_visual_direction_memory_id"].startswith("direction_2_38_")
        assert slide["run2_40_public_surface_machinery_hidden"] is False
        assert slide["run2_40_visible_machinery_terms"] >= 3
        assert slide["layout_metrics"]["trace_panel_visible"] is True

    assert_contains(
        result,
        [
            "Run 2.40",
            "same Run 2.38 and Run 2.39 data",
            "visual compiler",
            "trace_to_hidden_constraints_public_surface_composition",
            "drawRun240EditorialPoster",
            "drawRun240CinematicResult",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_40_visual_compiler() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "Run 2.40",
            "ppt-run2-40-prompt-only",
            "ppt-run2-40-run1-5-skill",
            "ppt-run2-40-full-vulca",
            "ppt-run2-40-bad-trace-visible-visual-compiler",
            "run2_40_visual_compiler_rerun_result.json",
            "visual_compiler_hidden_trace_public_composition",
            "same_run2_38_run2_39_data_no_database_expansion",
            "drawRun240EditorialPoster",
            "drawRun240CinematicResult",
        ],
    )


def test_run2_41_generator_thickens_content_and_visual_assets_without_database_expansion() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_41_content_visual_asset_arms.mjs"
    assert script_path.exists(), "missing Run 2.41 content/visual asset compiler generator"
    body = script_path.read_text(encoding="utf-8")

    assert_sequence(
        body,
        [
            'armId: "prompt_only"',
            'armId: "run1_5_skill"',
            'armId: "run2_41_full_content_visual_asset_compiler"',
            'armId: "bad_thin_content_visual_asset_compiler"',
        ],
    )
    assert_contains(
        body,
        [
            "run2_40_visual_compiler_rerun_result.json",
            "run2_39_public_video_visual_direction_rerun_result.json",
            "run2_38_public_video_slide_direction_memory.json",
            "commercial_usecase_bank.json",
            "sources.json",
            "contentVisualAssetCompilerTransform",
            "validateRun241ContentVisualAssetCompiler",
            "same_database_no_workflow_expansion_content_visual_asset_compiler",
            "visual_asset_surface_from_existing_sources_not_copied_media",
            "public_surface_machinery_hidden",
            "visible_business_detail_count",
            "visual_asset_surface_count",
            "editorial_scene_depth_score",
            "drawRun241MarketScenePoster",
            "drawRun241FailureStoryboard",
            "drawRun241BeforeAfterBusinessCase",
            "drawRun241ProductUiEvidenceScene",
            "drawRun241CinematicLaunchMoment",
            "drawRun241ReviewDecisionRoom",
        ],
    )


def test_run2_41_records_content_visual_asset_compiler_rerun_result() -> None:
    result = (PACK / "results" / "run2_41_content_visual_asset_compiler_rerun_result.md").read_text(
        encoding="utf-8"
    )
    result_json = load_json(PACK / "results" / "run2_41_content_visual_asset_compiler_rerun_result.json")
    presentations = ROOT / "outputs" / "019e7d9c-532a-70b3-8892-fa3ae42baef2" / "presentations"
    full_trace = load_json(presentations / "ppt-run2-41-full-vulca" / "trace_manifest.json")
    bad_trace = load_json(presentations / "ppt-run2-41-bad-thin-content-visual-asset-compiler" / "trace_manifest.json")

    assert result_json["status"] == "run2_41_content_visual_asset_compiler_rerun_public_blocked"
    assert result_json["source_generated_run_id"] == "2.40"
    assert result_json["database_expansion"] is False
    assert result_json["workflow_expansion"] is False
    assert result_json["input_chain"]["run2_40_rerun_result"].endswith("run2_40_visual_compiler_rerun_result.json")
    assert result_json["rerun"]["best_internal_arm"] == "run2_41_full_content_visual_asset_compiler"
    assert result_json["rerun"]["best_internal_arm_verdict"] == (
        "content_visual_asset_compiler_thickens_public_surface_without_database_or_workflow_expansion"
    )
    assert result_json["quality_delta"]["target_layer"] == "content_visual_asset_composition_thickness"
    assert result_json["quality_delta"]["source_data_status"] == "same_database_no_workflow_expansion_content_visual_asset_compiler"
    assert result_json["quality_delta"]["new_database_records_added"] == 0
    assert result_json["quality_delta"]["new_workflow_gates_added"] == 0
    assert result_json["quality_delta"]["full_slides_with_visible_business_detail_count_min_5"] == 6
    assert result_json["quality_delta"]["full_slides_with_visual_asset_surface_count_min_3"] == 6
    assert result_json["quality_delta"]["full_slides_with_machinery_hidden"] == 6
    assert result_json["rerun"]["combined_contact_sheet"].endswith("run2-41-four-arm-contact-sheet.png")
    assert result_json["rerun"]["full_skill_series_sheet"].endswith("run2-full-skill-series-horizontal.png")

    assert full_trace["arm_id"] == "run2_41_full_content_visual_asset_compiler"
    assert full_trace["run2_41_content_visual_asset_compiler_status"] == "same_database_content_visual_asset_compiler_applied"
    assert full_trace["run2_40_source_rerun_status"] == "run2_40_visual_compiler_rerun_public_blocked"
    assert len(full_trace["slides"]) == 6
    for slide in full_trace["slides"]:
        assert slide["run2_41_public_surface_machinery_hidden"] is True
        assert slide["run2_41_visible_business_detail_count"] >= 5
        assert slide["run2_41_visual_asset_surface_count"] >= 3
        assert slide["run2_41_editorial_scene_depth_score"] >= 0.7
        assert slide["run2_41_content_scene_payload"]
        assert slide["run2_41_visual_asset_surface_types"]
        assert slide["layout_metrics"]["presentation_surface_weight"] >= 0.72
        assert slide["layout_metrics"]["text_density_tier"] in {"medium", "rich"}
        assert slide["layout_metrics"]["trace_panel_visible"] is False
        assert len(slide["run2_41_code_module_ids"]) == 1
        assert slide["run2_41_code_module_ids"][0].startswith("drawRun241")

    assert bad_trace["arm_id"] == "bad_thin_content_visual_asset_compiler"
    assert bad_trace["workflow_expansion"] is False
    for slide in bad_trace["slides"]:
        assert slide["run2_41_public_surface_machinery_hidden"] is True
        assert slide["run2_41_visible_machinery_terms"] == 0
        assert slide["run2_41_visible_business_detail_count"] <= 2
        assert slide["run2_41_visual_asset_surface_count"] <= 1
        assert slide["run2_41_editorial_scene_depth_score"] < 0.5
        assert slide["layout_metrics"]["trace_panel_visible"] is False

    assert_contains(
        result,
        [
            "Run 2.41",
            "content visual asset compiler",
            "same database",
            "no workflow expansion",
            "drawRun241MarketScenePoster",
            "drawRun241CinematicLaunchMoment",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_41_content_visual_asset_compiler() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "Run 2.41",
            "ppt-run2-41-prompt-only",
            "ppt-run2-41-run1-5-skill",
            "ppt-run2-41-full-vulca",
            "ppt-run2-41-bad-thin-content-visual-asset-compiler",
            "run2_41_content_visual_asset_compiler_rerun_result.json",
            "content_visual_asset_composition_thickness",
            "same_database_no_workflow_expansion_content_visual_asset_compiler",
            "drawRun241MarketScenePoster",
            "drawRun241CinematicLaunchMoment",
        ],
    )


def test_run2_42_content_visual_asset_quality_audit_scores_run2_41_outputs(tmp_path: Path) -> None:
    script_path = ROOT / "scripts" / "audit_ppt_run2_42_content_visual_asset_quality.py"
    assert script_path.exists(), "missing Run 2.42 content/visual asset quality audit script"
    result_json = tmp_path / "run2_42_content_visual_asset_quality_audit.json"
    result_md = tmp_path / "run2_42_content_visual_asset_quality_audit.md"

    subprocess.run(
        [
            sys.executable,
            str(script_path),
            "--result-json",
            str(result_json),
            "--result-md",
            str(result_md),
        ],
        cwd=ROOT,
        check=True,
    )

    audit = load_json(result_json)
    assert audit["schema_version"] == "ppt_run2_content_visual_asset_quality_audit.v1"
    assert audit["run_id"] == "2.42"
    assert audit["status"] == "run2_42_content_visual_asset_quality_audit_public_blocked"
    assert audit["source_generated_run"] == "2.41"
    assert audit["source_prior_generated_run"] == "2.40"
    assert audit["creates_new_ppt_deck"] is False
    assert audit["public_ready"] is False
    assert audit["no_new_deck_proof"]["status"] == "pass"
    assert audit["no_new_deck_proof"]["new_pptx_created"] is False

    assert audit["input_chain"]["run2_41_full_trace_manifest"].endswith("ppt-run2-41-full-vulca/trace_manifest.json")
    assert audit["input_chain"]["run2_41_bad_trace_manifest"].endswith(
        "ppt-run2-41-bad-thin-content-visual-asset-compiler/trace_manifest.json"
    )
    assert audit["input_chain"]["run2_41_rerun_result"].endswith(
        "run2_41_content_visual_asset_compiler_rerun_result.json"
    )
    assert audit["input_chain"]["run2_40_rerun_result"].endswith("run2_40_visual_compiler_rerun_result.json")

    trace = audit["trace_closure"]
    assert trace["full_arm"]["arm_id"] == "run2_41_full_content_visual_asset_compiler"
    assert trace["full_arm"]["slide_count"] == 6
    assert trace["full_arm"]["content_visual_asset_compiler_applied_slides"] == 6
    assert trace["full_arm"]["visible_business_detail_min5_slides"] == 6
    assert trace["full_arm"]["visual_asset_surface_min3_slides"] == 6
    assert trace["full_arm"]["machinery_hidden_slides"] == 6
    assert trace["bad_control"]["arm_id"] == "bad_thin_content_visual_asset_compiler"
    assert trace["bad_control"]["thin_content_control_slides"] == 6
    assert trace["bad_control"]["machinery_leak_slides"] == 0

    assessment = audit["visual_quality_assessment"]
    assert assessment["content_visual_asset_gate"] == "pass_internal_only"
    assert assessment["same_database_control_gate"] == "pass_internal_only"
    assert assessment["design_quality_gate"] == "blocked"
    assert assessment["public_video_readiness"] == "blocked"
    assert assessment["root_cause_primary"] == (
        "visual_asset_surfaces_are_still_schematic_native_shapes_not_true_product_or_scene_assets"
    )
    assert assessment["top_next_layer_to_thicken"] == (
        "usecase_specific_visual_asset_semantics_editorial_composition_and_typography_hierarchy"
    )
    assert assessment["content_visual_asset_thickness_delta_from_bad_control"] == "proven"

    assert len(audit["role_records"]) == 6
    for record in audit["role_records"]:
        assert record["run2_41_data_consumed"] is True
        assert record["content_visual_asset_thickness_passed"] is True
        assert record["bad_control_boundary_passed"] is True
        assert record["public_video_grade"] is False
        assert record["full_arm_metrics"]["visible_business_detail_count"] >= 5
        assert record["full_arm_metrics"]["visual_asset_surface_count"] >= 3
        assert record["bad_control_metrics"]["visible_business_detail_count"] <= 2
        assert record["bad_control_metrics"]["visual_asset_surface_count"] <= 1
        assert "visual_asset_surface_is_named_but_not_semantically_rendered_enough" in record["aesthetic_failure_reasons"]


def test_run2_42_records_content_visual_asset_quality_audit_result() -> None:
    result = (PACK / "results" / "run2_42_content_visual_asset_quality_audit.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_42_content_visual_asset_quality_audit.json")

    assert result_json["status"] == "run2_42_content_visual_asset_quality_audit_public_blocked"
    assert result_json["source_generated_run"] == "2.41"
    assert result_json["source_prior_generated_run"] == "2.40"
    assert result_json["visual_quality_assessment"]["content_visual_asset_gate"] == "pass_internal_only"
    assert result_json["visual_quality_assessment"]["design_quality_gate"] == "blocked"
    assert result_json["visual_quality_assessment"]["top_next_layer_to_thicken"] == (
        "usecase_specific_visual_asset_semantics_editorial_composition_and_typography_hierarchy"
    )
    assert result_json["next_required_action"] == (
        "build_run2_43_visual_asset_semantics_editorial_composition_workflow"
    )
    assert_contains(
        result,
        [
            "Run 2.42 Content Visual Asset Quality Audit",
            "audit-only",
            "Run 2.41",
            "content visual asset thickness passes",
            "design quality is blocked",
            "visual_asset_surfaces_are_still_schematic_native_shapes_not_true_product_or_scene_assets",
            "usecase_specific_visual_asset_semantics_editorial_composition_and_typography_hierarchy",
            "Run 2.43 data/workflow",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_embeds_run2_42_content_visual_asset_quality_audit() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "run2_42_content_visual_asset_quality_audit.json",
            "Run 2.42 content visual asset quality audit",
            "visual_asset_surfaces_are_still_schematic_native_shapes_not_true_product_or_scene_assets",
            "usecase_specific_visual_asset_semantics_editorial_composition_and_typography_hierarchy",
            "content_visual_asset_thickness_delta_from_bad_control",
            "build_run2_43_visual_asset_semantics_editorial_composition_workflow",
        ],
    )


def test_run2_43_builder_writes_visual_asset_semantics_workflow_pack(tmp_path: Path) -> None:
    script_path = ROOT / "scripts" / "build_ppt_run2_43_visual_asset_semantics_workflow.py"
    assert script_path.exists(), "missing Run 2.43 visual asset semantics workflow script"
    result_json = tmp_path / "run2_43_visual_asset_semantics_workflow_result.json"
    result_md = tmp_path / "run2_43_visual_asset_semantics_workflow_result.md"

    subprocess.run(
        [
            sys.executable,
            str(script_path),
            "--out-dir",
            str(tmp_path),
            "--result-json",
            str(result_json),
            "--result-md",
            str(result_md),
        ],
        cwd=ROOT,
        check=True,
    )

    semantic_memory_path = tmp_path / "run2_43_semantic_visual_asset_memory.json"
    editorial_memory_path = tmp_path / "run2_43_editorial_composition_typography_memory.json"
    workflow_gates_path = tmp_path / "run2_43_visual_asset_semantics_workflow_gates.json"
    assert semantic_memory_path.exists()
    assert editorial_memory_path.exists()
    assert workflow_gates_path.exists()

    result = load_json(result_json)
    assert result["schema_version"] == "ppt_run2_visual_asset_semantics_workflow_result.v1"
    assert result["run_id"] == "2.43"
    assert result["status"] == "run2_43_visual_asset_semantics_workflow_ready_public_blocked"
    assert result["source_audit_run"] == "2.42"
    assert result["source_generated_run"] == "2.41"
    assert result["source_prior_generated_run"] == "2.40"
    assert result["creates_new_ppt_deck"] is False
    assert result["public_ready"] is False
    assert result["target_layer"] == (
        "usecase_specific_visual_asset_semantics_editorial_composition_and_typography_hierarchy"
    )
    assert result["next_required_action"] == (
        "consume_run2_43_visual_asset_semantics_workflow_before_run2_44_rerun"
    )
    assert result["output_chain"]["semantic_visual_asset_memory"].endswith(
        "run2_43_semantic_visual_asset_memory.json"
    )
    assert result["output_chain"]["editorial_composition_typography_memory"].endswith(
        "run2_43_editorial_composition_typography_memory.json"
    )
    assert result["output_chain"]["visual_asset_semantics_workflow_gates"].endswith(
        "run2_43_visual_asset_semantics_workflow_gates.json"
    )
    assert result["artifact_counts"]["semantic_visual_asset_records"] == 18
    assert result["artifact_counts"]["editorial_composition_typography_records"] == 6
    assert result["artifact_counts"]["workflow_gates"] == 6

    semantic_memory = load_json(semantic_memory_path)
    assert semantic_memory["status"] == "run2_43_semantic_visual_asset_memory_ready_public_blocked"
    assert semantic_memory["source_boundary"]["copied_screenshots"] == "forbidden"
    assert semantic_memory["source_boundary"]["raw_tutorial_or_video_media"] == "forbidden"
    semantic_records = semantic_memory["semantic_visual_asset_records"]
    assert len(semantic_records) == 18
    for record in semantic_records:
        assert record["semantic_asset_id"].startswith("semantic_asset_2_43_")
        assert record["source_audit_run"] == "2.42"
        assert record["source_generated_run"] == "2.41"
        assert record["source_run2_41_surface_type"]
        assert record["usecase_specific_semantic_object"]
        assert record["observable_scene_evidence"]
        assert record["allowed_native_ppt_representation"]
        assert "generic placeholder" in record["forbidden_schematic_substitutes"]
        assert "copied screenshots/source media" in record["forbidden_schematic_substitutes"]
        assert "no copied screenshots/source media" in record["source_boundary_note"]

    editorial_memory = load_json(editorial_memory_path)
    assert editorial_memory["status"] == (
        "run2_43_editorial_composition_typography_memory_ready_public_blocked"
    )
    editorial_records = editorial_memory["editorial_composition_typography_records"]
    assert len(editorial_records) == 6
    for record in editorial_records:
        assert record["editorial_typography_memory_id"].startswith("editorial_typography_2_43_")
        assert len(record["required_semantic_asset_ids"]) == 3
        assert len(record["typography_hierarchy_rules"]) >= 3
        assert len(record["editorial_composition_rules"]) >= 3
        assert record["climax_or_scene_depth_rule"]

    gates = load_json(workflow_gates_path)
    assert gates["status"] == "run2_43_visual_asset_semantics_workflow_gates_ready_public_blocked"
    assert gates["next_rerun_contract"] == "must_be_consumed_before_run2_44_four_arm_rerun"
    gate_records = gates["visual_asset_semantics_workflow_gates"]
    assert len(gate_records) == 6
    for gate in gate_records:
        assert len(gate["required_semantic_asset_ids"]) == 3
        assert gate["required_editorial_typography_memory_id"].startswith("editorial_typography_2_43_")
        assert gate["forbid_generic_placeholder_assets"] is True
        assert gate["forbid_schematic_native_shape_only_surface"] is True
        assert gate["forbid_copied_source_media"] is True
        assert set(gate["required_trace_fields"]) >= {
            "run2_43_semantic_visual_asset_ids",
            "run2_43_editorial_typography_memory_id",
            "run2_43_visual_asset_semantics_gate_id",
            "run2_43_visual_asset_semantics_execution_status",
        }


def test_run2_43_records_visual_asset_semantics_workflow_result() -> None:
    result = (PACK / "results" / "run2_43_visual_asset_semantics_workflow_result.md").read_text(
        encoding="utf-8"
    )
    result_json = load_json(PACK / "results" / "run2_43_visual_asset_semantics_workflow_result.json")
    workflow = load_json(PACK / "skill_workflow.json")

    assert (PACK / "run2_43_semantic_visual_asset_memory.json").exists()
    assert (PACK / "run2_43_editorial_composition_typography_memory.json").exists()
    assert (PACK / "run2_43_visual_asset_semantics_workflow_gates.json").exists()
    assert result_json["status"] == "run2_43_visual_asset_semantics_workflow_ready_public_blocked"
    assert result_json["source_audit_run"] == "2.42"
    assert result_json["source_generated_run"] == "2.41"
    assert result_json["source_prior_generated_run"] == "2.40"
    assert result_json["target_layer"] == (
        "usecase_specific_visual_asset_semantics_editorial_composition_and_typography_hierarchy"
    )
    assert result_json["next_required_action"] == (
        "consume_run2_43_visual_asset_semantics_workflow_before_run2_44_rerun"
    )
    assert workflow["status"] == "run2_53_product_surface_scene_workflow_directed_public_blocked"
    assert {stage["id"] for stage in workflow["stages"]} >= {
        "compile_run2_43_semantic_visual_asset_memory",
        "compile_run2_43_editorial_composition_typography_memory",
        "apply_run2_43_visual_asset_semantics_workflow_gates",
    }
    assert any(
        trigger["id"] == "run2_43_visual_asset_semantics_required_before_run2_44_rerun"
        for trigger in workflow["repair_triggers"]
    )
    assert_contains(
        result,
        [
            "Run 2.43 Visual Asset Semantics Workflow",
            "data/workflow-only",
            "Run 2.42",
            "semantic visual assets",
            "editorial composition",
            "typography hierarchy",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_embeds_run2_43_visual_asset_semantics_workflow() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "run2_43_visual_asset_semantics_workflow_result.json",
            "Run 2.43 visual asset semantics workflow",
            "run2_43_semantic_visual_asset_memory.json",
            "run2_43_editorial_composition_typography_memory.json",
            "run2_43_visual_asset_semantics_workflow_gates.json",
            "usecase_specific_visual_asset_semantics_editorial_composition_and_typography_hierarchy",
            "consume_run2_43_visual_asset_semantics_workflow_before_run2_44_rerun",
        ],
    )


def test_run2_44_dataflow_readiness_audit_flags_current_consumption_bug(tmp_path: Path) -> None:
    script_path = ROOT / "scripts" / "audit_ppt_run2_44_dataflow_readiness.py"
    assert script_path.exists(), "missing Run 2.44 dataflow readiness audit script"
    result_json = tmp_path / "run2_44_dataflow_readiness_audit.json"
    result_md = tmp_path / "run2_44_dataflow_readiness_audit.md"

    subprocess.run(
        [
            sys.executable,
            str(script_path),
            "--result-json",
            str(result_json),
            "--result-md",
            str(result_md),
        ],
        cwd=ROOT,
        check=True,
    )

    audit = load_json(result_json)
    assert audit["schema_version"] == "ppt_run2_dataflow_readiness_audit.v1"
    assert audit["run_id"] == "2.44-preflight"
    assert audit["status"] == "run2_44_dataflow_readiness_blocked"
    assert audit["bug_confirmed"] is True
    assert audit["not_a_run2_44_output"] is True
    assert audit["risk_priority"] == "high"
    assert audit["creates_new_ppt_deck"] is False
    assert audit["public_ready"] is False
    assert audit["latest_visible_ppt_run"] == "2.41"
    assert audit["latest_workflow_run"] == "2.43"
    assert audit["root_cause_primary"] == "latest_visible_ppt_does_not_consume_latest_run2_43_workflow"

    flow = audit["dataflow_findings"]
    assert flow["run2_43_consumed_by_latest_visible_ppt"] is False
    assert flow["run2_41_generator_reads_run2_38_data"] is True
    assert flow["run2_41_generator_reads_run2_43_data"] is False
    assert flow["run2_41_trace_has_run2_38_direction_and_recipe_ids"] is True
    assert flow["run2_41_trace_has_run2_38_gate_ids"] is False
    assert flow["run2_41_trace_has_run2_43_ids"] is False
    assert flow["run2_41_data_drives_text_and_trace"] is True
    assert flow["run2_41_visual_geometry_is_hardcoded"] is True
    assert flow["run2_43_semantic_memory_is_hardcoded_from_prior_trace"] is True
    assert flow["run2_43_builder_reads_multimodal_source_records"] is False

    next_gate = audit["next_rerun_gate"]
    assert next_gate["required_before_run2_44_generator"] == [
        "run2_43_semantic_visual_asset_memory.json",
        "run2_43_editorial_composition_typography_memory.json",
        "run2_43_visual_asset_semantics_workflow_gates.json",
    ]
    assert "run2_43_semantic_visual_asset_ids" in next_gate["required_trace_fields"]
    assert next_gate["generator_must_fail_if_run2_43_not_consumed"] is True
    assert audit["next_required_action"] == (
        "build_run2_44_generator_that_consumes_run2_43_memory_for_visual_geometry_before_render"
    )


def test_run2_44_dataflow_readiness_audit_result_is_recorded() -> None:
    result = (PACK / "results" / "run2_44_dataflow_readiness_audit.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_44_dataflow_readiness_audit.json")

    assert result_json["status"] == "run2_44_dataflow_readiness_blocked"
    assert result_json["bug_confirmed"] is True
    assert result_json["not_a_run2_44_output"] is True
    assert result_json["risk_priority"] == "high"
    assert result_json["root_cause_primary"] == "latest_visible_ppt_does_not_consume_latest_run2_43_workflow"
    assert result_json["dataflow_findings"]["run2_41_visual_geometry_is_hardcoded"] is True
    assert result_json["dataflow_findings"]["run2_43_builder_reads_multimodal_source_records"] is False
    assert_contains(
        result,
        [
            "Run 2.44 Dataflow Readiness Audit",
            "bug confirmed",
            "not a generated Run 2.44 output",
            "high-priority technical debt",
            "latest visible PPT does not consume Run 2.43",
            "data drives text and trace more than visual geometry",
            "Run 2.44 generator must fail if Run 2.43 memory is not consumed",
        ],
    )


def test_ppt_run_html_viewer_embeds_run2_44_dataflow_readiness_audit() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "run2_44_dataflow_readiness_audit.json",
            "Run 2.44 dataflow readiness audit",
            "latest_visible_ppt_does_not_consume_latest_run2_43_workflow",
            "run2_41_visual_geometry_is_hardcoded",
            "run2_43_builder_reads_multimodal_source_records",
            "build_run2_44_generator_that_consumes_run2_43_memory_for_visual_geometry_before_render",
        ],
    )


def test_run2_44_generator_consumes_run2_43_workflow_for_geometry() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_44_semantic_geometry_arms.mjs"
    assert script_path.exists(), "missing Run 2.44 semantic geometry generator"
    body = script_path.read_text(encoding="utf-8")

    assert_contains(
        body,
        [
            "run2_43_semantic_visual_asset_memory.json",
            "run2_43_editorial_composition_typography_memory.json",
            "run2_43_visual_asset_semantics_workflow_gates.json",
            "run2_44_dataflow_readiness_audit.json",
            "readRun244PackJsonForArm",
            "validateRun244SemanticGeometryCompiler",
            "semanticGeometryTransform",
            "layoutFromRun243Memory",
            "drawRun244SemanticGeometryRail",
            "run2_43_semantic_visual_asset_ids",
            "run2_43_editorial_typography_memory_id",
            "run2_43_visual_asset_semantics_gate_id",
            "run2_44_geometry_source",
            "run2_44_data_bound_geometry",
        ],
    )


def test_run2_44_records_semantic_geometry_rerun_result() -> None:
    result = (PACK / "results" / "run2_44_semantic_geometry_rerun_result.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_44_semantic_geometry_rerun_result.json")
    full_trace = load_json(
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run2-44-full-vulca"
        / "trace_manifest.json"
    )
    bad_trace = load_json(
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run2-44-bad-run2-43-name-only-geometry"
        / "trace_manifest.json"
    )

    assert result_json["status"] == "run2_44_semantic_geometry_rerun_public_blocked"
    assert result_json["source_data_workflow_run_id"] == "2.43"
    assert result_json["source_dataflow_audit_run_id"] == "2.44-preflight"
    assert result_json["rerun"]["best_internal_arm"] == "run2_44_full_semantic_geometry_compiler"
    assert result_json["quality_delta"]["target_layer"] == "semantic_visual_asset_geometry_binding"
    assert result_json["quality_delta"]["full_slides_with_run2_43_semantic_asset_ids"] == 6
    assert result_json["quality_delta"]["full_slides_with_data_bound_geometry"] == 6
    assert result_json["input_chain"]["semantic_visual_asset_memory"].endswith(
        "run2_43_semantic_visual_asset_memory.json"
    )
    assert result_json["input_chain"]["editorial_composition_typography_memory"].endswith(
        "run2_43_editorial_composition_typography_memory.json"
    )
    assert result_json["input_chain"]["visual_asset_semantics_workflow_gates"].endswith(
        "run2_43_visual_asset_semantics_workflow_gates.json"
    )

    assert full_trace["arm_id"] == "run2_44_full_semantic_geometry_compiler"
    assert full_trace["run2_44_semantic_geometry_status"] == (
        "run2_43_workflow_consumed_for_data_bound_geometry"
    )
    assert full_trace["source_data_workflow_run_id"] == "2.43"
    assert len(full_trace["slides"]) == 6
    for slide in full_trace["slides"]:
        assert len(slide["run2_43_semantic_visual_asset_ids"]) == 3
        assert all(asset_id.startswith("semantic_asset_2_43_") for asset_id in slide["run2_43_semantic_visual_asset_ids"])
        assert slide["run2_43_editorial_typography_memory_id"].startswith("editorial_typography_2_43_")
        assert slide["run2_43_visual_asset_semantics_gate_id"].startswith("gate_2_43_")
        assert slide["run2_43_visual_asset_semantics_execution_status"] == (
            "consumed_before_native_ppt_drawing"
        )
        assert slide["run2_43_source_boundary_status"] == "derived_only_no_copied_media"
        assert slide["run2_44_data_bound_geometry"] is True
        assert slide["run2_44_geometry_source"] == (
            "run2_43_semantic_visual_asset_memory+run2_43_editorial_composition_typography_memory+run2_43_visual_asset_semantics_workflow_gates"
        )
        assert slide["run2_44_layout_signature_target"]
        assert len(slide["run2_44_semantic_asset_geometry_slots"]) == 3
        assert all(slot["asset_id"].startswith("semantic_asset_2_43_") for slot in slide["run2_44_semantic_asset_geometry_slots"])
        assert len(slide["run2_44_code_module_ids"]) == 1
        assert slide["run2_44_code_module_ids"][0].startswith("drawRun244")

    assert bad_trace["arm_id"] == "bad_run2_43_name_only_geometry"
    assert bad_trace["run2_44_semantic_geometry_status"] == "run2_43_names_only_geometry_control"
    for slide in bad_trace["slides"]:
        assert slide["run2_43_semantic_visual_asset_ids"] == []
        assert slide["run2_43_editorial_typography_memory_id"] == ""
        assert slide["run2_43_visual_asset_semantics_gate_id"] == ""
        assert slide["run2_43_visual_asset_semantics_execution_status"] == "not_consumed_name_only_control"
        assert slide["run2_44_data_bound_geometry"] is False
        assert slide["run2_44_semantic_asset_geometry_slots"] == []

    assert_contains(
        result,
        [
            "Run 2.44 Semantic Geometry Rerun",
            "consumes Run 2.43",
            "semantic visual asset ids",
            "data-bound geometry",
            "bad_run2_43_name_only_geometry",
            "public blocked",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_44_semantic_geometry_rerun() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "Run 2.44",
            "ppt-run2-44-prompt-only",
            "ppt-run2-44-run1-5-skill",
            "ppt-run2-44-full-vulca",
            "ppt-run2-44-bad-run2-43-name-only-geometry",
            "run2_44_semantic_geometry_rerun_result.json",
            "semantic_visual_asset_geometry_binding",
            "run2_43_workflow_consumed_for_data_bound_geometry",
        ],
    )


def test_run2_45_semantic_geometry_effectiveness_audit_compares_2_44_to_controls(tmp_path: Path) -> None:
    script_path = ROOT / "scripts" / "audit_ppt_run2_45_semantic_geometry_effectiveness.py"
    assert script_path.exists(), "missing Run 2.45 semantic geometry effectiveness audit script"
    result_json = tmp_path / "run2_45_semantic_geometry_effectiveness_audit.json"
    result_md = tmp_path / "run2_45_semantic_geometry_effectiveness_audit.md"
    presentations_dir = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
    )
    before_pptx = sorted(path.relative_to(ROOT).as_posix() for path in presentations_dir.rglob("*2-45*.pptx"))

    subprocess.run(
        [
            sys.executable,
            str(script_path),
            "--result-json",
            str(result_json),
            "--result-md",
            str(result_md),
        ],
        cwd=ROOT,
        check=True,
        capture_output=True,
        text=True,
    )

    after_pptx = sorted(path.relative_to(ROOT).as_posix() for path in presentations_dir.rglob("*2-45*.pptx"))
    audit = load_json(result_json)
    report = result_md.read_text(encoding="utf-8")

    assert before_pptx == []
    assert after_pptx == []
    assert audit["run_id"] == "2.45"
    assert audit["status"] == "run2_45_semantic_geometry_effectiveness_audit_public_blocked"
    assert audit["source_generated_run"] == "2.44"
    assert audit["source_workflow_run"] == "2.43"
    assert audit["creates_new_ppt_deck"] is False
    assert audit["public_ready"] is False
    assert audit["input_chain"]["run2_44_full_trace_manifest"].endswith(
        "ppt-run2-44-full-vulca/trace_manifest.json"
    )
    assert audit["input_chain"]["run2_44_bad_trace_manifest"].endswith(
        "ppt-run2-44-bad-run2-43-name-only-geometry/trace_manifest.json"
    )
    assert audit["input_chain"]["run2_44_rerun_result"].endswith(
        "run2_44_semantic_geometry_rerun_result.json"
    )

    trace = audit["semantic_geometry_trace_effectiveness"]
    assert trace["dataflow_bug_fixed"] is True
    assert trace["full_arm"]["all_slides_have_semantic_ids_memory_gate_and_geometry"] is True
    assert trace["full_arm"]["slides_with_run2_43_semantic_asset_ids"] == 6
    assert trace["full_arm"]["slides_with_data_bound_geometry"] == 6
    assert trace["bad_control"]["name_only_boundary_passed"] is True
    assert trace["bad_control"]["slides_without_semantic_asset_ids"] == 6
    assert trace["bad_control"]["slides_without_data_bound_geometry"] == 6

    visual = audit["visual_effectiveness_assessment"]
    assert visual["composition_compiler_kind"] == "slot_based_semantic_geometry"
    assert visual["slot_based_geometry_slides"] == 6
    assert visual["dataflow_fix_visual_delta_from_bad_control"] == "proven_internal_only"
    assert visual["public_video_grade_visual_quality"] is False
    assert visual["visual_quality_gate"] == "blocked"
    assert visual["root_cause_primary"] == "run2_44_dataflow_fixed_but_composition_compiler_still_slot_based"
    assert visual["top_next_layer_to_thicken"] == "multimodal_composition_memory_and_visual_object_grammar"

    assert audit["gate_summary"]["dataflow_gate"] == "pass_internal_only"
    assert audit["gate_summary"]["composition_quality_gate"] == "blocked"
    assert audit["next_required_action"] == (
        "build_run2_46_multimodal_composition_memory_and_workflow_thickening"
    )
    assert_contains(
        report,
        [
            "Run 2.45 Semantic Geometry Effectiveness Audit",
            "dataflow bug is fixed",
            "slot-based semantic geometry",
            "not public-video-grade",
            "Run 2.46",
        ],
    )


def test_run2_45_records_semantic_geometry_effectiveness_audit_result() -> None:
    result = (PACK / "results" / "run2_45_semantic_geometry_effectiveness_audit.md").read_text(
        encoding="utf-8"
    )
    result_json = load_json(PACK / "results" / "run2_45_semantic_geometry_effectiveness_audit.json")

    assert result_json["status"] == "run2_45_semantic_geometry_effectiveness_audit_public_blocked"
    assert result_json["source_generated_run"] == "2.44"
    assert result_json["semantic_geometry_trace_effectiveness"]["dataflow_bug_fixed"] is True
    assert result_json["visual_effectiveness_assessment"]["visual_quality_gate"] == "blocked"
    assert result_json["visual_effectiveness_assessment"]["top_next_layer_to_thicken"] == (
        "multimodal_composition_memory_and_visual_object_grammar"
    )
    assert result_json["delivery_artifacts"]["pptx_paths"] == []
    assert_contains(
        result,
        [
            "Run 2.45 Semantic Geometry Effectiveness Audit",
            "audit-only",
            "dataflow bug is fixed",
            "composition compiler is still slot-based",
            "Run 2.46",
        ],
    )


def test_ppt_run_html_viewer_embeds_run2_45_semantic_geometry_effectiveness_audit() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "run2_45_semantic_geometry_effectiveness_audit.json",
            "Run 2.45 semantic geometry effectiveness audit",
            "dataflow_bug_fixed",
            "slot_based_semantic_geometry",
            "build_run2_46_multimodal_composition_memory_and_workflow_thickening",
        ],
    )
    assert_contains(
        viewer,
        [
            '"latestRunId": "2.77"',
            "Run 2.45 semantic geometry effectiveness audit",
            "run2_45_semantic_geometry_effectiveness_audit.json",
            "slot_based_semantic_geometry",
        ],
    )


def test_run2_46_builder_writes_multimodal_composition_memory_pack(tmp_path: Path) -> None:
    script_path = ROOT / "scripts" / "build_ppt_run2_46_multimodal_composition_memory.py"
    assert script_path.exists(), "missing Run 2.46 multimodal composition memory builder"
    result_json = tmp_path / "run2_46_multimodal_composition_memory_result.json"
    result_md = tmp_path / "run2_46_multimodal_composition_memory_result.md"

    subprocess.run(
        [
            sys.executable,
            str(script_path),
            "--out-dir",
            str(tmp_path),
            "--result-json",
            str(result_json),
            "--result-md",
            str(result_md),
        ],
        cwd=ROOT,
        check=True,
        capture_output=True,
        text=True,
    )

    result = load_json(result_json)
    decomposition = load_json(tmp_path / "run2_46_multimodal_composition_decomposition.json")
    grammar = load_json(tmp_path / "run2_46_visual_object_grammar_memory.json")
    gates = load_json(tmp_path / "run2_46_composition_workflow_gates.json")
    report = result_md.read_text(encoding="utf-8")

    assert result["run_id"] == "2.46"
    assert result["status"] == "run2_46_multimodal_composition_memory_ready_public_blocked"
    assert result["source_audit_run"] == "2.45"
    assert result["source_generated_run"] == "2.44"
    assert result["source_workflow_run"] == "2.43"
    assert result["creates_new_ppt_deck"] is False
    assert result["target_layer"] == "multimodal_composition_memory_and_visual_object_grammar"
    assert result["artifact_counts"] == {
        "multimodal_composition_decomposition_records": 6,
        "visual_object_grammar_records": 6,
        "composition_workflow_gates": 6,
    }
    assert result["next_required_action"] == "consume_run2_46_multimodal_composition_memory_before_run2_47_rerun"

    assert decomposition["schema_version"] == "ppt_run2_multimodal_composition_decomposition.v1"
    assert decomposition["status"] == "run2_46_multimodal_composition_decomposition_ready_public_blocked"
    assert len(decomposition["multimodal_composition_decomposition_records"]) == 6
    for record in decomposition["multimodal_composition_decomposition_records"]:
        assert record["decomposition_id"].startswith("composition_decomposition_2_46_")
        assert record["source_audit_root_cause"] == "run2_44_dataflow_fixed_but_composition_compiler_still_slot_based"
        assert record["slot_based_failure_mode"] == "semantic_objects_bound_to_geometry_slots"
        assert len(record["multimodal_source_record_ids"]) >= 1
        assert record["composition_move"] in {
            "poster_scene_depth",
            "failure_to_route_storyboard",
            "asymmetric_before_after_delta",
            "active_product_surface",
            "cinematic_result_object",
            "decision_room_handoff",
        }
        assert "do not copy" in " ".join(record["source_boundary_rules"]).lower()
        assert len(record["native_ppt_composition_implications"]) >= 3

    assert grammar["schema_version"] == "ppt_run2_visual_object_grammar_memory.v1"
    assert grammar["status"] == "run2_46_visual_object_grammar_memory_ready_public_blocked"
    assert len(grammar["visual_object_grammar_records"]) == 6
    for record in grammar["visual_object_grammar_records"]:
        assert record["visual_object_grammar_id"].startswith("visual_object_grammar_2_46_")
        assert record["required_decomposition_id"].startswith("composition_decomposition_2_46_")
        assert record["required_run2_43_editorial_typography_memory_id"].startswith("editorial_typography_2_43_")
        assert len(record["required_run2_43_semantic_visual_asset_ids"]) == 3
        assert record["replaces_run2_44_slot_based_geometry"] is True
        assert len(record["visual_object_grammar"]) >= 4
        assert len(record["composition_quality_checks"]) >= 4
        assert "run2_46_visual_object_grammar_id" in record["required_trace_fields"]

    assert gates["schema_version"] == "ppt_run2_composition_workflow_gates.v1"
    assert gates["status"] == "run2_46_composition_workflow_gates_ready_public_blocked"
    assert gates["next_rerun_contract"] == "must_be_consumed_before_run2_47_four_arm_rerun"
    assert len(gates["composition_workflow_gates"]) == 6
    for gate in gates["composition_workflow_gates"]:
        assert gate["gate_id"].startswith("gate_2_46_")
        assert gate["required_visual_object_grammar_id"].startswith("visual_object_grammar_2_46_")
        assert gate["forbid_slot_based_geometry_as_primary_surface"] is True
        assert gate["require_multimodal_composition_decomposition"] is True
        assert "run2_46_visual_object_grammar_id" in gate["required_trace_fields"]
        assert "run2_46_composition_gate_id" in gate["required_trace_fields"]

    assert_contains(
        report,
        [
            "Run 2.46 Multimodal Composition Memory",
            "audit-only",
            "multimodal composition decomposition",
            "visual object grammar",
            "must be consumed before Run 2.47",
        ],
    )


def test_run2_46_records_multimodal_composition_memory_result() -> None:
    result = (PACK / "results" / "run2_46_multimodal_composition_memory_result.md").read_text(
        encoding="utf-8"
    )
    result_json = load_json(PACK / "results" / "run2_46_multimodal_composition_memory_result.json")

    assert (PACK / "run2_46_multimodal_composition_decomposition.json").exists()
    assert (PACK / "run2_46_visual_object_grammar_memory.json").exists()
    assert (PACK / "run2_46_composition_workflow_gates.json").exists()
    assert result_json["status"] == "run2_46_multimodal_composition_memory_ready_public_blocked"
    assert result_json["creates_new_ppt_deck"] is False
    assert result_json["source_audit_run"] == "2.45"
    assert result_json["artifact_counts"]["visual_object_grammar_records"] == 6
    assert result_json["output_chain"]["visual_object_grammar_memory"].endswith(
        "run2_46_visual_object_grammar_memory.json"
    )
    assert_contains(
        result,
        [
            "Run 2.46 Multimodal Composition Memory",
            "slot-based",
            "visual object grammar",
            "Run 2.47",
        ],
    )


def test_ppt_run_html_viewer_embeds_run2_46_multimodal_composition_memory() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "run2_46_multimodal_composition_memory_result.json",
            "Run 2.46 multimodal composition memory",
            "run2_46_multimodal_composition_decomposition.json",
            "run2_46_visual_object_grammar_memory.json",
            "run2_46_composition_workflow_gates.json",
            "consume_run2_46_multimodal_composition_memory_before_run2_47_rerun",
            "generated proof advances through Run 2.47 and Run 2.50",
        ],
    )
    assert_contains(
        viewer,
        [
            '"latestRunId": "2.77"',
            "Run 2.46 multimodal composition memory",
            "run2_46_multimodal_composition_memory_result.json",
            "visual object grammar",
        ],
    )


def test_run2_47_generator_consumes_run2_46_composition_grammar() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_47_composition_grammar_arms.mjs"
    assert script_path.exists(), "missing Run 2.47 composition grammar generator"
    body = script_path.read_text(encoding="utf-8")

    assert_contains(
        body,
        [
            "run2_46_multimodal_composition_decomposition.json",
            "run2_46_visual_object_grammar_memory.json",
            "run2_46_composition_workflow_gates.json",
            "run2_46_multimodal_composition_memory_result.json",
            "readRun247PackJsonForArm",
            "validateRun247CompositionGrammarCompiler",
            "compositionGrammarTransform",
            "drawRun247ComposedObjectScene",
            "run2_46_visual_object_grammar_id",
            "run2_46_multimodal_composition_decomposition_id",
            "run2_46_composition_gate_id",
            "run2_46_slot_based_geometry_replaced",
            "run2_46_source_boundary_status",
            "bad_run2_46_missing_composition_grammar",
        ],
    )


def test_run2_47_records_composition_grammar_rerun_result() -> None:
    result = (PACK / "results" / "run2_47_composition_grammar_rerun_result.md").read_text(
        encoding="utf-8"
    )
    result_json = load_json(PACK / "results" / "run2_47_composition_grammar_rerun_result.json")
    full_trace = load_json(
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run2-47-full-vulca"
        / "trace_manifest.json"
    )
    bad_trace = load_json(
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run2-47-bad-missing-composition-grammar"
        / "trace_manifest.json"
    )

    assert result_json["status"] == "run2_47_composition_grammar_rerun_public_blocked"
    assert result_json["source_composition_memory_run_id"] == "2.46"
    assert result_json["source_generated_run_id"] == "2.44"
    assert result_json["rerun"]["best_internal_arm"] == "run2_47_full_composition_grammar_compiler"
    assert result_json["quality_delta"]["target_layer"] == "composition_grammar_binding"
    assert result_json["quality_delta"]["full_slides_with_run2_46_visual_object_grammar_id"] == 6
    assert result_json["quality_delta"]["full_slides_with_slot_based_geometry_replaced"] == 6
    assert result_json["quality_delta"]["bad_control_slides_without_run2_46_grammar"] == 6
    assert result_json["input_chain"]["visual_object_grammar_memory"].endswith(
        "run2_46_visual_object_grammar_memory.json"
    )
    assert result_json["input_chain"]["composition_workflow_gates"].endswith(
        "run2_46_composition_workflow_gates.json"
    )

    assert full_trace["arm_id"] == "run2_47_full_composition_grammar_compiler"
    assert full_trace["run2_47_composition_grammar_status"] == (
        "run2_46_composition_grammar_consumed_before_native_ppt_drawing"
    )
    assert full_trace["source_composition_memory_run_id"] == "2.46"
    assert len(full_trace["slides"]) == 6
    for slide in full_trace["slides"]:
        assert slide["run2_46_visual_object_grammar_id"].startswith("visual_object_grammar_2_46_")
        assert slide["run2_46_multimodal_composition_decomposition_id"].startswith(
            "composition_decomposition_2_46_"
        )
        assert slide["run2_46_composition_gate_id"].startswith("gate_2_46_")
        assert slide["run2_46_slot_based_geometry_replaced"] is True
        assert slide["run2_46_source_boundary_status"] == "derived_only_no_copied_media"
        assert slide["run2_47_primary_composition_kind"] != "slot_based_semantic_geometry"
        assert len(slide["run2_47_visual_object_scene_objects"]) >= 3
        assert slide["run2_44_semantic_asset_geometry_slots"] == []
        assert len(slide["run2_47_code_module_ids"]) == 1
        assert slide["run2_47_code_module_ids"][0].startswith("drawRun247")

    assert bad_trace["arm_id"] == "bad_run2_46_missing_composition_grammar"
    assert bad_trace["run2_47_composition_grammar_status"] == (
        "run2_44_geometry_present_but_run2_46_composition_grammar_missing"
    )
    for slide in bad_trace["slides"]:
        assert slide["run2_46_visual_object_grammar_id"] == ""
        assert slide["run2_46_multimodal_composition_decomposition_id"] == ""
        assert slide["run2_46_composition_gate_id"] == ""
        assert slide["run2_46_slot_based_geometry_replaced"] is False
        assert len(slide["run2_44_semantic_asset_geometry_slots"]) == 3

    assert_contains(
        result,
        [
            "Run 2.47 Composition Grammar Rerun",
            "consumes Run 2.46",
            "visual object grammar",
            "slot-based geometry replaced",
            "bad_run2_46_missing_composition_grammar",
            "public blocked",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_47_composition_grammar_rerun() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "Run 2.47",
            "ppt-run2-47-prompt-only",
            "ppt-run2-47-run1-5-skill",
            "ppt-run2-47-full-vulca",
            "ppt-run2-47-bad-missing-composition-grammar",
            "run2_47_composition_grammar_rerun_result.json",
            "composition_grammar_binding",
            "run2_46_composition_grammar_consumed_before_native_ppt_drawing",
        ],
    )
    assert_contains(
        viewer,
        [
            '"latestRunId": "2.77"',
            "Run 2.47",
            "run2_47_composition_grammar_rerun_result.json",
            "visual object grammar",
        ],
    )


def test_run2_48_composition_grammar_effectiveness_audit_compares_2_47_to_bad_and_2_44(
    tmp_path: Path,
) -> None:
    script_path = ROOT / "scripts" / "audit_ppt_run2_48_composition_grammar_effectiveness.py"
    result_json = tmp_path / "run2_48_composition_grammar_effectiveness_audit.json"
    result_md = tmp_path / "run2_48_composition_grammar_effectiveness_audit.md"
    presentations = ROOT / "outputs" / "019e7d9c-532a-70b3-8892-fa3ae42baef2" / "presentations"
    pptx_before = sorted(path.name for path in presentations.glob("*2-48*.pptx"))

    assert script_path.exists(), "missing Run 2.48 composition grammar effectiveness audit script"
    completed = subprocess.run(
        [
            sys.executable,
            str(script_path),
            "--result-json",
            str(result_json),
            "--result-md",
            str(result_md),
        ],
        cwd=ROOT,
        check=True,
        capture_output=True,
        text=True,
    )

    pptx_after = sorted(path.name for path in presentations.glob("*2-48*.pptx"))
    audit = load_json(result_json)
    report = result_md.read_text(encoding="utf-8")

    assert "run2_48_composition_grammar_effectiveness_audit_public_blocked" in completed.stdout
    assert pptx_before == pptx_after == []
    assert audit["run_id"] == "2.48"
    assert audit["status"] == "run2_48_composition_grammar_effectiveness_audit_public_blocked"
    assert audit["source_generated_run"] == "2.47"
    assert audit["source_composition_memory_run"] == "2.46"
    assert audit["comparison_prior_generated_run"] == "2.44"
    assert audit["creates_new_ppt_deck"] is False
    assert audit["public_ready"] is False
    assert audit["composition_grammar_trace_effectiveness"]["full_arm"][
        "all_slides_have_grammar_decomposition_gate_and_slot_replacement"
    ] is True
    assert audit["composition_grammar_trace_effectiveness"]["full_arm"][
        "slides_with_run2_46_visual_object_grammar_id"
    ] == 6
    assert audit["composition_grammar_trace_effectiveness"]["full_arm"][
        "slides_with_run2_47_scene_objects"
    ] == 6
    assert audit["composition_grammar_trace_effectiveness"]["full_arm"][
        "slides_without_run2_44_slots"
    ] == 6
    assert audit["composition_grammar_trace_effectiveness"]["bad_control"][
        "missing_composition_grammar_boundary_passed"
    ] is True
    assert audit["composition_grammar_trace_effectiveness"]["bad_control"][
        "slides_without_run2_46_grammar"
    ] == 6
    assert audit["composition_grammar_trace_effectiveness"]["bad_control"][
        "slides_with_run2_44_slots"
    ] == 6
    assert audit["visual_effectiveness_assessment"]["composition_grammar_delta_from_run2_44"] == (
        "proven_internal_only"
    )
    assert audit["visual_effectiveness_assessment"]["visual_quality_gate"] == "blocked"
    assert audit["visual_effectiveness_assessment"]["public_video_grade_visual_quality"] is False
    assert audit["delivery_artifacts"]["pptx_paths"] == []
    assert_contains(
        report,
        [
            "Run 2.48 Composition Grammar Effectiveness Audit",
            "No Run 2.48 PPTX",
            "consumes Run 2.46",
            "slot-based geometry replaced",
            "not public-video-grade",
            "Run 2.49",
        ],
    )


def test_run2_48_records_composition_grammar_effectiveness_audit_result() -> None:
    result = (PACK / "results" / "run2_48_composition_grammar_effectiveness_audit.md").read_text(
        encoding="utf-8"
    )
    result_json = load_json(
        PACK / "results" / "run2_48_composition_grammar_effectiveness_audit.json"
    )

    assert result_json["status"] == "run2_48_composition_grammar_effectiveness_audit_public_blocked"
    assert result_json["source_generated_run"] == "2.47"
    assert result_json["source_composition_memory_run"] == "2.46"
    assert result_json["comparison_prior_generated_run"] == "2.44"
    assert result_json["delivery_artifacts"]["pptx_paths"] == []
    assert result_json["gate_summary"]["grammar_consumption_gate"] == "pass_internal_only"
    assert result_json["gate_summary"]["bad_control_gate"] == "pass"
    assert result_json["gate_summary"]["public_release_gate"] == "blocked"
    assert result_json["visual_effectiveness_assessment"]["root_cause_primary"] == (
        "run2_47_composition_grammar_consumed_but_visual_editorial_quality_still_not_public_grade"
    )
    assert result_json["visual_effectiveness_assessment"]["top_next_layer_to_thicken"] == (
        "readability_content_density_and_editorial_renderer_repair"
    )
    assert_contains(
        result,
        [
            "Run 2.48 is audit-only",
            "full arm is structurally stronger",
            "public release gate",
            "readability",
        ],
    )


def test_ppt_run_html_viewer_embeds_run2_48_composition_grammar_effectiveness_audit() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "run2_48_composition_grammar_effectiveness_audit.json",
            "Run 2.48 composition grammar effectiveness audit",
            "run248AuditStatus",
            "visual object grammar",
            "readability_content_density_and_editorial_renderer_repair",
        ],
    )
    assert_contains(
        viewer,
        [
            '"latestRunId": "2.77"',
            "Run 2.48 composition grammar effectiveness audit",
            "run2_48_composition_grammar_effectiveness_audit.json",
            "visual object grammar",
            "No Run 2.48 PPTX/download",
            "public-video-grade",
        ],
    )
    assert "ppt-run2-48" not in viewer


def test_run2_49_builder_creates_readability_content_density_renderer_repair_pack(
    tmp_path: Path,
) -> None:
    script_path = ROOT / "scripts" / "build_ppt_run2_49_readability_content_density_renderer_repair.py"
    result_json = tmp_path / "run2_49_readability_content_density_renderer_repair_result.json"
    result_md = tmp_path / "run2_49_readability_content_density_renderer_repair_result.md"
    presentations = ROOT / "outputs" / "019e7d9c-532a-70b3-8892-fa3ae42baef2" / "presentations"
    pptx_before = sorted(path.name for path in presentations.glob("*2-49*.pptx"))

    assert script_path.exists(), "missing Run 2.49 readability/content-density builder"
    completed = subprocess.run(
        [
            sys.executable,
            str(script_path),
            "--out-dir",
            str(tmp_path),
            "--result-json",
            str(result_json),
            "--result-md",
            str(result_md),
        ],
        cwd=ROOT,
        check=True,
        capture_output=True,
        text=True,
    )

    pptx_after = sorted(path.name for path in presentations.glob("*2-49*.pptx"))
    result = load_json(result_json)
    readability = load_json(tmp_path / "run2_49_readability_memory.json")
    density = load_json(tmp_path / "run2_49_content_evidence_density_memory.json")
    gates = load_json(tmp_path / "run2_49_editorial_renderer_workflow_gates.json")
    report = result_md.read_text(encoding="utf-8")

    assert "run2_49_readability_content_density_renderer_repair_ready_public_blocked" in completed.stdout
    assert pptx_before == pptx_after == []
    assert result["run_id"] == "2.49"
    assert result["status"] == "run2_49_readability_content_density_renderer_repair_ready_public_blocked"
    assert result["source_audit_run"] == "2.48"
    assert result["source_generated_run"] == "2.47"
    assert result["source_composition_memory_run"] == "2.46"
    assert result["target_layer"] == "readability_content_density_and_editorial_renderer_repair"
    assert result["creates_new_ppt_deck"] is False
    assert result["public_ready"] is False
    assert result["next_required_action"] == "consume_run2_49_before_run2_50_four_arm_rerun"
    assert result["delivery_artifacts"]["pptx_paths"] == []

    assert readability["status"] == "run2_49_readability_memory_ready_public_blocked"
    assert density["status"] == "run2_49_content_evidence_density_memory_ready_public_blocked"
    assert gates["status"] == "run2_49_editorial_renderer_workflow_gates_ready_public_blocked"
    assert len(readability["readability_records"]) == 6
    assert len(density["content_evidence_density_records"]) == 6
    assert len(gates["editorial_renderer_workflow_gates"]) == 6

    for record in readability["readability_records"]:
        assert record["min_contact_sheet_title_px"] >= 22
        assert record["max_headline_words"] <= 9
        assert record["forbid_title_clipping"] is True
        assert "contact_sheet_scale" in record["readability_gate_id"]
    for record in density["content_evidence_density_records"]:
        assert record["min_specific_business_evidence_objects"] >= 3
        assert record["min_inspectable_visual_proof_objects"] >= 2
        assert "generic abstract proof" in record["forbidden_evidence_substitutes"]
    for gate in gates["editorial_renderer_workflow_gates"]:
        assert gate["next_rerun_contract"] == "must_be_consumed_before_run2_50_four_arm_rerun"
        assert gate["forbid_square_block_grid_as_primary_surface"] is True
        assert gate["min_non_square_surface_ratio_variants"] >= 2
        assert gate["require_contact_sheet_readability"] is True
        assert gate["require_inspectable_business_evidence"] is True
        assert {
            "run2_49_readability_memory_id",
            "run2_49_content_evidence_density_memory_id",
            "run2_49_editorial_renderer_gate_id",
            "run2_49_renderer_contract_id",
        } <= set(gate["required_trace_fields"])

    assert_contains(
        report,
        [
            "Run 2.49",
            "data/workflow-only",
            "readability",
            "content evidence density",
            "editorial renderer",
            "Run 2.50",
        ],
    )


def test_run2_49_records_result_and_extends_skill_workflow() -> None:
    result = load_json(PACK / "results" / "run2_49_readability_content_density_renderer_repair_result.json")
    readability = load_json(PACK / "run2_49_readability_memory.json")
    density = load_json(PACK / "run2_49_content_evidence_density_memory.json")
    gates = load_json(PACK / "run2_49_editorial_renderer_workflow_gates.json")
    workflow = load_json(PACK / "skill_workflow.json")
    workflow_stage_ids = [stage["id"] for stage in workflow["stages"]]

    assert result["status"] == "run2_49_readability_content_density_renderer_repair_ready_public_blocked"
    assert result["output_chain"]["readability_memory"].endswith("run2_49_readability_memory.json")
    assert result["output_chain"]["content_evidence_density_memory"].endswith(
        "run2_49_content_evidence_density_memory.json"
    )
    assert result["output_chain"]["editorial_renderer_workflow_gates"].endswith(
        "run2_49_editorial_renderer_workflow_gates.json"
    )
    assert result["artifact_counts"] == {
        "readability_records": 6,
        "content_evidence_density_records": 6,
        "editorial_renderer_workflow_gates": 6,
    }
    assert readability["source_audit_run"] == "2.48"
    assert density["source_audit_run"] == "2.48"
    assert gates["source_audit_run"] == "2.48"
    assert "compile_run2_49_readability_memory" in workflow_stage_ids
    assert "compile_run2_49_content_evidence_density_memory" in workflow_stage_ids
    assert "apply_run2_49_editorial_renderer_workflow_gates" in workflow_stage_ids
    assert "compile_run2_51_editorial_copy_memory" in workflow_stage_ids
    assert "compile_run2_51_shape_text_socket_memory" in workflow_stage_ids
    assert "apply_run2_51_renderer_archetype_workflow_gates" in workflow_stage_ids
    assert workflow_stage_ids.index("compile_run2_49_readability_memory") < workflow_stage_ids.index(
        "compile_run2_49_content_evidence_density_memory"
    )
    assert workflow_stage_ids.index("compile_run2_49_content_evidence_density_memory") < workflow_stage_ids.index(
        "apply_run2_49_editorial_renderer_workflow_gates"
    )
    assert workflow_stage_ids.index("compile_run2_51_editorial_copy_memory") < workflow_stage_ids.index(
        "compile_run2_51_shape_text_socket_memory"
    )
    assert workflow_stage_ids.index("compile_run2_51_shape_text_socket_memory") < workflow_stage_ids.index(
        "apply_run2_51_renderer_archetype_workflow_gates"
    )


def test_ppt_run_html_viewer_embeds_run2_49_data_only_repair_pack() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "run2_49_readability_content_density_renderer_repair_result.json",
            "run249ResultStatus",
            "Run 2.49 readability/content density/editorial renderer repair",
            "Data-only Run",
            "consume_run2_49_before_run2_50_four_arm_rerun",
        ],
    )
    assert_contains(
        viewer,
        [
            '"latestRunId": "2.77"',
            "Run 2.49 readability/content density/editorial renderer repair",
            "Data-only Run",
            "run2_49_readability_memory.json",
            "run2_49_content_evidence_density_memory.json",
            "run2_49_editorial_renderer_workflow_gates.json",
            "Run 2.50",
        ],
    )
    assert "ppt-run2-49" not in viewer


def test_ppt_run_html_viewer_embeds_run2_51_data_workflow_repair_pack() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "run2_51_editorial_copy_memory.json",
            "run2_51_shape_text_socket_memory.json",
            "run2_51_renderer_archetype_workflow_gates.json",
            "run251ResultStatus",
            "Run 2.51 remains the prior data/workflow repair layer",
            "visual validation waits for the next generated rerun",
        ],
    )
    assert_contains(
        viewer,
        [
            '"latestRunId": "2.77"',
            "Run 2.51 remains the prior data/workflow repair layer",
            "run2_51_editorial_shape_text_repair_result.json",
            "run2_51_editorial_copy_memory.json",
            "run2_51_shape_text_socket_memory.json",
            "run2_51_renderer_archetype_workflow_gates.json",
            "Data/workflow-only, no new PPT deck",
            "Visual validation",
            "Deferred to generated rerun",
        ],
    )
    assert "ppt-run2-51" not in viewer


def test_ppt_run_html_viewer_surfaces_latest_data_workflow_repair_first() -> None:
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    latest_heading = "Latest data/workflow repair"
    historical_heading = "Why 2.8 still looks close to 2.7"
    assert latest_heading in viewer
    assert viewer.index(latest_heading) < viewer.index(historical_heading)
    assert_contains(
        viewer,
        [
            "Run 2.53 product-surface scene repair",
            "consume_run2_53_before_run2_54_four_arm_rerun",
            "Data/workflow-only, no new PPT deck",
            "Run 2.51 remains the prior data/workflow repair layer",
            "Run 2.49 readability/content density repair",
            ".dataBandHead .pill",
            "text-overflow: ellipsis",
            "Scene renderer workflow gates",
        ],
    )
    assert 'title="${escapeHtml(refs.run253ResultStatus || "missing")}"' in viewer


def test_run2_50_generator_consumes_run2_49_repair_pack() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_50_readability_density_renderer_arms.mjs"
    assert script_path.exists(), "missing Run 2.50 readability/content-density renderer generator"
    body = script_path.read_text(encoding="utf-8")

    assert_contains(
        body,
        [
            "run2_49_readability_memory.json",
            "run2_49_content_evidence_density_memory.json",
            "run2_49_editorial_renderer_workflow_gates.json",
            "run2_49_readability_memory_id",
            "run2_49_content_evidence_density_memory_id",
            "run2_49_editorial_renderer_gate_id",
            "run2_49_renderer_contract_id",
            "run2_49_business_evidence_density_status",
            "bad_run2_49_missing_repair_pack",
            "drawRun250EditorialEvidenceScene",
        ],
    )


def test_run2_50_records_readability_density_renderer_rerun_result() -> None:
    result = (PACK / "results" / "run2_50_readability_density_renderer_rerun_result.md").read_text(
        encoding="utf-8"
    )
    result_json = load_json(PACK / "results" / "run2_50_readability_density_renderer_rerun_result.json")
    presentations = ROOT / "outputs" / "019e7d9c-532a-70b3-8892-fa3ae42baef2" / "presentations"
    full_trace = load_json(presentations / "ppt-run2-50-full-vulca" / "trace_manifest.json")
    bad_trace = load_json(presentations / "ppt-run2-50-bad-missing-run2-49-repair-pack" / "trace_manifest.json")

    assert result_json["status"] == "run2_50_readability_density_renderer_rerun_public_blocked"
    assert result_json["source_repair_run_id"] == "2.49"
    assert result_json["source_generated_run_id"] == "2.47"
    assert result_json["rerun"]["best_internal_arm"] == "run2_50_full_readability_density_renderer"
    assert result_json["quality_delta"]["target_layer"] == "readability_content_density_and_editorial_renderer_binding"
    assert result_json["quality_delta"]["full_slides_with_run2_49_readability_memory_id"] == 6
    assert result_json["quality_delta"]["full_slides_with_run2_49_content_evidence_density_memory_id"] == 6
    assert result_json["quality_delta"]["full_slides_with_run2_49_editorial_renderer_gate_id"] == 6
    assert result_json["quality_delta"]["full_slides_with_business_evidence_density_pass"] == 6
    assert result_json["quality_delta"]["bad_control_slides_without_run2_49_repair_pack"] == 6
    assert result_json["rerun"]["combined_contact_sheet"].endswith("run2-50-four-arm-contact-sheet.png")

    assert full_trace["arm_id"] == "run2_50_full_readability_density_renderer"
    assert full_trace["run2_50_readability_density_renderer_status"] == (
        "run2_49_repair_pack_consumed_before_native_ppt_drawing"
    )
    assert len(full_trace["slides"]) == 6
    for slide in full_trace["slides"]:
        assert slide["run2_49_readability_memory_id"].startswith("readability_memory_2_49_")
        assert slide["run2_49_content_evidence_density_memory_id"].startswith("content_evidence_density_2_49_")
        assert slide["run2_49_editorial_renderer_gate_id"].startswith("gate_2_49_")
        assert slide["run2_49_renderer_contract_id"].startswith("renderer_contract_2_49_")
        assert slide["run2_49_contact_sheet_readability_status"] == "pass_internal"
        assert slide["run2_49_business_evidence_density_status"] == "pass_internal"
        assert len(slide["run2_49_business_evidence_objects"]) >= 3
        assert len(slide["run2_49_inspectable_proof_objects"]) >= 2
        assert slide["run2_49_non_square_surface_ratio_variants"] >= 2
        assert slide["run2_50_primary_surface_kind"] != "square_block_grid"
        assert slide["layout_metrics"]["proof_objects"] >= 2
        assert slide["layout_metrics"]["visible_words"] >= 70
        assert slide["run2_50_code_module_ids"][0].startswith("drawRun250")

    assert bad_trace["arm_id"] == "bad_run2_49_missing_repair_pack"
    for slide in bad_trace["slides"]:
        assert slide["run2_49_readability_memory_id"] == ""
        assert slide["run2_49_content_evidence_density_memory_id"] == ""
        assert slide["run2_49_editorial_renderer_gate_id"] == ""
        assert slide["run2_49_business_evidence_density_status"] == "fail_missing_run2_49"

    assert_contains(
        result,
        [
            "Run 2.50 Readability Density Renderer Rerun",
            "consumes Run 2.49",
            "readability",
            "content evidence density",
            "editorial renderer",
            "bad_run2_49_missing_repair_pack",
            "public blocked",
        ],
    )


def test_run2_51_builder_creates_editorial_copy_shape_socket_repair_pack(tmp_path: Path) -> None:
    script_path = ROOT / "scripts" / "build_ppt_run2_51_editorial_copy_shape_sockets.py"
    result_json = tmp_path / "run2_51_editorial_shape_text_repair_result.json"
    result_md = tmp_path / "run2_51_editorial_shape_text_repair_result.md"
    presentations = ROOT / "outputs" / "019e7d9c-532a-70b3-8892-fa3ae42baef2" / "presentations"
    pptx_before = sorted(path.name for path in presentations.glob("*2-51*.pptx"))

    assert script_path.exists(), "missing Run 2.51 editorial copy/shape socket builder"
    body = script_path.read_text(encoding="utf-8")
    body_lower = body.lower()
    for forbidden in (
        "presentation_artifact_tool",
        "pptxgen",
        "slide.shapes",
        "from pptx",
        "import pptx",
        "Presentation(",
    ):
        if forbidden == "Presentation(":
            assert forbidden not in body
        else:
            assert forbidden.lower() not in body_lower
    completed = subprocess.run(
        [
            sys.executable,
            str(script_path),
            "--out-dir",
            str(tmp_path),
            "--result-json",
            str(result_json),
            "--result-md",
            str(result_md),
        ],
        cwd=ROOT,
        check=True,
        capture_output=True,
        text=True,
    )

    pptx_after = sorted(path.name for path in presentations.glob("*2-51*.pptx"))
    result = load_json(result_json)
    copy_memory = load_json(tmp_path / "run2_51_editorial_copy_memory.json")
    socket_memory = load_json(tmp_path / "run2_51_shape_text_socket_memory.json")
    gates = load_json(tmp_path / "run2_51_renderer_archetype_workflow_gates.json")
    report = result_md.read_text(encoding="utf-8")

    assert "run2_51_editorial_shape_text_repair_ready_public_blocked" in completed.stdout
    assert pptx_before == pptx_after == []
    assert result["run_id"] == "2.51"
    assert result["status"] == "run2_51_editorial_shape_text_repair_ready_public_blocked"
    assert result["source_data_workflow_run"] == "2.49"
    assert result["source_generated_run"] == "2.50"
    assert result["target_layer"] == "editorial_copy_and_shape_text_socket_repair"
    assert result["creates_new_ppt_deck"] is False
    assert result["visual_validation_deferred_to_generated_rerun"] is True
    assert result["public_ready"] is False
    assert result["next_required_action"] == "consume_run2_51_before_run2_52_four_arm_rerun"
    assert result["delivery_artifacts"]["pptx_paths"] == []
    assert result["artifact_counts"] == {
        "editorial_copy_records": 6,
        "shape_text_socket_records": 6,
        "renderer_archetype_workflow_gates": 6,
    }

    assert copy_memory["status"] == "run2_51_editorial_copy_memory_ready_public_blocked"
    assert socket_memory["status"] == "run2_51_shape_text_socket_memory_ready_public_blocked"
    assert gates["status"] == "run2_51_renderer_archetype_workflow_gates_ready_public_blocked"
    assert {record["role"] for record in copy_memory["editorial_copy_records"]} == EXPECTED_RUN2_51_ROLES
    assert {record["role"] for record in socket_memory["shape_text_socket_records"]} == EXPECTED_RUN2_51_ROLES
    assert {record["role"] for record in gates["renderer_archetype_workflow_gates"]} == EXPECTED_RUN2_51_ROLES
    for record in copy_memory["editorial_copy_records"]:
        bundle = record["public_surface_copy_bundle"]
        assert EXPECTED_RUN2_51_COPY_BUNDLE_KEYS <= set(bundle)
        assert word_count(bundle["headline"]) <= 7
        assert word_count(bundle["subline"]) <= 18
        public_values = public_text_values(bundle)
        assert public_values
        for value in public_values:
            normalized_value = normalize(value)
            for forbidden_term in EXPECTED_RUN2_51_FORBIDDEN_PUBLIC_TERMS:
                assert normalize(forbidden_term) not in normalized_value
    for gate in gates["renderer_archetype_workflow_gates"]:
        assert EXPECTED_RUN2_51_TRACE_FIELDS <= set(gate["required_trace_fields"])

    assert_contains(
        report,
        [
            "Run 2.51",
            "data/workflow-only",
            "Run 2.51 creates editorial copy, shape text sockets, and renderer archetype workflow gates",
            "editorial copy",
            "shape text sockets",
            "visual validation is deferred",
            "deferred visual validation",
            "Run 2.52",
        ],
    )


def test_run2_51_records_editorial_copy_shape_socket_repair_pack() -> None:
    result = load_json(PACK / "results" / "run2_51_editorial_shape_text_repair_result.json")
    copy_memory = load_json(PACK / "run2_51_editorial_copy_memory.json")
    socket_memory = load_json(PACK / "run2_51_shape_text_socket_memory.json")
    gates = load_json(PACK / "run2_51_renderer_archetype_workflow_gates.json")

    assert result["status"] == "run2_51_editorial_shape_text_repair_ready_public_blocked"
    assert result["visual_validation_deferred_to_generated_rerun"] is True
    assert result["output_chain"] == {
        "editorial_copy_memory": "docs/product/ppt-run2-data-skill-quality/run2_51_editorial_copy_memory.json",
        "shape_text_socket_memory": "docs/product/ppt-run2-data-skill-quality/run2_51_shape_text_socket_memory.json",
        "renderer_archetype_workflow_gates": "docs/product/ppt-run2-data-skill-quality/run2_51_renderer_archetype_workflow_gates.json",
    }
    assert copy_memory["schema_version"] == "ppt_run2_editorial_copy_memory.v1"
    assert socket_memory["schema_version"] == "ppt_run2_shape_text_socket_memory.v1"
    assert gates["schema_version"] == "ppt_run2_renderer_archetype_workflow_gates.v1"
    assert copy_memory["source_generated_run"] == "2.50"
    assert socket_memory["source_generated_run"] == "2.50"
    assert gates["source_generated_run"] == "2.50"
    assert len(copy_memory["editorial_copy_records"]) == 6
    assert len(socket_memory["shape_text_socket_records"]) == 6
    assert len(gates["renderer_archetype_workflow_gates"]) == 6

    for record in copy_memory["editorial_copy_records"]:
        bundle = record["public_surface_copy_bundle"]
        assert set(bundle) == EXPECTED_RUN2_51_COPY_BUNDLE_KEYS
        assert word_count(bundle["headline"]) <= 7
        assert len(bundle["headline"]) <= 48
        assert word_count(bundle["subline"]) <= 18
        assert len(bundle["subline"]) <= 120
        for value in bundle["proof_nuggets"]:
            assert word_count(value) <= 8
            assert len(value) <= 54
        for value in bundle["annotations"]:
            assert word_count(value) <= 6
            assert len(value) <= 42
        for value in bundle["state_labels"]:
            assert word_count(value) <= 4
            assert len(value) <= 28
        for text in public_text_values(bundle):
            lowered = normalize(text)
            assert not any(normalize(term) in lowered for term in EXPECTED_RUN2_51_FORBIDDEN_PUBLIC_TERMS)
        assert record["visual_validation_deferred_to_generated_rerun"] is True
        assert record["next_rerun_obligation"] == "must_be_consumed_before_run2_52_four_arm_rerun"

    for record in socket_memory["shape_text_socket_records"]:
        assert len(record["socket_contracts"]) >= 4
        assert len(record["shape_primitives"]) >= 3
        assert len(record["geometry_constraints"]) >= 3
        assert record["primary_archetype"]
        assert record["forbidden_layout_patterns"]
        assert record["proof_socket_ids"]
        for socket in record["socket_contracts"]:
            assert socket["owning_shape_id"]
            assert socket["placement_rule"]
            assert socket["character_budget"] > 0
            assert socket["max_lines"] >= 1
            assert socket["overflow_policy"] in {"shrink_to_fit", "wrap_within_socket", "reject_and_recompile"}

    for gate in gates["renderer_archetype_workflow_gates"]:
        assert set(gate["required_trace_fields"]) == EXPECTED_RUN2_51_TRACE_FIELDS
        assert gate["consumer_contract"]["next_generated_run"] == "2.52"
        assert gate["consumer_contract"]["must_bind_before_public_text"] is True
        assert set(gate["consumer_contract"]["required_trace_fields"]) == EXPECTED_RUN2_51_TRACE_FIELDS
        assert gate["next_rerun_contract"] == "must_be_consumed_before_run2_52_four_arm_rerun"
        assert gate["forbid_square_block_grid_as_primary_surface"] is True
        assert gate["max_equal_card_clusters"] <= 1
        assert gate["min_semantic_primitives"] >= 3
        assert gate["visual_validation_deferred_to_generated_rerun"] is True


def test_run2_51_extends_skill_workflow_without_claiming_generated_deck() -> None:
    workflow = load_json(PACK / "skill_workflow.json")
    result = load_json(PACK / "results" / "run2_51_editorial_shape_text_repair_result.json")
    workflow_stage_ids = [stage["id"] for stage in workflow["stages"]]
    stage_by_id = {stage["id"]: stage for stage in workflow["stages"]}

    assert result["creates_new_ppt_deck"] is False
    assert result["visual_validation_deferred_to_generated_rerun"] is True
    assert "compile_run2_51_editorial_copy_memory" in workflow_stage_ids
    assert "compile_run2_51_shape_text_socket_memory" in workflow_stage_ids
    assert "apply_run2_51_renderer_archetype_workflow_gates" in workflow_stage_ids
    assert workflow_stage_ids.index("compile_run2_51_editorial_copy_memory") < workflow_stage_ids.index(
        "compile_run2_51_shape_text_socket_memory"
    )
    assert workflow_stage_ids.index("compile_run2_51_shape_text_socket_memory") < workflow_stage_ids.index(
        "apply_run2_51_renderer_archetype_workflow_gates"
    )
    assert stage_by_id["compile_run2_51_editorial_copy_memory"]["order"] == 47
    assert stage_by_id["compile_run2_51_shape_text_socket_memory"]["order"] == 48
    assert stage_by_id["apply_run2_51_renderer_archetype_workflow_gates"]["order"] == 49


def test_run2_51_builder_rejects_malformed_run2_50_source() -> None:
    from scripts import build_ppt_run2_51_editorial_copy_shape_sockets as run251_builder

    run249_result = load_json(PACK / "results" / "run2_49_readability_content_density_renderer_repair_result.json")
    run249_readability = load_json(PACK / "run2_49_readability_memory.json")
    run249_density = load_json(PACK / "run2_49_content_evidence_density_memory.json")
    run249_gates = load_json(PACK / "run2_49_editorial_renderer_workflow_gates.json")
    run250_result = load_json(PACK / "results" / "run2_50_readability_density_renderer_rerun_result.json")

    wrong_layer = json.loads(json.dumps(run250_result))
    wrong_layer["quality_delta"]["target_layer"] = "wrong_layer"
    with pytest.raises(ValueError, match="target_layer"):
        run251_builder.validate_inputs(
            run249_result, run249_readability, run249_density, run249_gates, wrong_layer
        )

    wrong_modules = json.loads(json.dumps(run250_result))
    wrong_modules["quality_delta"]["repair_modules"] = ["only_one"]
    with pytest.raises(ValueError, match="repair_modules"):
        run251_builder.validate_inputs(
            run249_result, run249_readability, run249_density, run249_gates, wrong_modules
        )


def test_run2_52_generator_consumes_run2_51_editorial_socket_pack() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_52_editorial_socket_renderer_arms.mjs"
    assert script_path.exists(), "missing Run 2.52 editorial socket renderer generator"
    body = script_path.read_text(encoding="utf-8")

    assert_contains(
        body,
        [
            "run2_51_editorial_shape_text_repair_result.json",
            "run2_51_editorial_copy_memory.json",
            "run2_51_shape_text_socket_memory.json",
            "run2_51_renderer_archetype_workflow_gates.json",
            "run2_51_editorial_copy_memory_id",
            "run2_51_shape_text_socket_memory_id",
            "run2_51_renderer_archetype_gate_id",
            "run2_51_primary_archetype",
            "run2_51_public_surface_copy_status",
            "run2_51_text_socket_placement_status",
            "bad_run2_51_missing_editorial_socket_pack",
            "drawRun252EditorialSocketScene",
            "run2_51_editorial_socket_pack_consumed_before_native_ppt_drawing",
            "build_ppt_full_skill_series_sheet.py",
            "Run 2.27",
            "ppt-run2-27-full-vulca",
            "Run 2.28",
            "ppt-run2-28-full-vulca",
            "Run 2.29",
            "ppt-run2-29-full-vulca",
            "Run 2.31",
            "ppt-run2-31-full-vulca",
            "Run 2.40",
            "ppt-run2-40-full-vulca",
            "Run 2.52",
            "ppt-run2-52-full-vulca",
        ],
    )


def test_run2_52_records_editorial_socket_renderer_rerun_result() -> None:
    result_md = (PACK / "results" / "run2_52_editorial_socket_renderer_rerun_result.md").read_text(
        encoding="utf-8"
    )
    result_json = load_json(PACK / "results" / "run2_52_editorial_socket_renderer_rerun_result.json")
    presentations = ROOT / "outputs" / "019e7d9c-532a-70b3-8892-fa3ae42baef2" / "presentations"
    full_trace = load_json(presentations / "ppt-run2-52-full-vulca" / "trace_manifest.json")
    bad_trace = load_json(
        presentations / "ppt-run2-52-bad-missing-run2-51-editorial-socket-pack" / "trace_manifest.json"
    )

    assert result_json["run_id"] == "2.52"
    assert result_json["status"] == "run2_52_editorial_socket_renderer_rerun_public_blocked"
    assert result_json["source_repair_run_id"] == "2.51"
    assert result_json["source_generated_run_id"] == "2.50"
    assert result_json["rerun"]["best_internal_arm"] == "run2_52_full_editorial_socket_renderer"
    assert result_json["quality_delta"]["target_layer"] == (
        "editorial_copy_shape_text_socket_and_renderer_archetype_binding"
    )
    assert result_json["quality_delta"]["source_data_status"] == (
        "run2_51_editorial_socket_pack_consumed_before_native_ppt_drawing"
    )
    assert result_json["quality_delta"]["full_slides_with_run2_51_editorial_copy_memory_id"] == 6
    assert result_json["quality_delta"]["full_slides_with_run2_51_shape_text_socket_memory_id"] == 6
    assert result_json["quality_delta"]["full_slides_with_run2_51_renderer_archetype_gate_id"] == 6
    assert result_json["quality_delta"]["full_slides_with_socket_bound_public_text"] == 6
    assert result_json["quality_delta"]["bad_control_slides_without_run2_51_pack"] == 6
    assert result_json["rerun"]["combined_contact_sheet"].endswith("run2-52-four-arm-contact-sheet.png")

    assert full_trace["arm_id"] == "run2_52_full_editorial_socket_renderer"
    assert full_trace["run2_52_editorial_socket_renderer_status"] == (
        "run2_51_editorial_socket_pack_consumed_before_native_ppt_drawing"
    )
    assert len(full_trace["slides"]) == 6
    for slide in full_trace["slides"]:
        assert EXPECTED_RUN2_52_TRACE_FIELDS <= set(slide)
        assert slide["run2_51_editorial_copy_memory_id"].startswith("editorial_copy_2_51_")
        assert slide["run2_51_shape_text_socket_memory_id"].startswith("shape_text_socket_2_51_")
        assert slide["run2_51_renderer_archetype_gate_id"].startswith("gate_2_51_")
        assert slide["run2_51_public_surface_copy_status"] == "pass_internal"
        assert slide["run2_51_text_socket_placement_status"] == "pass_internal"
        assert slide["run2_51_shape_vocabulary_status"] == "pass_internal"
        assert slide["run2_51_character_fit_status"] == "pass_internal"
        assert slide["run2_51_forbidden_surface_terms_count"] == 0
        assert slide["run2_51_equal_card_cluster_count"] <= 1
        assert slide["run2_51_semantic_primitive_count"] >= 3
        assert slide["run2_52_socket_bound_public_text_elements"] >= 4
        assert slide["run2_52_shape_primitive_count"] >= 3
        assert slide["run2_52_primary_surface_kind"] != "square_block_grid"
        assert slide["layout_metrics"]["visible_words"] >= 48
        assert slide["layout_metrics"]["proof_objects"] >= 2
        assert slide["run2_52_code_module_ids"][0].startswith("drawRun252")

    assert bad_trace["arm_id"] == "bad_run2_51_missing_editorial_socket_pack"
    for slide in bad_trace["slides"]:
        assert slide["run2_51_editorial_copy_memory_id"] == ""
        assert slide["run2_51_shape_text_socket_memory_id"] == ""
        assert slide["run2_51_renderer_archetype_gate_id"] == ""
        assert slide["run2_51_public_surface_copy_status"] == "fail_missing_run2_51"

    assert_contains(
        result_md,
        [
            "Run 2.52 Editorial Socket Renderer Rerun",
            "consumes Run 2.51",
            "editorial copy",
            "shape text sockets",
            "renderer archetype",
            "bad_run2_51_missing_editorial_socket_pack",
            "public blocked",
        ],
    )


def test_run2_52_generator_rejects_malformed_run2_51_source() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_52_editorial_socket_renderer_arms.mjs"
    body = script_path.read_text(encoding="utf-8")

    assert "validateRun252RepairPack" in body
    assert "Run 2.52 must consume Run 2.51 repair result" in body
    assert "Run 2.52 editorial copy status mismatch" in body
    assert "Run 2.52 socket status mismatch" in body
    assert "Run 2.52 renderer archetype gate status mismatch" in body
    assert "consumer_contract.next_generated_run" in body
    assert "must_bind_before_public_text" in body
    assert "missing role contract" in body


def test_run2_52_bad_control_trace_does_not_leak_run2_51_fields() -> None:
    presentations = ROOT / "outputs" / "019e7d9c-532a-70b3-8892-fa3ae42baef2" / "presentations"
    bad_trace = load_json(
        presentations / "ppt-run2-52-bad-missing-run2-51-editorial-socket-pack" / "trace_manifest.json"
    )

    for slide in bad_trace["slides"]:
        leaked_values = [
            value
            for key, value in slide.items()
            if key.startswith("run2_51_")
            and key not in {
                "run2_51_editorial_copy_memory_id",
                "run2_51_shape_text_socket_memory_id",
                "run2_51_renderer_archetype_gate_id",
                "run2_51_public_surface_copy_status",
            }
            and value not in ("", 0, [], False, None, "fail_missing_run2_51")
        ]
        assert leaked_values == []


def test_run2_53_builds_product_surface_scene_repair_pack() -> None:
    script_path = ROOT / "scripts" / "build_ppt_run2_53_product_surface_scene_repair.py"
    assert script_path.exists()

    result_md = (PACK / "results" / "run2_53_product_surface_scene_repair_result.md").read_text(
        encoding="utf-8"
    )
    result = load_json(PACK / "results" / "run2_53_product_surface_scene_repair_result.json")
    scene_memory = load_json(PACK / "run2_53_product_surface_scene_memory.json")
    evidence_memory = load_json(PACK / "run2_53_business_visual_evidence_memory.json")
    gates = load_json(PACK / "run2_53_scene_renderer_workflow_gates.json")
    run251_copy = load_json(PACK / "run2_51_editorial_copy_memory.json")
    run251_gates = load_json(PACK / "run2_51_renderer_archetype_workflow_gates.json")

    assert result["run_id"] == "2.53"
    assert result["status"] == "run2_53_product_surface_scene_repair_ready_public_blocked"
    assert result["public_ready"] is False
    assert result["creates_new_ppt_deck"] is False
    assert result["selected_usecase_id"] == "usecase_design_to_production_platform_launch"
    assert result["source_repair_run_id"] == "2.51"
    assert result["source_generated_run_id"] == "2.52"
    assert result["target_layer"] == "product_surface_scene_and_business_visual_evidence_repair"
    assert result["next_required_action"] == "consume_run2_53_before_run2_54_four_arm_rerun"
    assert result["artifact_counts"] == {
        "product_surface_scene_records": 6,
        "business_visual_evidence_records": 6,
        "scene_renderer_workflow_gates": 6,
    }

    assert scene_memory["status"] == "run2_53_product_surface_scene_memory_ready_public_blocked"
    assert evidence_memory["status"] == "run2_53_business_visual_evidence_memory_ready_public_blocked"
    assert gates["status"] == "run2_53_scene_renderer_workflow_gates_ready_public_blocked"

    scene_records = scene_memory["product_surface_scene_records"]
    evidence_records = evidence_memory["business_visual_evidence_records"]
    gate_records = gates["scene_renderer_workflow_gates"]
    assert {record["role"] for record in scene_records} == EXPECTED_RUN2_51_ROLES
    assert {record["role"] for record in evidence_records} == EXPECTED_RUN2_51_ROLES
    assert {record["role"] for record in gate_records} == EXPECTED_RUN2_51_ROLES

    run251_copy_by_role = {record["role"]: record for record in run251_copy["editorial_copy_records"]}
    run251_gate_by_role = {gate["role"]: gate for gate in run251_gates["renderer_archetype_workflow_gates"]}
    scene_by_role = {record["role"]: record for record in scene_records}

    for scene in scene_records:
        role = scene["role"]
        assert scene["id"] == f"product_surface_scene_2_53_{role}"
        assert scene["selected_usecase_id"] == "usecase_design_to_production_platform_launch"
        assert scene["source_run_ids"] == ["2.51", "2.52"]
        assert scene["must_render_product_or_business_object"] is True
        assert len(scene["surface_slots"]) >= 3
        assert "generic geometric diagram" in scene["forbidden_patterns"]
        assert "evidence card only" in scene["forbidden_patterns"]
        assert "floating annotation card" in scene["forbidden_patterns"]
        assert scene["next_rerun_obligation"] == "must_be_consumed_before_run2_54_four_arm_rerun"

    for evidence in evidence_records:
        role = evidence["role"]
        assert evidence["id"] == f"business_visual_evidence_2_53_{role}"
        assert evidence["required_product_surface_scene_id"] == scene_by_role[role]["id"]
        assert evidence["required_editorial_copy_memory_id"] == run251_copy_by_role[role]["copy_memory_id"]
        assert evidence["observable_business_object"]
        assert evidence["reader_question_answered"]
        assert len(evidence["minimum_visual_specificity_checks"]) >= 3
        assert evidence["source_boundary"]["raw_media_copied"] is False
        assert evidence["source_boundary"]["source_layout_copied"] is False

    for gate in gate_records:
        role = gate["role"]
        assert gate["id"] == f"scene_renderer_gate_2_53_{role}"
        assert gate["required_product_surface_scene_id"] == scene_by_role[role]["id"]
        assert gate["required_business_visual_evidence_id"] == f"business_visual_evidence_2_53_{role}"
        assert gate["required_run2_51_renderer_archetype_gate_id"] == run251_gate_by_role[role]["gate_id"]
        assert gate["consumer_contract"]["next_generated_run"] == "2.54"
        assert gate["consumer_contract"]["must_bind_before_native_drawing"] is True
        assert set(gate["required_trace_fields"]) == EXPECTED_RUN2_53_TRACE_FIELDS

    assert_contains(
        result_md,
        [
            "Run 2.53 Product Surface Scene Repair",
            "data/workflow-only",
            "product surface scene",
            "business visual evidence",
            "Run 2.54",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_run2_53_extends_skill_workflow_without_claiming_generated_deck() -> None:
    workflow = load_json(PACK / "skill_workflow.json")
    result = load_json(PACK / "results" / "run2_53_product_surface_scene_repair_result.json")
    stage_ids = [stage["id"] for stage in workflow["stages"]]
    stage_by_id = {stage["id"]: stage for stage in workflow["stages"]}

    assert result["creates_new_ppt_deck"] is False
    assert result["visual_validation_deferred_to_generated_rerun"] is True
    assert "compile_run2_53_product_surface_scene_memory" in stage_ids
    assert "compile_run2_53_business_visual_evidence_memory" in stage_ids
    assert "apply_run2_53_scene_renderer_workflow_gates" in stage_ids
    assert stage_ids.index("compile_run2_53_product_surface_scene_memory") < stage_ids.index(
        "compile_run2_53_business_visual_evidence_memory"
    )
    assert stage_ids.index("compile_run2_53_business_visual_evidence_memory") < stage_ids.index(
        "apply_run2_53_scene_renderer_workflow_gates"
    )
    assert stage_by_id["compile_run2_53_product_surface_scene_memory"]["order"] == 50
    assert stage_by_id["compile_run2_53_business_visual_evidence_memory"]["order"] == 51
    assert stage_by_id["apply_run2_53_scene_renderer_workflow_gates"]["order"] == 52


def test_run2_53_builder_rejects_malformed_run2_52_source() -> None:
    from scripts import build_ppt_run2_53_product_surface_scene_repair as run253_builder

    run252_result = load_json(PACK / "results" / "run2_52_editorial_socket_renderer_rerun_result.json")
    run252_trace = load_json(
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run2-52-full-vulca"
        / "trace_manifest.json"
    )
    run251_copy = load_json(PACK / "run2_51_editorial_copy_memory.json")
    run251_socket = load_json(PACK / "run2_51_shape_text_socket_memory.json")
    run251_gates = load_json(PACK / "run2_51_renderer_archetype_workflow_gates.json")

    malformed_result = json.loads(json.dumps(run252_result))
    malformed_result["quality_delta"]["source_data_status"] = "wrong_source"
    with pytest.raises(ValueError, match="source_data_status"):
        run253_builder.validate_inputs(
            malformed_result,
            run252_trace,
            run251_copy,
            run251_socket,
            run251_gates,
        )

    malformed_result = json.loads(json.dumps(run252_result))
    malformed_result["quality_delta"]["full_slides_with_run2_51_editorial_copy_memory_id"] = 5
    with pytest.raises(ValueError, match="editorial copy"):
        run253_builder.validate_inputs(
            malformed_result,
            run252_trace,
            run251_copy,
            run251_socket,
            run251_gates,
        )


def test_ppt_run_html_viewer_mentions_run2_53_product_surface_scene_repair() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "run2_53_product_surface_scene_memory.json",
            "run2_53_business_visual_evidence_memory.json",
            "run2_53_scene_renderer_workflow_gates.json",
            "run253ResultStatus",
            "Run 2.53 product-surface scene repair",
            "data/workflow-only",
            "consume_run2_53_before_run2_54_four_arm_rerun",
        ],
    )
    assert_contains(
        viewer,
        [
            '"latestRunId": "2.77"',
            "Run 2.53 product-surface scene repair",
            "run2_53_product_surface_scene_memory.json",
            "run2_53_business_visual_evidence_memory.json",
            "run2_53_scene_renderer_workflow_gates.json",
            "Run 2.54",
        ],
    )
    assert "Run 2.53" not in script.split("RUN_SPECS", 1)[1].split("def rel", 1)[0]
    assert "ppt-run2-53-" not in script


def test_run2_54_generator_consumes_run2_53_product_surface_scene_pack() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_54_product_surface_scene_arms.mjs"
    assert script_path.exists(), "missing Run 2.54 product-surface-scene generator"
    body = script_path.read_text(encoding="utf-8")

    assert_contains(
        body,
        [
            "run2_53_product_surface_scene_repair_result.json",
            "run2_53_product_surface_scene_memory.json",
            "run2_53_business_visual_evidence_memory.json",
            "run2_53_scene_renderer_workflow_gates.json",
            "run2_53_product_surface_scene_id",
            "run2_53_business_visual_evidence_id",
            "run2_53_scene_renderer_gate_id",
            "bad_run2_53_missing_product_surface_scene_pack",
            "drawRun254ProductSurfaceScene",
            "run2_53_product_surface_scene_pack_consumed_before_native_ppt_drawing",
        ],
    )
    assert "Run 2.54 must consume Run 2.53 repair result" in body
    assert "must_bind_before_native_drawing" in body


def test_run2_54_records_product_surface_scene_rerun_result() -> None:
    result_md = (PACK / "results" / "run2_54_product_surface_scene_rerun_result.md").read_text(
        encoding="utf-8"
    )
    result = load_json(PACK / "results" / "run2_54_product_surface_scene_rerun_result.json")
    presentations = ROOT / "outputs" / "019e7d9c-532a-70b3-8892-fa3ae42baef2" / "presentations"
    full_trace = load_json(presentations / "ppt-run2-54-full-vulca" / "trace_manifest.json")
    bad_trace = load_json(
        presentations / "ppt-run2-54-bad-missing-run2-53-product-surface-scene-pack" / "trace_manifest.json"
    )

    assert result["run_id"] == "2.54"
    assert result["status"] == "run2_54_product_surface_scene_rerun_public_blocked"
    assert result["selected_usecase_id"] == "usecase_design_to_production_platform_launch"
    assert result["source_repair_run_id"] == "2.53"
    assert result["source_generated_run_id"] == "2.52"
    assert result["rerun"]["best_internal_arm"] == "run2_54_full_product_surface_scene"
    assert result["quality_delta"]["target_layer"] == "product_surface_scene_and_business_visual_evidence_binding"
    assert result["quality_delta"]["source_data_status"] == (
        "run2_53_product_surface_scene_pack_consumed_before_native_ppt_drawing"
    )
    assert result["quality_delta"]["full_slides_with_run2_53_product_surface_scene_id"] == 6
    assert result["quality_delta"]["full_slides_with_run2_53_business_visual_evidence_id"] == 6
    assert result["quality_delta"]["full_slides_with_run2_53_scene_renderer_gate_id"] == 6
    assert result["quality_delta"]["full_slides_without_forbidden_generic_geometry"] == 6
    assert result["quality_delta"]["bad_control_slides_without_run2_53_pack"] == 6
    assert result["rerun"]["combined_contact_sheet"].endswith("run2-54-four-arm-contact-sheet.png")
    assert result["rerun"]["full_skill_series_sheet"].endswith("run2-full-skill-series-horizontal.png")

    assert full_trace["arm_id"] == "run2_54_full_product_surface_scene"
    assert full_trace["run2_54_product_surface_scene_status"] == (
        "run2_53_product_surface_scene_pack_consumed_before_native_ppt_drawing"
    )
    assert full_trace["source_repair_run_id"] == "2.53"
    assert full_trace["source_generated_run_id"] == "2.52"
    assert len(full_trace["slides"]) == 6
    for slide in full_trace["slides"]:
        assert set(EXPECTED_RUN2_54_TRACE_FIELDS) <= set(slide)
        assert slide["run2_53_product_surface_scene_id"].startswith("product_surface_scene_2_53_")
        assert slide["run2_53_business_visual_evidence_id"].startswith("business_visual_evidence_2_53_")
        assert slide["run2_53_scene_renderer_gate_id"].startswith("scene_renderer_gate_2_53_")
        assert slide["run2_53_primary_product_or_business_object"]
        assert slide["run2_53_visual_specificity_status"] == "pass_internal"
        assert slide["run2_53_forbidden_generic_geometry_count"] == 0
        assert slide["run2_54_product_surface_slots_rendered"] >= 3
        assert slide["run2_54_business_visual_evidence_objects"] >= 1
        assert slide["run2_54_primary_surface_kind"] != "square_block_grid"
        assert "drawRun254ProductSurfaceScene" in slide["run2_54_code_module_ids"]
        assert slide["layout_metrics"]["visible_words"] >= 50

    assert bad_trace["arm_id"] == "bad_run2_53_missing_product_surface_scene_pack"
    for slide in bad_trace["slides"]:
        assert slide["run2_53_product_surface_scene_id"] == ""
        assert slide["run2_53_business_visual_evidence_id"] == ""
        assert slide["run2_53_scene_renderer_gate_id"] == ""
        assert slide["run2_53_visual_specificity_status"] == "fail_missing_run2_53"

    assert_contains(
        result_md,
        [
            "Run 2.54 Product Surface Scene Rerun",
            "consumes Run 2.53",
            "product surface scene",
            "business visual evidence",
            "bad_run2_53_missing_product_surface_scene_pack",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_run2_54_generator_rejects_malformed_run2_53_source() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_54_product_surface_scene_arms.mjs"
    body = script_path.read_text(encoding="utf-8")

    assert "validateRun254RepairPack" in body
    assert "Run 2.54 must consume Run 2.53 repair result" in body
    assert "Run 2.54 product surface scene status mismatch" in body
    assert "Run 2.54 business visual evidence status mismatch" in body
    assert "Run 2.54 scene renderer gate status mismatch" in body
    assert "must_bind_before_native_drawing" in body
    assert "missing role contract" in body


def test_run2_54_bad_control_trace_does_not_leak_run2_53_fields() -> None:
    presentations = ROOT / "outputs" / "019e7d9c-532a-70b3-8892-fa3ae42baef2" / "presentations"
    bad_trace = load_json(
        presentations / "ppt-run2-54-bad-missing-run2-53-product-surface-scene-pack" / "trace_manifest.json"
    )

    for slide in bad_trace["slides"]:
        leaked_values = [
            value
            for key, value in slide.items()
            if key.startswith("run2_53_")
            and key
            not in {
                "run2_53_product_surface_scene_id",
                "run2_53_business_visual_evidence_id",
                "run2_53_scene_renderer_gate_id",
                "run2_53_primary_product_or_business_object",
                "run2_53_visual_specificity_status",
                "run2_53_forbidden_generic_geometry_count",
            }
            and value not in ("", 0, [], False, None, "fail_missing_run2_53")
        ]
        assert leaked_values == []


def test_ppt_run_html_viewer_mentions_run2_54_product_surface_scene_rerun() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "Run 2.54",
            "ppt-run2-54-prompt-only",
            "ppt-run2-54-run1-5-skill",
            "ppt-run2-54-full-vulca",
            "ppt-run2-54-bad-missing-run2-53-product-surface-scene-pack",
            "run2-54-four-arm-contact-sheet.png",
            "run2_54_product_surface_scene_rerun_result.json",
            "run2_53_product_surface_scene_pack_consumed_before_native_ppt_drawing",
        ],
    )
    assert_contains(
        viewer,
        [
            '"latestRunId": "2.77"',
            "Run 2.54",
            "run2-54-four-arm-contact-sheet.png",
            "ppt-run2-54-full-vulca",
            "Run 2.53 product-surface scene repair",
            "run2_54_product_surface_scene_rerun_result.json",
        ],
    )


def test_run2_55_generator_consumes_run2_54_and_repairs_text_shape_integration() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_55_text_shape_integration_arms.mjs"
    assert script_path.exists(), "missing Run 2.55 text-shape integration generator"
    body = script_path.read_text(encoding="utf-8")

    assert_contains(
        body,
        [
            "run2_54_product_surface_scene_rerun_result.json",
            "ppt-run2-54-full-vulca/trace_manifest.json",
            "run2_53_product_surface_scene_memory.json",
            "run2_53_business_visual_evidence_memory.json",
            "drawRun255TextShapeIntegration",
            "text_shape_integration_and_shape_vocabulary_repair",
            "bad_run2_54_without_text_shape_integration",
            "run2_54_product_surface_scene_status",
            "run2_55_text_shape_integration_status",
            "run2_55_non_rectangular_shape_families_rendered",
            "run2_55_text_shape_binding_pairs",
        ],
    )
    assert "Run 2.55 must consume Run 2.54 generated result" in body
    assert "Run 2.55 must compare against Run 2.54 full trace" in body
    assert "equal rectangle cluster" in body


def test_run2_55_records_text_shape_integration_rerun_result() -> None:
    result_md = (PACK / "results" / "run2_55_text_shape_integration_rerun_result.md").read_text(
        encoding="utf-8"
    )
    result = load_json(PACK / "results" / "run2_55_text_shape_integration_rerun_result.json")
    presentations = ROOT / "outputs" / "019e7d9c-532a-70b3-8892-fa3ae42baef2" / "presentations"
    full_trace = load_json(presentations / "ppt-run2-55-full-vulca" / "trace_manifest.json")
    bad_trace = load_json(
        presentations / "ppt-run2-55-bad-without-text-shape-integration" / "trace_manifest.json"
    )

    assert result["run_id"] == "2.55"
    assert result["status"] == "run2_55_text_shape_integration_rerun_public_blocked"
    assert result["selected_usecase_id"] == "usecase_design_to_production_platform_launch"
    assert result["source_repair_run_id"] == "2.53"
    assert result["source_generated_run_id"] == "2.54"
    assert result["rerun"]["best_internal_arm"] == "run2_55_full_text_shape_integration"
    assert result["quality_delta"]["target_layer"] == "text_shape_integration_and_shape_vocabulary_repair"
    assert result["quality_delta"]["source_data_status"] == (
        "run2_54_product_surface_scene_rerun_consumed_before_text_shape_redraw"
    )
    assert result["quality_delta"]["full_slides_with_named_text_containers"] == 6
    assert result["quality_delta"]["full_slides_with_non_rectangular_shape_families"] == 6
    assert result["quality_delta"]["full_slides_with_text_shape_binding_pairs"] == 6
    assert result["quality_delta"]["full_slides_without_text_overflow_risk"] == 6
    assert result["quality_delta"]["full_slides_without_equal_rectangle_clusters"] == 6
    assert result["quality_delta"]["bad_control_slides_without_text_shape_integration"] == 6
    assert result["rerun"]["combined_contact_sheet"].endswith("run2-55-four-arm-contact-sheet.png")
    assert result["rerun"]["full_skill_series_sheet"].endswith("run2-full-skill-series-horizontal.png")

    assert full_trace["arm_id"] == "run2_55_full_text_shape_integration"
    assert full_trace["run2_55_text_shape_integration_status"] == (
        "run2_54_product_surface_scene_rerun_consumed_before_text_shape_redraw"
    )
    assert full_trace["source_repair_run_id"] == "2.53"
    assert full_trace["source_generated_run_id"] == "2.54"
    assert len(full_trace["slides"]) == 6
    layout_strategies = {slide["run2_55_primary_layout_strategy"] for slide in full_trace["slides"]}
    assert len(layout_strategies) >= 4
    for slide in full_trace["slides"]:
        assert set(EXPECTED_RUN2_54_TRACE_FIELDS) <= set(slide)
        assert set(EXPECTED_RUN2_55_TRACE_FIELDS) <= set(slide)
        assert slide["run2_53_product_surface_scene_id"].startswith("product_surface_scene_2_53_")
        assert slide["run2_54_product_surface_slots_rendered"] >= 3
        assert slide["run2_55_text_shape_integration_status"] == "pass_internal"
        assert slide["run2_55_named_text_containers_rendered"] >= 4
        assert slide["run2_55_non_rectangular_shape_families_rendered"] >= 3
        assert slide["run2_55_text_shape_binding_pairs"] >= 4
        assert slide["run2_55_text_overflow_risk_count"] == 0
        assert slide["run2_55_equal_rectangle_cluster_count"] == 0
        assert slide["run2_55_editorial_hierarchy_levels"] >= 4
        assert "drawRun255TextShapeIntegration" in slide["run2_55_code_module_ids"]
        assert slide["layout_metrics"]["text_box_count"] >= 7
        assert slide["layout_metrics"]["visible_words"] >= 64

    assert bad_trace["arm_id"] == "bad_run2_54_without_text_shape_integration"
    for slide in bad_trace["slides"]:
        assert slide["run2_53_product_surface_scene_id"].startswith("product_surface_scene_2_53_")
        assert slide["run2_55_text_shape_integration_status"] == "fail_missing_text_shape_integration"
        assert slide["run2_55_non_rectangular_shape_families_rendered"] <= 1
        assert slide["run2_55_equal_rectangle_cluster_count"] >= 1

    assert_contains(
        result_md,
        [
            "Run 2.55 Text Shape Integration Rerun",
            "consumes Run 2.54",
            "named text containers",
            "non-rectangular shape families",
            "bad_run2_54_without_text_shape_integration",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_run2_55_generator_rejects_malformed_run2_54_source() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_55_text_shape_integration_arms.mjs"
    body = script_path.read_text(encoding="utf-8")

    assert "validateRun255TextShapeInputs" in body
    assert "Run 2.55 must consume Run 2.54 generated result" in body
    assert "Run 2.55 must compare against Run 2.54 full trace" in body
    assert "Run 2.55 requires six Run 2.54 full trace slides" in body
    assert "Run 2.55 missing Run 2.53 product surface scene id" in body
    assert "equal rectangle cluster" in body


def test_ppt_run_html_viewer_mentions_run2_55_text_shape_integration_rerun() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "Run 2.55",
            "ppt-run2-55-prompt-only",
            "ppt-run2-55-run1-5-skill",
            "ppt-run2-55-full-vulca",
            "ppt-run2-55-bad-without-text-shape-integration",
            "run2-55-four-arm-contact-sheet.png",
            "run2_55_text_shape_integration_rerun_result.json",
            "run2_54_product_surface_scene_rerun_consumed_before_text_shape_redraw",
        ],
    )
    assert_contains(
        viewer,
        [
            '"latestRunId": "2.77"',
            "Run 2.55",
            "run2-55-four-arm-contact-sheet.png",
            "ppt-run2-55-full-vulca",
            "Run 2.55 generated result",
            "run2_55_text_shape_integration_rerun_result.json",
        ],
    )


def test_run2_56_generator_consumes_run2_55_and_splits_role_renderers() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_56_role_renderer_split_arms.mjs"
    assert script_path.exists(), "missing Run 2.56 role-renderer split generator"
    body = script_path.read_text(encoding="utf-8")

    assert_contains(
        body,
        [
            "run2_55_text_shape_integration_rerun_result.json",
            "ppt-run2-55-full-vulca/trace_manifest.json",
            "run2_51_shape_text_socket_memory.json",
            "drawRun256CoverPosterStage",
            "drawRun256SetupRouteMap",
            "drawRun256ContrastBeforeAfterLens",
            "drawRun256ProofWorkspaceSurface",
            "drawRun256ClimaxEditorialHero",
            "drawRun256CloseDecisionWall",
            "role_specific_renderer_variation_and_layout_qa",
            "bad_run2_55_reused_single_template",
            "run2_56_layout_signature",
            "run2_56_text_collision_risk_count",
        ],
    )
    assert "Run 2.56 must consume Run 2.55 generated result" in body
    assert "Run 2.56 must compare against Run 2.55 full trace" in body
    assert "six unique role renderer ids" in body


def test_run2_56_records_role_renderer_split_rerun_result() -> None:
    result_md = (PACK / "results" / "run2_56_role_renderer_split_rerun_result.md").read_text(
        encoding="utf-8"
    )
    result = load_json(PACK / "results" / "run2_56_role_renderer_split_rerun_result.json")
    presentations = ROOT / "outputs" / "019e7d9c-532a-70b3-8892-fa3ae42baef2" / "presentations"
    full_trace = load_json(presentations / "ppt-run2-56-full-vulca" / "trace_manifest.json")
    bad_trace = load_json(
        presentations / "ppt-run2-56-bad-reused-single-template" / "trace_manifest.json"
    )

    assert result["run_id"] == "2.56"
    assert result["status"] == "run2_56_role_renderer_split_rerun_public_blocked"
    assert result["selected_usecase_id"] == "usecase_design_to_production_platform_launch"
    assert result["source_repair_run_id"] == "2.53"
    assert result["source_generated_run_id"] == "2.55"
    assert result["rerun"]["best_internal_arm"] == "run2_56_full_role_renderer_split"
    assert result["quality_delta"]["target_layer"] == "role_specific_renderer_variation_and_layout_qa"
    assert result["quality_delta"]["source_data_status"] == (
        "run2_55_text_shape_integration_consumed_before_role_renderer_redraw"
    )
    assert result["quality_delta"]["full_slides_with_unique_role_renderer"] == 6
    assert result["quality_delta"]["full_slides_with_unique_layout_signature"] == 6
    assert result["quality_delta"]["full_slides_with_role_archetype_binding"] == 6
    assert result["quality_delta"]["full_slides_without_text_collision_risk"] == 6
    assert result["quality_delta"]["full_slides_without_text_overflow_risk"] == 6
    assert result["quality_delta"]["bad_control_slides_with_reused_run2_55_template"] == 6
    assert result["rerun"]["combined_contact_sheet"].endswith("run2-56-four-arm-contact-sheet.png")
    assert result["rerun"]["full_skill_series_sheet"].endswith("run2-full-skill-series-horizontal.png")

    assert full_trace["arm_id"] == "run2_56_full_role_renderer_split"
    assert full_trace["run2_56_role_renderer_split_status"] == (
        "run2_55_text_shape_integration_consumed_before_role_renderer_redraw"
    )
    assert full_trace["source_repair_run_id"] == "2.53"
    assert full_trace["source_generated_run_id"] == "2.55"
    assert len(full_trace["slides"]) == 6
    renderer_ids = {slide["run2_56_role_renderer_id"] for slide in full_trace["slides"]}
    layout_signatures = {slide["run2_56_layout_signature"] for slide in full_trace["slides"]}
    sameness_buckets = {slide["run2_56_visual_sameness_bucket"] for slide in full_trace["slides"]}
    anchor_regions = {slide["run2_56_primary_anchor_region"] for slide in full_trace["slides"]}
    assert len(renderer_ids) == 6
    assert len(layout_signatures) == 6
    assert len(sameness_buckets) == 6
    assert len(anchor_regions) >= 5

    for slide in full_trace["slides"]:
        assert set(EXPECTED_RUN2_54_TRACE_FIELDS) <= set(slide)
        assert set(EXPECTED_RUN2_55_TRACE_FIELDS) <= set(slide)
        assert set(EXPECTED_RUN2_56_TRACE_FIELDS) <= set(slide)
        assert slide["run2_55_text_shape_integration_status"] == "pass_internal"
        assert slide["run2_56_distinct_role_surface_status"] == "pass_internal"
        assert slide["run2_56_role_archetype_binding_status"] == "pass_internal"
        assert slide["run2_56_role_specific_geometry_count"] >= 5
        assert slide["run2_56_text_collision_risk_count"] == 0
        assert slide["run2_56_text_overflow_risk_count"] == 0
        assert slide["run2_56_role_renderer_id"] in slide["run2_56_code_module_ids"]
        assert slide["layout_metrics"]["zones"] >= 6

    assert bad_trace["arm_id"] == "bad_run2_55_reused_single_template"
    for slide in bad_trace["slides"]:
        assert slide["run2_55_text_shape_integration_status"] == "pass_internal"
        assert slide["run2_56_distinct_role_surface_status"] == "fail_reused_single_template"
        assert slide["run2_56_layout_signature"] == "reused_run2_55_stage_side_template"
        assert slide["run2_56_role_renderer_id"] == "drawRun255TextShapeIntegration"

    assert_contains(
        result_md,
        [
            "Run 2.56 Role Renderer Split Rerun",
            "consumes Run 2.55",
            "six unique role renderer ids",
            "six unique layout signatures",
            "bad_run2_55_reused_single_template",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_56_role_renderer_split_rerun() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "Run 2.56",
            "ppt-run2-56-prompt-only",
            "ppt-run2-56-run1-5-skill",
            "ppt-run2-56-full-vulca",
            "ppt-run2-56-bad-reused-single-template",
            "run2-56-four-arm-contact-sheet.png",
            "run2_56_role_renderer_split_rerun_result.json",
            "run2_55_text_shape_integration_consumed_before_role_renderer_redraw",
        ],
    )
    assert_contains(
        viewer,
        [
            '"latestRunId": "2.77"',
            "Run 2.56",
            "run2-56-four-arm-contact-sheet.png",
            "ppt-run2-56-full-vulca",
            "Run 2.56 generated result",
            "run2_56_role_renderer_split_rerun_result.json",
        ],
    )


def test_run2_57_records_product_capability_content_layer() -> None:
    result_md = (PACK / "results" / "run2_57_product_capability_content_result.md").read_text(
        encoding="utf-8"
    )
    result = load_json(PACK / "results" / "run2_57_product_capability_content_result.json")
    capability_memory = load_json(PACK / "run2_57_product_capability_memory.json")
    message_contracts = load_json(PACK / "run2_57_slide_message_contracts.json")
    workflow_gates = load_json(PACK / "run2_57_content_workflow_gates.json")

    assert result["run_id"] == "2.57"
    assert result["status"] == "run2_57_product_capability_content_ready_public_blocked"
    assert result["public_ready"] is False
    assert result["selected_usecase_id"] == "usecase_design_to_production_platform_launch"
    assert result["source_generated_run_id"] == "2.56"
    assert result["database_expansion"] is True
    assert result["workflow_expansion"] is True
    assert result["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert result["generation_boundary"]["creates_new_ppt_deck"] is False
    assert result["generation_boundary"]["latest_generated_run_id"] == "2.56"
    assert result["quality_delta"]["target_layer"] == "product_capability_narrative_and_content_specificity"
    assert result["quality_delta"]["source_data_status"] == "run2_56_role_renderer_split_rerun_public_blocked"
    assert result["quality_delta"]["slides_requiring_product_content_repair"] >= 4
    assert result["next_required_action"] == (
        "run2_58_generate_four_arm_ppt_consuming_run2_57_product_content_layer"
    )
    assert set(result["output_chain"]) == {
        "product_capability_memory",
        "slide_message_contracts",
        "content_workflow_gates",
    }

    assert capability_memory["status"] == "run2_57_product_capability_memory_ready_public_blocked"
    capability_records = capability_memory["product_capability_records"]
    capability_ids = {record["id"] for record in capability_records}
    assert {record["capability_layer"] for record in capability_records} == EXPECTED_RUN2_57_CAPABILITY_LAYERS
    assert len(capability_records) == 7
    for record in capability_records:
        assert record["plain_language_claim"]
        assert record["visible_ppt_obligation"]
        assert record["why_it_matters"]
        assert record["workflow_contract"]
        assert record["next_renderer_obligation"]
        assert record["trace_obligation"] in EXPECTED_RUN2_57_TRACE_FIELDS

    logic_records = capability_memory["product_logic_relation_records"]
    assert len(logic_records) >= 6
    for relation in logic_records:
        assert relation["from_capability_id"] in capability_ids
        assert relation["to_capability_id"] in capability_ids
        assert relation["relation_type"]
        assert relation["must_be_visible_as"]

    competitor_records = capability_memory["competitor_boundary_records"]
    assert {record["competitor_boundary_id"] for record in competitor_records} == (
        EXPECTED_RUN2_57_COMPETITOR_BOUNDARIES
    )
    for record in competitor_records:
        assert "code-generated editable PPT" in record["vulca_difference"]
        assert record["slide_obligation"]

    assert message_contracts["status"] == "run2_57_slide_message_contracts_ready_public_blocked"
    slide_contracts = message_contracts["slide_message_contracts"]
    assert {contract["role"] for contract in slide_contracts} == EXPECTED_RUN2_51_ROLES
    for contract in slide_contracts:
        assert contract["contract_id"].startswith("message_contract_2_57_")
        assert contract["reader_question"]
        assert contract["required_answer"]
        assert contract["visible_product_object"]
        assert len(contract["required_product_terms"]) >= 3
        assert contract["minimum_specific_content_units"] >= 5
        assert set(contract["source_capability_ids"]) <= capability_ids
        assert contract["forbidden_empty_claims"]
        assert contract["next_renderer_obligation"]["required_trace_fields"] == sorted(
            EXPECTED_RUN2_57_TRACE_FIELDS
        )

    assert workflow_gates["status"] == "run2_57_content_workflow_gates_ready_public_blocked"
    gates = workflow_gates["content_workflow_gates"]
    assert len(gates) >= 5
    for gate in gates:
        assert gate["gate_id"].startswith("gate_2_57_")
        assert gate["gate_type"]
        assert gate["pass_fail_checks"]
        assert gate["bad_control_probe"]
        assert set(gate["required_trace_fields"]) <= EXPECTED_RUN2_57_TRACE_FIELDS
        assert gate["consumer_contract"]["next_generated_run"] == "2.58"
    assert set(workflow_gates["next_generated_run_contract"]["required_trace_fields"]) == (
        EXPECTED_RUN2_57_TRACE_FIELDS
    )
    assert workflow_gates["next_generated_run_contract"]["bad_control_arm"] == (
        "bad_run2_56_without_product_capability_content"
    )

    assert_contains(
        result_md,
        [
            "Run 2.57 Product Capability Content Layer",
            "data/workflow-only",
            "code-generated editable PPT",
            "product capability",
            "slide message contracts",
            "content workflow gates",
            "Run 2.58",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_57_product_capability_content_layer() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "Run 2.57",
            "run2_57_product_capability_content_result.json",
            "run2_57_product_capability_memory.json",
            "run2_57_slide_message_contracts.json",
            "run2_57_content_workflow_gates.json",
            "run257ProductCapabilityStatus",
            "run257SlideMessageContractStatus",
            "run257ContentWorkflowGateStatus",
            "product_capability_narrative_and_content_specificity",
        ],
    )
    assert_contains(
        viewer,
        [
            '"latestRunId": "2.77"',
            "Run 2.57 product capability content layer",
            "run2_57_product_capability_content_result.json",
            "run2_57_product_capability_memory.json",
            "run2_57_slide_message_contracts.json",
            "run2_57_content_workflow_gates.json",
            "No PPT deck generated in Run 2.57",
            "code-generated editable PPT",
            "Next generated run",
            "2.58",
        ],
    )
    assert "ppt-run2-57-" not in script
    assert "ppt-run2-57-" not in viewer


def test_run2_58_generator_consumes_run2_57_product_content_layer() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_58_product_content_contract_arms.mjs"
    assert script_path.exists(), "missing Run 2.58 product content contract generator"
    body = script_path.read_text(encoding="utf-8")

    assert_contains(
        body,
        [
            "run2_57_product_capability_content_result.json",
            "run2_57_product_capability_memory.json",
            "run2_57_slide_message_contracts.json",
            "run2_57_content_workflow_gates.json",
            "ppt-run2-56-full-vulca/trace_manifest.json",
            "drawRun258CapabilityLaunchSurface",
            "drawRun258SourceToMemoryPipeline",
            "drawRun258CompetitorBoundaryMatrix",
            "drawRun258TraceQaInspectionBoard",
            "drawRun258OperatingLoopClimax",
            "drawRun258ReleaseDecisionWall",
            "product_capability_content_contract_binding",
            "bad_run2_56_without_product_capability_content",
            "run2_57_product_capability_content_consumed_before_native_ppt_drawing",
        ],
    )
    assert "Run 2.58 must consume Run 2.57 product capability content result" in body
    assert "Run 2.58 must compare against Run 2.56 full trace" in body
    assert "generic claim count must be zero" in body


def test_run2_58_records_product_content_contract_rerun_result() -> None:
    result_md = (PACK / "results" / "run2_58_product_content_contract_rerun_result.md").read_text(
        encoding="utf-8"
    )
    result = load_json(PACK / "results" / "run2_58_product_content_contract_rerun_result.json")
    presentations = ROOT / "outputs" / "019e7d9c-532a-70b3-8892-fa3ae42baef2" / "presentations"
    full_trace = load_json(presentations / "ppt-run2-58-full-vulca" / "trace_manifest.json")
    bad_trace = load_json(
        presentations / "ppt-run2-58-bad-without-product-capability-content" / "trace_manifest.json"
    )

    assert result["run_id"] == "2.58"
    assert result["status"] == "run2_58_product_content_contract_rerun_public_blocked"
    assert result["public_ready"] is False
    assert result["selected_usecase_id"] == "usecase_design_to_production_platform_launch"
    assert result["source_repair_run_id"] == "2.57"
    assert result["source_generated_run_id"] == "2.56"
    assert result["rerun"]["best_internal_arm"] == "run2_58_full_product_content_contract"
    assert result["quality_delta"]["target_layer"] == "product_capability_content_contract_binding"
    assert result["quality_delta"]["source_data_status"] == (
        "run2_57_product_capability_content_consumed_before_native_ppt_drawing"
    )
    assert result["quality_delta"]["full_slides_with_run2_57_capability_ids"] == 6
    assert result["quality_delta"]["full_slides_with_message_contracts"] == 6
    assert result["quality_delta"]["full_slides_with_content_workflow_gates"] == 6
    assert result["quality_delta"]["full_slides_with_reader_question_answered"] == 6
    assert result["quality_delta"]["full_slides_with_zero_generic_claims"] == 6
    assert result["quality_delta"]["full_slides_with_code_generated_ppt_claim"] >= 4
    assert result["quality_delta"]["bad_control_slides_without_run2_57_content"] == 6
    assert "human_release_approval" in result["remaining_public_release_gates"]
    assert "public_blocked" in result["release_boundary"]
    assert result["rerun"]["combined_contact_sheet"].endswith("run2-58-four-arm-contact-sheet.png")
    assert result["rerun"]["full_skill_series_sheet"].endswith("run2-full-skill-series-horizontal.png")

    assert full_trace["arm_id"] == "run2_58_full_product_content_contract"
    assert full_trace["run2_58_product_content_contract_status"] == (
        "run2_57_product_capability_content_consumed_before_native_ppt_drawing"
    )
    assert full_trace["source_repair_run_id"] == "2.57"
    assert full_trace["source_generated_run_id"] == "2.56"
    assert len(full_trace["slides"]) == 6
    for slide in full_trace["slides"]:
        assert set(EXPECTED_RUN2_54_TRACE_FIELDS) <= set(slide)
        assert set(EXPECTED_RUN2_55_TRACE_FIELDS) <= set(slide)
        assert set(EXPECTED_RUN2_56_TRACE_FIELDS) <= set(slide)
        assert set(EXPECTED_RUN2_58_TRACE_FIELDS) <= set(slide)
        assert slide["run2_56_distinct_role_surface_status"] == "pass_internal"
        assert len(slide["run2_57_product_capability_ids"]) >= 2
        assert slide["run2_57_slide_message_contract_id"].startswith("message_contract_2_57_")
        assert slide["run2_57_content_workflow_gate_id"].startswith("gate_2_57_")
        assert slide["run2_57_content_specificity_status"] == "pass_internal"
        assert slide["run2_57_reader_question_answered_status"] == "pass_internal"
        assert slide["run2_57_generic_claim_count"] == 0
        assert slide["run2_57_required_product_terms_rendered"] >= 3
        assert slide["run2_58_product_content_contract_status"] == "pass_internal"
        assert slide["run2_58_reader_question_visible"] is True
        assert slide["run2_58_product_terms_visible_count"] >= 3
        assert slide["run2_58_proof_object_count"] >= 4
        assert "drawRun258ProductContentContract" in slide["run2_58_code_module_ids"]
        assert slide["layout_metrics"]["visible_words"] >= 82

    assert bad_trace["arm_id"] == "bad_run2_56_without_product_capability_content"
    for slide in bad_trace["slides"]:
        assert slide["run2_56_distinct_role_surface_status"] == "pass_internal"
        assert slide["run2_57_product_capability_ids"] == []
        assert slide["run2_57_slide_message_contract_id"] == ""
        assert slide["run2_57_content_specificity_status"] == "fail_missing_run2_57"
        assert slide["run2_57_generic_claim_count"] >= 1
        assert slide["run2_58_bad_control_boundary_status"] == "fail_missing_run2_57_product_content"

    assert_contains(
        result_md,
        [
            "Run 2.58 Product Content Contract Rerun",
            "consumes Run 2.57",
            "product capability content",
            "slide message contracts",
            "content workflow gates",
            "bad_run2_56_without_product_capability_content",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_58_product_content_contract_rerun() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "Run 2.58",
            "ppt-run2-58-prompt-only",
            "ppt-run2-58-run1-5-skill",
            "ppt-run2-58-full-vulca",
            "ppt-run2-58-bad-without-product-capability-content",
            "run2-58-four-arm-contact-sheet.png",
            "run2_58_product_content_contract_rerun_result.json",
            "run2_57_product_capability_content_consumed_before_native_ppt_drawing",
        ],
    )
    assert_contains(
        viewer,
        [
            '"latestRunId": "2.77"',
            "Run 2.58",
            "run2-58-four-arm-contact-sheet.png",
            "ppt-run2-58-full-vulca",
            "Run 2.58 generated result",
            "run2_58_product_content_contract_rerun_result.json",
            "code-generated editable PPT",
        ],
    )


def test_run2_58_records_open_source_slide_code_learning_map() -> None:
    research = load_json(PACK / "run2_58_open_source_slide_code_research.json")

    assert research["status"] == "run2_58_open_source_slide_code_research_ready_public_blocked"
    assert research["selected_usecase_id"] == "usecase_design_to_production_platform_launch"
    assert research["usage_boundary"]["learning_mode"] == "runtime_learning_only"
    assert research["usage_boundary"]["do_not_copy_source_assets"] is True
    assert research["usage_boundary"]["do_not_copy_brand_or_template_layouts"] is True
    assert research["usage_boundary"]["rewrite_as_vulca_native_artifact_tool_modules"] is True

    records = research["open_source_reference_records"]
    assert {record["id"] for record in records} == EXPECTED_RUN2_58_OPEN_SOURCE_REFERENCE_IDS
    for record in records:
        assert record["label"]
        assert record["source_url"].startswith(("https://github.com/", "https://sli.dev/", "https://marp.app/"))
        assert record["source_type"] in {
            "pptx_generation_library",
            "slide_as_code_framework",
            "html_slide_runtime",
            "research_codebase",
        }
        assert record["what_to_learn"]
        assert len(record["learnable_code_patterns"]) >= 3
        assert record["integration_boundary"]
        assert record["license_or_usage_note"]
        assert record["product_relevance"]
        assert record["do_not_copy_source_assets"] is True

    serialized = json.dumps(research, ensure_ascii=False)
    assert_contains(
        serialized,
        [
            "PptxGenJS",
            "Slidev",
            "Marp",
            "reveal.js",
            "SlideCoder",
            "code-generated editable PPT",
            "do_not_copy_source_assets",
            "runtime_learning_only",
        ],
    )


def test_ppt_run_html_viewer_surfaces_run2_58_experiment_lab_and_open_source_map() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    required_terms = [
        "Run 2.58 experiment lab",
        "run2_58_product_content_contract_rerun_result.json",
        "generate_ppt_run2_58_product_content_contract_arms.mjs",
        "ppt-run2-58-full-vulca/trace_manifest.json",
        "run2-58-four-arm-contact-sheet.png",
        "Open-source slide-code learning map",
        "run2_58_open_source_slide_code_research.json",
        "PptxGenJS",
        "Slidev",
        "Marp",
        "reveal.js",
        "SlideCoder",
        "do not copy source assets",
        "code-generated editable PPT",
    ]
    assert_contains(script, required_terms)
    assert_contains(viewer, required_terms)


def test_run2_59_records_content_aware_composition_compiler() -> None:
    contracts = load_json(PACK / "run2_59_content_composition_contracts.json")
    capacity = load_json(PACK / "run2_59_layout_capacity_model.json")
    selector = load_json(PACK / "run2_59_content_to_layout_selector.json")
    trace_policy = load_json(PACK / "run2_59_public_surface_trace_policy.json")
    gates = load_json(PACK / "run2_59_composition_workflow_gates.json")
    result_md = (PACK / "results" / "run2_59_content_aware_composition_compiler_result.md").read_text(
        encoding="utf-8"
    )
    result = load_json(PACK / "results" / "run2_59_content_aware_composition_compiler_result.json")

    assert result["run_id"] == "2.59"
    assert result["status"] == "run2_59_content_aware_composition_compiler_ready_public_blocked"
    assert result["public_ready"] is False
    assert result["selected_usecase_id"] == "usecase_design_to_production_platform_launch"
    assert result["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert result["generation_boundary"]["creates_new_ppt_deck"] is False
    assert result["generation_boundary"]["latest_generated_run_id"] == "2.58"
    assert result["quality_delta"]["target_layer"] == "content_aware_composition_compiler"
    assert result["quality_delta"]["fixes_failure_mode"] == (
        "content_thickness_and_layout_module_memory_not_bound_before_native_ppt_drawing"
    )
    assert result["next_required_action"] == (
        "run2_60_generate_four_arm_ppt_consuming_run2_59_composition_compiler"
    )
    assert set(result["output_chain"]) == {
        "content_composition_contracts",
        "layout_capacity_model",
        "content_to_layout_selector",
        "public_surface_trace_policy",
        "composition_workflow_gates",
    }

    contract_records = contracts["content_composition_contracts"]
    capacity_records = capacity["layout_capacity_records"]
    selector_records = selector["content_to_layout_selection_records"]
    gate_records = gates["composition_workflow_gates"]
    capacity_by_module = {record["layout_module_id"]: record for record in capacity_records}
    contract_by_id = {record["content_contract_id"]: record for record in contract_records}

    assert contracts["status"] == "run2_59_content_composition_contracts_ready_public_blocked"
    assert {record["role"] for record in contract_records} == EXPECTED_RUN2_51_ROLES
    for record in contract_records:
        assert record["content_contract_id"].startswith("content_contract_2_59_")
        assert record["source_message_contract_id"].startswith("message_contract_2_57_")
        assert record["public_claim"]
        assert record["primary_proof_object"]
        assert 1 <= len(record["evidence_chips"]) <= 3
        assert len(record["trace_only_details"]) >= 2
        assert record["max_public_visible_words"] <= 74
        assert set(record["required_trace_fields_for_run2_60"]) == EXPECTED_RUN2_59_TRACE_FIELDS

    assert capacity["status"] == "run2_59_layout_capacity_model_ready_public_blocked"
    assert {record["layout_module_id"] for record in capacity_records} == EXPECTED_RUN2_59_LAYOUT_MODULE_IDS
    for record in capacity_records:
        assert record["source_layout_module_memory"] == "run2_15_layout_module_memory.json"
        assert 1 <= record["max_title_lines"] <= 3
        assert 44 <= record["max_visible_words"] <= 96
        assert 1 <= record["max_evidence_chips"] <= 4
        assert record["primary_object_area_min_pct"] >= 20
        assert record["spacing_rule"]
        assert record["forbidden_patterns"]
        assert record["fallback_if_over_capacity"]
        assert record["overflow_failure_behavior"]["if_public_words_exceed_capacity"] in {
            "move_detail_to_trace_viewer",
            "switch_to_dense_evidence_compression",
            "split_content_before_generation",
        }
        assert record["overflow_failure_behavior"]["run2_60_must_fail_without_fallback"] is True

    assert selector["status"] == "run2_59_content_to_layout_selector_ready_public_blocked"
    assert {record["role"] for record in selector_records} == EXPECTED_RUN2_51_ROLES
    assert {record["selected_layout_module_id"] for record in selector_records} <= EXPECTED_RUN2_59_LAYOUT_MODULE_IDS
    assert any(record["capacity_fit_status"] == "requires_dense_evidence_compression" for record in selector_records)
    for record in selector_records:
        contract = contract_by_id[record["content_contract_id"]]
        selected_capacity = capacity_by_module[record["selected_layout_module_id"]]
        assert record["role"] == contract["role"]
        assert record["content_burden_level"] in {"low", "medium", "high", "climax"}
        assert record["capacity_fit_status"] in {
            "pass_internal",
            "pass_after_trace_separation",
            "requires_dense_evidence_compression",
        }
        assert record["public_visible_word_budget"] <= selected_capacity["max_visible_words"]
        assert record["public_visible_word_budget"] <= contract["max_public_visible_words"]
        assert record["selection_reason"]
        assert record["fallback_route"]
        assert set(record["required_trace_fields_for_run2_60"]) == EXPECTED_RUN2_59_TRACE_FIELDS
        assert record["over_capacity_behavior"]["trigger"] == (
            "public_visible_word_budget_exceeds_selected_layout_capacity"
        )
        assert record["over_capacity_behavior"]["required_action"] in {
            "move_trace_only_details_to_viewer",
            "switch_to_dense_evidence_compression",
            "split_content_before_generation",
        }
        assert record["over_capacity_behavior"]["run2_60_trace_must_record_fallback"] is True

    assert trace_policy["status"] == "run2_59_public_surface_trace_policy_ready_public_blocked"
    assert trace_policy["policy_id"] == "run2_59_public_slide_trace_viewer_split"
    assert "reader_question" in trace_policy["public_slide_surface"]["allowed_content"]
    assert "workflow_gate_raw_checklist" in trace_policy["public_slide_surface"]["forbidden_content"]
    assert "workflow_gate_raw_checklist" in trace_policy["trace_viewer_surface"]["allowed_content"]
    assert trace_policy["run2_60_consumer_obligation"]["must_keep_latest_generated_run_id"] == "2.58_until_2.60_exists"

    assert gates["status"] == "run2_59_composition_workflow_gates_ready_public_blocked"
    assert {record["role"] for record in gate_records} == EXPECTED_RUN2_51_ROLES
    for gate in gate_records:
        assert gate["gate_id"].startswith("gate_2_59_")
        assert gate["content_contract_id"] in contract_by_id
        assert gate["selected_layout_module_id"] in capacity_by_module
        assert set(gate["required_trace_fields"]) == EXPECTED_RUN2_59_TRACE_FIELDS
        assert gate["pass_fail_checks"]
        assert gate["bad_control_probe"]
    assert set(gates["next_generated_run_contract"]["required_trace_fields"]) == EXPECTED_RUN2_59_TRACE_FIELDS
    assert gates["next_generated_run_contract"]["run_id"] == "2.60"
    assert gates["next_generated_run_contract"]["bad_control_arm"] == (
        "bad_run2_58_without_run2_59_composition_compiler"
    )

    assert_contains(
        result_md,
        [
            "Run 2.59 Content-Aware Composition Compiler",
            "data/workflow-only",
            "content-aware composition",
            "layout capacity model",
            "public slide / trace viewer split",
            "Run 2.60",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_59_content_aware_composition_compiler() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    script_terms = [
        "Run 2.59 content-aware composition compiler",
        "run2_59_content_composition_contracts.json",
        "run2_59_layout_capacity_model.json",
        "run2_59_content_to_layout_selector.json",
        "run2_59_public_surface_trace_policy.json",
        "run2_59_composition_workflow_gates.json",
        "run2_59_content_aware_composition_compiler_result.json",
        "content_thickness_and_layout_module_memory_not_bound_before_native_ppt_drawing",
        "public slide / trace viewer split",
        "run2_60_generate_four_arm_ppt_consuming_run2_59_composition_compiler",
    ]
    viewer_terms = [
        *script_terms,
        '"latestRunId": "2.77"',
    ]
    assert_contains(script, script_terms)
    assert_contains(viewer, viewer_terms)
    assert "ppt-run2-59-" not in script
    assert "ppt-run2-59-" not in viewer


def test_run2_60_generator_consumes_run2_59_composition_compiler() -> None:
    script_path = ROOT / "scripts" / "generate_ppt_run2_60_content_aware_composition_arms.mjs"
    assert script_path.exists(), "missing Run 2.60 content-aware composition generator"
    body = script_path.read_text(encoding="utf-8")

    assert_contains(
        body,
        [
            "run2_59_content_composition_contracts.json",
            "run2_59_layout_capacity_model.json",
            "run2_59_content_to_layout_selector.json",
            "run2_59_public_surface_trace_policy.json",
            "run2_59_composition_workflow_gates.json",
            "run2_58_product_content_contract_rerun_result.json",
            "ppt-run2-58-full-vulca/trace_manifest.json",
            "drawRun260ContentAwareComposition",
            "drawRun260EditorialCoverField",
            "drawRun260ProductTheaterStage",
            "drawRun260BeforeAfterRoute",
            "drawRun260DenseEvidenceCompression",
            "drawRun260MetricRevealStage",
            "drawRun260QuietReleaseHandoff",
            "bad_run2_58_without_run2_59_composition_compiler",
            "run2_59_composition_compiler_consumed_before_native_ppt_drawing",
            "public_visible_word_budget_exceeds_selected_layout_capacity",
        ],
    )
    assert "Run 2.60 must consume Run 2.59 content composition contracts" in body
    assert "Run 2.60 must compare against Run 2.58 full trace" in body


def test_run2_60_records_content_aware_composition_rerun_result() -> None:
    result_md = (PACK / "results" / "run2_60_content_aware_composition_rerun_result.md").read_text(
        encoding="utf-8"
    )
    result = load_json(PACK / "results" / "run2_60_content_aware_composition_rerun_result.json")
    presentations = ROOT / "outputs" / "019e7d9c-532a-70b3-8892-fa3ae42baef2" / "presentations"
    full_trace = load_json(presentations / "ppt-run2-60-full-vulca" / "trace_manifest.json")
    bad_trace = load_json(
        presentations / "ppt-run2-60-bad-without-composition-compiler" / "trace_manifest.json"
    )

    assert result["run_id"] == "2.60"
    assert result["status"] == "run2_60_content_aware_composition_rerun_public_blocked"
    assert result["public_ready"] is False
    assert result["selected_usecase_id"] == "usecase_design_to_production_platform_launch"
    assert result["source_repair_run_id"] == "2.59"
    assert result["source_generated_run_id"] == "2.58"
    assert result["rerun"]["best_internal_arm"] == "run2_60_full_content_aware_composition"
    assert result["quality_delta"]["target_layer"] == "content_aware_composition_compiler_consumed"
    assert result["quality_delta"]["source_data_status"] == (
        "run2_59_composition_compiler_consumed_before_native_ppt_drawing"
    )
    assert result["quality_delta"]["full_slides_with_run2_59_content_contracts"] == 6
    assert result["quality_delta"]["full_slides_with_layout_capacity_fit"] == 6
    assert result["quality_delta"]["full_slides_with_public_trace_split"] == 6
    assert result["quality_delta"]["full_slides_without_layout_collision"] == 6
    assert result["quality_delta"]["bad_control_slides_without_run2_59_compiler"] == 6
    assert result["rerun"]["combined_contact_sheet"].endswith("run2-60-four-arm-contact-sheet.png")
    assert result["rerun"]["full_skill_series_sheet"].endswith("run2-full-skill-series-horizontal.png")
    assert "public_blocked" in result["release_boundary"]

    assert full_trace["arm_id"] == "run2_60_full_content_aware_composition"
    assert full_trace["run2_60_content_aware_composition_status"] == (
        "run2_59_composition_compiler_consumed_before_native_ppt_drawing"
    )
    assert full_trace["source_repair_run_id"] == "2.59"
    assert full_trace["source_generated_run_id"] == "2.58"
    assert len(full_trace["slides"]) == 6
    for slide in full_trace["slides"]:
        assert set(EXPECTED_RUN2_59_TRACE_FIELDS) <= set(slide)
        assert set(EXPECTED_RUN2_60_TRACE_FIELDS) <= set(slide)
        assert slide["run2_59_content_contract_id"].startswith("content_contract_2_59_")
        assert slide["run2_59_layout_module_id"] in EXPECTED_RUN2_59_LAYOUT_MODULE_IDS
        assert slide["run2_59_capacity_fit_status"] in {
            "pass_after_trace_separation",
            "requires_dense_evidence_compression",
        }
        assert slide["run2_59_public_visible_word_budget"] <= 74
        assert slide["run2_59_trace_only_detail_count"] >= 2
        assert slide["run2_60_public_trace_split_status"] == "pass_internal"
        assert slide["run2_60_over_capacity_fallback_status"] in {
            "pass_no_over_capacity",
            "pass_fallback_recorded",
        }
        assert slide["run2_60_layout_collision_status"] == "pass_internal"
        assert slide["run2_60_content_aware_composition_status"] == "pass_internal"
        assert "drawRun260ContentAwareComposition" in slide["run2_60_code_module_ids"]
        assert slide["layout_metrics"]["visible_words"] <= 96

    assert bad_trace["arm_id"] == "bad_run2_58_without_run2_59_composition_compiler"
    for slide in bad_trace["slides"]:
        assert slide["run2_59_content_contract_id"] == ""
        assert slide["run2_59_layout_module_id"] == ""
        assert slide["run2_59_capacity_fit_status"] == "fail_missing_run2_59"
        assert slide["run2_60_public_trace_split_status"] == "fail_missing_run2_59_composition_compiler"
        assert slide["run2_60_layout_collision_status"] == "risk_reuses_run2_58_renderer"

    assert_contains(
        result_md,
        [
            "Run 2.60 Content-Aware Composition Rerun",
            "consumes Run 2.59",
            "layout capacity model",
            "content-to-layout selector",
            "public slide / trace viewer split",
            "bad_run2_58_without_run2_59_composition_compiler",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_run2_61_records_narrative_proof_dataset() -> None:
    narrative = load_json(PACK / "run2_61_narrative_proof_dataset.json")
    selector = load_json(PACK / "run2_61_story_to_visual_carrier_selector.json")
    fusion = load_json(PACK / "run2_61_text_socket_fusion_contracts.json")
    source_policy = load_json(PACK / "run2_61_source_to_public_proof_policy.json")
    gates = load_json(PACK / "run2_61_narrative_workflow_gates.json")
    result = load_json(PACK / "results" / "run2_61_narrative_proof_dataset_result.json")
    result_md = (PACK / "results" / "run2_61_narrative_proof_dataset_result.md").read_text(
        encoding="utf-8"
    )

    assert narrative["status"] == "run2_61_narrative_proof_dataset_ready_public_blocked"
    assert selector["status"] == "run2_61_story_to_visual_carrier_selector_ready_public_blocked"
    assert fusion["status"] == "run2_61_text_socket_fusion_contracts_ready_public_blocked"
    assert source_policy["status"] == "run2_61_source_to_public_proof_policy_ready_public_blocked"
    assert gates["status"] == "run2_61_narrative_workflow_gates_ready_public_blocked"
    assert result["run_id"] == "2.61"
    assert result["generation_boundary"]["creates_new_ppt_deck"] is False

    narrative_records = narrative["narrative_proof_records"]
    selector_records = selector["story_to_visual_carrier_records"]
    fusion_records = fusion["text_socket_fusion_contracts"]
    gate_records = gates["narrative_workflow_gates"]

    assert len(narrative_records) == len(EXPECTED_RUN2_61_ROLES)
    assert len(selector_records) == len(EXPECTED_RUN2_61_ROLES)
    assert len(fusion_records) == len(EXPECTED_RUN2_61_ROLES)
    assert len(gate_records) == len(EXPECTED_RUN2_61_ROLES)

    assert {record["role"] for record in narrative_records} == EXPECTED_RUN2_61_ROLES
    assert {record["role"] for record in selector_records} == EXPECTED_RUN2_61_ROLES
    assert {record["role"] for record in fusion_records} == EXPECTED_RUN2_61_ROLES
    assert {record["role"] for record in gate_records} == EXPECTED_RUN2_61_ROLES

    selector_by_role = {record["role"]: record for record in selector_records}
    fusion_by_role = {record["role"]: record for record in fusion_records}
    gate_by_role = {record["role"]: record for record in gate_records}

    for record in narrative_records:
        assert record["narrative_proof_id"].startswith("narrative_proof_2_61_")
        assert record["selected_usecase_id"] == "usecase_design_to_production_platform_launch"
        assert record["reader_question"]
        assert record["required_answer"]
        assert record["business_action"]
        assert record["public_takeaway"]
        assert set(record["copy_units"]) == EXPECTED_RUN2_61_REQUIRED_COPY_UNITS
        assert len(record["source_refs"]) >= 2
        assert len(record["proof_payload"]["secondary_evidence_objects"]) >= 2
        assert record["proof_payload"]["product_mechanism"]
        assert record["proof_payload"]["business_consequence"]
        assert record["public_proof_replacement"]["replacement_type"] in {
            "source_pack_object",
            "inspection_board",
            "before_after_route_break",
            "native_editable_proxy",
            "operating_loop_proxy",
            "release_gate_proxy",
        }
        assert record["density_budget"]["minimum_public_proof_objects"] >= 2
        assert record["density_budget"]["minimum_socket_bound_copy_units"] >= 4
        assert record["density_budget"]["maximum_free_floating_labels"] <= 2
        assert record["bad_control_probe"].startswith("fail_if_")
        assert record["next_rerun_obligation"] == "must_be_consumed_before_run2_62_four_arm_rerun"

        selected = selector_by_role[record["role"]]
        fused = fusion_by_role[record["role"]]
        gate = gate_by_role[record["role"]]
        assert selected["source_narrative_proof_id"] == record["narrative_proof_id"]
        assert selected["selected_layout_module_id"].startswith("module_2_15_")
        assert selected["selected_socket_memory_id"].startswith("shape_text_socket_2_51_")
        assert selected["visual_carrier_type"] in EXPECTED_RUN2_61_CARRIER_TYPES
        assert selected["visual_weight_requirement"]["minimum_primary_carrier_area_pct"] >= 30
        assert fused["source_narrative_proof_id"] == record["narrative_proof_id"]
        assert fused["source_socket_memory_id"] == selected["selected_socket_memory_id"]
        assert {binding["copy_unit_key"] for binding in fused["socket_bindings"]} >= {
            "headline",
            "subhead",
            "proof_badges",
            "annotations",
            "state_labels",
        }
        assert gate["source_narrative_proof_id"] == record["narrative_proof_id"]
        assert gate["next_run_required_trace_fields"] == [
            "run2_61_narrative_proof_id",
            "run2_61_visual_carrier_selector_id",
            "run2_61_text_socket_fusion_contract_id",
            "run2_61_public_proof_replacement_id",
            "run2_61_narrative_workflow_gate_id",
        ]

    assert source_policy["policy_id"] == "run2_61_source_to_public_proof_policy"
    assert "copy_source_slide_or_video_frame" in source_policy["forbidden_source_copying_behaviors"]
    assert "native_editable_proxy_object" in source_policy["allowed_source_abstraction_types"]
    assert result["next_required_action"] == "run2_62_generate_four_arm_ppt_consuming_run2_61_narrative_proof_dataset"
    assert "Run 2.61 Narrative Proof Dataset" in result_md
    assert "Do not advance to Run 3.0" in result_md


def test_ppt_run_html_viewer_mentions_run2_60_content_aware_composition_rerun() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    required_terms = [
        "Run 2.60",
        "ppt-run2-60-prompt-only",
        "ppt-run2-60-run1-5-skill",
        "ppt-run2-60-full-vulca",
        "ppt-run2-60-bad-without-composition-compiler",
        "run2-60-four-arm-contact-sheet.png",
        "run2_60_content_aware_composition_rerun_result.json",
        "run2_59_composition_compiler_consumed_before_native_ppt_drawing",
        '"latestRunId": "2.77"',
        "Run 2.60 generated result",
        "content-aware composition compiler consumed",
    ]
    assert_contains(script, required_terms[:-2])
    assert_contains(viewer, required_terms)


def test_ppt_run_html_viewer_mentions_run2_61_narrative_proof_dataset() -> None:
    viewer_script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer_html = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")
    for body in (viewer_script, viewer_html):
        assert_contains(
            body,
            [
                "Run 2.61 narrative proof dataset",
                "Why 2.60 felt thin",
                "Per-slide narrative proof table",
                "Visual carrier selector",
                "Text socket fusion",
                "Run 2.61 expanded role contracts",
                "2.61 consumption contract",
                "Socket plan",
                "Trace gate",
                "maximum_public_visible_words",
                "Public proof replacement",
                "run2_61_narrative_proof_dataset.json",
                "run2_61_story_to_visual_carrier_selector.json",
                "run2_61_text_socket_fusion_contracts.json",
                "run2_61_source_to_public_proof_policy.json",
                "run2_61_narrative_workflow_gates.json",
                "run2_62_generate_four_arm_ppt_consuming_run2_61_narrative_proof_dataset",
                '"latestRunId": "2.77"',
            ],
        )
    assert_contains(
        viewer_html,
        [
            "input-to-output product surface",
            "socket_2_61_cover_headline",
            "gate_2_61_cover_narrative_proof_fusion",
            "fail_if_cover_renders_claim_without_narrative_proof_and_socket_fusion",
        ],
    )


def test_run2_62_records_narrative_proof_rerun_result() -> None:
    result_md = (PACK / "results" / "run2_62_narrative_proof_rerun_result.md").read_text(
        encoding="utf-8"
    )
    result = load_json(PACK / "results" / "run2_62_narrative_proof_rerun_result.json")
    presentations = ROOT / "outputs" / "019e7d9c-532a-70b3-8892-fa3ae42baef2" / "presentations"
    full_trace = load_json(presentations / "ppt-run2-62-full-vulca" / "trace_manifest.json")
    bad_trace = load_json(
        presentations / "ppt-run2-62-bad-without-narrative-proof" / "trace_manifest.json"
    )

    assert result["run_id"] == "2.62"
    assert result["status"] == "run2_62_narrative_proof_rerun_public_blocked"
    assert result["public_ready"] is False
    assert result["selected_usecase_id"] == "usecase_design_to_production_platform_launch"
    assert result["source_repair_run_id"] == "2.61"
    assert result["source_generated_run_id"] == "2.60"
    assert result["rerun"]["best_internal_arm"] == "run2_62_full_narrative_proof"
    assert result["quality_delta"]["target_layer"] == "run2_61_narrative_proof_dataset_consumed"
    assert result["quality_delta"]["source_data_status"] == (
        "run2_61_narrative_proof_dataset_consumed_before_native_ppt_drawing"
    )
    assert result["quality_delta"]["full_slides_with_run2_61_contracts"] == 6
    assert result["quality_delta"]["full_slides_with_socket_bindings"] == 6
    assert result["quality_delta"]["full_slides_with_public_proof_replacements"] == 6
    assert result["quality_delta"]["bad_control_slides_without_run2_61_contracts"] == 6
    assert result["rerun"]["combined_contact_sheet"].endswith("run2-62-four-arm-contact-sheet.png")
    assert result["rerun"]["full_skill_series_sheet"].endswith("run2-full-skill-series-horizontal.png")
    assert "public_blocked" in result["release_boundary"]

    assert full_trace["arm_id"] == "run2_62_full_narrative_proof"
    assert full_trace["run2_62_narrative_proof_consumption_status"] == (
        "run2_61_narrative_proof_dataset_consumed_before_native_ppt_drawing"
    )
    assert full_trace["source_repair_run_id"] == "2.61"
    assert full_trace["source_generated_run_id"] == "2.60"
    assert len(full_trace["slides"]) == 6
    for slide in full_trace["slides"]:
        assert set(EXPECTED_RUN2_62_TRACE_FIELDS) <= set(slide)
        assert slide["run2_61_narrative_proof_id"].startswith("narrative_proof_2_61_")
        assert slide["run2_61_visual_carrier_selector_id"].startswith("visual_carrier_selector_2_61_")
        assert slide["run2_61_text_socket_fusion_contract_id"].startswith("text_socket_fusion_2_61_")
        assert slide["run2_61_public_proof_replacement_id"].startswith("public_proof_replacement_2_61_")
        assert slide["run2_61_narrative_workflow_gate_id"].startswith("gate_2_61_")
        assert slide["run2_62_narrative_proof_consumption_status"] == "pass_internal"
        assert slide["run2_62_socket_binding_count"] >= 5
        assert slide["run2_62_public_proof_object_count"] >= 2
        assert slide["run2_62_required_answer_visible"] == "pass_internal"
        assert "drawRun262NarrativeProofComposition" in slide["run2_62_code_module_ids"]
        assert slide["layout_metrics"]["visible_words"] <= 120

    assert bad_trace["arm_id"] == "bad_run2_60_without_run2_61_narrative_proof_dataset"
    for slide in bad_trace["slides"]:
        for field in EXPECTED_RUN2_61_CONSUMER_TRACE_FIELDS:
            assert slide[field] == ""
        assert slide["run2_62_narrative_proof_consumption_status"] == "fail_missing_run2_61"
        assert slide["run2_62_socket_binding_count"] == 0
        assert slide["run2_62_required_answer_visible"] == "fail_missing_run2_61"

    assert_contains(
        result_md,
        [
            "Run 2.62 Narrative Proof Rerun",
            "consumes Run 2.61",
            "bad_run2_60_without_run2_61_narrative_proof_dataset",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_62_narrative_proof_rerun() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    required_terms = [
        "Run 2.62",
        "ppt-run2-62-prompt-only",
        "ppt-run2-62-run1-5-skill",
        "ppt-run2-62-full-vulca",
        "ppt-run2-62-bad-without-narrative-proof",
        "run2-62-four-arm-contact-sheet.png",
        "run2_62_narrative_proof_rerun_result.json",
        "run2_61_narrative_proof_dataset_consumed_before_native_ppt_drawing",
        '"latestRunId": "2.77"',
        "Run 2.62 generated narrative proof consumption",
    ]
    assert_contains(script, required_terms[:-2])
    assert_contains(viewer, required_terms)


def test_run2_63_audits_narrative_proof_consumption_effectiveness() -> None:
    audit_md = (
        PACK / "results" / "run2_63_narrative_proof_consumption_effectiveness_audit.md"
    ).read_text(encoding="utf-8")
    audit = load_json(
        PACK / "results" / "run2_63_narrative_proof_consumption_effectiveness_audit.json"
    )

    assert audit["run_id"] == "2.63"
    assert audit["status"] == (
        "run2_63_narrative_proof_consumption_effectiveness_audit_public_blocked"
    )
    assert audit["creates_new_ppt_deck"] is False
    assert audit["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert audit["source_generated_run"] == "2.62"
    assert audit["source_data_run"] == "2.61"
    assert audit["data_consumption_assessment"]["run2_61_data_consumed"] is True
    assert audit["data_consumption_assessment"]["full_slides_with_run2_61_contracts"] == 6
    assert audit["data_consumption_assessment"]["full_slides_with_socket_bindings"] == 6
    assert audit["data_consumption_assessment"]["bad_control_without_run2_61_contracts"] == 6
    assert audit["root_cause_assessment"]["primary_layer"] == "renderer_composition_grammar"
    assert audit["root_cause_assessment"]["not_primary_layer"] == "raw_data_or_workflow_consumption"
    assert set(audit["root_cause_assessment"]["renderer_blockers"]) >= (
        EXPECTED_RUN2_63_RENDERER_BLOCKERS
    )
    assert audit["gate_summary"]["data_consumption_gate"] == "pass_internal"
    assert audit["gate_summary"]["renderer_effectiveness_gate"] == "blocked"
    assert audit["gate_summary"]["public_release_gate"] == "blocked"
    assert audit["next_required_action"] == (
        "build_run2_64_renderer_composition_repair_for_dynamic_sockets_and_semantic_diagrams_before_rerun"
    )

    assert_contains(
        audit_md,
        [
            "Run 2.63 Narrative Proof Consumption Effectiveness Audit",
            "2.62 consumes 2.61",
            "renderer/composition grammar",
            "static socket plan",
            "text fit",
            "Run 2.64",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_63_consumption_effectiveness_audit() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    required_terms = [
        "Run 2.63 narrative proof effectiveness audit",
        "run2_63_narrative_proof_consumption_effectiveness_audit.json",
        "renderer_composition_grammar",
        "static_socket_plan_repeated_on_every_slide",
        "build_run2_64_renderer_composition_repair_for_dynamic_sockets_and_semantic_diagrams_before_rerun",
    ]
    assert_contains(script, required_terms)
    assert_contains(viewer, required_terms)


def test_run2_64_records_renderer_composition_repair_contract() -> None:
    dynamic_socket = load_json(PACK / "run2_64_dynamic_socket_renderer_memory.json")
    semantic_diagram = load_json(PACK / "run2_64_semantic_diagram_renderer_memory.json")
    text_fit = load_json(PACK / "run2_64_text_fit_renderer_gates.json")
    dry_run = load_json(PACK / "run2_64_renderer_dry_run_binding_matrix.json")
    result = load_json(PACK / "results" / "run2_64_renderer_composition_repair_result.json")
    result_md = (PACK / "results" / "run2_64_renderer_composition_repair_result.md").read_text(
        encoding="utf-8"
    )

    assert dynamic_socket["status"] == "run2_64_dynamic_socket_renderer_memory_ready_public_blocked"
    assert semantic_diagram["status"] == "run2_64_semantic_diagram_renderer_memory_ready_public_blocked"
    assert text_fit["status"] == "run2_64_text_fit_renderer_gates_ready_public_blocked"
    assert dry_run["status"] == "run2_64_renderer_dry_run_binding_matrix_ready_public_blocked"
    assert result["run_id"] == "2.64"
    assert result["status"] == "run2_64_renderer_composition_repair_ready_public_blocked"
    assert result["creates_new_ppt_deck"] is False
    assert result["source_audit_run"] == "2.63"
    assert result["source_generated_run"] == "2.62"
    assert result["quality_delta"]["target_layer"] == "renderer_composition_grammar_repair_contract"
    assert set(result["quality_delta"]["repairs_blockers"]) >= EXPECTED_RUN2_63_RENDERER_BLOCKERS
    assert result["next_required_action"] == (
        "run2_65_generate_four_arm_ppt_consuming_run2_64_renderer_composition_repair"
    )

    dynamic_records = dynamic_socket["dynamic_socket_renderer_records"]
    diagram_records = semantic_diagram["semantic_diagram_renderer_records"]
    text_gate_records = text_fit["text_fit_renderer_gates"]
    dry_records = dry_run["dry_run_binding_records"]
    assert {record["role"] for record in dynamic_records} == EXPECTED_RUN2_61_ROLES
    assert {record["role"] for record in diagram_records} == EXPECTED_RUN2_61_ROLES
    assert {record["role"] for record in text_gate_records} == EXPECTED_RUN2_61_ROLES
    assert {record["role"] for record in dry_records} == EXPECTED_RUN2_61_ROLES

    for record in dynamic_records:
        assert record["source_run2_63_blockers"] == sorted(EXPECTED_RUN2_63_RENDERER_BLOCKERS)
        assert record["static_socket_plan_replaced"] is True
        assert set(record["active_copy_unit_keys"]) == EXPECTED_RUN2_61_REQUIRED_COPY_UNITS
        assert len(record["active_socket_bindings"]) == len(EXPECTED_RUN2_61_REQUIRED_COPY_UNITS)
        assert set(record["required_trace_fields_for_run2_65"]) == EXPECTED_RUN2_64_TRACE_FIELDS

    for record in diagram_records:
        assert record["source_visual_carrier_type"] in EXPECTED_RUN2_61_CARRIER_TYPES
        assert record["forbid_generic_repeated_shape_system"] is True
        assert len(record["proof_object_bindings"]) >= 3
        assert set(record["required_trace_fields_for_run2_65"]) == EXPECTED_RUN2_64_TRACE_FIELDS

    for record in text_gate_records:
        assert record["runtime_claim_boundary"] == "must_be_verified_by_run2_65_render_trace"
        assert record["preflight_gate_status"] == "metadata_gate_only_not_runtime_proof"
        assert record["forbid_ellipsis_or_clipped_public_copy"] is True
        assert set(record["required_trace_fields_for_run2_65"]) == EXPECTED_RUN2_64_TRACE_FIELDS

    for record in dry_records:
        assert record["all_required_copy_units_bound"] is True
        assert record["orphan_socket_count"] == 0
        assert record["unbound_copy_unit_count"] == 0
        assert record["ready_for_run2_65_consumption"] is True

    assert_contains(
        result_md,
        [
            "Run 2.64 Renderer Composition Repair",
            "dynamic socket",
            "semantic diagram",
            "text-fit",
            "dry-run binding matrix",
            "Run 2.65",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_64_renderer_composition_repair() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    required_terms = [
        "Run 2.64 renderer composition repair",
        "run2_64_dynamic_socket_renderer_memory.json",
        "run2_64_semantic_diagram_renderer_memory.json",
        "run2_64_text_fit_renderer_gates.json",
        "run2_64_renderer_dry_run_binding_matrix.json",
        "run2_65_generate_four_arm_ppt_consuming_run2_64_renderer_composition_repair",
    ]
    assert_contains(script, required_terms)
    assert_contains(viewer, required_terms)


def test_run2_65_records_renderer_composition_rerun_result() -> None:
    result_md = (PACK / "results" / "run2_65_renderer_composition_rerun_result.md").read_text(
        encoding="utf-8"
    )
    result = load_json(PACK / "results" / "run2_65_renderer_composition_rerun_result.json")
    presentations = ROOT / "outputs" / "019e7d9c-532a-70b3-8892-fa3ae42baef2" / "presentations"
    full_trace = load_json(presentations / "ppt-run2-65-full-vulca" / "trace_manifest.json")
    bad_trace = load_json(
        presentations / "ppt-run2-65-bad-without-renderer-composition-repair" / "trace_manifest.json"
    )

    assert result["run_id"] == "2.65"
    assert result["status"] == "run2_65_renderer_composition_rerun_public_blocked"
    assert result["public_ready"] is False
    assert result["selected_usecase_id"] == "usecase_design_to_production_platform_launch"
    assert result["source_repair_run_id"] == "2.64"
    assert result["source_generated_run_id"] == "2.62"
    assert result["rerun"]["best_internal_arm"] == "run2_65_full_renderer_composition_repair"
    assert result["quality_delta"]["target_layer"] == "run2_64_renderer_composition_repair_consumed"
    assert result["quality_delta"]["source_data_status"] == (
        "run2_64_renderer_composition_repair_consumed_before_native_ppt_drawing"
    )
    assert result["quality_delta"]["full_slides_with_run2_64_contracts"] == 6
    assert result["quality_delta"]["full_slides_with_dynamic_sockets"] == 6
    assert result["quality_delta"]["full_slides_with_semantic_diagrams"] == 6
    assert result["quality_delta"]["full_slides_with_text_fit_runtime_pass"] == 6
    assert result["quality_delta"]["bad_control_slides_without_run2_64_contracts"] == 6
    assert result["rerun"]["combined_contact_sheet"].endswith("run2-65-four-arm-contact-sheet.png")
    assert result["rerun"]["full_skill_series_sheet"].endswith("run2-full-skill-series-horizontal.png")

    assert full_trace["arm_id"] == "run2_65_full_renderer_composition_repair"
    assert full_trace["run2_65_renderer_composition_repair_status"] == (
        "run2_64_renderer_composition_repair_consumed_before_native_ppt_drawing"
    )
    assert full_trace["source_repair_run_id"] == "2.64"
    assert full_trace["source_generated_run_id"] == "2.62"
    assert len(full_trace["slides"]) == 6
    for slide in full_trace["slides"]:
        assert set(EXPECTED_RUN2_65_TRACE_FIELDS) <= set(slide)
        assert slide["run2_64_dynamic_socket_renderer_id"].startswith("dynamic_socket_renderer_2_64_")
        assert slide["run2_64_semantic_diagram_renderer_id"].startswith("semantic_diagram_renderer_2_64_")
        assert slide["run2_64_text_fit_gate_id"].startswith("text_fit_gate_2_64_")
        assert slide["run2_64_renderer_dry_run_binding_id"].startswith("renderer_dry_run_binding_2_64_")
        assert slide["run2_64_dynamic_socket_plan_status"] == "pass_internal"
        assert slide["run2_64_semantic_diagram_binding_status"] == "pass_internal"
        assert slide["run2_64_text_fit_preflight_status"] == "pass_internal"
        assert slide["run2_65_renderer_composition_repair_status"] == "pass_internal"
        assert slide["run2_65_dynamic_socket_count"] >= 6
        assert slide["run2_65_semantic_diagram_type"]
        assert slide["run2_65_text_fit_runtime_status"] == "pass_internal"
        assert slide["run2_65_visual_delta_from_run2_62"] == "dynamic_socket_semantic_diagram_text_fit_delta"
        assert "drawRun265RendererCompositionRepair" in slide["run2_65_code_module_ids"]
        assert slide["layout_metrics"]["visible_words"] <= 125
        assert slide["layout_metrics"]["text_overflow_risk_count"] == 0

    full_layout_dir = presentations / "ppt-run2-65-full-vulca" / "layout" / "final"
    for layout_path in sorted(full_layout_dir.glob("slide-*.layout.json")):
        layout = load_json(layout_path)
        text_values: list[str] = []

        def collect_text_values(value: object) -> None:
            if isinstance(value, dict):
                for key, nested in value.items():
                    if key in {"text", "textPreview"} and isinstance(nested, str):
                        text_values.append(nested)
                    else:
                        collect_text_values(nested)
            elif isinstance(value, list):
                for item in value:
                    collect_text_values(item)

        collect_text_values(layout)
        assert text_values
        assert all("..." not in text_value and "…" not in text_value for text_value in text_values)

    assert bad_trace["arm_id"] == "bad_run2_64_without_renderer_composition_repair"
    for slide in bad_trace["slides"]:
        for field in EXPECTED_RUN2_64_TRACE_FIELDS:
            assert slide[field] == "" or slide[field] == "fail_missing_run2_64"
        assert slide["run2_65_renderer_composition_repair_status"] == "fail_missing_run2_64"
        assert slide["run2_65_dynamic_socket_count"] == 0
        assert slide["run2_65_text_fit_runtime_status"] == "fail_missing_run2_64"

    assert_contains(
        result_md,
        [
            "Run 2.65 Renderer Composition Rerun",
            "consumes Run 2.64",
            "dynamic socket",
            "semantic diagram",
            "text-fit",
            "bad_run2_64_without_renderer_composition_repair",
            "public blocked",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_65_renderer_composition_rerun() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    required_terms = [
        "Run 2.65",
        "ppt-run2-65-prompt-only",
        "ppt-run2-65-run1-5-skill",
        "ppt-run2-65-full-vulca",
        "ppt-run2-65-bad-without-renderer-composition-repair",
        "run2-65-four-arm-contact-sheet.png",
        "run2_65_renderer_composition_rerun_result.json",
        "run2_64_renderer_composition_repair_consumed_before_native_ppt_drawing",
        '"latestRunId": "2.77"',
    ]
    assert_contains(script, required_terms[:-1])
    assert_contains(viewer, required_terms)


def test_run2_66_records_reference_first_redesign_contract() -> None:
    failure_audit = load_json(PACK / "run2_66_visual_failure_audit.json")
    design_grammar = load_json(PACK / "run2_66_reference_first_design_grammar.json")
    art_direction = load_json(PACK / "run2_66_slide_art_direction_contracts.json")
    workflow_gates = load_json(PACK / "run2_66_reference_first_workflow_gates.json")
    result = load_json(PACK / "results" / "run2_66_reference_first_redesign_result.json")
    result_md = (PACK / "results" / "run2_66_reference_first_redesign_result.md").read_text(
        encoding="utf-8"
    )

    assert failure_audit["status"] == "run2_66_visual_failure_audit_ready_public_blocked"
    assert failure_audit["source_generated_run"] == "2.65"
    assert failure_audit["failure_layer"] == "public_surface_art_direction_not_renderer_contracts"
    assert {
        "engineering_proof_surface_dominates_public_surface",
        "reference_data_reduced_to_trace_contracts",
        "renderer_shape_vocabulary_too_primitive",
        "public_story_lacks_scene_specific_business_evidence",
        "tests_measure_trace_not_public_video_aesthetic_grade",
    } <= set(failure_audit["root_causes"])
    assert failure_audit["blocked_next_action"] == "do_not_tune_run2_65_rectangles"

    assert design_grammar["status"] == "run2_66_reference_first_design_grammar_ready_public_blocked"
    assert design_grammar["run_id"] == "2.66"
    assert design_grammar["source_generated_run"] == "2.65"
    assert design_grammar["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert design_grammar["raw_media_policy"] == "forbidden_reference_analysis_only"
    assert len(design_grammar["reference_records"]) >= 5
    assert len(design_grammar["role_design_grammar_records"]) == 6
    assert {
        "stripe_sessions_2025_product_keynote",
        "figma_demo_platform_overview",
        "figma_config_2026_opening_keynote",
        "figma_config_2024_slides_keynote",
    } <= {record["reference_id"] for record in design_grammar["reference_records"]}

    for record in design_grammar["role_design_grammar_records"]:
        assert record["reference_archetype_id"].startswith("reference_first_archetype_2_66_")
        assert record["role"] in EXPECTED_RUN2_61_ROLES
        assert len(record["source_reference_ids"]) >= 3
        assert record["public_first_read_object"]
        assert record["layout_archetype"] in {
            "cinematic_product_reveal",
            "operating_loop_stage",
            "before_after_theater",
            "inspectable_product_workspace",
            "hero_product_surface_demo",
            "release_decision_map",
        }
        assert record["composition_contract"]["forbid_rect_ellipse_only_primary_surface"] is True
        assert record["composition_contract"]["forbid_bottom_socket_legend_as_primary_information"] is True
        assert record["composition_contract"]["min_depth_layers"] >= 4
        assert record["composition_contract"]["min_scene_specific_business_objects"] >= 3
        assert record["composition_contract"]["max_visible_trace_terms"] == 0
        assert len(record["native_ppt_required_modules"]) >= 4

    assert art_direction["status"] == "run2_66_slide_art_direction_contracts_ready_public_blocked"
    assert len(art_direction["slide_art_direction_contracts"]) == 6
    for contract in art_direction["slide_art_direction_contracts"]:
        assert contract["role"] in EXPECTED_RUN2_61_ROLES
        assert contract["content_strategy"]["min_scene_specific_content_units"] >= 4
        assert contract["copy_strategy"]["public_copy_mode"] == "reader_outcome_first"
        assert contract["visual_strategy"]["must_show_product_or_business_scene"] is True
        assert contract["qa_contract"]["human_review_question"] == (
            "Would this read as a public product/keynote slide without seeing trace metadata?"
        )

    assert workflow_gates["status"] == "run2_66_reference_first_workflow_gates_ready_public_blocked"
    assert workflow_gates["next_generated_run_contract"]["run_id"] == "2.67"
    assert workflow_gates["next_generated_run_contract"]["must_consume_before_native_ppt_drawing"] is True
    assert workflow_gates["bad_control_arm"] == "bad_run2_65_without_run2_66_reference_first_grammar"
    assert set(workflow_gates["required_trace_fields_for_run2_67"]) >= {
        "run2_66_reference_archetype_id",
        "run2_66_public_first_read_object",
        "run2_66_layout_archetype",
        "run2_66_art_direction_contract_id",
        "run2_66_public_surface_aesthetic_gate_status",
    }

    assert result["run_id"] == "2.66"
    assert result["status"] == "run2_66_reference_first_redesign_ready_public_blocked"
    assert result["creates_new_ppt_deck"] is False
    assert result["source_generated_run"] == "2.65"
    assert result["quality_delta"]["target_layer"] == "reference_first_public_surface_art_direction"
    assert result["quality_delta"]["role_design_grammar_records"] == 6
    assert result["next_required_action"] == (
        "run2_67_generate_four_arm_ppt_consuming_run2_66_reference_first_design_grammar"
    )
    assert_contains(
        result_md,
        [
            "Run 2.66 Reference-First Redesign",
            "does not generate a new PPT",
            "reference-first design grammar",
            "public surface art direction",
            "Run 2.67",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_66_reference_first_redesign() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    required_terms = [
        "Run 2.66 reference-first redesign",
        "run2_66_visual_failure_audit.json",
        "run2_66_reference_first_design_grammar.json",
        "run2_66_slide_art_direction_contracts.json",
        "run2_66_reference_first_workflow_gates.json",
        "run2_67_generate_four_arm_ppt_consuming_run2_66_reference_first_design_grammar",
        '"latestRunId": "2.77"',
    ]
    assert_contains(script, required_terms[:-1])
    assert_contains(viewer, required_terms)


def test_run2_67_generates_reference_first_public_surface_rerun() -> None:
    result = load_json(PACK / "results" / "run2_67_reference_first_rerun_result.json")
    result_md = (PACK / "results" / "run2_67_reference_first_rerun_result.md").read_text(
        encoding="utf-8"
    )
    presentations = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
    )
    full_trace = load_json(presentations / "ppt-run2-67-full-vulca" / "trace_manifest.json")
    bad_trace = load_json(
        presentations / "ppt-run2-67-bad-without-reference-first-grammar" / "trace_manifest.json"
    )

    assert result["run_id"] == "2.67"
    assert result["status"] == "run2_67_reference_first_rerun_public_blocked"
    assert result["source_design_run_id"] == "2.66"
    assert result["source_generated_run_id"] == "2.65"
    assert result["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert result["rerun"]["best_internal_arm"] == "run2_67_full_reference_first_design_grammar"
    assert result["rerun"]["combined_contact_sheet"].endswith("run2-67-four-arm-contact-sheet.png")
    assert result["quality_delta"]["full_slides_with_run2_66_reference_grammar"] == 6
    assert result["quality_delta"]["full_slides_with_reference_archetypes"] == 6
    assert result["quality_delta"]["full_slides_with_public_surface_aesthetic_gate"] == 6
    assert result["quality_delta"]["full_slides_with_visible_trace_terms_removed"] == 6
    assert result["quality_delta"]["bad_control_slides_without_run2_66_reference_grammar"] == 6
    assert result["control_boundary"]["bad_run2_65_without_run2_66_reference_first_grammar"]

    assert full_trace["arm_id"] == "run2_67_full_reference_first_design_grammar"
    assert full_trace["source_design_run_id"] == "2.66"
    assert full_trace["source_generated_run_id"] == "2.65"
    assert full_trace["run2_67_reference_first_consumption_status"] == (
        "run2_66_reference_first_design_grammar_consumed_before_native_ppt_drawing"
    )
    assert len(full_trace["slides"]) == 6
    for slide in full_trace["slides"]:
        assert slide["role"] in EXPECTED_RUN2_61_ROLES
        assert slide["run2_66_reference_archetype_id"].startswith(
            "reference_first_archetype_2_66_"
        )
        assert slide["run2_66_public_first_read_object"]
        assert slide["run2_66_layout_archetype"] in {
            "cinematic_product_reveal",
            "operating_loop_stage",
            "before_after_theater",
            "inspectable_product_workspace",
            "hero_product_surface_demo",
            "release_decision_map",
        }
        assert slide["run2_66_art_direction_contract_id"].startswith(
            "art_direction_contract_2_66_"
        )
        assert slide["run2_66_public_surface_aesthetic_gate_status"] == "pass_internal"
        assert slide["run2_67_reference_first_code_module_id"].startswith(
            "drawRun267ReferenceFirst"
        )
        assert slide["run2_67_scene_specific_business_object_count"] >= 3
        assert slide["run2_67_depth_layer_count"] >= 4
        assert slide["run2_67_visible_trace_terms"] == 0
        assert slide["layout_metrics"]["trace_panel_visible"] is False
        assert slide["layout_metrics"]["visible_words"] >= 42
        assert slide["layout_metrics"]["visible_words"] <= 132

    assert bad_trace["arm_id"] == "bad_run2_65_without_run2_66_reference_first_grammar"
    for slide in bad_trace["slides"]:
        assert slide["run2_66_reference_archetype_id"] == ""
        assert slide["run2_66_art_direction_contract_id"] == ""
        assert slide["run2_66_public_surface_aesthetic_gate_status"] == (
            "fail_missing_run2_66"
        )

    assert_contains(
        result_md,
        [
            "Run 2.67 Reference-First Rerun",
            "consumes Run 2.66",
            "reference-first design grammar",
            "public surface art direction",
            "bad control",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_67_reference_first_rerun() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    required_terms = [
        "Run 2.67",
        "run2_67_reference_first_rerun_result.json",
        "ppt-run2-67-full-vulca",
        "ppt-run2-67-bad-without-reference-first-grammar",
        "run2-67-four-arm-contact-sheet.png",
        '"latestRunId": "2.77"',
    ]
    assert_contains(script, required_terms[:-1])
    assert_contains(viewer, required_terms)


def test_run2_68_generates_targeted_renderer_debug_rerun() -> None:
    result = load_json(PACK / "results" / "run2_68_targeted_debug_rerun_result.json")
    result_md = (PACK / "results" / "run2_68_targeted_debug_rerun_result.md").read_text(
        encoding="utf-8"
    )
    presentations = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
    )
    full_trace = load_json(presentations / "ppt-run2-68-full-vulca" / "trace_manifest.json")
    bad_trace = load_json(
        presentations / "ppt-run2-68-bad-without-targeted-debug" / "trace_manifest.json"
    )

    assert result["run_id"] == "2.68"
    assert result["status"] == "run2_68_targeted_debug_rerun_public_blocked"
    assert result["source_generated_run_id"] == "2.67"
    assert result["source_design_run_id"] == "2.66"
    assert result["debug_scope"]["targeted_roles"] == ["setup", "proof", "close"]
    assert result["rerun"]["best_internal_arm"] == "run2_68_full_targeted_debug_repair"
    assert result["quality_delta"]["targeted_debug_roles_fixed"] == 3
    assert result["quality_delta"]["full_slides_with_text_overlap_risk_removed"] == 6
    assert result["quality_delta"]["bad_control_slides_without_targeted_debug_repair"] == 3

    targeted_modules = {
        "setup": "drawRun268LayeredOperatingStageDebug",
        "proof": "drawRun268InspectableWorkspaceDebug",
        "close": "drawRun268ReleaseDecisionBoardDebug",
    }
    assert full_trace["arm_id"] == "run2_68_full_targeted_debug_repair"
    for slide in full_trace["slides"]:
        assert slide["run2_68_text_overlap_risk_status"] == "pass_internal"
        assert slide["layout_metrics"]["text_overflow_risk_count"] == 0
        if slide["role"] in targeted_modules:
            assert slide["run2_68_targeted_debug_role"] is True
            assert slide["run2_68_debug_renderer_module_id"] == targeted_modules[slide["role"]]
            assert slide["run2_68_bug_fix_status"] == "pass_internal"
            assert slide["run2_68_renderer_bug_fixed"] in {
                "operating_stage_was_node_diagram",
                "proof_title_body_overlap_and_wireframe_workspace",
                "release_map_was_random_node_graph",
            }
        else:
            assert slide["run2_68_targeted_debug_role"] is False

    assert bad_trace["arm_id"] == "bad_run2_67_without_targeted_debug_repair"
    for slide in bad_trace["slides"]:
        if slide["role"] in targeted_modules:
            assert slide["run2_68_bug_fix_status"] == "fail_missing_run2_68_debug"

    assert_contains(
        result_md,
        [
            "Run 2.68 Targeted Renderer Debug Rerun",
            "S02 setup",
            "S04 proof",
            "S06 close",
            "text overlap",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_68_targeted_debug_rerun() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    required_terms = [
        "Run 2.68",
        "run2_68_targeted_debug_rerun_result.json",
        "ppt-run2-68-full-vulca",
        "ppt-run2-68-bad-without-targeted-debug",
        "run2-68-four-arm-contact-sheet.png",
        '"latestRunId": "2.77"',
    ]
    assert_contains(script, required_terms[:-1])
    assert_contains(viewer, required_terms)


def test_run2_69_generates_public_content_fill_rerun() -> None:
    result = load_json(PACK / "results" / "run2_69_public_content_fill_rerun_result.json")
    result_md = (PACK / "results" / "run2_69_public_content_fill_rerun_result.md").read_text(
        encoding="utf-8"
    )
    presentations = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
    )
    full_trace = load_json(presentations / "ppt-run2-69-full-vulca" / "trace_manifest.json")
    bad_trace = load_json(
        presentations / "ppt-run2-69-bad-without-public-content-fill" / "trace_manifest.json"
    )

    assert result["run_id"] == "2.69"
    assert result["status"] == "run2_69_public_content_fill_rerun_public_blocked"
    assert result["source_generated_run_id"] == "2.68"
    assert result["source_design_run_id"] == "2.66"
    assert result["rerun"]["best_internal_arm"] == "run2_69_full_public_content_fill"
    assert result["quality_delta"]["full_slides_with_public_content_slots"] == 6
    assert result["quality_delta"]["full_slides_without_visible_debug_terms"] == 6
    assert result["quality_delta"]["full_slides_with_visual_slot_labels"] == 6
    assert result["quality_delta"]["bad_control_slides_without_public_content_fill"] == 6

    assert full_trace["arm_id"] == "run2_69_full_public_content_fill"
    for slide in full_trace["slides"]:
        assert slide["run2_69_public_content_fill_status"] == "pass_internal"
        assert slide["run2_69_box_fill_status"] == "pass_internal"
        assert slide["run2_69_internal_debug_terms_visible"] == 0
        assert slide["run2_69_public_content_slot_count"] >= 5
        assert len(slide["run2_69_visual_content_slot_texts"]) >= 5
        assert slide["layout_metrics"]["text_overflow_risk_count"] == 0
        public_text = " ".join(slide["run2_69_visual_content_slot_texts"]).lower()
        assert "debug fix" not in public_text
        assert "targeted debug" not in public_text

    assert bad_trace["arm_id"] == "bad_run2_68_without_public_content_fill"
    for slide in bad_trace["slides"]:
        assert slide["run2_69_public_content_fill_status"] == "fail_missing_run2_69_public_content_fill"

    assert_contains(
        result_md,
        [
            "Run 2.69 Public Content Fill Rerun",
            "fill visual boxes",
            "remove debug fix copy",
            "public-facing slide copy",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_69_public_content_fill_rerun() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    required_terms = [
        "Run 2.69",
        "run2_69_public_content_fill_rerun_result.json",
        "ppt-run2-69-full-vulca",
        "ppt-run2-69-bad-without-public-content-fill",
        "run2-69-four-arm-contact-sheet.png",
        '"latestRunId": "2.77"',
    ]
    assert_contains(script, required_terms[:-1])
    assert_contains(viewer, required_terms)


def test_run2_70_generates_high_fidelity_mock_content_rerun() -> None:
    result = load_json(PACK / "results" / "run2_70_high_fidelity_mock_content_rerun_result.json")
    result_md = (PACK / "results" / "run2_70_high_fidelity_mock_content_rerun_result.md").read_text(
        encoding="utf-8"
    )
    presentations = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
    )
    full_trace = load_json(presentations / "ppt-run2-70-full-vulca" / "trace_manifest.json")
    bad_trace = load_json(
        presentations / "ppt-run2-70-bad-without-high-fidelity-mock-content" / "trace_manifest.json"
    )

    assert result["run_id"] == "2.70"
    assert result["status"] == "run2_70_high_fidelity_mock_content_rerun_public_blocked"
    assert result["source_generated_run_id"] == "2.69"
    assert result["source_design_run_id"] == "2.66"
    assert result["rerun"]["best_internal_arm"] == "run2_70_full_high_fidelity_mock_content"
    assert result["mock_content_scope"]["target_roles"] == ["contrast", "proof", "climax"]
    assert result["quality_delta"]["full_targeted_slides_with_high_fidelity_mock_content"] == 3
    assert result["quality_delta"]["full_targeted_slides_without_empty_wireframes"] == 3
    assert result["quality_delta"]["full_targeted_slides_with_realistic_mock_objects"] == 3
    assert result["quality_delta"]["bad_control_targeted_slides_without_high_fidelity_mock_content"] == 3

    expected_surface_by_role = {
        "contrast": "product_scene_mock",
        "proof": "inspectable_generated_slide_mock",
        "climax": "editable_presentation_surface_mock",
    }
    assert full_trace["arm_id"] == "run2_70_full_high_fidelity_mock_content"
    targeted_full_slides = [slide for slide in full_trace["slides"] if slide["role"] in expected_surface_by_role]
    assert len(targeted_full_slides) == 3
    for slide in targeted_full_slides:
        assert slide["run2_70_high_fidelity_mock_content_status"] == "pass_internal"
        assert slide["run2_70_empty_wireframe_replaced"] is True
        assert slide["run2_70_realistic_mock_object_count"] >= 4
        assert slide["run2_70_product_mock_surface_kind"] == expected_surface_by_role[slide["role"]]
        assert slide["layout_metrics"]["text_overflow_risk_count"] == 0

    assert bad_trace["arm_id"] == "bad_run2_69_without_high_fidelity_mock_content"
    targeted_bad_slides = [slide for slide in bad_trace["slides"] if slide["role"] in expected_surface_by_role]
    assert len(targeted_bad_slides) == 3
    for slide in targeted_bad_slides:
        assert (
            slide["run2_70_high_fidelity_mock_content_status"]
            == "fail_missing_run2_70_high_fidelity_mock_content"
        )

    assert_contains(
        result_md,
        [
            "Run 2.70 High-Fidelity Mock Content Rerun",
            "replace abstract wireframes",
            "product scene mock",
            "editable presentation surface mock",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_70_high_fidelity_mock_content_rerun() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    required_terms = [
        "Run 2.70",
        "run2_70_high_fidelity_mock_content_rerun_result.json",
        "ppt-run2-70-full-vulca",
        "ppt-run2-70-bad-without-high-fidelity-mock-content",
        "run2-70-four-arm-contact-sheet.png",
        '"latestRunId": "2.77"',
    ]
    assert_contains(script, required_terms[:-1])
    assert_contains(viewer, required_terms)


def test_run2_71_generates_component_semantics_content_binding_rerun() -> None:
    result = load_json(PACK / "results" / "run2_71_component_semantics_rerun_result.json")
    result_md = (PACK / "results" / "run2_71_component_semantics_rerun_result.md").read_text(
        encoding="utf-8"
    )
    presentations = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
    )
    full_trace = load_json(presentations / "ppt-run2-71-full-vulca" / "trace_manifest.json")
    bad_trace = load_json(
        presentations / "ppt-run2-71-bad-without-component-semantics" / "trace_manifest.json"
    )

    assert result["run_id"] == "2.71"
    assert result["status"] == "run2_71_component_semantics_rerun_public_blocked"
    assert result["source_generated_run_id"] == "2.70"
    assert result["source_design_run_id"] == "2.66"
    assert result["rerun"]["best_internal_arm"] == "run2_71_full_component_semantics"
    assert result["component_semantics_scope"]["target_roles"] == ["contrast", "proof", "climax"]
    assert result["quality_delta"]["full_targeted_slides_with_component_manifests"] == 3
    assert result["quality_delta"]["full_targeted_slides_with_text_component_bindings"] == 3
    assert result["quality_delta"]["full_targeted_slides_with_non_rectangular_primitives"] == 3
    assert result["quality_delta"]["full_targeted_slides_with_distinct_component_archetypes"] == 3
    assert result["quality_delta"]["bad_control_targeted_slides_without_component_semantics"] == 3

    expected_archetype_by_role = {
        "contrast": "before_after_product_scene",
        "proof": "inspectable_evidence_workspace",
        "climax": "editable_presentation_surface",
    }
    assert full_trace["arm_id"] == "run2_71_full_component_semantics"
    targeted_full_slides = [
        slide for slide in full_trace["slides"] if slide["role"] in expected_archetype_by_role
    ]
    assert len(targeted_full_slides) == 3
    archetypes = set()
    for slide in targeted_full_slides:
        archetypes.add(slide["run2_71_component_archetype"])
        assert slide["run2_71_component_semantics_status"] == "pass_internal"
        assert slide["run2_71_component_archetype"] == expected_archetype_by_role[slide["role"]]
        assert len(slide["run2_71_component_manifest"]) >= 5
        assert len(slide["run2_71_text_component_bindings"]) >= 5
        assert len(slide["run2_71_non_rectangular_primitives"]) >= 2
        assert slide["run2_71_weakness_repair_status"] == "pass_internal"
        assert slide["layout_metrics"]["text_overflow_risk_count"] == 0
    assert len(archetypes) == 3

    assert bad_trace["arm_id"] == "bad_run2_70_without_component_semantics"
    targeted_bad_slides = [
        slide for slide in bad_trace["slides"] if slide["role"] in expected_archetype_by_role
    ]
    assert len(targeted_bad_slides) == 3
    for slide in targeted_bad_slides:
        assert slide["run2_71_component_semantics_status"] == "fail_missing_run2_71_component_semantics"

    assert_contains(
        result_md,
        [
            "Run 2.71 Component Semantics Rerun",
            "component manifest",
            "text-to-component binding",
            "non-rectangular primitives",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_71_component_semantics_rerun() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    required_terms = [
        "Run 2.71",
        "run2_71_component_semantics_rerun_result.json",
        "ppt-run2-71-full-vulca",
        "ppt-run2-71-bad-without-component-semantics",
        "run2-71-four-arm-contact-sheet.png",
        '"latestRunId": "2.77"',
    ]
    assert_contains(script, required_terms[:-1])
    assert_contains(viewer, required_terms)


def _bbox_contains(outer: dict, inner: dict) -> bool:
    return (
        inner["x"] >= outer["x"]
        and inner["y"] >= outer["y"]
        and inner["x"] + inner["w"] <= outer["x"] + outer["w"]
        and inner["y"] + inner["h"] <= outer["y"] + outer["h"]
    )


def test_run2_72_generates_shape_bound_text_rerun() -> None:
    result = load_json(PACK / "results" / "run2_72_shape_bound_text_rerun_result.json")
    result_md = (PACK / "results" / "run2_72_shape_bound_text_rerun_result.md").read_text(
        encoding="utf-8"
    )
    presentations = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
    )
    full_trace = load_json(presentations / "ppt-run2-72-full-vulca" / "trace_manifest.json")
    bad_trace = load_json(
        presentations / "ppt-run2-72-bad-without-shape-bound-text" / "trace_manifest.json"
    )

    assert result["run_id"] == "2.72"
    assert result["status"] == "run2_72_shape_bound_text_rerun_public_blocked"
    assert result["source_generated_run_id"] == "2.71"
    assert result["source_design_run_id"] == "2.66"
    assert result["rerun"]["best_internal_arm"] == "run2_72_full_shape_bound_text"
    assert result["shape_bound_text_scope"]["target_roles"] == ["contrast", "proof", "climax"]
    assert result["quality_delta"]["target_layer"] == "run2_72_shape_bound_text_geometry"
    assert result["quality_delta"]["full_targeted_slides_with_shape_bound_text"] == 3
    assert result["quality_delta"]["full_targeted_slides_with_binding_geometry_records"] == 3
    assert result["quality_delta"]["full_targeted_slides_with_text_inside_component_bounds"] == 3
    assert result["quality_delta"]["full_targeted_slides_with_zero_unbound_component_labels"] == 3
    assert result["quality_delta"]["bad_control_targeted_slides_without_shape_bound_text"] == 3

    expected_archetype_by_role = {
        "contrast": "before_after_product_scene",
        "proof": "inspectable_evidence_workspace",
        "climax": "editable_presentation_surface",
    }
    assert full_trace["arm_id"] == "run2_72_full_shape_bound_text"
    targeted_full_slides = [
        slide for slide in full_trace["slides"] if slide["role"] in expected_archetype_by_role
    ]
    assert len(targeted_full_slides) == 3
    for slide in targeted_full_slides:
        assert slide["run2_71_component_semantics_status"] == "pass_internal"
        assert slide["run2_72_shape_bound_text_status"] == "pass_internal"
        assert slide["run2_72_component_archetype"] == expected_archetype_by_role[slide["role"]]
        assert slide["run2_72_bound_text_component_count"] >= 5
        assert slide["run2_72_unbound_text_component_count"] == 0
        assert slide["run2_72_text_inside_component_bounds"] is True
        assert len(slide["run2_72_binding_geometry_records"]) >= 5
        for record in slide["run2_72_binding_geometry_records"]:
            assert record["text_fits_component"] is True
            assert _bbox_contains(record["component_bbox"], record["text_bbox"])
            assert record["text_bbox"]["w"] > 0
            assert record["text_bbox"]["h"] > 0
        assert slide["layout_metrics"]["text_overflow_risk_count"] == 0

    assert bad_trace["arm_id"] == "bad_run2_72_without_shape_bound_text"
    targeted_bad_slides = [
        slide for slide in bad_trace["slides"] if slide["role"] in expected_archetype_by_role
    ]
    assert len(targeted_bad_slides) == 3
    for slide in targeted_bad_slides:
        assert slide["run2_72_shape_bound_text_status"] == "fail_missing_run2_72_shape_bound_text"
        assert slide["run2_72_binding_geometry_records"] == []

    assert_contains(
        result_md,
        [
            "Run 2.72 Shape-Bound Text Rerun",
            "component bbox",
            "text bbox",
            "inside the component bounds",
            "Do not advance to Run 3.0",
        ],
    )


def test_ppt_run_html_viewer_mentions_run2_72_shape_bound_text_rerun() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    required_terms = [
        "Run 2.72",
        "run2_72_shape_bound_text_rerun_result.json",
        "ppt-run2-72-full-vulca",
        "ppt-run2-72-bad-without-shape-bound-text",
        "run2-72-four-arm-contact-sheet.png",
        '"latestRunId": "2.77"',
    ]
    assert_contains(script, required_terms[:-1])
    assert_contains(viewer, required_terms)


def test_ppt_run_html_viewer_mentions_run2_50_readability_density_renderer_rerun() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")
    viewer = (
        ROOT
        / "outputs"
        / "019e7d9c-532a-70b3-8892-fa3ae42baef2"
        / "presentations"
        / "ppt-run-viewer.html"
    ).read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "Run 2.50",
            "ppt-run2-50-prompt-only",
            "ppt-run2-50-run1-5-skill",
            "ppt-run2-50-full-vulca",
            "ppt-run2-50-bad-missing-run2-49-repair-pack",
            "run2_50_readability_density_renderer_rerun_result.json",
            "Generated four-arm rerun over Run 2.49",
            "run250ResultStatus",
            "readability_content_density_and_editorial_renderer_binding",
        ],
    )
    assert_contains(
        viewer,
        [
            '"latestRunId": "2.77"',
            "Run 2.50",
            "run2-50-four-arm-contact-sheet.png",
            "ppt-run2-50-prompt-only",
            "ppt-run2-50-run1-5-skill",
            "ppt-run2-50-full-vulca",
            "ppt-run2-50-bad-missing-run2-49-repair-pack",
            "Generated four-arm rerun over Run 2.49",
            "Run 2.50 readability density renderer rerun",
            "run2_50_readability_density_renderer_rerun_result.json",
        ],
    )


def test_ppt_layout_quality_checker_flags_geometry_failures(tmp_path: Path) -> None:
    layout_dir = tmp_path / "layout"
    layout_dir.mkdir()
    ok_path = layout_dir / "slide-01.layout.json"
    bad_path = layout_dir / "slide-02.layout.json"
    out_path = tmp_path / "layout-report.txt"

    ok_path.write_text(
        json.dumps(
            {
                "slide": {"frame": {"width": 1280, "height": 720}},
                "elements": [
                    {
                        "id": "headline",
                        "text": "One clear headline",
                        "bbox": [80, 80, 520, 64],
                        "resolvedFontSize": 28,
                        "textLayout": {"lineCount": 1},
                    },
                    {"id": "proof-object", "bbox": [760, 120, 320, 240]},
                ],
            }
        ),
        encoding="utf-8",
    )
    bad_path.write_text(
        json.dumps(
            {
                "slide": {"frame": {"width": 1280, "height": 720}},
                "elements": [
                    {
                        "id": "escaped",
                        "text": "Escaped box",
                        "bbox": [-10, 20, 200, 40],
                        "resolvedFontSize": 18,
                        "textLayout": {"lineCount": 1},
                    },
                    {
                        "id": "overlap-a",
                        "text": "Overlap A",
                        "bbox": [50, 50, 200, 40],
                        "resolvedFontSize": 18,
                        "textLayout": {"lineCount": 1},
                    },
                    {
                        "id": "overlap-b",
                        "text": "Overlap B",
                        "bbox": [60, 55, 190, 38],
                        "resolvedFontSize": 18,
                        "textLayout": {"lineCount": 1},
                    },
                    {
                        "id": "microtype",
                        "text": "Fine print",
                        "bbox": [300, 300, 100, 5],
                        "resolvedFontSize": 6,
                        "textLayout": {"lineCount": 1},
                    },
                ],
            }
        ),
        encoding="utf-8",
    )

    issue_codes = {issue.code for issue in check_layout(bad_path)}
    assert {"out-of-bounds", "text-overlap", "microtype", "tight-text"} <= issue_codes

    report = write_report(layout_dir, out_path)

    assert report == {"layout_files": 2, "layout_errors": 2, "layout_warnings": 2}
    assert "text-overlap" in out_path.read_text(encoding="utf-8")


def test_ppt_contact_sheet_supports_custom_labels(tmp_path: Path) -> None:
    first = tmp_path / "first.png"
    second = tmp_path / "second.png"
    out = tmp_path / "sheet.png"
    Image.new("RGB", (160, 90), "#ffffff").save(first)
    Image.new("RGB", (160, 90), "#111318").save(second)

    build_contact_sheet([first, second], out, "Two arm sheet", cols=2, labels=["Prompt-only", "Full arm"])

    assert out.exists()
    sheet = Image.open(out)
    assert sheet.size[0] > 0
    assert sheet.size[1] > 0


def test_run2_bad_aesthetic_memory_has_structured_replacement() -> None:
    replacement = load_json(PACK / "generation_briefs" / "bad_aesthetic_memory_replacement.json")

    assert replacement["replacement_for"] == "aesthetic_memory.json"
    assert replacement["status"] == "negative-control-only"
    assert EXPECTED_MOVES <= {move["aesthetic_move"] for move in replacement["moves"]}
    for move in replacement["moves"]:
        assert EXPECTED_BAD_MEMORY_FIELDS <= set(move)
        assert move["rhythm_role"] in EXPECTED_RHYTHM_ROLES
        assert move["negative_rules"]
        assert {"max_claims", "max_panels", "max_words"} <= set(move["density_budget"])
    assert_contains(json.dumps(replacement), ["dashboard", "dense", "small labels", "no visual climax"])


def test_run2_results_reviewed_and_public_blocked() -> None:
    comparison = (PACK / "results" / "comparison_report.md").read_text(encoding="utf-8")
    delivery = (PACK / "results" / "delivery_gate.md").read_text(encoding="utf-8")
    trace_contract = load_json(PACK / "results" / "trace_manifest_contract.json")

    assert_contains(comparison, ["Status", "run2_18-thickness-pack-public-blocked"])
    assert_contains(comparison, ["Prior delivery status", "motion-renderer-proof-public-blocked"])
    assert_contains(
        comparison,
        [
            "Run 2.17 Motion Proof",
            "run2_17_motion_renderer_proof_result.json",
            "separate HTML motion renderer",
            "cover_attention_reset",
            "before_after_reveal",
            "climax_scale_emphasis",
        ],
    )
    assert_contains(
        comparison,
        [
            "Run 2.17",
            "run2_17_motion_delivery_audit.json",
            "HTML viewer is static",
            "Keynote readout",
        ],
    )
    assert_contains(comparison, ["Run 2.16", "latest reviewed generated PPT result"])
    assert_contains(
        comparison,
        [
            "Run 2.13",
            "run2_13_full_skill",
            "run2_12_thick_multimodal_evidence.json",
            "bad_thick_data_memory",
        ],
    )
    assert_contains(
        comparison,
        [
            "Run 2.12",
            "run2_12_thick_multimodal_evidence.json",
            "run2_12_design_memory_seed.json",
            "run2_12_workflow_gate_seed.json",
        ],
    )
    assert_contains(comparison, ["prompt_only", "run1_5_skill", "run2_skill", "bad_aesthetic_memory"])
    assert_contains(
        comparison,
        ["Run 2.3", "run2_3_full_skill", "visual_component_execution", "native visual components"],
    )
    assert_contains(
        comparison,
        ["Run 2.4", "motion grammar", "video_demo_beat_map", "presentation_sequence_components"],
    )
    assert_contains(
        comparison,
        ["Run 2.6", "run2_6_full_skill", "data_workflow_policy_execution", "workflow_decision_policy"],
    )
    assert_contains(
        comparison,
        [
            "Run 2.10",
            "run2_10_full_skill",
            "run2_10_visual_system_memory.json",
            "actual code module ids",
        ],
    )
    assert_contains(
        comparison,
        [
            "Run 2.8",
            "run2_8_full_skill",
            "run2_8_executable_design_memory.json",
            "actual native module calls",
        ],
    )
    assert_contains(
        comparison,
        ["Run 2.2", "run2_2_full_skill", "multimodal_learning", "visual_learning_target_execution"],
    )
    assert_contains(comparison, ["Run 2.1", "run2_1_full_skill", "product learning", "not public-release claims"])
    assert "0.00" not in comparison
    assert_contains(delivery, ["public publishing", "blocked", "native render", "human approval", "trace manifest"])
    assert_contains(delivery, ["Run 2.3", "pass for local Run 2.3 arms", "native visual components"])
    assert_contains(delivery, ["Run 2.6", "run2-6-four-arm-contact-sheet", "workflow_decision_policy"])
    assert_contains(
        delivery,
        ["Run 2.8", "run2-8-four-arm-contact-sheet", "actual native module calls"],
    )
    assert_contains(
        delivery,
        ["Run 2.10", "run2-10-four-arm-contact-sheet", "actual native visual-system module calls"],
    )
    assert_contains(
        delivery,
        ["Run 2.13", "run2-13-four-arm-contact-sheet", "Run 2.13 full-arm trace includes Run 2.12 evidence"],
    )
    assert_contains(delivery, ["Run 2.2", "run2-2-four-arm-contact-sheet", "public-video-grade visual proof"])
    assert trace_contract["required_output_name"] == "trace_manifest.json"
    assert "aesthetic_move_ids" in trace_contract["per_slide_required_fields"]
    assert "visual_component_ids" in trace_contract["per_slide_required_fields"]
    assert "video_beat_ids" in trace_contract["per_slide_required_fields"]
    assert "motion_target_ids" in trace_contract["per_slide_required_fields"]
    assert "sequence_component_ids" in trace_contract["per_slide_required_fields"]
    assert "runtime_isolation" in trace_contract["arm_required_fields"]
    assert "native_ppt_checks" in trace_contract["per_slide_required_fields"]
    assert "layout_geometry_checks" in trace_contract["per_slide_required_fields"]
    assert "commercial_usecase_id" in trace_contract["per_slide_required_fields"]
    assert "aesthetic_benchmark_ids" in trace_contract["per_slide_required_fields"]
    assert "theme_policy_id" in trace_contract["per_slide_required_fields"]
    assert EXPECTED_RUN2_13_TRACE_FIELDS <= set(trace_contract["per_slide_required_fields"])
    assert EXPECTED_RUN2_14_TRACE_FIELDS <= set(trace_contract["per_slide_required_fields"])
    assert EXPECTED_RUN2_15_TRACE_FIELDS <= set(trace_contract["per_slide_required_fields"])
    assert EXPECTED_RUN2_16_TRACE_FIELDS <= set(trace_contract["per_slide_required_fields"])
    assert trace_contract["native_ppt_thresholds"]["full_slide_rasterized_allowed"] is False


def test_run2_3_records_native_component_rerun_result() -> None:
    result = (PACK / "results" / "run2_3_rerun_result.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_3_rerun_result.json")

    assert result_json["status"] == "rerun_completed_public_blocked"
    assert result_json["public_ready"] is False
    assert result_json["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert result_json["rerun"]["status"] == "completed"
    assert result_json["rerun"]["best_internal_arm"] == "run2_3_full_skill"
    assert result_json["rerun"]["generated_outputs_committed"] is False
    assert (
        result_json["rerun"]["best_internal_arm_verdict"]
        == "native_visual_components_visible_but_not_public_release_ready"
    )
    assert result_json["next_required_action"] == "turn_component_execution_into_public_grade_visual_system"
    assert_contains(
        json.dumps(result_json["native_component_learning"]),
        ["native_component_visible", "layout_errors_zero", "still_internal_demo_grade"],
    )
    assert_contains(
        result,
        [
            "Run 2.3",
            "visual_target_components.json",
            "native visual components",
            "before/after thumbnail",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_run2_4_records_motion_sequence_rerun_result() -> None:
    result = (PACK / "results" / "run2_4_rerun_result.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_4_rerun_result.json")

    assert result_json["status"] == "rerun_completed_public_blocked"
    assert result_json["public_ready"] is False
    assert result_json["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert result_json["rerun"]["status"] == "completed"
    assert result_json["rerun"]["best_internal_arm"] == "run2_4_full_skill"
    assert result_json["rerun"]["generated_outputs_committed"] is False
    assert result_json["rerun"]["best_internal_arm_verdict"] == "motion_sequence_visible_but_not_public_video_grade"
    assert (
        result_json["next_required_action"]
        == "thicken_public_grade_visual_system_and_motion_rendering_then_rerun_same_five_layers"
    )
    assert_contains(
        json.dumps(result_json["motion_sequence_learning"]),
        [
            "native_sequence_visible",
            "video_beat_ids_motion_target_ids_sequence_component_ids_ordered_reveal_steps_present",
            "still_schematic_internal_demo_grade",
        ],
    )
    assert_contains(
        result,
        [
            "Run 2.4",
            "motion",
            "sequence",
            "run2_4_full_skill",
            "public-video-grade",
            "Do not advance to Run 3.0",
        ],
    )


def test_run2_5_records_production_design_rerun_result() -> None:
    result = (PACK / "results" / "run2_5_rerun_result.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_5_rerun_result.json")

    assert result_json["status"] == "rerun_completed_public_blocked"
    assert result_json["public_ready"] is False
    assert result_json["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert result_json["rerun"]["status"] == "completed"
    assert result_json["rerun"]["best_internal_arm"] == "run2_5_full_skill"
    assert result_json["rerun"]["generated_outputs_committed"] is False
    assert (
        result_json["rerun"]["best_internal_arm_verdict"]
        == "production_modules_visible_but_not_public_release_ready"
    )
    assert result_json["next_required_action"] == "native_render_and_human_review_then_deepen_visual_polish_if_needed"
    assert_contains(
        json.dumps(result_json["production_design_learning"]),
        [
            "production_reference_ids",
            "aesthetic_memory_v2_ids",
            "visual_production_module_ids",
            "fallback_policy_recorded",
            "still_internal_demo_grade",
        ],
    )
    assert_contains(
        result,
        [
            "Run 2.5",
            "production_reference_decompositions.json",
            "aesthetic_memory_v2.json",
            "visual_production_modules.json",
            "run2_5_full_skill",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_run2_6_records_data_workflow_rerun_result() -> None:
    result = (PACK / "results" / "run2_6_rerun_result.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_6_rerun_result.json")

    assert result_json["status"] == "rerun_completed_public_blocked"
    assert result_json["public_ready"] is False
    assert result_json["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert result_json["rerun"]["status"] == "completed"
    assert result_json["rerun"]["best_internal_arm"] == "run2_6_full_skill"
    assert result_json["rerun"]["generated_outputs_committed"] is False
    assert (
        result_json["rerun"]["best_internal_arm_verdict"]
        == "data_workflow_policy_visible_but_not_public_release_ready"
    )
    assert result_json["next_required_action"] == "deepen_same_five_layers_with_visual_quality_not_run3"
    assert_contains(
        json.dumps(result_json["data_workflow_learning"]),
        [
            "commercial_usecase_id",
            "aesthetic_benchmark_ids",
            "theme_policy_id",
            "typography_system_id",
            "spacing_token_set_id",
            "source_brand_sanitization",
        ],
    )
    assert_contains(
        result,
        [
            "Run 2.6",
            "commercial_usecase_bank.json",
            "aesthetic_benchmark_bank.json",
            "workflow_decision_policy.json",
            "run2_6_full_skill",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_run2_6r_records_visual_repair_rerun_result() -> None:
    result = (PACK / "results" / "run2_6r_visual_repair_result.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_6r_visual_repair_result.json")

    assert result_json["status"] == "rerun_completed_public_blocked"
    assert result_json["public_ready"] is False
    assert result_json["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert result_json["rerun"]["status"] == "completed"
    assert result_json["rerun"]["best_internal_arm"] == "run2_6r_visual_repair_full_skill"
    assert result_json["rerun"]["generated_outputs_committed"] is False
    assert result_json["rerun"]["best_internal_arm_verdict"] in {
        "visual_repair_visible_but_not_public_release_ready",
        "visual_repair_insufficient_same_stage_repeat_required",
    }
    assert result_json["next_required_action"] in {
        "native_render_and_human_review_then_continue_same_five_layers_if_needed",
        "repeat_same_stage_visual_repair_before_native_review",
    }
    assert_contains(
        json.dumps(result_json["visual_repair_learning"]),
        [
            "visual_repair_policy_ids",
            "visual_delta_from_run2_5",
            "visual_repair_validation_probe",
            "typography",
            "spacing",
            "climax",
            "theme",
        ],
    )
    assert result_json["qa_summary"]["arm_isolation_guard"] == "passed"
    assert result_json["qa_summary"]["run2_6r_full_media_entries"] == 0
    assert result_json["qa_summary"]["run2_6r_full_picture_shapes"] == 0
    assert result_json["qa_summary"]["run2_6r_full_native_shape_minimum_passed"] is True
    assert_contains(
        result,
        [
            "Run 2.6R",
            "visual_repair_policy.json",
            "run2_6r_visual_repair_full_skill",
            "full-skill-series",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_run2_8_records_executable_design_memory_rerun_result() -> None:
    result = (PACK / "results" / "run2_8_executable_design_memory_result.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_8_executable_design_memory_result.json")
    gate_json = load_json(PACK / "results" / "run2_8_visual_qa_gate.json")

    assert result_json["status"] == "rerun_completed_public_blocked"
    assert result_json["public_ready"] is False
    assert result_json["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert result_json["rerun"]["status"] == "completed"
    assert result_json["rerun"]["best_internal_arm"] == "run2_8_full_skill"
    assert result_json["rerun"]["generated_outputs_committed"] is False
    assert result_json["qa_summary"]["arm_isolation_guard"] == "passed"
    assert result_json["qa_summary"]["trace_truthfulness_guard"] == "passed"
    assert result_json["qa_summary"]["run2_8_full_media_entries"] == 0
    assert result_json["qa_summary"]["run2_8_full_picture_shapes"] == 0
    assert result_json["control_boundary"]["bad_memory_schema"].startswith("commercial_case_plus_decomposition_only")
    assert "run2-8-four-arm-contact-sheet.png" in result_json["rerun"]["combined_contact_sheet"]
    assert "run2-full-skill-series-horizontal.png" in result_json["rerun"]["full_skill_series_sheet"]
    assert result_json["rerun"]["html_viewer"].endswith("/ppt-run-viewer.html")
    assert gate_json["status"] == "run2_8_visual_qa_gate_public_blocked"
    assert gate_json["observed"]["full_arm_trace_origin"] == "actual_native_module_calls"
    assert gate_json["observed"]["full_arm_required_code_bindings_missing"] == 0
    assert gate_json["observed"]["full_arm_budget_failures"] == 0
    assert gate_json["decision"] == "public_blocked"
    assert_contains(
        json.dumps(result_json["five_layer_learning"]),
        [
            "run2_8_tutorial_decomposition.json",
            "run2_8_executable_design_memory.json",
            "run2_8_workflow_gate_matrix.json",
            "actual native module calls",
        ],
    )
    assert_contains(
        result,
        [
            "Run 2.8",
            "executable design memory",
            "HTML viewer",
            "workflow gate matrix",
            "actual native module calls",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_run2_9_records_visual_primitive_rerun_result() -> None:
    result = (PACK / "results" / "run2_9_visual_primitive_rerun_result.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_9_visual_primitive_rerun_result.json")

    assert result_json["status"] == "rerun_completed_public_blocked"
    assert result_json["public_ready"] is False
    assert result_json["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert result_json["rerun"]["status"] == "completed"
    assert result_json["rerun"]["best_internal_arm"] == "run2_9_full_skill"
    assert result_json["rerun"]["generated_outputs_committed"] is False
    assert "run2-9-four-arm-contact-sheet.png" in result_json["rerun"]["combined_contact_sheet"]
    assert "run2-full-skill-series-horizontal.png" in result_json["rerun"]["full_skill_series_sheet"]
    assert result_json["rerun"]["html_viewer"].endswith("/ppt-run-viewer.html")
    assert result_json["qa_summary"]["trace_truthfulness_guard"].startswith("passed")
    assert result_json["qa_summary"]["case_pack_validator"] == "passed with --profile run2"
    assert result_json["control_boundary"]["bad_visual_primitive_memory"].startswith("commercial_case_plus_visual_primitive_repair_only")
    assert_contains(
        json.dumps(result_json["actual_full_arm_modules_by_role"]),
        [
            "drawRun29TypographicField",
            "drawRun29EditorialSpread",
            "drawRun29LayeredProductSurface",
            "drawRun29MotionStoryboard",
            "drawRun29ClimaxStage",
        ],
    )
    assert_contains(
        json.dumps(result_json["five_layer_learning"]),
        [
            "run2_9_visual_primitive_repair.json",
            "run2_9_executable_visual_modules.json",
            "run2_9_visual_gate_matrix.json",
            "actual drawRun29 visual module calls",
        ],
    )
    assert_contains(
        result,
        [
            "Run 2.9",
            "visual primitive",
            "HTML viewer",
            "drawRun29ClimaxStage",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_run2_10_records_visual_system_rerun_result() -> None:
    result = (PACK / "results" / "run2_10_visual_system_rerun_result.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_10_visual_system_rerun_result.json")

    assert result_json["status"] == "rerun_completed_public_blocked"
    assert result_json["public_ready"] is False
    assert result_json["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert result_json["rerun"]["status"] == "completed"
    assert result_json["rerun"]["best_internal_arm"] == "run2_10_full_skill"
    assert result_json["rerun"]["generated_outputs_committed"] is False
    assert "run2-10-four-arm-contact-sheet.png" in result_json["rerun"]["combined_contact_sheet"]
    assert "run2-full-skill-series-horizontal.png" in result_json["rerun"]["full_skill_series_sheet"]
    assert result_json["rerun"]["html_viewer"].endswith("/ppt-run-viewer.html")
    assert result_json["qa_summary"]["trace_truthfulness_guard"].startswith("passed")
    assert result_json["qa_summary"]["case_pack_validator"] == "passed with --profile run2"
    assert result_json["qa_summary"]["browser_viewer_check"].startswith("passed")
    assert_contains(
        json.dumps(result_json["actual_full_arm_modules_by_role"]),
        [
            "drawRun210FullBleedVisualField",
            "drawRun210ProductTheater",
            "drawRun210KineticSequence",
            "drawRun210EditorialTypeSystem",
            "drawRun210NonRectangularProofPath",
            "drawRun210CinematicClimax",
        ],
    )
    assert_contains(
        json.dumps(result_json["visual_learning"]),
        [
            "structural asymmetry",
            "not merely recolored",
            "public blocked",
        ],
    )
    assert_contains(
        result,
        [
            "Run 2.10",
            "visual system",
            "HTML viewer",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_run2_13_records_thick_data_rerun_result() -> None:
    result = (PACK / "results" / "run2_13_thick_data_rerun_result.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_13_thick_data_rerun_result.json")

    assert result_json["status"] == "rerun_completed_public_blocked"
    assert result_json["public_ready"] is False
    assert result_json["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert result_json["rerun"]["status"] == "completed"
    assert result_json["rerun"]["best_internal_arm"] == "run2_13_full_skill"
    assert result_json["rerun"]["generated_outputs_committed"] is False
    assert "run2-13-four-arm-contact-sheet.png" in result_json["rerun"]["combined_contact_sheet"]
    assert "run2-full-skill-series-horizontal.png" in result_json["rerun"]["full_skill_series_sheet"]
    assert result_json["rerun"]["html_viewer"].endswith("/ppt-run-viewer.html")
    assert result_json["qa_summary"]["trace_truthfulness_guard"].startswith("passed")
    assert result_json["qa_summary"]["case_pack_validator"] == "passed with --profile run2"
    assert result_json["control_boundary"]["bad_thick_data_memory"].startswith(
        "commercial_case_plus_thick_evidence_only"
    )
    assert_contains(
        json.dumps(result_json["actual_full_arm_modules_by_role"]),
        [
            "drawRun213LaunchArcRoute",
            "drawRun213TypeWhitespaceSystem",
            "drawRun213ProductDemoSequence",
            "drawRun213MetricClimaxObject",
            "drawRun213WorkflowGateRail",
        ],
    )
    assert_contains(
        json.dumps(result_json["thick_data_learning"]),
        [
            "run2_12_thick_multimodal_evidence.json",
            "run2_12_design_memory_seed.json",
            "run2_12_workflow_gate_seed.json",
            "evidence_only_negative_control",
        ],
    )
    assert_contains(
        result,
        [
            "Run 2.13",
            "thick data",
            "run2_12_thick_multimodal_evidence.json",
            "HTML viewer",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_run2_14_records_aesthetic_trace_rerun_result() -> None:
    result = (PACK / "results" / "run2_14_aesthetic_trace_rerun_result.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_14_aesthetic_trace_rerun_result.json")

    assert result_json["status"] == "rerun_completed_public_blocked"
    assert result_json["public_ready"] is False
    assert result_json["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert result_json["rerun"]["status"] == "completed"
    assert result_json["rerun"]["best_internal_arm"] == "run2_14_full_skill"
    assert result_json["rerun"]["generated_outputs_committed"] is False
    assert "run2-14-four-arm-contact-sheet.png" in result_json["rerun"]["combined_contact_sheet"]
    assert "run2-full-skill-series-horizontal.png" in result_json["rerun"]["full_skill_series_sheet"]
    assert result_json["rerun"]["html_viewer"].endswith("/ppt-run-viewer.html")
    assert result_json["qa_summary"]["trace_truthfulness_guard"].startswith("passed")
    assert result_json["qa_summary"]["case_pack_validator"] == "passed with --profile run2"
    assert result_json["control_boundary"]["bad_visible_workflow_memory"].startswith(
        "run2_12_data_workflow_without_run2_10_aesthetic_shell"
    )
    assert_contains(
        json.dumps(result_json["actual_full_arm_modules_by_role"]),
        [
            "drawRun214PresentationShell",
            "drawRun214CinematicProductTheater",
            "drawRun214HiddenWorkflowRoute",
            "drawRun214MetricRevealStage",
            "drawRun214QuietReleaseHandoff",
        ],
    )
    assert_contains(
        json.dumps(result_json["aesthetic_trace_learning"]),
        [
            "2.10 visual system aesthetic shell",
            "2.13 thick data trace core",
            "manifest_only_trace_public_surface",
            "workflow_hidden_from_slide_surface",
        ],
    )
    assert_contains(
        result,
        [
            "Run 2.14",
            "2.10 visual system aesthetic",
            "2.13 thick data trace",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_run2_15_records_layout_selector_data_workflow_result() -> None:
    result = (PACK / "results" / "run2_15_layout_selector_result.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_15_layout_selector_result.json")
    readme = (PACK / "results" / "README.md").read_text(encoding="utf-8")
    comparison = (PACK / "results" / "comparison_report.md").read_text(encoding="utf-8")
    delivery = (PACK / "results" / "delivery_gate.md").read_text(encoding="utf-8")

    assert result_json["status"] == "selector_data_workflow_ready_public_blocked"
    assert result_json["public_ready"] is False
    assert result_json["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert result_json["run_id"] == "2.15"
    assert result_json["creates_new_ppt_outputs"] is False
    assert result_json["next_generation_gate"] == "run2_15_selector_gate_matrix_must_drive_next_four_arm_rerun"
    assert_contains(
        json.dumps(result_json["selector_artifacts"]),
        [
            "run2_15_layout_selector_sources.json",
            "run2_15_layout_module_memory.json",
            "run2_15_layout_selector_gate_matrix.json",
        ],
    )
    assert_contains(
        json.dumps(result_json["learning"]),
        [
            "layout module selector",
            "text resilience",
            "trace hidden from public slide surface",
            "product theater realism",
        ],
    )
    assert_contains(
        result,
        [
            "Run 2.15",
            "layout module selector",
            "data/workflow-only",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )
    assert "ppt-run2-15-" not in json.dumps(result_json)
    assert_contains(
        readme,
        [
            "Run 2.15 intentionally adds no new `ppt-run2-15-*` output folders",
            "Run 2.16 is the selector-driven generated rerun that consumes the Run 2.15 artifacts",
        ],
    )
    assert_contains(
        comparison,
        [
            "It does not generate a new deck",
            "latest reviewed generated PPT result",
        ],
    )
    assert_contains(
        delivery,
        [
            "Run 2.15 is data/workflow-only and does not create PPTX artifacts",
            "Run 2.16 is the generated rerun that consumes it",
        ],
    )


def test_run2_16_records_selector_driven_rerun_result() -> None:
    result = (PACK / "results" / "run2_16_selector_rerun_result.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_16_selector_rerun_result.json")

    assert result_json["status"] == "rerun_completed_public_blocked"
    assert result_json["public_ready"] is False
    assert result_json["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert result_json["rerun"]["status"] == "completed"
    assert result_json["rerun"]["best_internal_arm"] == "run2_16_full_skill"
    assert result_json["rerun"]["generated_outputs_committed"] is False
    assert "run2-16-four-arm-contact-sheet.png" in result_json["rerun"]["combined_contact_sheet"]
    assert "run2-full-skill-series-horizontal.png" in result_json["rerun"]["full_skill_series_sheet"]
    assert result_json["rerun"]["html_viewer"].endswith("/ppt-run-viewer.html")
    assert result_json["qa_summary"]["trace_truthfulness_guard"].startswith("passed")
    assert result_json["qa_summary"]["case_pack_validator"] == "passed with --profile run2"
    assert result_json["control_boundary"]["bad_selector_memory"].startswith(
        "selector_sources_without_module_memory_or_gate_matrix"
    )
    assert_contains(
        json.dumps(result_json["actual_full_arm_modules_by_role"]),
        [
            "drawRun216EditorialCoverField",
            "drawRun216ProductTheaterStage",
            "drawRun216BeforeAfterRoute",
            "drawRun216MetricRevealStage",
            "drawRun216QuietReleaseHandoff",
            "drawRun216DenseEvidenceCompression",
        ],
    )
    assert_contains(
        json.dumps(result_json["selector_learning"]),
        [
            "run2_15_layout_selector_sources.json",
            "run2_15_layout_module_memory.json",
            "run2_15_layout_selector_gate_matrix.json",
            "selector_gate_matrix_executed_before_native_ppt_generation",
            "bad_selector_memory",
        ],
    )
    assert_contains(
        result,
        [
            "Run 2.16",
            "selector-driven",
            "run2_15_layout_selector_gate_matrix.json",
            "HTML viewer",
            "public blocked",
            "Do not advance to Run 3.0",
        ],
    )


def test_run2_17_records_motion_delivery_audit() -> None:
    result = (PACK / "results" / "run2_17_motion_delivery_audit.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_17_motion_delivery_audit.json")

    assert result_json["status"] == "motion_delivery_audit_public_blocked"
    assert result_json["public_ready"] is False
    assert result_json["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert result_json["source_generated_run_id"] == "2.16"
    assert result_json["delivery_truth"]["html_viewer_mode"] == "static_slide_preview_only"
    assert result_json["delivery_truth"]["keynote_expected_behavior"] == "editable_static_slides_no_native_animation"
    assert result_json["delivery_truth"]["native_ppt_animation_status"] == "absent_in_current_pptx"
    assert result_json["delivery_truth"]["motion_xml_scan_scope"] == "tag_presence_only_not_playback_verification"
    assert result_json["motion_renderer_gap"]["next_run_recommendation"] == "run2_17_motion_renderer_proof"
    assert result_json["motion_renderer_gap"]["keep_static_ppt_as"] == "editable_product_output"
    assert result_json["motion_renderer_gap"]["public_video_path"] == "separate_html_or_video_motion_renderer_until_pptx_animation_is_verified"

    arm_audits = result_json["arm_audits"]
    assert {arm["arm_id"] for arm in arm_audits} == {
        "prompt_only",
        "run1_5_skill",
        "run2_16_full_skill",
        "bad_selector_memory",
    }
    for arm in arm_audits:
        assert arm["pptx_path"].endswith(".pptx")
        assert "ppt-run2-16" in arm["pptx_path"]
        assert arm["slide_count"] == 6
        assert arm["media_entry_count"] == 0
        assert arm["motion"]["has_motion_xml"] is False
        assert arm["motion"]["scan_scope"] == "motion_xml_tag_presence_only_not_playback_verification"
        assert arm["motion"]["slides_with_motion_xml"] == []
        assert arm["keynote_readout"] == "static_editable_slides_only"

    assert_contains(
        result,
        [
            "Run 2.17",
            "HTML viewer is static",
            "Keynote",
            "no transition/timing/animation",
            "Run 2.16",
            "Do not advance to Run 3.0",
        ],
    )


def test_run2_17_records_motion_renderer_proof_contract() -> None:
    result = (PACK / "results" / "run2_17_motion_renderer_proof_result.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_17_motion_renderer_proof_result.json")

    assert result_json["status"] == "motion_renderer_proof_created_public_blocked"
    assert result_json["public_ready"] is False
    assert result_json["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert result_json["source_generated_run_id"] == "2.16"
    assert result_json["source_audit"] == "run2_17_motion_delivery_audit.json"
    assert result_json["delivery_boundary"]["static_ppt_role"] == "editable_product_output"
    assert result_json["delivery_boundary"]["motion_proof_role"] == "separate_html_motion_renderer"
    assert result_json["delivery_boundary"]["native_pptx_animation_claim"] == "not_claimed"
    assert result_json["delivery_boundary"]["keynote_animation_claim"] == "not_claimed"
    assert result_json["local_outputs"]["html"].endswith("run2-17-motion-renderer-proof.html")
    assert result_json["local_outputs"]["manifest"].endswith("run2-17-motion-renderer-proof-manifest.json")

    scenes = result_json["scenes"]
    assert [scene["scene_id"] for scene in scenes] == [
        "cover_attention_reset",
        "before_after_reveal",
        "climax_scale_emphasis",
    ]
    for scene in scenes:
        assert scene["source_motion_contract_ids"]
        assert scene["animation_steps"]
        assert scene["reduced_motion_fallback"]
        assert scene["public_release_gate"] == "blocked_until_human_review_and_video_export_gate"

    assert_contains(
        result,
        [
            "Run 2.17 motion renderer proof",
            "separate HTML motion renderer",
            "not Keynote animation",
            "cover_attention_reset",
            "before_after_reveal",
            "climax_scale_emphasis",
            "public blocked",
        ],
    )


def test_run2_17_motion_renderer_proof_script_is_code_generated_and_bounded() -> None:
    script = (ROOT / "scripts" / "build_ppt_run2_17_motion_renderer_proof.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "run2-17-motion-renderer-proof.html",
            "run2-17-motion-renderer-proof-manifest.json",
            "cover_attention_reset",
            "before_after_reveal",
            "climax_scale_emphasis",
            "not_pptx_not_keynote_animation",
            "prefers-reduced-motion",
            "Play",
            "Pause",
            "Restart",
        ],
    )
    assert ".pptx" not in script.lower()


def test_run2_18_records_data_memory_workflow_thickness_pack() -> None:
    evidence = load_json(PACK / "run2_18_multimodal_evidence_expansion.json")
    memory = load_json(PACK / "run2_18_design_memory_expansion.json")
    workflow_gates = load_json(PACK / "run2_18_workflow_gate_expansion.json")

    assert evidence["status"] == "run2_18_multimodal_evidence_expansion_ready"
    assert memory["status"] == "run2_18_design_memory_expansion_ready"
    assert workflow_gates["status"] == "run2_18_workflow_gate_expansion_ready"
    for payload in [evidence, memory, workflow_gates]:
        assert payload["stage_policy"] == "repeat_same_five_layers_not_run3"
        assert payload["public_ready"] is False
        assert payload["creates_new_ppt_deck"] is False

    records = evidence["records"]
    assert len(records) >= 8
    assert {record["record_id"] for record in records} >= {
        "thick_2_18_figma_launch_identity_system",
        "thick_2_18_stripe_keynote_business_demo",
        "thick_2_18_apple_liquid_glass_motion_surface",
        "thick_2_18_duarte_slide_design_method",
    }
    for record in records:
        assert RUN2_18_EVIDENCE_FIELDS <= set(record)
        assert record["commercial_usecase_ids"]
        assert record["source_ids"]
        assert record["modality_mix"]
        assert record["memory_targets"]
        assert record["workflow_gate_targets"]
        assert_contains(record["anti_copy_boundary"], ["derived", "no copied"])
        assert "raw_transcript" not in record
        for marker in RUN2_8_FORBIDDEN_MEDIA_MARKERS:
            assert marker not in " ".join(iter_string_values(record)).lower()

    memories = memory["memory_expansions"]
    assert len(memories) >= 6
    evidence_ids = {record["record_id"] for record in records}
    for item in memories:
        assert RUN2_18_MEMORY_FIELDS <= set(item)
        assert set(item["evidence_record_ids"]) <= evidence_ids
        assert item["trace_fields_required"]
        assert_contains(item["code_generation_binding"], ["native PPT", "code"])
        assert_contains(item["negative_control_failure"], ["bad control", "fails"])

    gates = workflow_gates["gates"]
    assert len(gates) >= 6
    memory_ids = {item["memory_id"] for item in memories}
    for gate in gates:
        assert RUN2_18_WORKFLOW_GATE_FIELDS <= set(gate)
        assert gate["required_before_next_rerun"] is True
        assert set(gate["evidence_record_ids"]) <= evidence_ids
        assert set(gate["memory_ids"]) <= memory_ids
        assert gate["trace_fields"]
        assert_contains(gate["bad_control_probe"], ["bad control", "fail"])
        assert_contains(gate["release_boundary"], ["public", "blocked"])


def test_run2_18_records_thickness_result_and_no_new_ppt_output() -> None:
    result = (PACK / "results" / "run2_18_thickness_result.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_18_thickness_result.json")
    workflow = load_json(PACK / "skill_workflow.json")

    assert result_json["status"] == "run2_18_thickness_pack_ready_public_blocked"
    assert result_json["public_ready"] is False
    assert result_json["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert result_json["creates_new_ppt_deck"] is False
    assert result_json["latest_generated_ppt_run"] == "2.16"
    assert result_json["next_required_action"] == "consume_run2_18_thickness_pack_in_next_four_arm_rerun"
    assert result_json["artifact_counts"] == {
        "evidence_records": 8,
        "memory_expansions": 6,
        "workflow_gates": 6,
    }
    assert "ppt-run2-18-" not in result
    assert_contains(
        result,
        [
            "Run 2.18",
            "thickness pack",
            "data",
            "design memory",
            "workflow gate",
            "not generate a new PPT",
            "Do not advance to Run 3.0",
        ],
    )

    assert workflow["status"] == "run2_53_product_surface_scene_workflow_directed_public_blocked"
    assert {stage["id"] for stage in workflow["stages"]} >= {
        "expand_run2_18_multimodal_evidence",
        "expand_run2_18_design_memory",
        "apply_run2_18_workflow_gate_expansion",
    }
    assert any(trigger["id"] == "run2_18_thickness_pack_required_before_next_rerun" for trigger in workflow["repair_triggers"])


def test_ppt_run_html_viewer_mentions_run2_18_thickness_pack() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "run2_18_multimodal_evidence_expansion.json",
            "run2_18_design_memory_expansion.json",
            "run2_18_workflow_gate_expansion.json",
            "Run 2.18 thickness pack",
            "data/workflow thickness",
            "not a new PPT",
        ],
    )
    assert "Run 2.18" not in script.split("RUN_SPECS", 1)[1].split("def rel", 1)[0]
    assert "ppt-run2-18-" not in script


def test_ppt_run_html_viewer_mentions_run2_15_selector_artifacts() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "run2_15_layout_selector_sources.json",
            "run2_15_layout_module_memory.json",
            "run2_15_layout_selector_gate_matrix.json",
            "Run 2.15 selector",
            "layout module selector",
        ],
    )
    assert "ppt-run2-15-" not in script


def test_ppt_run_html_viewer_mentions_run2_17_motion_delivery_audit() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "run2_17_motion_delivery_audit.json",
            "Run 2.17 motion delivery audit",
            "HTML viewer is static",
            "Keynote readout",
            "static_slide_preview_only",
        ],
    )
    assert "ppt-run2-17-" not in script


def test_ppt_run_html_viewer_mentions_run2_17_motion_renderer_proof() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "run2_17_motion_renderer_proof_result.json",
            "Run 2.17 motion renderer proof",
            "run2-17-motion-renderer-proof.html",
            "separate_html_motion_renderer",
            "not_claimed",
        ],
    )
    assert "Run 2.17" not in script.split("RUN_SPECS", 1)[1].split("def rel", 1)[0]


def test_ppt_run_html_viewer_builder_tracks_latest_outputs() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "ppt-run-viewer.html",
            "Run 2.8",
            "ppt-run2-8-prompt-only",
            "ppt-run2-8-run1-5-skill",
            "ppt-run2-8-full-vulca",
            "ppt-run2-8-bad-memory-schema",
            "Run 2.9",
            "ppt-run2-9-prompt-only",
            "ppt-run2-9-run1-5-skill",
            "ppt-run2-9-full-vulca",
            "ppt-run2-9-bad-visual-primitive-memory",
            "Run 2.10",
            "ppt-run2-10-prompt-only",
            "ppt-run2-10-run1-5-skill",
            "ppt-run2-10-full-vulca",
            "ppt-run2-10-bad-visual-system-memory",
            "Run 2.13",
            "ppt-run2-13-prompt-only",
            "ppt-run2-13-run1-5-skill",
            "ppt-run2-13-full-vulca",
            "ppt-run2-13-bad-thick-data-memory",
            "Run 2.14",
            "ppt-run2-14-prompt-only",
            "ppt-run2-14-run1-5-skill",
            "ppt-run2-14-full-vulca",
            "ppt-run2-14-bad-visible-workflow-memory",
            "Run 2.16",
            "ppt-run2-16-prompt-only",
            "ppt-run2-16-run1-5-skill",
            "ppt-run2-16-full-vulca",
            "ppt-run2-16-bad-selector-memory",
            "Full skill series",
            "Data / Skill",
            "renderData",
            "safeHref",
            "sources.json",
            "run2_7_multimodal_source_records.json",
            "run2_8_tutorial_decomposition.json",
            "run2_8_executable_design_memory.json",
            "run2_8_workflow_gate_matrix.json",
            "run2_9_visual_primitive_repair.json",
            "run2_9_executable_visual_modules.json",
            "run2_9_visual_gate_matrix.json",
            "run2_10_visual_system_sources.json",
            "run2_10_visual_system_memory.json",
            "run2_10_visual_system_gate_matrix.json",
            "run2_12_thick_multimodal_evidence.json",
            "run2_12_design_memory_seed.json",
            "run2_12_workflow_gate_seed.json",
            "Run 2.9 visual primitive repair",
            "Run 2.10 visual-system sources",
            "Run 2.10 visual-system memory",
            "Run 2.10 visual-system gate matrix",
            "Run 2.12 thick multimodal evidence",
            "Run 2.12 design memory seeds",
            "Run 2.12 workflow gate seeds",
            "vulca_ppt_skill.md",
            "Why 2.8 still looks close to 2.7",
        ],
    )


def test_ppt_run_html_viewer_cli_defaults_are_repo_relative(tmp_path: Path) -> None:
    out = tmp_path / "viewer.html"
    result = subprocess.run(
        [sys.executable, str(ROOT / "scripts" / "build_ppt_run_html_viewer.py"), "--out", str(out)],
        cwd=tmp_path,
        capture_output=True,
        text=True,
        check=False,
    )

    assert result.returncode == 0, result.stderr
    payload = json.loads(result.stdout)
    assert payload["latest"] == "2.77"
    assert payload["runs"] >= 30
    assert out.exists()
    assert_contains(
        out.read_text(encoding="utf-8"),
        [
            '"latestRunId": "2.77"',
            "Run 2.58 experiment lab",
            "Open-source slide-code learning map",
        ],
    )
    assert not (tmp_path / "outputs").exists()


def test_ppt_run_html_viewer_mentions_run2_51_to_run2_52_consumption_chain() -> None:
    script = (ROOT / "scripts" / "build_ppt_run_html_viewer.py").read_text(encoding="utf-8")

    assert_contains(
        script,
        [
            "Run 2.52",
            "ppt-run2-52-prompt-only",
            "ppt-run2-52-run1-5-skill",
            "ppt-run2-52-full-vulca",
            "ppt-run2-52-bad-missing-run2-51-editorial-socket-pack",
            "run2-52-four-arm-contact-sheet.png",
            "run2_52_editorial_socket_renderer_rerun_result.json",
            "run2_51_editorial_copy_memory.json",
            "run2_51_shape_text_socket_memory.json",
            "run2_51_renderer_archetype_workflow_gates.json",
            "run2_51_editorial_shape_text_repair_result.json",
            "Run 2.51 is data/workflow-only",
            "run2_51_editorial_socket_pack_consumed_before_native_ppt_drawing",
            "bad_run2_51_missing_editorial_socket_pack",
        ],
    )
    assert "Run 2.51" not in script.split("RUN_SPECS", 1)[1].split("def rel", 1)[0]
    assert "ppt-run2-51-" not in script


def test_run2_audit_records_trace_and_native_render_blockers() -> None:
    audit = (PACK / "results" / "audit_review.md").read_text(encoding="utf-8")
    audit_json = load_json(PACK / "results" / "audit_review.json")
    delivery = (PACK / "results" / "delivery_gate.md").read_text(encoding="utf-8")
    comparison = (PACK / "results" / "comparison_report.md").read_text(encoding="utf-8")

    assert audit_json["status"] == "reviewed-public-blocked"
    assert audit_json["run2_skill_verdict"] == "best_internal_arm_not_public_ready"
    assert "trace_qa_outcome_refresh_required" in audit_json["blocking_issues"]
    assert "native_render_inspection_blocked" in audit_json["blocking_issues"]
    assert_contains(audit, ["trace QA outcome refresh required", "native render inspection blocked"])
    assert_contains(audit, ["Keynote -609 recovery", "manual export"])
    assert_contains(delivery, ["trace QA outcome fields refreshed", "blocked"])
    assert_contains(comparison, ["Trace QA outcome refresh", "not public-release claims"])


def test_run2_2_records_visual_target_rerun_result() -> None:
    result = (PACK / "results" / "run2_2_rerun_result.md").read_text(encoding="utf-8")
    result_json = load_json(PACK / "results" / "run2_2_rerun_result.json")

    assert result_json["status"] == "rerun_completed_public_blocked"
    assert result_json["public_ready"] is False
    assert result_json["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert result_json["rerun"]["status"] == "completed"
    assert result_json["rerun"]["best_internal_arm"] == "run2_2_full_skill"
    assert result_json["rerun"]["generated_outputs_committed"] is False
    assert (
        result_json["rerun"]["best_internal_arm_verdict"]
        == "visible_learning_targets_present_but_not_public_video_grade"
    )
    assert result_json["next_required_action"] == "turn_visual_targets_into_real_native_slide_thumbnails_and_rerun"
    assert_contains(
        json.dumps(result_json["visible_learning"]),
        ["schematic", "not_high_fidelity", "traceable", "placeholder"],
    )
    assert_contains(
        result,
        [
            "Run 2.2",
            "multimodal database",
            "visual learning targets",
            "schematic",
            "not public-video-grade",
            "Do not advance to Run 3.0",
            "native-PPT components",
        ],
    )


def test_run2_1_records_rerun_result_and_next_visual_depth_gate() -> None:
    readiness = (PACK / "results" / "run2_1_readiness.md").read_text(encoding="utf-8")
    readiness_json = load_json(PACK / "results" / "run2_1_readiness.json")

    assert readiness_json["status"] == "rerun_completed_public_blocked"
    assert readiness_json["next_required_action"] == "thicken_visual_taste_evidence_and_rerun"
    assert readiness_json["public_ready"] is False
    assert "extraction_units" in readiness_json["completed_depth_layers"]
    assert "skill_workflow" in readiness_json["completed_depth_layers"]
    assert "trace_refresh_utility" in readiness_json["completed_depth_layers"]
    assert readiness_json["rerun"]["status"] == "completed"
    assert readiness_json["rerun"]["best_internal_arm"] == "run2_1_full_skill"
    assert readiness_json["rerun"]["generated_outputs_committed"] is False
    assert readiness_json["rerun"]["best_internal_arm_verdict"] == "stronger_product_learning_not_public_video_grade"
    assert readiness_json["stage_policy"] == "repeat_same_five_layers_not_run3"
    assert_contains(
        readiness,
        [
            "Run 2.1",
            "rerun four arms",
            "public blocked",
            "not a new stage",
            "does not advance",
            "realistic slide diffs",
            "visual transformation",
        ],
    )


def test_run2_generation_protocol_blocks_leakage_and_rasterized_report_failures() -> None:
    body = (PACK / "generation_protocol.md").read_text(encoding="utf-8")

    assert_contains(
        body,
        [
            "runtime isolation",
            "fresh generation prompt",
            "cache directory",
            "full-slide raster",
            "image-to-native-object ratio",
            "visible text overlap",
            "clipped text",
            "default-styled tables",
        ],
    )


def test_run2_pack_does_not_commit_generated_artifacts() -> None:
    blocked_suffixes = {".jpeg", ".jpg", ".mp4", ".pdf", ".png", ".pptx"}
    blocked = [
        str(path.relative_to(PACK))
        for path in PACK.rglob("*")
        if path.is_file() and path.suffix.lower() in blocked_suffixes
    ]

    assert blocked == []

    tracked = tracked_files()
    pack_prefix = "docs/product/ppt-run2-data-skill-quality/"
    tracked_pack_artifact_suffixes = (
        ".ppt",
        ".pptx",
        ".png",
        ".jpg",
        ".jpeg",
        ".mp4",
        ".pdf",
        ".layout.json",
    )
    tracked_generated = [
        path
        for path in tracked
        if path.startswith("outputs/")
        or (path.startswith(pack_prefix) and path.lower().endswith(tracked_pack_artifact_suffixes))
        or "contact-sheet" in path
        or "/preview/" in path
    ]

    assert tracked_generated == []
